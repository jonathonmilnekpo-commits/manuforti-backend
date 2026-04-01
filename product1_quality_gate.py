#!/usr/bin/env python3
"""
Product 1 Quality Gate — Pre-delivery validation
Run this before emailing ANY Product 1 report
"""

from pptx import Presentation
import sys

def validate_product1(filepath, supplier_name):
    """Strict validation against canonical standard"""
    
    print(f"\n{'='*60}")
    print(f"PRODUCT 1 QUALITY GATE: {supplier_name}")
    print(f"{'='*60}\n")
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"✗ FAILED: Cannot open file: {e}")
        return False
    
    errors = []
    warnings = []
    
    # Check 1: Slide count
    if len(prs.slides) != 9:
        errors.append(f"Slide count: {len(prs.slides)} (expected 9)")
    else:
        print("✓ Slide count: 9")
    
    # Check 2: Slide dimensions
    width = float(prs.slide_width) / 914400
    height = float(prs.slide_height) / 914400
    if abs(width - 13.333) > 0.01 or abs(height - 7.5) > 0.01:
        errors.append(f"Dimensions: {width:.3f}\" x {height:.3f}\" (expected 13.333\" x 7.5\")")
    else:
        print(f"✓ Dimensions: {width:.3f}\" x {height:.3f}\"")
    
    # Check 3: Manu Forti logo on every slide
    logo_errors = []
    for i, slide in enumerate(prs.slides, 1):
        logo_found = False
        for shape in slide.shapes:
            if type(shape).__name__ == 'Picture':
                left = float(shape.left) / 914400
                top = float(shape.top) / 914400
                # Logo should be at bottom-right (left > 12, top > 6.5)
                if left > 12 and top > 6.5:
                    if abs(left - 12.533) < 0.05 and abs(top - 6.613) < 0.05:
                        logo_found = True
        if not logo_found:
            logo_errors.append(f"Slide {i}")
    
    if logo_errors:
        errors.append(f"Manu Forti logo missing on: {', '.join(logo_errors)}")
    else:
        print("✓ Manu Forti logo: All 9 slides")
    
    # Check 4: Source line on every slide
    source_errors = []
    for i, slide in enumerate(prs.slides, 1):
        source_found = False
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                if 'Manu Forti Intelligence' in shape.text and 'Confidential' in shape.text:
                    source_found = True
        if not source_found:
            source_errors.append(f"Slide {i}")
    
    if source_errors:
        errors.append(f"Source line missing on: {', '.join(source_errors)}")
    else:
        print("✓ Source line: All 9 slides")
    
    # Check 5: Risk gauge on Slide 2 (if exists)
    if len(prs.slides) >= 2:
        slide2 = prs.slides[1]
        gauge_found = False
        for shape in slide2.shapes:
            if type(shape).__name__ == 'Picture':
                left = float(shape.left) / 914400
                top = float(shape.top) / 914400
                width = float(shape.width) / 914400
                height = float(shape.height) / 914400
                
                # Risk gauge should be at left=0.3, top=1.05, size=7.2x3.1
                if abs(left - 0.3) < 0.1 and abs(top - 1.05) < 0.1:
                    if abs(width - 7.2) < 0.1:
                        gauge_found = True
                        if abs(height - 3.1) > 0.1:
                            errors.append(f"Risk gauge height: {height:.3f}\" (expected 3.1\")")
        
        if not gauge_found:
            warnings.append("Risk gauge not found on Slide 2 (may be placeholder)")
        else:
            print("✓ Risk gauge: Slide 2 present")
    
    # Report results
    print(f"\n{'='*60}")
    if errors:
        print("VALIDATION FAILED — DO NOT DELIVER")
        print(f"{'='*60}")
        for error in errors:
            print(f"  ✗ {error}")
        if warnings:
            for warning in warnings:
                print(f"  ⚠ {warning}")
        return False
    else:
        print("✓ ALL CHECKS PASSED — READY FOR DELIVERY")
        print(f"{'='*60}")
        if warnings:
            for warning in warnings:
                print(f"  ⚠ {warning}")
        return True

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 quality_gate.py <pptx_file> <supplier_name>")
        sys.exit(1)
    
    filepath = sys.argv[1]
    supplier_name = sys.argv[2]
    
    passed = validate_product1(filepath, supplier_name)
    sys.exit(0 if passed else 1)
