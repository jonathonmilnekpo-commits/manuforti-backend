#!/usr/bin/env python3
"""
Product 1 v15 Report Validator
Checks generated PPTX decks against the canonical v15 schema
"""

import json
import sys
from pathlib import Path
from typing import Dict, List, Tuple, Optional
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")

class Product1Validator:
    """Validates Product 1 reports against v15 schema"""
    
    def __init__(self, schema_path: str = "product1_v15_schema.json"):
        self.schema = self._load_schema(schema_path)
        self.errors: List[str] = []
        self.warnings: List[str] = []
        self.passed: List[str] = []
        
    def _load_schema(self, path: str) -> Dict:
        """Load the v15 schema from JSON"""
        with open(path, 'r') as f:
            return json.load(f)
    
    def validate_deck(self, pptx_path: str) -> Tuple[bool, List[str], List[str], List[str]]:
        """
        Validate a PPTX file against the schema
        Returns: (is_valid, errors, warnings, passed_checks)
        """
        self.errors = []
        self.warnings = []
        self.passed = []
        
        if not PPTX_AVAILABLE:
            self.errors.append("python-pptx library required for validation")
            return False, self.errors, self.warnings, self.passed
        
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            self.errors.append(f"File not found: {pptx_path}")
            return False, self.errors, self.warnings, self.passed
        
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            self.errors.append(f"Failed to open PPTX: {e}")
            return False, self.errors, self.warnings, self.passed
        
        # Run all validations
        self._validate_slide_count(prs)
        self._validate_structure(prs)
        self._validate_branding(prs)
        self._validate_design_elements(prs)
        
        is_valid = len(self.errors) == 0
        return is_valid, self.errors, self.warnings, self.passed
    
    def _validate_slide_count(self, prs):
        """Check correct number of slides"""
        expected = self.schema['slides']['count']
        actual = len(prs.slides)
        
        if actual == expected:
            self.passed.append(f"✓ Slide count: {actual} slides (expected {expected})")
        else:
            self.errors.append(f"✗ Slide count: {actual} slides (expected {expected})")
    
    def _validate_structure(self, prs):
        """Validate each slide has required elements"""
        slides_schema = self.schema['slides']['items']
        
        for i, (slide, slide_schema) in enumerate(zip(prs.slides, slides_schema)):
            slide_num = i + 1
            slide_name = slide_schema['name']
            required = slide_schema['required_elements']
            
            # Check for text-based indicators of required elements
            slide_text = self._extract_slide_text(slide).lower()
            
            # Slide-specific validations
            if slide_schema['type'] == 'title':
                self._validate_title_slide(slide, slide_num, slide_name)
            elif slide_schema['type'] == 'executive_summary':
                self._validate_executive_summary(slide, slide_num, slide_name, slide_text)
            elif slide_schema['type'] == 'recommendation':
                self._validate_recommendation(slide, slide_num, slide_name, slide_text)
            elif slide_schema['type'] == 'supplier_profile':
                self._validate_supplier_profile(slide, slide_num, slide_name, slide_text)
            elif slide_schema['type'] == 'financial_health':
                self._validate_financial_health(slide, slide_num, slide_name, slide_text)
            elif slide_schema['type'] == 'market_position':
                self._validate_market_position(slide, slide_num, slide_name)
            elif slide_schema['type'] == 'operational_risk':
                self._validate_operational_risk(slide, slide_num, slide_name, slide_text)
            elif slide_schema['type'] == 'commercial_intelligence':
                self._validate_commercial_intelligence(slide, slide_num, slide_name)
            elif slide_schema['type'] == 'esg_assessment':
                self._validate_esg_assessment(slide, slide_num, slide_name, slide_text)
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide"""
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        return " ".join(texts)
    
    def _count_shapes_with_text(self, slide, keywords: List[str]) -> int:
        """Count shapes containing any of the keywords"""
        count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_lower = shape.text.lower()
                if any(kw.lower() in text_lower for kw in keywords):
                    count += 1
        return count
    
    def _validate_title_slide(self, slide, num: int, name: str):
        """Validate slide 1 - Title"""
        text = self._extract_slide_text(slide)
        
        checks = [
            ('supplier name', ['supplier', 'analysis', 'report']),
            ('report date', ['2025', '2026', 'january', 'february', 'march']),
            ('confidentiality', ['confidential', 'private'])
        ]
        
        for check_name, keywords in checks:
            if any(kw in text for kw in keywords):
                self.passed.append(f"✓ Slide {num} ({name}): Contains {check_name}")
            else:
                self.warnings.append(f"⚠ Slide {num} ({name}): Missing {check_name}")
    
    def _validate_executive_summary(self, slide, num: int, name: str, text: str):
        """Validate slide 2 - Executive Summary"""
        required_indicators = [
            ('risk rating', ['low', 'medium', 'high', 'risk']),
            ('risk score', ['/', '100', 'score']),
            ('recommendation', ['approve', 'approved', 'reject']),
        ]
        
        for indicator, keywords in required_indicators:
            if any(kw in text for kw in keywords):
                self.passed.append(f"✓ Slide {num} ({name}): Has {indicator}")
            else:
                self.warnings.append(f"⚠ Slide {num} ({name}): Missing {indicator}")
    
    def _validate_recommendation(self, slide, num: int, name: str, text: str):
        """Validate slide 3 - Recommendation"""
        # Check for status indicators
        has_approved = 'approved' in text or 'approve' in text
        has_conditions = 'condition' in text or 'with' in text
        
        if has_approved:
            self.passed.append(f"✓ Slide {num} ({name}): Has approval status")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing approval status")
        
        if has_conditions:
            self.passed.append(f"✓ Slide {num} ({name}): Has conditions")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing conditions")
        
        # Check for prohibited green on recommendation slide
        self._check_no_green_on_recommendation(slide, num)
    
    def _check_no_green_on_recommendation(self, slide, num: int):
        """Ensure no green color is used on recommendation slide"""
        green_found = False
        for shape in slide.shapes:
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type == 1:  # SOLID fill
                        rgb = shape.fill.fore_color.rgb
                        if rgb and rgb[1] > 150 and rgb[0] < 100 and rgb[2] < 100:
                            green_found = True
                except:
                    pass
        
        if green_found:
            self.errors.append(f"✗ Slide {num}: GREEN detected on Recommendation slide (prohibited by v15)")
        else:
            self.passed.append(f"✓ Slide {num}: No green color (compliant with v15)")
    
    def _validate_supplier_profile(self, slide, num: int, name: str, text: str):
        """Validate slide 4 - Supplier Profile"""
        indicators = [
            ('company overview', ['founded', 'employees', 'headquarters', 'revenue']),
            ('leadership', ['ceo', 'cfo', 'founder', 'director', 'president']),
            ('org structure', ['subsidiary', 'parent', 'group', 'division'])
        ]
        
        for indicator, keywords in indicators:
            if any(kw in text for kw in keywords):
                self.passed.append(f"✓ Slide {num} ({name}): Has {indicator}")
            else:
                self.warnings.append(f"⚠ Slide {num} ({name}): Missing {indicator}")
    
    def _validate_financial_health(self, slide, num: int, name: str, text: str):
        """Validate slide 5 - Financial Health"""
        required_metrics = self.schema['slides']['items'][4]['required_metrics']
        
        # Check for financial keywords
        financial_keywords = [
            'revenue', 'ebitda', 'profit', 'debt', 'cash', 'cagr', 'order book'
        ]
        found_keywords = [kw for kw in financial_keywords if kw in text]
        
        if len(found_keywords) >= 5:
            self.passed.append(f"✓ Slide {num} ({name}): Contains {len(found_keywords)}/8 required financial metrics")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Only {len(found_keywords)}/8 financial metrics found")
        
        # Specifically check for debt metrics (required by Jonathon)
        debt_keywords = ['gross debt', 'net cash', 'net debt', 'debt/ebitda']
        found_debt = [kw for kw in debt_keywords if kw.replace('/', ' ').replace('  ', ' ') in text.replace('/', ' ')]
        
        if found_debt:
            self.passed.append(f"✓ Slide {num} ({name}): Has debt metrics ({', '.join(found_debt)})")
        else:
            self.errors.append(f"✗ Slide {num} ({name}): Missing required debt metrics (gross debt, net cash/debt, debt/EBITDA)")
    
    def _validate_market_position(self, slide, num: int, name: str):
        """Validate slide 6 - Market Position"""
        # Check for chart shapes
        chart_count = sum(1 for s in slide.shapes if s.shape_type == 19)  # CHART type
        
        if chart_count > 0:
            self.passed.append(f"✓ Slide {num} ({name}): Contains chart")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): No chart detected")
        
        # Check for peer/competitor mentions
        text = self._extract_slide_text(slide).lower()
        if any(kw in text for kw in ['competitor', 'peer', 'vs', 'versus', 'market share', 'position']):
            self.passed.append(f"✓ Slide {num} ({name}): Has competitive context")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing competitive context")
    
    def _validate_operational_risk(self, slide, num: int, name: str, text: str):
        """Validate slide 7 - Operational Capability + Risk"""
        # Check for risk table and matrix indicators
        has_timeline = any(kw in text for kw in ['investment', 'timeline', 'facility', 'expansion'])
        has_risk_table = any(kw in text for kw in ['risk', 'impact', 'probability', 'mitigation'])
        has_capabilities = any(kw in text for kw in ['capability', 'certification', 'iso', 'capacity'])
        
        if has_timeline:
            self.passed.append(f"✓ Slide {num} ({name}): Has investment timeline")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing investment timeline")
        
        if has_risk_table:
            self.passed.append(f"✓ Slide {num} ({name}): Has risk assessment")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing risk assessment")
        
        if has_capabilities:
            self.passed.append(f"✓ Slide {num} ({name}): Has capabilities section")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing capabilities section")
    
    def _validate_commercial_intelligence(self, slide, num: int, name: str):
        """Validate slide 8 - Commercial Intelligence"""
        # Check for charts
        chart_count = sum(1 for s in slide.shapes if s.shape_type == 19)
        
        if chart_count >= 2:
            self.passed.append(f"✓ Slide {num} ({name}): Contains both required charts")
        elif chart_count == 1:
            self.warnings.append(f"⚠ Slide {num} ({name}): Only 1 chart (expected 2: radar + lollipop)")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): No charts detected")
        
        # Check for negotiation notes
        text = self._extract_slide_text(slide).lower()
        if any(kw in text for kw in ['negotiation', 'bargaining', 'leverage', 'position', 'commercial']):
            self.passed.append(f"✓ Slide {num} ({name}): Has commercial/negotiation content")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing negotiation notes")
    
    def _validate_esg_assessment(self, slide, num: int, name: str, text: str):
        """Validate slide 9 - ESG Assessment"""
        # Check for ESG pillars
        has_env = any(kw in text for kw in ['environmental', 'iso 14001', 'scope', 'carbon'])
        has_soc = any(kw in text for kw in ['social', 'labour', 'labor', 'workforce', 'grievance'])
        has_gov = any(kw in text for kw in ['governance', 'anti-corruption', 'sanctions', 'ownership'])
        
        if has_env:
            self.passed.append(f"✓ Slide {num} ({name}): Has Environmental section")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing Environmental section")
        
        if has_soc:
            self.passed.append(f"✓ Slide {num} ({name}): Has Social section")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing Social section")
        
        if has_gov:
            self.passed.append(f"✓ Slide {num} ({name}): Has Governance section")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): Missing Governance section")
        
        # Check for conditions
        if 'condition' in text:
            conditions = text.count('condition')
            if conditions >= 6:
                self.passed.append(f"✓ Slide {num} ({name}): Has 6+ ESG conditions")
            else:
                self.warnings.append(f"⚠ Slide {num} ({name}): Only {conditions} conditions (expected 6)")
        else:
            self.warnings.append(f"⚠ Slide {num} ({name}): No ESG conditions found")
    
    def _validate_branding(self, prs):
        """Validate branding elements across all slides"""
        source_format = self.schema['branding']['source_line']['format']
        footer_text = self.schema['branding']['footer']
        prohibited = self.schema['branding']['prohibited_phrases']
        
        # Check first slide for source line format
        first_slide_text = self._extract_slide_text(prs.slides[0]).lower()
        
        if 'manu forti' in first_slide_text or 'confidential' in first_slide_text:
            self.passed.append("✓ Branding: Source line format present")
        else:
            self.warnings.append("⚠ Branding: Source line format not detected")
        
        # Check for prohibited phrases across all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            text = self._extract_slide_text(slide).lower()
            for phrase in prohibited:
                if phrase.lower() in text:
                    self.errors.append(f"✗ Slide {slide_num}: Prohibited phrase found: '{phrase}'")
        
        if not any('Prohibited phrase' in e for e in self.errors):
            self.passed.append("✓ Branding: No prohibited phrases found")
    
    def _validate_design_elements(self, prs):
        """Validate design system compliance"""
        # This is a simplified check - full color validation would require more detailed analysis
        self.passed.append("✓ Design: Deck loaded successfully")
    
    def generate_report(self, output_path: str = None) -> str:
        """Generate a validation report"""
        lines = [
            "=" * 60,
            "PRODUCT 1 v15 REPORT VALIDATION",
            "=" * 60,
            f"Schema Version: {self.schema.get('schema_version', 'unknown')}",
            f"Template: {self.schema.get('template_version', 'unknown')} ({self.schema.get('template_status', '')})",
            "-" * 60,
            ""
        ]
        
        if self.passed:
            lines.extend(["✓ PASSED CHECKS:", ""])
            lines.extend([f"  {p}" for p in self.passed])
            lines.append("")
        
        if self.warnings:
            lines.extend(["⚠ WARNINGS:", ""])
            lines.extend([f"  {w}" for w in self.warnings])
            lines.append("")
        
        if self.errors:
            lines.extend(["✗ ERRORS:", ""])
            lines.extend([f"  {e}" for e in self.errors])
            lines.append("")
        
        # Summary
        total = len(self.passed) + len(self.warnings) + len(self.errors)
        lines.extend([
            "-" * 60,
            f"SUMMARY: {len(self.passed)} passed, {len(self.warnings)} warnings, {len(self.errors)} errors",
            "-" * 60,
            f"STATUS: {'✓ VALID' if len(self.errors) == 0 else '✗ INVALID'}"
        ])
        
        report = "\n".join(lines)
        
        if output_path:
            with open(output_path, 'w') as f:
                f.write(report)
        
        return report


def main():
    """CLI entry point"""
    if len(sys.argv) < 2:
        print("Usage: python validate_report.py <pptx_file> [output_report.txt]")
        print("\nExample:")
        print("  python validate_report.py Boskalis_Product1_v15_Final.pptx")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    validator = Product1Validator()
    is_valid, errors, warnings, passed = validator.validate_deck(pptx_file)
    
    report = validator.generate_report(output_file)
    print(report)
    
    sys.exit(0 if is_valid else 1)


if __name__ == "__main__":
    main()
