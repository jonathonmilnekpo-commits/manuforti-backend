#!/usr/bin/env python3
"""Extract Jarotech Product 1 formatting specifications"""

from pptx import Presentation
import json

prs = Presentation('/Users/jonathonmilne/Documents/Jarotech_AS_Product1_v15_PIPELINE.pptx')

def emu_to_inches(emu):
    return float(emu) / 914400

print('=== JAROTECH PRODUCT 1 V15 - CANONICAL FORMAT ===')
print()
print('Slide Dimensions: %.3f" x %.3f"' % (emu_to_inches(prs.slide_width), emu_to_inches(prs.slide_height)))
print()

# Analyze each slide
for slide_idx, slide in enumerate(prs.slides, 1):
    print(f'=== SLIDE {slide_idx} ===')
    
    for shape in slide.shapes:
        shape_type = type(shape).__name__
        
        # Get position and size
        left = emu_to_inches(shape.left)
        top = emu_to_inches(shape.top)
        width = emu_to_inches(shape.width)
        height = emu_to_inches(shape.height)
        
        # Get text if available
        text = ""
        if hasattr(shape, 'text'):
            text = shape.text.strip()[:100]
        
        print(f"  {shape_type}:")
        print(f"    Position: left={left:.3f}\", top={top:.3f}\"")
        print(f"    Size: width={width:.3f}\", height={height:.3f}\"")
        if text:
            print(f"    Text: '{text}'")
    print()
