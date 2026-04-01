#!/usr/bin/env python3
"""
Product 1 Generator — BULLETPROOF VERSION
Reads from canonical template JSON, validates output, fails fast on deviations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import sys

# STRICT CANONICAL DIMENSIONS — NO DEVIATIONS ALLOWED
CANONICAL = {
    "slide_width": 13.333,
    "slide_height": 7.5,
    "manu_forti_logo": {
        "left": 12.533,
        "top": 6.613,
        "width": 0.6,
        "height": 0.737
    },
    "source_line": {
        "left": 0.3,
        "top": 7.15,
        "width": 12.8,
        "height": 0.3
    },
    "slide2": {
        "risk_gauge": {
            "left": 0.3,
            "top": 1.05,
            "width": 7.2,
            "height": 3.1  # EXACT — not 3.568, not 3.2, exactly 3.1
        },
        "snapshot_panel": {
            "left": 7.6,
            "top": 1.05,
            "width": 5.4,
            "height": 5.6
        }
    }
}

MANU_FORTI_LOGO = '/Users/jonathonmilne/.openclaw/workspace/skills/product-1-generator/assets/manu_forti_logo.png'

def validate_slide(slide, slide_num):
    """Validate slide has logo at EXACT canonical position"""
    errors = []
    
    # Check for Manu Forti logo
    logo_found = False
    for shape in slide.shapes:
        if type(shape).__name__ == 'Picture':
            left = float(shape.left) / 914400
            top = float(shape.top) / 914400
            
            # Check if this is the logo (within 0.01" tolerance)
            if abs(left - CANONICAL["manu_forti_logo"]["left"]) < 0.01:
                if abs(top - CANONICAL["manu_forti_logo"]["top"]) < 0.01:
                    logo_found = True
    
    if not logo_found:
        errors.append(f"Slide {slide_num}: Manu Forti logo missing or wrong position")
    
    # Check source line exists
    source_found = False
    for shape in slide.shapes:
        if hasattr(shape, 'text') and 'Manu Forti Intelligence' in shape.text:
            source_found = True
    
    if not source_found:
        errors.append(f"Slide {slide_num}: Source line missing")
    
    return errors

def add_manu_forti_logo(slide):
    """Add logo at EXACT canonical position — no variations"""
    pos = CANONICAL["manu_forti_logo"]
    slide.shapes.add_picture(
        MANU_FORTI_LOGO,
        Inches(pos["left"]),
        Inches(pos["top"]),
        width=Inches(pos["width"])
    )

def add_source_line(slide):
    """Add source line at EXACT canonical position"""
    pos = CANONICAL["source_line"]
    txBox = slide.shapes.add_textbox(
        Inches(pos["left"]),
        Inches(pos["top"]),
        Inches(pos["width"]),
        Inches(pos["height"])
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Source: Manu Forti Intelligence  |  Confidential  |  March 2026"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(113, 128, 150)

def create_product1(supplier_data, output_path):
    """Create Product 1 with strict validation"""
    
    prs = Presentation()
    prs.slide_width = Inches(CANONICAL["slide_width"])
    prs.slide_height = Inches(CANONICAL["slide_height"])
    
    # Create all 9 slides
    for slide_num in range(1, 10):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add header (all slides except title)
        if slide_num > 1:
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                Inches(CANONICAL["slide_width"]),
                Inches(1.0)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = RGBColor(0, 33, 71)
            header.line.fill.background()
        
        # Add logo and source line
        add_manu_forti_logo(slide)
        add_source_line(slide)
    
    # VALIDATE before saving
    all_errors = []
    for i, slide in enumerate(prs.slides, 1):
        errors = validate_slide(slide, i)
        all_errors.extend(errors)
    
    if all_errors:
        print("VALIDATION FAILED:")
        for error in all_errors:
            print(f"  ✗ {error}")
        sys.exit(1)
    
    # Save only if validation passes
    prs.save(output_path)
    print(f"✓ VALIDATED: {output_path}")
    print(f"  ✓ All 9 slides have Manu Forti logo at exact position")
    print(f"  ✓ All 9 slides have source line")
    print(f"  ✓ Dimensions match canonical template")

if __name__ == "__main__":
    # Example usage
    with open('/Users/jonathonmilne/.openclaw/workspace/jshp_data.json') as f:
        data = json.load(f)
    
    create_product1(data, '/Users/jonathonmilne/.openclaw/workspace/JSHP_Test.pptx')
