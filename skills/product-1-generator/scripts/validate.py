#!/usr/bin/env python3
"""
Product 1 v16 QC Validator
Validates PPTX output against v16 standards before delivery.
"""

import sys
from pathlib import Path
from pptx import Presentation

def validate_product1(pptx_path):
    """
    Validates Product 1 report against v16 checklist.
    Returns (passed, errors, warnings)
    """
    errors = []
    warnings = []
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return False, [f"Cannot open PPTX: {e}"], []
    
    # Check 1: File size (should be >500KB to indicate embedded charts)
    file_size = Path(pptx_path).stat().st_size / 1024
    if file_size < 500:
        errors.append(f"File size {file_size:.0f}KB too small - charts may not be embedded")
    elif file_size < 1000:
        warnings.append(f"File size {file_size:.0f}KB - verify all charts present")
    
    # Check 2: Slide count (must be exactly 9)
    slide_count = len(prs.slides)
    if slide_count != 9:
        errors.append(f"Expected 9 slides, found {slide_count}")
    
    # Check 3: Each slide has content
    for i, slide in enumerate(prs.slides, 1):
        if len(slide.shapes) < 2:
            errors.append(f"Slide {i} has insufficient content ({len(slide.shapes)} shapes)")
    
    # Check 4: Look for chart images (indicates visual content)
    chart_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                chart_count += 1
    
    if chart_count < 5:
        errors.append(f"Only {chart_count} images found - expected 9 charts")
    elif chart_count < 9:
        warnings.append(f"Found {chart_count} images - verify all 9 charts present")
    
    # Check 5: Text content validation (no placeholder text)
    placeholder_indicators = ['[Enter', '[Placeholder', '[Insert', 'TODO', 'XXX']
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.upper()
                for indicator in placeholder_indicators:
                    if indicator.upper() in text:
                        errors.append(f"Slide {i} contains placeholder text: '{indicator}'")
    
    # Check 6: Source line present on slides
    source_found = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and 'Source:' in shape.text:
                source_found += 1
                break
    
    if source_found < 9:
        warnings.append(f"Source line missing on {9 - source_found} slides")
    
    # Summary
    passed = len(errors) == 0
    
    print(f"\n{'='*60}")
    print(f"Product 1 v16 QC Validation: {pptx_path.name}")
    print(f"{'='*60}")
    print(f"File size: {file_size:.1f} KB")
    print(f"Slides: {slide_count}")
    print(f"Charts/images: {chart_count}")
    print(f"Source lines: {source_found}/9")
    print(f"{'='*60}")
    
    if errors:
        print("\n❌ ERRORS (must fix before delivery):")
        for e in errors:
            print(f"  • {e}")
    
    if warnings:
        print("\n⚠️  WARNINGS (review recommended):")
        for w in warnings:
            print(f"  • {w}")
    
    if passed and not warnings:
        print("\n✅ ALL CHECKS PASSED - Ready for delivery")
    elif passed:
        print("\n✅ PASSED with warnings - Review before delivery")
    else:
        print("\n❌ FAILED - Fix errors before delivery")
    
    print(f"{'='*60}\n")
    
    return passed, errors, warnings

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python validate.py <pptx_file>")
        sys.exit(1)
    
    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    
    passed, errors, warnings = validate_product1(pptx_path)
    sys.exit(0 if passed else 1)
