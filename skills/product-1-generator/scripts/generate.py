#!/usr/bin/env python3
"""
Product 1 v15 Generator
Creates PPTX from structured JSON input.
"""

import json
import sys
from pathlib import Path
from datetime import datetime

# Try importing with fallback paths
import sys
from pathlib import Path

# Add common site-packages paths
for path in [
    '/opt/homebrew/lib/python3.14/site-packages',
    '/opt/homebrew/lib/python3.13/site-packages', 
    '/usr/local/lib/python3.11/site-packages',
    str(Path.home() / '.local/lib/python3.11/site-packages'),
    str(Path.home() / 'Library/Python/3.11/lib/python/site-packages'),
]:
    if Path(path).exists():
        sys.path.insert(0, path)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"Error: python-pptx not installed ({e})", file=sys.stderr)
    sys.exit(1)

# Color constants (RGB tuples)
NAVY = (0x00, 0x21, 0x47)
STEEL_BLUE = (0x2B, 0x6C, 0xB0)
MID_GREY = (0x71, 0x80, 0x96)

# V15 Color Palette (RGB tuples)
NAVY = (0x00, 0x21, 0x47)
STEEL_BLUE = (0x2B, 0x6C, 0xB0)
MID_GREY = (0x71, 0x80, 0x96)
GREEN = (0x48, 0xBB, 0x78)
AMBER = (0xD6, 0x9E, 0x2E)
RED = (0xE5, 0x3E, 0x3E)

def create_title_slide(prs, data):
    """Create slide 1: Title"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = f"SUPPLIER EVALUATION REPORT\n{data['supplier']}"
    
    # Add subtitle with stats
    subtitle = slide.placeholders[1]
    subtitle.text = f"{data.get('sector', 'Industrial Services')}\n{data.get('stats', '')}"
    
    return slide

def create_exec_summary_slide(prs, data):
    """Create slide 2: Executive Summary"""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "EXECUTIVE SUMMARY"
    
    # Add risk gauge and summary
    summary = data.get('executive_summary', {})
    risk = summary.get('overall_risk', 'MEDIUM')
    
    content = slide.placeholders[1]
    content.text = f"""{data['supplier']} — Risk & Suitability Overview

OVERALL RISK: {risk}
Risk Score: {summary.get('risk_score', 'N/A')}/100

Key Insight:
{summary.get('key_insight', 'No insight provided')}

Risk Factors:
{chr(10).join(['• ' + f for f in summary.get('risk_factors', [])])}
"""
    return slide

def create_recommendation_slide(prs, data):
    """Create slide 3: Recommendation"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "RECOMMENDATION"
    
    rec = data.get('recommendation', {})
    verdict = rec.get('verdict', 'FLAG FOR REVIEW')
    
    conditions = chr(10).join(['• ' + c for c in rec.get('conditions', [])])
    
    content = slide.placeholders[1]
    content.text = f"""RECOMMENDATION: {verdict}

{data.get('supplier')} is recommended for {data.get('sector', 'relevant contracts')}.

Commercial Conditions:
{conditions}

Risk Summary:
• Financial: {rec.get('risk_summary', {}).get('financial', 'N/A')}
• Operational: {rec.get('risk_summary', {}).get('operational', 'N/A')}
• Geopolitical: {rec.get('risk_summary', {}).get('geopolitical', 'N/A')}
• ESG: {rec.get('risk_summary', {}).get('esg', 'N/A')}
"""
    return slide

def create_profile_slide(prs, data):
    """Create slide 4: Supplier Profile"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "SUPPLIER PROFILE"
    
    profile = data.get('profile', {})
    leadership = profile.get('leadership', [])
    
    leaders_text = chr(10).join([f"{l['name']} — {l['role']}\n  {l.get('tenure', '')}" for l in leadership[:3]])
    
    content = slide.placeholders[1]
    content.text = f"""Corporate Structure & Global Footprint

Company Overview:
{profile.get('description', 'No description provided')}

Leadership:
{leaders_text}

Headquarters: {profile.get('headquarters', 'N/A')}
Employees: {profile.get('employees', 'N/A')}
"""
    return slide

def create_financial_slide(prs, data):
    """Create slide 5: Financial Health"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "FINANCIAL HEALTH"
    
    fin = data.get('financials', {})
    
    content = slide.placeholders[1]
    content.text = f"""Revenue & EBITDA Trajectory

Financial Highlights:
• Revenue 2024: {fin.get('revenue_2024', 'N/A')} ({fin.get('revenue_yoy', 'N/A')})
• EBITDA: {fin.get('ebitda', 'N/A')} ({fin.get('ebitda_margin', 'N/A')} margin)
• Net Profit: {fin.get('net_profit', 'N/A')}
• 3-Year CAGR: {fin.get('cagr_3yr', 'N/A')}
• Order Book: {fin.get('order_book', 'N/A')}
• Gross Debt: {fin.get('gross_debt', 'N/A')}
• Net Cash: {fin.get('net_cash', 'N/A')}
• Debt/EBITDA: {fin.get('debt_ebitda', 'N/A')}

Trend: {fin.get('trend', 'Stable')}
"""
    return slide

def create_market_slide(prs, data):
    """Create slide 6: Market Position"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "MARKET POSITION"
    
    market = data.get('market_position', {})
    competitors = market.get('competitors', [])
    
    comp_text = chr(10).join([f"• {c['name']}: {c.get('revenue', 'N/A')}" for c in competitors[:5]])
    advantages = chr(10).join([f"• {a}" for a in market.get('advantages', [])])
    
    content = slide.placeholders[1]
    content.text = f"""Marine Contractor Revenue Comparison — 2024

Competitive Landscape:
{data['supplier']} ranks #{market.get('rank', 'N/A')} globally.

Key Competitors:
{comp_text}

Competitive Advantages:
{advantages}
"""
    return slide

def create_risk_slide(prs, data):
    """Create slide 7: Operational Capability & Risk"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "OPERATIONAL CAPABILITY & RISK ASSESSMENT"
    
    risks = data.get('risks', [])
    risk_text = chr(10).join([f"• {r['category']}: {r['description'][:60]}... (Impact: {r.get('impact', 'N/A')})" for r in risks[:5]])
    
    content = slide.placeholders[1]
    content.text = f"""Milestones & Strategic Investment | Risk Matrix

Risk Summary:
• Financial: {data.get('recommendation', {}).get('risk_summary', {}).get('financial', 'LOW')}
• Operational: {data.get('recommendation', {}).get('risk_summary', {}).get('operational', 'LOW')}
• Geopolitical: {data.get('recommendation', {}).get('risk_summary', {}).get('geopolitical', 'MEDIUM')}
• ESG: {data.get('recommendation', {}).get('risk_summary', {}).get('esg', 'MEDIUM')}

Key Risks:
{risk_text}
"""
    return slide

def create_commercial_slide(prs, data):
    """Create slide 8: Commercial Intelligence"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "COMMERCIAL INTELLIGENCE & PEER RISK COMPARISON"
    
    content = slide.placeholders[1]
    content.text = f"""Benchmarking Radar | Peer Risk Profile

Commercial Terms:
• Standard payment terms: Net-30 to Net-60 days
• Milestone/progress payments standard
• Performance bonds typically required

Key Watch Points:
⚠ Monitor project-country risk individually
⚠ Review order book concentration
⚠ Track working capital requirements
"""
    return slide

def create_esg_slide(prs, data):
    """Create slide 9: ESG Assessment"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "ESG ASSESSMENT"
    
    esg = data.get('esg', {})
    controversies = chr(10).join([f"⚠ {c}" for c in esg.get('controversies', [])])
    
    content = slide.placeholders[1]
    content.text = f"""Environmental, Social & Governance Screening

Overall Rating: {esg.get('overall', 'MEDIUM')}

Pillar Ratings:
• Environmental: {esg.get('environmental', 'MEDIUM')}
• Social: {esg.get('social', 'LOW')}
• Governance: {esg.get('governance', 'LOW')}

Controversies:
{controversies if controversies else 'None identified'}

Source: Manu Forti Intelligence | Confidential | {datetime.now().strftime('%B %Y')}
"""
    return slide

def generate_pptx(input_json: str, output_pptx: str):
    """Generate v15 PPTX from JSON input."""
    
    # Load input data
    with open(input_json) as f:
        data = json.load(f)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Generate all 9 slides
    create_title_slide(prs, data)
    create_exec_summary_slide(prs, data)
    create_recommendation_slide(prs, data)
    create_profile_slide(prs, data)
    create_financial_slide(prs, data)
    create_market_slide(prs, data)
    create_risk_slide(prs, data)
    create_commercial_slide(prs, data)
    create_esg_slide(prs, data)
    
    # Save
    prs.save(output_pptx)
    print(f"Generated: {output_pptx}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: generate.py <input.json> <output.pptx>")
        sys.exit(1)
    
    generate_pptx(sys.argv[1], sys.argv[2])
