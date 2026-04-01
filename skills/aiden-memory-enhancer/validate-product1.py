#!/usr/bin/env python3
"""
Product 1 Validator - Checks output against canonical v15 template
ENHANCED: Now validates graphics/charts on each slide
"""

import json
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict

class Product1Validator:
    """Validates Product 1 PPTX against canonical v15 requirements with graphics checking"""
    
    REQUIRED_SLIDES = 9
    REQUIRED_METRICS = [
        "Revenue",
        "EBITDA", 
        "Net Profit",
        "CAGR",
        "Order Book",
        "Gross Debt",
        "Net Cash",
        "Debt/EBITDA"
    ]
    
    # Graphics requirements per slide
    GRAPHICS_REQUIREMENTS = {
        2: {
            "name": "Executive Summary",
            "required": ["risk_gauge"],
            "check": lambda shapes, text: any(
                s.shape_type == MSO_SHAPE_TYPE.PICTURE or 
                ("low" in text and "medium" in text and "high" in text)
                for s in shapes
            ),
            "description": "Risk gauge dial/visual (LOW/MEDIUM/HIGH)"
        },
        5: {
            "name": "Financial Health",
            "required": ["dual_axis_chart"],
            "check": lambda shapes, text: any(
                hasattr(s, "chart") or "revenue" in text and "ebitda" in text
                for s in shapes
            ),
            "description": "Dual-axis chart (revenue + EBITDA)"
        },
        6: {
            "name": "Market Position",
            "required": ["horizontal_bar_chart"],
            "check": lambda shapes, text: any(
                hasattr(s, "chart") and "competitor" in text
                for s in shapes
            ),
            "description": "Horizontal bar chart vs competitors"
        },
        7: {
            "name": "Operational + Risk",
            "required": ["timeline", "risk_matrix"],
            "check_timeline": lambda shapes, text: any(
                str(year) in text for year in range(2015, 2027)
            ),
            "check_matrix": lambda shapes, text: any(
                s.shape_type == MSO_SHAPE_TYPE.TABLE or 
                ("impact" in text and "probability" in text)
                for s in shapes
            ),
            "description": "Timeline + 2x2 risk matrix"
        },
        8: {
            "name": "Commercial Intelligence",
            "required": ["radar_chart", "peer_risk_chart"],
            "check": lambda shapes, text: any(
                hasattr(s, "chart") and ("benchmark" in text or "peer" in text)
                for s in shapes
            ),
            "description": "Radar chart + peer risk comparison"
        },
        9: {
            "name": "ESG Assessment",
            "required": ["esg_columns"],
            "check": lambda shapes, text: any(
                s.shape_type == MSO_SHAPE_TYPE.TABLE or
                all(esg in text for esg in ["environmental", "social", "governance"])
                for s in shapes
            ),
            "description": "E/S/G three-column layout"
        }
    }
    
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
        self.errors = []
        self.warnings = []
        self.graphics_errors = []
        
    def validate(self) -> Dict:
        """Run all validation checks including graphics"""
        self._check_slide_count()
        self._check_graphics_on_all_slides()
        self._check_slide_5_financials()
        self._check_branding()
        self._check_logos()
        
        return {
            "valid": len(self.errors) == 0 and len(self.graphics_errors) == 0,
            "errors": self.errors,
            "warnings": self.warnings,
            "graphics_errors": self.graphics_errors,
            "slide_count": len(self.prs.slides),
            "graphics_check": self._get_graphics_summary()
        }
    
    def _check_slide_count(self):
        """Verify exactly 9 slides"""
        if len(self.prs.slides) != self.REQUIRED_SLIDES:
            self.errors.append(
                f"Slide count: {len(self.prs.slides)}, expected {self.REQUIRED_SLIDES}"
            )
    
    def _check_graphics_on_all_slides(self):
        """Check required graphics on each slide"""
        for slide_num, requirements in self.GRAPHICS_REQUIREMENTS.items():
            if slide_num > len(self.prs.slides):
                continue
                
            slide = self.prs.slides[slide_num - 1]
            slide_text = " ".join([s.text.lower() for s in slide.shapes if hasattr(s, "text")])
            
            # Check using the custom check function
            if "check" in requirements:
                if not requirements["check"](slide.shapes, slide_text):
                    self.graphics_errors.append(
                        f"Slide {slide_num} ({requirements['name']}): Missing {requirements['description']}"
                    )
            
            # Special handling for slide 7 (two checks)
            if slide_num == 7:
                if not requirements["check_timeline"](slide.shapes, slide_text):
                    self.graphics_errors.append(
                        f"Slide 7: Missing timeline (no year markers found)"
                    )
                if not requirements["check_matrix"](slide.shapes, slide_text):
                    self.graphics_errors.append(
                        f"Slide 7: Missing 2x2 risk matrix (no table or impact/probability text)"
                    )
    
    def _get_graphics_summary(self) -> Dict:
        """Get summary of graphics found on each slide"""
        summary = {}
        for i, slide in enumerate(self.prs.slides, 1):
            images = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
            charts = sum(1 for s in slide.shapes if hasattr(s, "chart"))
            tables = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE)
            
            summary[f"slide_{i}"] = {
                "images": images,
                "charts": charts,
                "tables": tables,
                "has_visual": images + charts + tables > 0
            }
        return summary
    
    def _check_slide_5_financials(self):
        """Check Slide 5 has all required financial metrics"""
        if len(self.prs.slides) < 5:
            return
            
        slide = self.prs.slides[4]
        slide_text = ""
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        
        missing_metrics = []
        for metric in self.REQUIRED_METRICS:
            if metric.lower() not in slide_text.lower():
                missing_metrics.append(metric)
        
        if missing_metrics:
            self.errors.append(
                f"Slide 5: Missing required metrics: {', '.join(missing_metrics)}"
            )
    
    def _check_branding(self):
        """Check for Manu Forti branding"""
        found_branding = False
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "Manu Forti" in shape.text:
                        found_branding = True
                        break
        
        if not found_branding:
            self.warnings.append("Missing Manu Forti branding")
    
    def _check_logos(self):
        """Check logos present on slides"""
        slide_1 = self.prs.slides[0]
        has_logo = False
        
        for shape in slide_1.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_logo = True
                break
        
        if not has_logo:
            self.warnings.append("Slide 1: May be missing supplier logo")


def main():
    if len(sys.argv) < 2:
        print("Usage: validate-product1.py <pptx_file>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    validator = Product1Validator(pptx_path)
    result = validator.validate()
    
    print(f"\n{'='*70}")
    print(f"Product 1 v15 Validation with Graphics Check")
    print(f"{'='*70}")
    print(f"File: {pptx_path}")
    print(f"Status: {'✓ VALID' if result['valid'] else '✗ INVALID'}")
    print(f"Slides: {result['slide_count']}")
    
    if result['errors']:
        print(f"\n❌ Structure Errors ({len(result['errors'])}):")
        for error in result['errors']:
            print(f"   • {error}")
    
    if result['graphics_errors']:
        print(f"\n🎨 Graphics Errors ({len(result['graphics_errors'])}):")
        for error in result['graphics_errors']:
            print(f"   • {error}")
    
    if result['warnings']:
        print(f"\n⚠️  Warnings ({len(result['warnings'])}):")
        for warning in result['warnings']:
            print(f"   • {warning}")
    
    print(f"\n📊 Graphics Summary:")
    for slide, gfx in result['graphics_check'].items():
        status = "✓" if gfx['has_visual'] else "✗"
        print(f"   {status} {slide}: {gfx['images']} images, {gfx['charts']} charts, {gfx['tables']} tables")
    
    print(f"\n{'='*70}\n")
    
    # Update skill usage
    update_skill_usage(result['valid'], result['errors'] + result['graphics_errors'])
    
    sys.exit(0 if result['valid'] else 1)


def update_skill_usage(success: bool, errors: List[str]):
    """Update skill-usage.json with validation results"""
    import os
    from datetime import datetime
    
    usage_file = os.path.expanduser(
        "~/.openclaw/workspace/memory/skill-usage.json"
    )
    
    try:
        with open(usage_file, 'r') as f:
            data = json.load(f)
    except:
        data = {"skills": {}}
    
    if "product-1-generator" not in data["skills"]:
        data["skills"]["product-1-generator"] = {
            "times_used": 0,
            "successes": 0,
            "failures": 0,
            "validation_failures": []
        }
    
    skill_data = data["skills"]["product-1-generator"]
    skill_data["times_used"] = skill_data.get("times_used", 0) + 1
    skill_data["last_validated"] = datetime.now().isoformat()
    
    if success:
        skill_data["successes"] = skill_data.get("successes", 0) + 1
    else:
        skill_data["failures"] = skill_data.get("failures", 0) + 1
        skill_data["validation_failures"] = skill_data.get("validation_failures", [])
        skill_data["validation_failures"].extend(errors)
        skill_data["validation_failures"] = skill_data["validation_failures"][-10:]
    
    with open(usage_file, 'w') as f:
        json.dump(data, f, indent=2)
    
    print(f"✓ Skill usage updated: {'success' if success else 'failure'} recorded")


if __name__ == "__main__":
    main()