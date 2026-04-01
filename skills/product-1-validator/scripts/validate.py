#!/usr/bin/env python3
"""
Product 1 v15 Validator - Enhanced with Graphics Validation
Validates PPTX against locked v15 template structure AND checks all visuals are present.
"""

import sys
import json
from pathlib import Path
from typing import List, Dict, Tuple

# Try importing with fallback paths
for path in [
    '/opt/homebrew/lib/python3.14/site-packages',
    '/opt/homebrew/lib/python3.13/site-packages',
    '/usr/local/lib/python3.11/site-packages',
    str(Path.home() / '.local/lib/python3.11/site-packages'),
    str(Path.home() / 'Library/Python/3.11/lib/python/site-packages'),
]:
    if Path(path).exists():
        sys.path.insert(0, path)

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    print(json.dumps({
        "valid": False,
        "error": f"python-pptx not installed: {e}"
    }))
    sys.exit(1)

# V15 Template Requirements
REQUIRED_SLIDES = 9
REQUIRED_METRICS = [
    "revenue", "ebitda", "net profit", "cagr", "order book",
    "gross debt", "net cash", "debt/ebitda"
]
REQUIRED_SECTIONS = [
    "executive summary",
    "recommendation", 
    "supplier profile",
    "financial health",
    "market position",
    "operational capability",
    "commercial intelligence",
    "esg assessment"
]

# Graphics requirements per slide
GRAPHICS_REQUIREMENTS = {
    2: {
        "name": "Executive Summary",
        "required": ["risk_gauge"],
        "descriptions": {
            "risk_gauge": "Risk gauge dial/visual (LOW/MEDIUM/HIGH indicator)"
        }
    },
    5: {
        "name": "Financial Health",
        "required": ["dual_axis_chart"],
        "descriptions": {
            "dual_axis_chart": "Dual-axis chart (revenue bars + EBITDA line)"
        }
    },
    6: {
        "name": "Market Position",
        "required": ["horizontal_bar_chart"],
        "descriptions": {
            "horizontal_bar_chart": "Horizontal bar chart vs named competitors"
        }
    },
    7: {
        "name": "Operational + Risk",
        "required": ["timeline", "risk_matrix"],
        "descriptions": {
            "timeline": "Strategic investment timeline",
            "risk_matrix": "2x2 risk matrix (Impact vs Probability)"
        }
    },
    8: {
        "name": "Commercial Intelligence",
        "required": ["radar_chart", "peer_risk_chart"],
        "descriptions": {
            "radar_chart": "Spider/radar benchmarking chart",
            "peer_risk_chart": "Peer risk comparison (lollipop or bar)"
        }
    },
    9: {
        "name": "ESG Assessment",
        "required": ["esg_columns"],
        "descriptions": {
            "esg_columns": "E/S/G three-column layout with ratings"
        }
    }
}


def detect_chart_type(shape) -> str:
    """Detect what type of chart/visual a shape represents."""
    
    # Check for pictures/images (risk gauge, logos)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    
    # Check for charts
    if hasattr(shape, "chart"):
        chart = shape.chart
        if hasattr(chart, "chart_type"):
            chart_type = str(chart.chart_type)
            
            # Detect radar/spider
            if "RADAR" in chart_type or "SPIDER" in chart_type:
                return "radar_chart"
            
            # Detect bar charts
            if "BAR" in chart_type:
                # Horizontal vs vertical
                if "HORIZONTAL" in chart_type or "CLUSTERED" in chart_type:
                    return "horizontal_bar_chart"
                return "bar_chart"
            
            # Detect line charts
            if "LINE" in chart_type:
                return "line_chart"
            
            # Detect combo/dual-axis
            if "COMBO" in chart_type:
                return "dual_axis_chart"
        
        return "chart"
    
    # Check for tables (ESG columns, risk matrix)
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"
    
    # Check for grouped shapes (complex visuals)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    
    return "unknown"


def analyze_slide_graphics(slide, slide_num: int) -> Dict:
    """Analyze what graphics are present on a slide."""
    
    graphics_found = {
        "images": 0,
        "charts": 0,
        "tables": 0,
        "groups": 0,
        "chart_types": [],
        "has_visual": False
    }
    
    for shape in slide.shapes:
        shape_type = detect_chart_type(shape)
        
        if shape_type == "image":
            graphics_found["images"] += 1
            graphics_found["has_visual"] = True
        elif shape_type in ["chart", "radar_chart", "horizontal_bar_chart", 
                           "bar_chart", "line_chart", "dual_axis_chart"]:
            graphics_found["charts"] += 1
            graphics_found["chart_types"].append(shape_type)
            graphics_found["has_visual"] = True
        elif shape_type == "table":
            graphics_found["tables"] += 1
            graphics_found["has_visual"] = True
        elif shape_type == "group":
            graphics_found["groups"] += 1
            graphics_found["has_visual"] = True
    
    return graphics_found


def check_specific_graphics(slide_num: int, graphics_found: Dict, slide_text: str) -> List[str]:
    """Check for specific required graphics on each slide."""
    
    missing = []
    
    if slide_num not in GRAPHICS_REQUIREMENTS:
        return missing
    
    req = GRAPHICS_REQUIREMENTS[slide_num]
    
    for required in req["required"]:
        found = False
        
        # Check by chart type
        if required == "risk_gauge":
            # Risk gauge is typically an image or has specific text
            if graphics_found["images"] > 0 or "low" in slide_text and "medium" in slide_text:
                found = True
        
        elif required == "dual_axis_chart":
            if "dual_axis_chart" in graphics_found["chart_types"] or graphics_found["charts"] > 0:
                found = True
        
        elif required == "horizontal_bar_chart":
            if "horizontal_bar_chart" in graphics_found["chart_types"]:
                found = True
            elif graphics_found["charts"] > 0 and "competitor" in slide_text:
                found = True
        
        elif required == "timeline":
            # Timeline often a group or has year markers
            if graphics_found["groups"] > 0 or any(year in slide_text for year in ["2020", "2021", "2022", "2023", "2024", "2025"]):
                found = True
        
        elif required == "risk_matrix":
            # Risk matrix is a table or 2x2 grid
            if graphics_found["tables"] > 0 or graphics_found["groups"] > 0:
                found = True
            elif "impact" in slide_text and "probability" in slide_text:
                found = True
        
        elif required == "radar_chart":
            if "radar_chart" in graphics_found["chart_types"]:
                found = True
            elif "spider" in slide_text or "benchmark" in slide_text:
                found = True
        
        elif required == "peer_risk_chart":
            if graphics_found["charts"] > 0 and "peer" in slide_text:
                found = True
        
        elif required == "esg_columns":
            # ESG should have table or three distinct sections
            if graphics_found["tables"] > 0:
                found = True
            elif all(esg in slide_text for esg in ["environmental", "social", "governance"]):
                found = True
        
        if not found:
            missing.append(req["descriptions"][required])
    
    return missing


def validate_pptx(pptx_path: str) -> dict:
    """Validate a Product 1 v15 PPTX file with graphics checking."""
    
    errors = []
    warnings = []
    graphics_errors = []
    score = 100
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"valid": False, "error": f"Cannot open PPTX: {e}"}
    
    # Check slide count
    slide_count = len(prs.slides)
    if slide_count != REQUIRED_SLIDES:
        errors.append(f"Slide count: {slide_count} (required: {REQUIRED_SLIDES})")
        score -= 20
    
    # Extract all text content and analyze graphics per slide
    all_text = []
    slide_texts = []
    slide_graphics = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.lower())
        
        slide_text_combined = " ".join(slide_text)
        slide_texts.append(slide_text_combined)
        all_text.extend(slide_text)
        
        # Analyze graphics on this slide
        graphics = analyze_slide_graphics(slide, i)
        slide_graphics.append(graphics)
        
        # Check for required graphics on this slide
        missing_graphics = check_specific_graphics(i, graphics, slide_text_combined)
        if missing_graphics:
            for mg in missing_graphics:
                graphics_errors.append(f"Slide {i}: Missing {mg}")
            score -= len(missing_graphics) * 8
    
    full_text = " ".join(all_text)
    
    # Check required sections
    missing_sections = []
    for section in REQUIRED_SECTIONS:
        if section not in full_text:
            missing_sections.append(section)
    
    if missing_sections:
        errors.append(f"Missing sections: {', '.join(missing_sections)}")
        score -= len(missing_sections) * 10
    
    # Check financial metrics
    missing_metrics = []
    for metric in REQUIRED_METRICS:
        if metric not in full_text:
            missing_metrics.append(metric)
    
    if missing_metrics:
        errors.append(f"Slide 5 missing metrics: {', '.join(missing_metrics)}")
        score -= len(missing_metrics) * 5
    
    # Check branding
    if "manu forti intelligence" not in full_text:
        warnings.append("Missing 'Manu Forti Intelligence' branding")
        score -= 5
    
    # Check recommendation slide (slide 3)
    if len(slide_texts) >= 3:
        rec_text = slide_texts[2]
        if "approve" not in rec_text:
            warnings.append("Slide 3: No recommendation verdict found")
        if "condition" not in rec_text and "decline" not in rec_text:
            warnings.append("Slide 3: No conditions or decline rationale found")
    
    # Check ESG slide (slide 9)
    if len(slide_texts) >= 9:
        esg_text = slide_texts[8]
        esg_elements = ["environmental", "social", "governance"]
        missing_esg = [e for e in esg_elements if e not in esg_text]
        if missing_esg:
            errors.append(f"Slide 9 missing ESG elements: {', '.join(missing_esg)}")
            score -= 10
    
    # Add graphics errors to main errors list
    if graphics_errors:
        errors.extend(graphics_errors)
    
    # Normalize score
    score = max(0, min(100, score))
    
    # Determine validity
    valid = score >= 90 and len(errors) == 0
    
    # Build graphics summary
    graphics_summary = {}
    for i, gfx in enumerate(slide_graphics, 1):
        graphics_summary[f"slide_{i}"] = {
            "images": gfx["images"],
            "charts": gfx["charts"],
            "tables": gfx["tables"],
            "has_visual": gfx["has_visual"]
        }
    
    return {
        "valid": valid,
        "score": score,
        "errors": errors,
        "warnings": warnings,
        "slide_count": slide_count,
        "sections_found": [s for s in REQUIRED_SECTIONS if s in full_text],
        "metrics_found": [m for m in REQUIRED_METRICS if m in full_text],
        "graphics_summary": graphics_summary,
        "graphics_errors": graphics_errors
    }


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: validate.py <pptx_file>")
        sys.exit(1)
    
    result = validate_pptx(sys.argv[1])
    print(json.dumps(result, indent=2))
    
    # Exit code
    if result.get("valid"):
        sys.exit(0)
    elif result.get("errors"):
        sys.exit(2)
    else:
        sys.exit(1)