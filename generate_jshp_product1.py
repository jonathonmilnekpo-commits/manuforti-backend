#!/usr/bin/env python3
"""
JSHP Transformer Product 1 — v15 Final
Complete 9-slide supplier evaluation report following canonical v15 template
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Palette (v15 locked) ──────────────────────────────────────────────────────
NAVY      = RGBColor(  0,  33,  71)  # #002147
COBALT    = RGBColor( 43, 108, 176)  # #2B6CB0
GREEN     = RGBColor( 72, 187, 120)  # #48BB78
AMBER     = RGBColor(237, 137,  54)  # #ED8936
RED       = RGBColor(229,  62,  62)  # #E53E3E
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(113, 128, 150)  # #718096
LGRAY     = RGBColor(226, 232, 240)  # #E2E8F0
BLACK     = RGBColor(  0,   0,   0)
PALE_BLUE = RGBColor(160, 200, 240)
PALE_GRAY = RGBColor(200, 200, 200)

VISUALS = "/Users/jonathonmilne/.openclaw/workspace/jshp_v15_visuals"
LOGO    = "/Users/jonathonmilne/.openclaw/workspace/manu_forti_logo.png"

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def header(slide, title, subtitle=None):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = WHITE
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
        sp = s.text_frame.paragraphs[0]
        sp.text = subtitle; sp.font.size = Pt(10); sp.font.italic = True
        sp.font.color.rgb = PALE_BLUE

def manu_forti_logo(slide):
    """Add Manu Forti logo bottom-right every slide — exact v15 positioning."""
    try:
        slide.shapes.add_picture(LOGO, Inches(12.533), Inches(6.613), width=Inches(0.6))
    except:
        # Fallback if logo not found — create placeholder
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(12.533), Inches(6.613), Inches(0.6), Inches(0.737))
        box.fill.solid(); box.fill.fore_color.rgb = COBALT
        box.line.fill.background()

def source(slide):
    s = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    p = s.text_frame.paragraphs[0]
    p.text = "Source: Manu Forti Intelligence  |  Confidential  |  March 2026"
    p.font.size = Pt(8); p.font.color.rgb = GRAY

def txt(slide, t, x, y, w, h, size=11, bold=False, col=None, wrap=True, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.text = t; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = col if col else BLACK
    if align: p.alignment = align

def blist(slide, title, items, x, y, w, h, title_col=None, size=10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title; p.font.bold = True; p.font.size = Pt(12)
    p.font.color.rgb = title_col or NAVY
    for item in items:
        p = tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = BLACK
        p.space_after = Pt(2)

def box(slide, x, y, w, h, fill, lc=None, rnd=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def img(slide, fname, x, y, w):
    slide.shapes.add_picture(f'{VISUALS}/{fname}', Inches(x), Inches(y), width=Inches(w))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
box(slide, 0, 0, 13.333, 7.5, NAVY)
box(slide, 0, 5.55, 13.333, 0.08, AMBER)

txt(slide, "SUPPLIER EVALUATION REPORT", 0.7, 1.4, 12, 0.5, size=15, col=PALE_BLUE)
txt(slide, "JiangSu HuaPeng Transformer Co., Ltd.", 0.7, 2.05, 12, 1.2, size=36, bold=True, col=WHITE)
txt(slide, "Power Transformers  |  Medium Power Transformers  |  Global Export",
    0.7, 3.5, 12, 0.6, size=19, col=PALE_GRAY)
box(slide, 0.7, 4.35, 4.5, 0.05, COBALT)
txt(slide, "US$1.0B Revenue  |  2,500 Employees  |  57 Years  |  200,000 MVA Capacity",
    0.7, 4.55, 12, 0.4, size=12, col=PALE_GRAY)
txt(slide, "Confidential  |  March 2026  |  Manu Forti Intelligence",
    0.7, 6.2, 12, 0.4, size=11, col=RGBColor(150, 150, 150))

manu_forti_logo(slide)

manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY (MEDIUM 48/100)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "EXECUTIVE SUMMARY",
       "JiangSu HuaPeng Transformer Co., Ltd. — Risk & Suitability Overview")

# Risk gauge — exact positioning from canonical template
# position: left=0.3", top=1.05", width=7.2", height=3.1"
img(slide, "01_risk_gauge.png", 0.3, 1.05, 7.2)

# Supplier snapshot panel — exact positioning
# position: left=7.6", top=1.05", width=5.4", height=5.6"
box(slide, 7.6, 1.05, 5.4, 5.6, LGRAY, rnd=True)
txt(slide, "SUPPLIER SNAPSHOT", 7.8, 1.22, 5.2, 0.32, size=11, bold=True, col=NAVY)

snapshot = [
    ("Supplier",   "JiangSu HuaPeng Transformer Co., Ltd."),
    ("Type",       "Family-Owned (Private)"),
    ("Sector",     "Power Transformer Manufacturing"),
    ("HQ",         "Liyang, Jiangsu, China"),
    ("Founded",    "1967 (57 years)"),
    ("Revenue",    "US$1.0B (2022)"),
    ("Employees",  "2,500 globally"),
    ("Capacity",   "200,000 MVA, up to 850kV"),
]
y = 1.64
for label, val in snapshot:
    txt(slide, label + ":", 7.8, y, 2.3, 0.27, size=9.5, bold=True, col=NAVY)
    txt(slide, val, 10.1, y, 2.7, 0.27, size=9.5)
    y += 0.30

# Key findings strip
box(slide, 0.3, 5.22, 7.2, 1.43, LGRAY, rnd=True)
txt(slide, "Key Findings", 0.5, 5.33, 7.0, 0.30, size=11, bold=True, col=NAVY)
findings = [
    "• World's largest Medium Power Transformer producer — market leader position",
    "• 57-year track record with zero catastrophic failures — exceptional quality record",
    "• Only Chinese company in top 10 North American transformer brands — unique position",
    "• 100% China-based manufacturing creates geopolitical exposure — conditions required",
]
y = 5.65
for f in findings:
    txt(slide, f, 0.5, y, 7.0, 0.26, size=9.5)
    y += 0.27

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: RECOMMENDATION (AMBER / APPROVE WITH CONDITIONS)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "RECOMMENDATION", "Decision Summary & Commercial Conditions")

# Amber recommendation banner
box(slide, 0.5, 1.1, 12.333, 1.35, AMBER, rnd=True)
txt(slide, "⚠  RECOMMENDATION: APPROVE WITH CONDITIONS",
    0.7, 1.18, 12, 0.55, size=22, bold=True, col=WHITE)
txt(slide,
    "JSHP is recommended for Medium Power Transformer supply and select High Power projects. "
    "World's largest MPT capacity and zero failure record support approval. "
    "Conditions apply due to 100% China manufacturing and geopolitical exposure.",
    0.7, 1.78, 12, 0.58, size=11, col=WHITE)

# Commercial Conditions
blist(slide, "Commercial Conditions", [
    "• Fixed pricing or hedging for multi-year contracts",
    "• Parent company guarantee for orders >US$5M",
    "• Performance bonds: 10% for large projects",
    "• Payment terms: Letter of Credit for initial orders",
    "• Confirm US sales office support for NA projects",
], 0.5, 2.65, 6.0, 2.5)

# Approval Conditions
blist(slide, "Approval Conditions (Mandatory)", [
    "• Quality inspection rights at Liyang facility",
    "• Quarterly production capacity confirmation",
    "• Geopolitical risk monitoring and mitigation plan",
    "• Alternative sourcing plan for critical projects",
    "• Currency hedging for CNY/USD/EUR exposure",
], 6.7, 2.65, 6.0, 2.5, title_col=AMBER)

# Risk summary strip — 4 pillars
box(slide, 0.5, 5.35, 12.333, 1.5, LGRAY, rnd=True)
txt(slide, "Overall Risk Summary — MEDIUM (48/100)", 0.7, 5.50, 12, 0.30, size=11, bold=True, col=NAVY)
risk_items = [
    ("Financial",    "LOW",    "US$1.0B revenue, debt-free family ownership"),
    ("Operational",  "LOW",    "200k MVA capacity, zero failure record"),
    ("Geopolitical", "HIGH",   "100% China manufacturing, trade policy risk"),
    ("ESG",          "MEDIUM", "Family governance, limited transparency"),
]
x = 0.6
for cat, rating, note in risk_items:
    col = GREEN if rating == "LOW" else (AMBER if rating == "MEDIUM" else RED)
    txt(slide, f"{cat}: {rating}", x, 5.88, 3.1, 0.27, size=9.5, bold=True, col=col)
    txt(slide, note,               x, 6.17, 3.1, 0.27, size=8.5, col=GRAY)
    x += 3.1

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: SUPPLIER PROFILE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "SUPPLIER PROFILE", "Corporate Structure & Global Footprint")

img(slide, "02_org_structure.png", 0.3, 1.05, 8.2)

blist(slide, "Company Overview", [
    "JiangSu HuaPeng Transformer Co., Ltd. (JSHP)",
    "is a leading Chinese manufacturer of power",
    "transformers, specializing in Medium Power",
    "Transformers where it holds the position of",
    "world's largest producer.",
    "",
    "Founded in 1967, the family-owned company",
    "has built a 57-year track record with zero",
    "catastrophic failures across thousands of",
    "installed units globally.",
], 8.7, 1.05, 4.3, 3.5)

blist(slide, "Manufacturing Excellence", [
    "• Annual capacity: 200,000 MVA",
    "• Voltage range: Up to 850kV, 1000MVA",
    "• 50-day avg fabrication (110kV-500kV)",
    "• 2022 production: 628 units (110-500kV)",
    "• 8,764 units (69kV and below)",
    "• ISO 9001 certified",
], 8.7, 4.5, 4.3, 2.4)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: FINANCIAL HEALTH
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "FINANCIAL HEALTH",
       "Revenue Trajectory 2019–2022  |  Family-Owned: Limited Public Financial Data")

img(slide, "03_financial_trend.png", 0.3, 1.05, 7.5)

R = 8.1   # right panel x-start
txt(slide, "Financial Highlights", R, 1.10, 4.9, 0.35, size=12, bold=True, col=NAVY)

metrics = [
    ("2022 Revenue",     "US$1.0B  (est.)"),
    ("Revenue Growth",   "+33% (2019–2022)"),
    ("Ownership",        "Family-owned, private"),
    ("Debt Status",      "Minimal / Debt-free"),
    ("Delivered 2022",   "111,500 MVA"),
    ("Liquidity",        "Strong (private reserves)"),
    ("Transparency",     "Limited (private co.)"),
    ("Credit Risk",      "LOW — stable ownership"),
]
y = 1.56
for label, val in metrics:
    val_col = GREEN if label in ("Debt Status", "Credit Risk") else BLACK
    txt(slide, label + ":", R,     y, 2.55, 0.29, size=9.5, bold=True, col=NAVY)
    txt(slide, val,          R+2.55, y, 2.25, 0.29, size=9.5, col=val_col)
    y += 0.30

# GREEN LOW box
box(slide, R, 4.05, 4.8, 0.44, GREEN, rnd=True)
txt(slide, "Financial Risk:  LOW  ✓", R+0.15, 4.12, 4.5, 0.30, size=12, bold=True, col=WHITE)

blist(slide, "Exposure Guidance", [
    "• Debt-free family ownership — low credit risk",
    "• Parent guarantee recommended for large orders",
    "• Limited financial transparency — private company",
    "• Payment: LC recommended for initial orders",
], R, 4.60, 4.8, 2.30)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: MARKET POSITION
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "MARKET POSITION",
       "Transformer Manufacturing — Revenue Comparison")

img(slide, "04_market_position.png", 0.3, 1.05, 8.0)

blist(slide, "Competitive Landscape", [
    "JSHP occupies a unique position as the only",
    "Chinese manufacturer in the top 10 North",
    "American transformer brands. While smaller",
    "than global giants like Siemens and ABB,",
    "JSHP dominates the Medium Power Transformer",
    "segment with unmatched specialization.",
], 8.5, 1.1, 4.5, 2.0)

blist(slide, "JSHP Competitive Advantages", [
    "• World's largest MPT producer by volume",
    "• Only Chinese company in top 10 NA brands",
    "• Zero catastrophic failure record",
    "• 50-day fabrication time (competitive)",
    "• Strong value proposition vs. European peers",
    "• 200,000 MVA manufacturing capacity",
], 8.5, 3.2, 4.5, 2.6)

txt(slide, "Key Competitors", 0.5, 5.5, 12, 0.30, size=12, bold=True, col=NAVY)
comps = [
    "Siemens Energy — Germany — €30B+ — Global leader, premium pricing",
    "ABB — Switzerland — €30B+ — Broad portfolio, strong service network",
    "Hitachi Energy — Japan — €10B+ — Technology leader, high reliability",
    "TBEA — China — €5B+ — Domestic China leader, expanding globally",
]
y = 5.84
for c in comps:
    txt(slide, "• " + c, 0.6, y, 12, 0.27, size=9.5)
    y += 0.29

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: OPERATIONAL CAPABILITY + RISK ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "OPERATIONAL CAPABILITY  &  RISK ASSESSMENT",
       "Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary")

img(slide, "05_investment_timeline.png", 0.3,  1.05, 7.0)
img(slide, "06_risk_matrix.png",         7.65, 1.05, 5.5)

# Bottom section
BOTY = 5.58
txt(slide, "Risk Summary", 0.3, BOTY, 7.0, 0.28, size=11, bold=True, col=NAVY)

risk_rows = [
    ("Geopolitical",       "🔴 HIGH",   "100% China manufacturing — tariff/trade exposure"),
    ("Trade Policy",       "🟡 MEDIUM", "US/EU tariffs on Chinese electrical equipment"),
    ("Supply Chain",       "🔴 HIGH",   "Concentrated China supply base"),
    ("Operational",        "🟢 LOW",    "Zero failures, ISO 9001 certified"),
    ("Financial",          "🟢 LOW",    "Debt-free family ownership"),
]
y = BOTY + 0.31
for cat, rating, note in risk_rows:
    if "🔴" in rating:
        col = RED
    elif "🟡" in rating:
        col = AMBER
    else:
        col = GREEN
    txt(slide, cat + ":",  0.3, y, 2.1, 0.24, size=8.5, bold=True, col=NAVY)
    txt(slide, rating,     2.4, y, 1.4, 0.24, size=8.5, bold=True, col=col)
    txt(slide, note,       3.8, y, 3.7, 0.24, size=8.0, col=GRAY)
    y += 0.25

# Key Capabilities
blist(slide, "Key Capabilities", [
    "• ISO 9001 certified quality management",
    "• 57 years continuous operation",
    "• 50+ countries served globally",
    "• US sales office for NA support",
], 7.65, BOTY, 5.5, 1.85)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: COMMERCIAL INTELLIGENCE + PEER RISK
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON",
       "Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms")

img(slide, "07_benchmarking_radar.png", 0.2, 1.0, 5.8)
img(slide, "08_peer_risk.png",          6.2, 1.0, 6.9)

box(slide, 0.3, 5.8, 12.7, 1.50, LGRAY, rnd=True)

blist(slide, "Commercial Terms & Negotiation", [
    "• Pricing: 20-30% below European competitors",
    "• Lead time: 50 days average (110kV-500kV)",
    "• MOQ: Flexible for established relationships",
    "• Warranty: Standard 2-year, extended available",
], 0.5, 5.85, 6.5, 1.4, size=9)

blist(slide, "Key Watch Points", [
    "⚠ Geopolitical: Monitor US-China trade policy developments",
    "⚠ Supply chain: Confirm raw material sourcing stability",
    "⚠ Currency: Hedge CNY/USD exposure on long-term contracts",
], 7.0, 5.85, 5.8, 1.4, title_col=AMBER, size=9)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: ESG ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "ESG ASSESSMENT",
       "Environmental, Social & Governance Screening  |  Overall Rating: MEDIUM")

img(slide, "09_esg_assessment.png", 1.67, 1.05, 10.0)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/jonathonmilne/.openclaw/workspace/JSHP_Transformer_Product1_v15_Final.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print("  9 slides | v15: MEDIUM (48/100) | Manu Forti branding | March 2026")
