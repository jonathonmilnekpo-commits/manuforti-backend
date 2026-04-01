#!/usr/bin/env python3
"""
OX2 Product 1 — v15 Final
9-slide supplier evaluation report for OX2 AB (Swedish renewable energy developer)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Palette (v15 locked) ──────────────────────────────────────────────────────
NAVY      = RGBColor(  0,  33,  71)   # #002147 - Wood Mackenzie navy
STEEL_BLUE= RGBColor( 43, 108, 176)   # #2B6CB0
GREEN     = RGBColor( 72, 187, 120)   # #48BB78
AMBER     = RGBColor(214, 158,  46)   # #D69E2E
RED       = RGBColor(229,  62,  62)   # #E53E3E
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(127, 140, 141)
LGRAY     = RGBColor(236, 240, 241)
BLACK     = RGBColor(  0,   0,   0)
PALE_BLUE = RGBColor(160, 200, 240)
PALE_GRAY = RGBColor(200, 200, 200)
OX2_BLUE  = RGBColor(  0,  90, 140)   # OX2 corporate blue

VISUALS = "/Users/jonathonmilne/.openclaw/workspace/ox2_v15_visuals"

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def header(slide, title, subtitle=None):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = WHITE
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
        sp = s.text_frame.paragraphs[0]
        sp.text = subtitle; sp.font.size = Pt(10); sp.font.italic = True
        sp.font.color.rgb = PALE_BLUE

def logo_tag(slide):
    """Add OX2 logo placeholder — text-based top-right of header."""
    # White rounded backing
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(10.92), Inches(0.06), Inches(2.26), Inches(0.87))
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    # OX2 text logo
    tx = slide.shapes.add_textbox(Inches(11.0), Inches(0.15), Inches(2.1), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    p.text = "OX2"; p.font.bold = True; p.font.size = Pt(28)
    p.font.color.rgb = OX2_BLUE; p.alignment = PP_ALIGN.CENTER

def manu_forti_logo(slide):
    """Add Manu Forti logo at bottom right of slide."""
    tx = slide.shapes.add_textbox(Inches(11.0), Inches(6.9), Inches(2.0), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "Manu Forti Intelligence"; p.font.size = Pt(8)
    p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.RIGHT

def source(slide):
    s = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(10), Inches(0.3))
    p = s.text_frame.paragraphs[0]
    p.text = "Source: Manu Forti Intelligence  |  Confidential  |  March 2026"
    p.font.size = Pt(8); p.font.color.rgb = GRAY

def txt(slide, t, x, y, w, h, size=11, bold=False, col=None, wrap=True, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.text = t; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = col if col else BLACK
    if align: p.alignment = align

def blist(slide, title, items, x, y, w, h, title_col=None, size=10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title; p.font.bold = True; p.font.size = Pt(12)
    p.font.color.rgb = title_col or NAVY
    for item in items:
        p = tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = BLACK
        p.space_after = Pt(2)

def box(slide, x, y, w, h, fill, lc=None, rnd=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def img(slide, fname, x, y, w):
    slide.shapes.add_picture(f'{VISUALS}/{fname}', Inches(x), Inches(y), width=Inches(w))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
box(slide, 0, 0, 13.333, 7.5, NAVY)
box(slide, 0, 5.55, 13.333, 0.08, GREEN)

txt(slide, "SUPPLIER EVALUATION REPORT", 0.7, 1.4, 12, 0.5, size=15, col=PALE_BLUE)
txt(slide, "OX2 AB", 0.7, 2.05, 12, 1.2, size=48, bold=True, col=WHITE)
txt(slide, "Renewable Energy Development  |  Onshore & Offshore Wind  |  Solar PV & Storage",
    0.7, 3.5, 12, 0.6, size=18, col=PALE_GRAY)
box(slide, 0.7, 4.35, 4.5, 0.05, STEEL_BLUE)
txt(slide, "SEK 5.4B Revenue  |  ~612 Employees  |  10+ Markets  |  EQT Owned",
    0.7, 4.55, 12, 0.4, size=12, col=PALE_GRAY)
txt(slide, "Confidential  |  March 2026  |  Manu Forti Intelligence",
    0.7, 6.2, 12, 0.4, size=11, col=RGBColor(150, 150, 150))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "EXECUTIVE SUMMARY",
       "OX2 AB — Risk & Suitability Overview")
logo_tag(slide)

# Risk gauge
img(slide, "01_risk_gauge.png", 0.3, 1.05, 7.2)

# Supplier snapshot panel
box(slide, 7.6, 1.05, 5.4, 5.6, LGRAY, rnd=True)
txt(slide, "SUPPLIER SNAPSHOT", 7.8, 1.22, 5.2, 0.32, size=11, bold=True, col=NAVY)

snapshot = [
    ("Supplier",   "OX2 AB"),
    ("Ownership",  "EQT Private Equity (Oct 2024)"),
    ("Sector",     "Renewable Energy Development"),
    ("HQ",         "Stockholm, Sweden"),
    ("Founded",    "2004 (19+ years experience)"),
    ("Revenue",    "SEK 5.4B ($520M est. 2024)"),
    ("EBITDA",     "SEK 1,000M+ ($96M+, ~18% margin)"),
    ("Pipeline",   "Multi-GW development portfolio"),
    ("Employees",  "~612 globally"),
]
y = 1.64
for label, val in snapshot:
    txt(slide, label + ":", 7.8, y, 2.3, 0.27, size=9.5, bold=True, col=NAVY)
    txt(slide, val, 10.1, y, 2.7, 0.27, size=9.5)
    y += 0.30

# Key findings strip
box(slide, 0.3, 5.22, 7.2, 1.43, LGRAY, rnd=True)
txt(slide, "Key Findings", 0.5, 5.33, 7.0, 0.30, size=11, bold=True, col=NAVY)
findings = [
    "• EQT acquisition (Oct 2024) provides financial stability and long-term capital",
    "• Pure-play renewable developer with 19+ years track record in Nordics/Baltics",
    "• Revenue volatility due to project timing — typical for development business",
    "• Multi-GW pipeline across onshore wind, offshore wind, solar PV, and storage",
]
y = 5.65
for f in findings:
    txt(slide, f, 0.5, y, 7.0, 0.26, size=9.5)
    y += 0.27

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: RECOMMENDATION (GREEN / APPROVE)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "RECOMMENDATION", "Decision Summary & Commercial Conditions")
logo_tag(slide)

# Green recommendation banner
box(slide, 0.5, 1.1, 12.333, 1.35, GREEN, rnd=True)
txt(slide, "✓  RECOMMENDATION: APPROVE",
    0.7, 1.28, 12, 0.55, size=24, bold=True, col=WHITE)
txt(slide,
    "OX2 is recommended for renewable energy project development partnerships. Strong EQT backing, "
    "proven track record, and pure-play renewable focus make them a suitable counterparty. "
    "Standard project finance structures and milestone-based payments recommended.",
    0.7, 1.88, 12, 0.48, size=11, col=WHITE)

# Commercial Conditions
blist(slide, "Commercial Conditions", [
    "• Milestone-based payments tied to project approvals and milestones",
    "• Standard development fee structures for project origination",
    "• EQT parent company guarantee for large-scale commitments",
    "• Performance security via escrow for upfront development costs",
    "• Regular financial reporting (quarterly) given private ownership",
], 0.5, 2.65, 6.0, 2.5)

# Risk Mitigation
blist(slide, "Risk Mitigation Measures", [
    "• Verify EQT commitment letter for long-term capital support",
    "• Secure step-in rights if development milestones missed",
    "• Contingency planning for grid connection delays",
    "• Diversify across multiple OX2 projects to reduce concentration",
    "• Monitor regulatory changes in key markets (Sweden, Finland, Poland)",
], 6.7, 2.65, 6.0, 2.5, title_col=STEEL_BLUE)

# Risk summary strip — 4 pillars
box(slide, 0.5, 5.35, 12.333, 1.5, LGRAY, rnd=True)
txt(slide, "Overall Risk Summary — LOW (35/100)", 0.7, 5.50, 12, 0.30, size=11, bold=True, col=NAVY)
risk_items = [
    ("Financial",    "LOW",    "SEK 5.4B revenue, PE-backed, strong liquidity"),
    ("Operational",  "LOW",    "19+ years experience, multi-GW pipeline"),
    ("Geopolitical", "LOW",    "Sweden HQ, EU/NATO aligned, diversified markets"),
    ("ESG",          "LOW",    "Pure-play renewable, strong sustainability focus"),
]
x = 0.6
for cat, rating, note in risk_items:
    col = GREEN if rating == "LOW" else AMBER
    txt(slide, f"{cat}: {rating}", x, 5.88, 3.1, 0.27, size=9.5, bold=True, col=col)
    txt(slide, note,               x, 6.17, 3.1, 0.27, size=8.5, col=GRAY)
    x += 3.1

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: SUPPLIER PROFILE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "SUPPLIER PROFILE", "Corporate Structure & Global Footprint")
logo_tag(slide)

img(slide, "02_org_structure.png", 0.3, 1.05, 8.2)

blist(slide, "Company Overview", [
    "OX2 is a leading European developer of renewable",
    "energy projects, specializing in onshore wind,",
    "offshore wind, solar PV, and energy storage.",
    "Founded in 2004, the company has developed",
    "a multi-gigawatt project pipeline across",
    "Northern Europe and Australia.",
    "",
    "Acquired by EQT Private Equity in October 2024,",
    "providing long-term capital and strategic backing",
    "for continued growth and project development.",
], 8.7, 1.05, 4.3, 3.5)

blist(slide, "Leadership", [
    "Paul Stormoen — CEO",
    "  Founder and industry veteran",
    "  20+ years in renewables",
    "",
    "HQ: Stockholm, Sweden",
    "Founded: 2004",
    "Delisted: October 2024",
], 8.7, 4.5, 4.3, 2.4)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: FINANCIAL HEALTH
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "FINANCIAL HEALTH",
       "Revenue & EBITDA Trajectory 2021–2024  |  EQT Backed  |  Strong Liquidity")
logo_tag(slide)

img(slide, "03_financial_trend.png", 0.3, 1.05, 7.5)

R = 8.1   # right panel x-start
txt(slide, "Financial Highlights", R, 1.10, 4.9, 0.35, size=12, bold=True, col=NAVY)

metrics = [
    ("2024 Revenue",     "SEK 5.4B  (-20% YoY)"),
    ("2024 EBITDA",      "SEK 1,000M+  (~18% margin)"),
    ("2023 Net Profit",  "SEK 800M+  ($77M+)"),
    ("3yr Revenue CAGR", "+9%  (2021–2023)"),
    ("Order Book",       "Multi-GW pipeline"),
    ("Gross Debt",       "Moderate (PE managed)"),
    ("Net Cash Position","Positive (EQT backed)"),
    ("Debt / EBITDA",    "<1.0x  ✓  (conservative)"),
]
y = 1.56
for label, val in metrics:
    val_col = GREEN if label in ("Liquidity", "Debt Profile") else BLACK
    txt(slide, label + ":", R,     y, 2.55, 0.29, size=9.5, bold=True, col=NAVY)
    txt(slide, val,          R+2.55, y, 2.25, 0.29, size=9.5, col=val_col)
    y += 0.30

# GREEN LOW box
box(slide, R, 4.05, 4.8, 0.44, GREEN, rnd=True)
txt(slide, "Financial Risk:  LOW  ✓", R+0.15, 4.12, 4.5, 0.30, size=12, bold=True, col=WHITE)

blist(slide, "Exposure Guidance", [
    "• EQT backing provides financial stability",
    "• Milestone-based payments reduce counterparty risk",
    "• No parent guarantee required for standard projects",
    "• Escrow recommended for large upfront commitments",
], R, 4.60, 4.8, 2.30)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: MARKET POSITION
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "MARKET POSITION",
       "Renewable Developer Revenue Comparison — 2024 Estimates")
logo_tag(slide)

img(slide, "04_market_position.png", 0.3, 1.05, 8.0)

blist(slide, "Competitive Landscape", [
    "OX2 operates as a focused pure-play developer",
    "with deep expertise in Nordic and Baltic markets.",
    "While smaller than utility-scale competitors,",
    "OX2's specialized focus and local relationships",
    "provide competitive advantages in target markets.",
], 8.5, 1.1, 4.5, 2.0)

blist(slide, "OX2 Competitive Advantages", [
    "• 19+ years Nordic/Baltic market expertise",
    "• Multi-GW development pipeline",
    "• EQT backing for long-term growth capital",
    "• Integrated development capabilities",
    "• Strong track record of project sales",
    "• Pure-play renewable focus (no fossil legacy)",
], 8.5, 3.2, 4.5, 2.6)

txt(slide, "Key Competitors", 0.5, 5.5, 12, 0.30, size=12, bold=True, col=NAVY)
comps = [
    "Ørsted — Denmark — $9.6B — Global offshore wind leader",
    "Vattenfall — Sweden — $15B+ — Swedish state-owned utility",
    "RWE — Germany — $20B+ — Major German utility/renewables",
    "Iberdrola — Spain — $25B+ — Spanish renewable giant",
    "Enel Green Power — Italy — $20B+ — Italian global renewables",
]
y = 5.84
for c in comps:
    txt(slide, "• " + c, 0.6, y, 12, 0.27, size=9.5)
    y += 0.29

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: OPERATIONAL CAPABILITY + RISK ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "OPERATIONAL CAPABILITY  &  RISK ASSESSMENT",
       "Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary")
logo_tag(slide)

img(slide, "05_investment_timeline.png", 0.3,  1.05, 7.0)
img(slide, "06_risk_matrix.png",         7.65, 1.05, 5.5)

# Bottom section
BOTY = 5.58
txt(slide, "Risk Summary", 0.3, BOTY, 7.0, 0.28, size=11, bold=True, col=NAVY)

risk_rows = [
    ("Financial",          "🟢 LOW",    "SEK 5.4B revenue; PE-backed stability"),
    ("Geopolitical",       "🟢 LOW",    "Sweden HQ; EU/NATO aligned markets"),
    ("Project Timing",     "🟡 MEDIUM", "Revenue volatility from project cycles"),
    ("Policy/Regulatory",  "🟡 MEDIUM", "Dependent on renewable energy support"),
    ("Competition",        "🟢 LOW",    "Differentiated niche positioning"),
]
y = BOTY + 0.31
for cat, rating, note in risk_rows:
    col = GREEN if "🟢" in rating else AMBER
    txt(slide, cat + ":",  0.3, y, 2.1, 0.24, size=8.5, bold=True, col=NAVY)
    txt(slide, rating,     2.4, y, 1.4, 0.24, size=8.5, bold=True, col=col)
    txt(slide, note,       3.8, y, 3.7, 0.24, size=8.0, col=GRAY)
    y += 0.25

# Key Capabilities
blist(slide, "Key Capabilities", [
    "• ISO certifications for project management",
    "• 19 years continuous renewable development",
    "• 10+ markets — diversified geographic presence",
    "• Full chain: origination → development → sale",
], 7.65, BOTY, 5.5, 1.85)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: COMMERCIAL INTELLIGENCE + PEER RISK
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON",
       "Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms")
logo_tag(slide)

img(slide, "07_benchmarking_radar.png", 0.2, 1.0, 5.8)
img(slide, "08_peer_risk.png",          6.2, 1.0, 6.9)

box(slide, 0.3, 5.8, 12.7, 1.50, LGRAY, rnd=True)

blist(slide, "Commercial Terms & Negotiation", [
    "• Pricing: Development fees based on project capacity/MW",
    "• Payment: Milestones for permitting, PPA, construction ready",
    "• Leverage: Volume commitment; framework agreements",
    "• IP: Client owns project rights; OX2 retains method IP",
], 0.5, 5.85, 6.5, 1.4, size=9)

blist(slide, "Key Watch Points", [
    "⚠ Project timing: Revenue recognition varies by project stage",
    "⚠ Grid delays: Monitor connection timelines in key markets",
    "⚠ Policy: Track renewable support scheme changes",
], 7.0, 5.85, 5.8, 1.4, title_col=AMBER, size=9)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: ESG ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "ESG ASSESSMENT",
       "Environmental, Social & Governance Screening  |  Overall Rating: LOW")
logo_tag(slide)

img(slide, "09_esg_assessment.png", 1.67, 1.05, 10.0)

source(slide)
manu_forti_logo(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/jonathonmilne/.openclaw/workspace/OX2_Product1_v15_Final.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print("  9 slides | v15: OX2 logo top-right every slide | LOW (35/100) | Manu Forti branding")
