"""
Procurement Intelligence Report Generator
Generates a professional PPTX + HTML/PDF from report data.
Usage: python3 generate_report.py --type supplier --data report_data.json --out /path/to/output
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import json, sys, os, argparse
from datetime import datetime

# ─── Brand colours ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy — primary
STEEL      = RGBColor(0x1C, 0x3A, 0x55)   # steel blue — secondary
ACCENT     = RGBColor(0x00, 0x9B, 0xD9)   # bright blue — accent / highlights
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY   = RGBColor(0x8A, 0x94, 0xA6)
DARK_GREY  = RGBColor(0x2D, 0x3A, 0x4A)

GREEN_RAG  = RGBColor(0x10, 0xB9, 0x81)
AMBER_RAG  = RGBColor(0xF5, 0x9E, 0x0B)
RED_RAG    = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def rag_colour(rating: str) -> RGBColor:
    r = rating.upper()
    if 'LOW' in r or 'GREEN' in r or '🟢' in r:   return GREEN_RAG
    if 'HIGH' in r or 'RED' in r or '🔴' in r:     return RED_RAG
    return AMBER_RAG


def add_rect(slide, l, t, w, h, fill_colour, line_colour=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=12, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return txb


def add_slide_header(slide, title, subtitle=None, dark=True):
    """Standard header bar across top of slide."""
    bg_col = NAVY if dark else STEEL
    add_rect(slide, 0, 0, 13.33, 1.1, bg_col)
    add_rect(slide, 0, 1.05, 13.33, 0.07, ACCENT)
    add_text(slide, title, 0.4, 0.15, 10, 0.6, size=20, bold=True, colour=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.65, 10, 0.4, size=11, colour=RGBColor(0xB0, 0xC4, 0xDE))


def add_footer(slide, client_name, date_str, page_num):
    add_rect(slide, 0, 7.15, 13.33, 0.35, NAVY)
    add_text(slide, f"Prepared for: {client_name}", 0.3, 7.18, 5, 0.28,
             size=8, colour=MID_GREY)
    add_text(slide, "CONFIDENTIAL", 5.5, 7.18, 3, 0.28,
             size=8, colour=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, f"{date_str}  |  {page_num}", 9.5, 7.18, 3.5, 0.28,
             size=8, colour=MID_GREY, align=PP_ALIGN.RIGHT)


def make_cover_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title if hasattr(slide.shapes, 'title') else None

    # Full navy background
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)

    # Accent bar left
    add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

    # Top accent line
    add_rect(slide, 0.18, 1.8, 13.15, 0.05, ACCENT)

    # Report type badge
    rtype = "SUPPLIER EVALUATION" if data.get('report_type') == 'supplier' else "CATEGORY MARKET INTELLIGENCE"
    add_rect(slide, 0.5, 1.1, 3.8, 0.5, STEEL)
    add_text(slide, rtype, 0.55, 1.15, 3.7, 0.4,
             size=10, bold=True, colour=ACCENT, align=PP_ALIGN.LEFT)

    # Main title
    add_text(slide, data.get('subject_name', '[SUBJECT]'), 0.5, 1.95, 9, 1.0,
             size=36, bold=True, colour=WHITE)

    # Subtitle
    add_text(slide, data.get('category', ''), 0.5, 2.95, 9, 0.5,
             size=18, colour=RGBColor(0xB0, 0xC4, 0xDE))

    # Divider
    add_rect(slide, 0.5, 3.55, 6, 0.04, ACCENT)

    # Meta block
    meta_y = 3.75
    meta = [
        ("Prepared for",  data.get('client_name', '')),
        ("Prepared by",   "Jonathon Milne  |  AI-Powered Procurement Intelligence"),
        ("Date",          data.get('date', datetime.today().strftime('%B %Y'))),
        ("Classification","Confidential"),
    ]
    for label, value in meta:
        add_text(slide, label.upper(), 0.5, meta_y, 2.8, 0.35,
                 size=8, colour=MID_GREY, bold=True)
        add_text(slide, value, 3.4, meta_y, 7, 0.35,
                 size=10, colour=WHITE)
        meta_y += 0.42

    # Bottom branding
    add_rect(slide, 0, 6.9, 13.33, 0.6, STEEL)
    add_text(slide, "AI-POWERED PROCUREMENT INTELLIGENCE", 0.5, 6.95, 8, 0.45,
             size=10, bold=True, colour=ACCENT)


def make_exec_summary_slide(prs, data, client_name, date_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    add_slide_header(slide, "Executive Summary", data.get('subject_name', ''))
    add_footer(slide, client_name, date_str, "2")

    es = data.get('executive_summary', {})

    # Overall RAG box
    rag = es.get('overall_risk', 'MEDIUM')
    rag_col = rag_colour(rag)
    add_rect(slide, 0.4, 1.3, 2.8, 1.2, rag_col)
    add_text(slide, "OVERALL RISK", 0.4, 1.3, 2.8, 0.45,
             size=9, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, rag.replace('🟢','').replace('🟡','').replace('🔴','').strip(),
             0.4, 1.72, 2.8, 0.6,
             size=22, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    # Recommendation box
    rec = es.get('recommendation', 'APPROVE WITH CONDITIONS')
    add_rect(slide, 3.4, 1.3, 5.5, 1.2, NAVY)
    add_text(slide, "RECOMMENDATION", 3.4, 1.3, 5.5, 0.45,
             size=9, bold=True, colour=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, rec, 3.4, 1.68, 5.5, 0.7,
             size=14, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

    # Key findings
    add_rect(slide, 0.4, 2.7, 8.5, 0.38, NAVY)
    add_text(slide, "KEY FINDINGS", 0.5, 2.75, 8, 0.28,
             size=10, bold=True, colour=ACCENT)

    findings = es.get('key_findings', [])
    for i, finding in enumerate(findings[:4]):
        y = 3.2 + i * 0.62
        add_rect(slide, 0.4, y, 0.35, 0.35, ACCENT)
        add_text(slide, str(i+1), 0.4, y, 0.35, 0.35,
                 size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, finding, 0.85, y, 8.0, 0.5,
                 size=10, colour=DARK_GREY)

    # Conditions panel (right side)
    conditions = es.get('conditions', [])
    if conditions:
        add_rect(slide, 9.2, 1.3, 3.9, 5.5, NAVY)
        add_text(slide, "CONDITIONS / WATCH POINTS", 9.3, 1.35, 3.7, 0.38,
                 size=9, bold=True, colour=ACCENT)
        cy = 1.85
        for c in conditions[:5]:
            add_rect(slide, 9.3, cy, 0.22, 0.22, AMBER_RAG)
            add_text(slide, c, 9.6, cy - 0.05, 3.3, 0.55,
                     size=9, colour=WHITE)
            cy += 0.65


def make_risk_summary_slide(prs, data, client_name, date_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    add_slide_header(slide, "Risk Assessment", data.get('subject_name', ''))
    add_footer(slide, client_name, date_str, "3")

    risks = data.get('risks', [
        {"category": "Financial",           "rating": "MEDIUM", "detail": "Revenue growing but limited cash reserves."},
        {"category": "Operational/Delivery","rating": "LOW",    "detail": "Strong track record, certified processes."},
        {"category": "Geopolitical/Country","rating": "LOW",    "detail": "Norwegian entity, EU-aligned."},
        {"category": "Concentration",       "rating": "HIGH",   "detail": "Single source — no qualified alternative."},
        {"category": "ESG & Compliance",    "rating": "MEDIUM", "detail": "No sustainability report published yet."},
        {"category": "Reputational",        "rating": "LOW",    "detail": "No adverse media findings."},
    ])

    add_text(slide, "Risk Category", 0.5, 1.25, 3.5, 0.35,
             size=9, bold=True, colour=MID_GREY)
    add_text(slide, "Rating", 4.1, 1.25, 1.5, 0.35,
             size=9, bold=True, colour=MID_GREY)
    add_text(slide, "Key Driver", 5.8, 1.25, 7.0, 0.35,
             size=9, bold=True, colour=MID_GREY)
    add_rect(slide, 0.4, 1.58, 12.5, 0.04, ACCENT)

    for i, risk in enumerate(risks):
        y = 1.72 + i * 0.75
        bg = NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A)
        add_rect(slide, 0.4, y, 12.5, 0.7, bg)

        # Category
        add_text(slide, risk['category'], 0.55, y + 0.18, 3.3, 0.4,
                 size=11, bold=True, colour=WHITE)

        # RAG chip
        rc = rag_colour(risk['rating'])
        add_rect(slide, 4.1, y + 0.18, 1.4, 0.35, rc)
        label = risk['rating'].replace('🟢','').replace('🟡','').replace('🔴','').strip()
        add_text(slide, label, 4.1, y + 0.18, 1.4, 0.35,
                 size=9, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

        # Detail
        add_text(slide, risk.get('detail', ''), 5.8, y + 0.12, 6.8, 0.5,
                 size=10, colour=RGBColor(0xB0, 0xC4, 0xDE))


def make_commercial_slide(prs, data, client_name, date_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    add_slide_header(slide, "Commercial Intelligence & Negotiation", data.get('subject_name', ''))
    add_footer(slide, client_name, date_str, "4")

    commercial = data.get('commercial', {})

    # Market pricing
    add_rect(slide, 0.4, 1.25, 5.9, 0.4, NAVY)
    add_text(slide, "MARKET PRICING", 0.55, 1.3, 5.7, 0.3,
             size=10, bold=True, colour=ACCENT)

    pricing = commercial.get('pricing', [
        {"item": "Unit price range", "low": "", "mid": "", "high": ""},
    ])
    headers = ["Item", "Low", "Mid", "High"]
    col_x   = [0.5, 3.2, 4.2, 5.2]
    header_y = 1.75
    for h, x in zip(headers, col_x):
        add_text(slide, h, x, header_y, 1.0, 0.3, size=9, bold=True, colour=MID_GREY)

    for i, row in enumerate(pricing[:3]):
        y = 2.1 + i * 0.45
        bg = NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A)
        add_rect(slide, 0.4, y, 5.9, 0.42, bg)
        add_text(slide, row.get('item',''), 0.55, y + 0.08, 2.6, 0.3, size=9, colour=WHITE)
        add_text(slide, row.get('low',''), 3.2, y + 0.08, 0.9, 0.3, size=9, colour=GREEN_RAG, align=PP_ALIGN.CENTER)
        add_text(slide, row.get('mid',''), 4.2, y + 0.08, 0.9, 0.3, size=9, colour=AMBER_RAG, align=PP_ALIGN.CENTER)
        add_text(slide, row.get('high',''), 5.2, y + 0.08, 0.9, 0.3, size=9, colour=RED_RAG, align=PP_ALIGN.CENTER)

    # Negotiation leverage
    add_rect(slide, 0.4, 3.45, 5.9, 0.4, NAVY)
    add_text(slide, "NEGOTIATION LEVERAGE POINTS", 0.55, 3.5, 5.7, 0.3,
             size=10, bold=True, colour=ACCENT)

    leverage = commercial.get('leverage', ["[Leverage point 1]", "[Leverage point 2]", "[Leverage point 3]"])
    for i, point in enumerate(leverage[:4]):
        y = 3.95 + i * 0.55
        add_rect(slide, 0.4, y, 0.3, 0.3, ACCENT)
        add_text(slide, "→", 0.43, y + 0.02, 0.25, 0.28,
                 size=10, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, point, 0.8, y, 5.4, 0.45, size=10, colour=DARK_GREY)

    # Red lines
    add_rect(slide, 6.8, 1.25, 6.1, 0.4, NAVY)
    add_text(slide, "RECOMMENDED CONTRACT TERMS", 6.95, 1.3, 5.9, 0.3,
             size=10, bold=True, colour=ACCENT)

    terms = commercial.get('contract_terms', [
        {"term": "Payment terms", "market_norm": "Net 30-60"},
        {"term": "Contract length", "market_norm": "12-24 months"},
        {"term": "Price escalation", "market_norm": "CPI-linked annual review"},
        {"term": "Liability cap", "market_norm": "1× annual contract value"},
        {"term": "Termination", "market_norm": "90 days notice, for cause immediate"},
    ])
    ty = 1.75
    for t in terms[:6]:
        add_rect(slide, 6.8, ty, 6.1, 0.5, NAVY if terms.index(t) % 2 == 0 else RGBColor(0x16, 0x28, 0x3A))
        add_text(slide, t['term'], 6.95, ty + 0.1, 2.8, 0.32, size=9, bold=True, colour=WHITE)
        add_text(slide, t['market_norm'], 9.85, ty + 0.1, 2.9, 0.32, size=9, colour=RGBColor(0xB0, 0xC4, 0xDE))
        ty += 0.52


def make_financial_capacity_slide(prs, data, client_name, date_str, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GREY)
    add_slide_header(slide, "Financial Health & Contract Capacity", data.get('subject_name', ''))
    add_footer(slide, client_name, date_str, str(page))

    fin = data.get('financial_capacity', {})

    # Left panel — contract size limits
    add_rect(slide, 0.4, 1.25, 6.1, 0.4, NAVY)
    add_text(slide, "RECOMMENDED CONTRACT SIZE LIMITS", 0.55, 1.3, 5.9, 0.3,
             size=10, bold=True, colour=ACCENT)

    limits = fin.get('contract_limits', [
        {"threshold": "Single contract value",       "limit": "NOK 8–10M",   "rationale": "≤25% of annual revenue"},
        {"threshold": "Annual spend (all contracts)", "limit": "NOK 15M",     "rationale": "Avoid >35% client dependency"},
        {"threshold": "Advance / mobilisation",       "limit": "NOK 500K",    "rationale": "Limited cash reserves"},
        {"threshold": "Unsecured exposure",           "limit": "NOK 5M",      "rationale": "Require bond above this level"},
    ])
    headers_x = [0.5, 3.2, 4.8]
    header_labels = ["Threshold", "Limit", "Rationale"]
    for h, x in zip(header_labels, headers_x):
        add_text(slide, h, x, 1.73, 1.5, 0.28, size=8, bold=True, colour=MID_GREY)
    add_rect(slide, 0.4, 1.98, 6.1, 0.03, ACCENT)

    for i, row in enumerate(limits):
        y = 2.05 + i * 0.55
        bg = NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A)
        add_rect(slide, 0.4, y, 6.1, 0.52, bg)
        add_text(slide, row['threshold'], 0.55, y + 0.1, 2.6, 0.32, size=9, colour=WHITE)
        add_text(slide, row['limit'],     3.2,  y + 0.1, 1.5, 0.32, size=10, bold=True, colour=ACCENT)
        add_text(slide, row['rationale'], 4.8,  y + 0.1, 1.6, 0.32, size=8,  colour=MID_GREY)

    # Protections
    add_rect(slide, 0.4, 4.35, 6.1, 0.38, NAVY)
    add_text(slide, "PROTECTIONS FOR LARGER CONTRACTS", 0.55, 4.4, 5.9, 0.28,
             size=9, bold=True, colour=AMBER_RAG)
    protections = fin.get('protections', [
        "Parent company / bank guarantee",
        "Performance bond (5–10% of contract value)",
        "Stage-gate milestone payment structure",
        "5% retention, released 12 months post-completion",
    ])
    for i, p in enumerate(protections[:4]):
        y = 4.83 + i * 0.44
        add_rect(slide, 0.42, y, 0.25, 0.25, AMBER_RAG)
        add_text(slide, "!", 0.42, y, 0.25, 0.25, size=9, bold=True,
                 colour=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, p, 0.76, y - 0.02, 5.6, 0.38, size=9, colour=DARK_GREY)

    # Right panel — operational capacity fit
    add_rect(slide, 6.8, 1.25, 6.1, 0.4, NAVY)
    add_text(slide, "OPERATIONAL CAPACITY FIT", 6.95, 1.3, 5.9, 0.3,
             size=10, bold=True, colour=ACCENT)

    capacity = data.get('capacity_fit', [
        {"factor": "Facility / floor space",      "supplier": "2,400 m²",          "requirement": "~800 m² laydown", "fit": "✅"},
        {"factor": "Max. single order value",      "supplier": "NOK 20M proven",    "requirement": "NOK 12M scope",   "fit": "✅"},
        {"factor": "Machining capability",         "supplier": "CNC to 6m, class X","requirement": "4m CNC + cert",   "fit": "✅"},
        {"factor": "Concurrent project load",      "supplier": "3 projects max",    "requirement": "1 project",       "fit": "✅"},
        {"factor": "Workforce / surge capacity",   "supplier": "45 FTE + 20% flex", "requirement": "~30 FTE equiv.",  "fit": "⚠️"},
    ])

    for h, x, w in [("Factor", 6.85, 2.2), ("Supplier", 9.1, 1.6), ("Req.", 10.75, 1.3), ("Fit", 12.1, 0.6)]:
        add_text(slide, h, x, 1.73, w, 0.28, size=8, bold=True, colour=MID_GREY)
    add_rect(slide, 6.8, 1.98, 6.1, 0.03, ACCENT)

    for i, row in enumerate(capacity):
        y = 2.05 + i * 0.55
        bg = NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A)
        add_rect(slide, 6.8, y, 6.1, 0.52, bg)
        add_text(slide, row['factor'],      6.85, y + 0.1, 2.2, 0.32, size=8,  colour=WHITE)
        add_text(slide, row['supplier'],    9.1,  y + 0.1, 1.6, 0.32, size=8,  colour=MID_GREY)
        add_text(slide, row['requirement'], 10.75,y + 0.1, 1.3, 0.32, size=8,  colour=MID_GREY)
        fit_col = GREEN_RAG if '✅' in row['fit'] else (AMBER_RAG if '⚠️' in row['fit'] else RED_RAG)
        add_rect(slide, 12.1, y + 0.08, 0.6, 0.35, fit_col)
        fit_label = 'OK' if '✅' in row['fit'] else ('WATCH' if '⚠️' in row['fit'] else 'FAIL')
        add_text(slide, fit_label, 12.1, y + 0.08, 0.6, 0.35,
                 size=6.5, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    # Lock-in warning
    lockins = data.get('technology_lockIn', [])
    if lockins:
        add_rect(slide, 6.8, 4.92, 6.1, 0.35, RGBColor(0x3D, 0x2E, 0x10))
        add_text(slide, "⚠  TECHNOLOGY LOCK-IN FLAGS", 6.95, 4.96, 5.9, 0.26,
                 size=9, bold=True, colour=AMBER_RAG)
        for i, lock in enumerate(lockins[:2]):
            y = 5.35 + i * 0.5
            add_text(slide, f"• {lock}", 6.9, y, 5.9, 0.42, size=9, colour=DARK_GREY)


def make_recommendation_slide(prs, data, client_name, date_str, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)
    add_rect(slide, 0, 6.9, 13.33, 0.6, STEEL)

    es = data.get('executive_summary', {})
    rec = es.get('recommendation', 'APPROVE WITH CONDITIONS')

    add_text(slide, "RECOMMENDATION", 0.5, 0.5, 12, 0.5,
             size=13, bold=True, colour=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, rec, 0.5, 1.1, 12, 1.0,
             size=32, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)
    add_rect(slide, 3, 2.2, 7.33, 0.05, ACCENT)

    rationale = data.get('recommendation_rationale',
        'Based on the analysis conducted, the supplier demonstrates sufficient capability and financial stability to proceed, subject to the conditions noted. The primary risks are manageable with appropriate contractual protections.')
    add_text(slide, rationale, 0.8, 2.4, 11.7, 1.2,
             size=12, colour=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

    # Next steps
    add_text(slide, "RECOMMENDED NEXT STEPS", 0.8, 3.8, 11.7, 0.4,
             size=11, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

    steps = data.get('next_steps', [
        {"action": "Request latest audited financials", "timeline": "Before contract"},
        {"action": "Conduct reference call with 2 existing customers", "timeline": "Within 30 days"},
        {"action": "Negotiate payment terms to net-60", "timeline": "During negotiation"},
        {"action": "Include 6-month performance review clause", "timeline": "Contract drafting"},
    ])

    sx = [0.5, 3.6, 6.7, 9.8]
    for i, step in enumerate(steps[:4]):
        x = sx[i]
        add_rect(slide, x, 4.3, 2.8, 1.8, STEEL)
        add_rect(slide, x, 4.3, 2.8, 0.12, ACCENT)
        add_text(slide, step['action'], x + 0.1, 4.5, 2.6, 1.0,
                 size=9, colour=WHITE)
        add_rect(slide, x, 5.85, 2.8, 0.35, RGBColor(0x0D, 0x1B, 0x2A))
        add_text(slide, step['timeline'], x + 0.1, 5.88, 2.6, 0.3,
                 size=8, colour=ACCENT, bold=True)

    add_text(slide, "AI-POWERED PROCUREMENT INTELLIGENCE  |  JONATHON MILNE",
             0.5, 6.95, 12, 0.4,
             size=9, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

    add_footer(slide, client_name, date_str, str(page))


def generate_pptx(data: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    client   = data.get('client_name', 'Client')
    date_str = data.get('date', datetime.today().strftime('%d %B %Y'))

    make_cover_slide(prs, data)
    make_exec_summary_slide(prs, data, client, date_str)
    make_risk_summary_slide(prs, data, client, date_str)
    make_financial_capacity_slide(prs, data, client, date_str, page=4)
    make_commercial_slide(prs, data, client, date_str)
    make_recommendation_slide(prs, data, client, date_str, page=6)

    prs.save(output_path)
    print(f"✅ PPTX saved: {output_path}")


# ─── Sample data for demo ─────────────────────────────────────────────────────

SAMPLE_DATA = {
    "report_type":   "supplier",
    "subject_name":  "Sample Supplier AS",
    "category":      "Engineering Services — Offshore Renewables",
    "client_name":   "Seabased AS",
    "date":          datetime.today().strftime("%d %B %Y"),
    "executive_summary": {
        "overall_risk":    "MEDIUM",
        "recommendation":  "APPROVE WITH CONDITIONS",
        "key_findings": [
            "Supplier has 8-year track record in offshore cable installation with DNV GL certification.",
            "Financial health is stable but thin cash reserves (~3 months runway) create delivery risk.",
            "No qualified alternative supplier identified in the Nordic market — concentration risk is HIGH.",
            "ESG credentials are limited — no published sustainability report, though no adverse findings.",
        ],
        "conditions": [
            "Request audited financials (FY2024) prior to contract execution.",
            "Cap initial contract value at NOK 5M pending performance review.",
            "Include 6-month break clause tied to delivery KPIs.",
            "Require key-person clause for lead engineer.",
        ]
    },
    "risks": [
        {"category": "Financial",            "rating": "MEDIUM", "detail": "Stable revenue, but limited cash reserves. Monitor closely."},
        {"category": "Operational/Delivery", "rating": "LOW",    "detail": "Strong 8-year track record, certified to DNV GL standards."},
        {"category": "Geopolitical/Country", "rating": "LOW",    "detail": "Norwegian-registered entity, fully EU/EEA compliant."},
        {"category": "Concentration",        "rating": "HIGH",   "detail": "Only qualified supplier in Nordic market. No alternative identified."},
        {"category": "ESG & Compliance",     "rating": "MEDIUM", "detail": "No sustainability report. Sanctions screening: CLEAR."},
        {"category": "Reputational",         "rating": "LOW",    "detail": "No adverse media findings across 3-year monitoring window."},
    ],
    "commercial": {
        "pricing": [
            {"item": "Day rate (senior engineer)",  "low": "12,000 NOK", "mid": "15,000 NOK", "high": "18,500 NOK"},
            {"item": "Mobilisation fee",            "low": "50,000 NOK", "mid": "85,000 NOK", "high": "120,000 NOK"},
            {"item": "Typical project value",       "low": "500K NOK",   "mid": "1.2M NOK",   "high": "3M+ NOK"},
        ],
        "leverage": [
            "Multi-year commitment (2–3 year frame agreement) unlocks 10–15% rate discount.",
            "Reference value: Seabased's profile gives the supplier a marquee client for marketing.",
            "Early engagement before peak season (Q1) gives scheduling leverage over competitors.",
            "Volume guarantee across multiple projects reduces their mobilisation cost per project.",
        ],
        "contract_terms": [
            {"term": "Payment terms",        "market_norm": "Net 30–45 days"},
            {"term": "Contract length",      "market_norm": "12–24 month frame"},
            {"term": "Price escalation",     "market_norm": "CPI-linked, annual review"},
            {"term": "Liability cap",        "market_norm": "1× annual contract value"},
            {"term": "Termination",          "market_norm": "90 days notice"},
            {"term": "Performance KPIs",     "market_norm": "Negotiable — push for monthly scorecard"},
        ]
    },
    "financial_capacity": {
        "contract_limits": [
            {"threshold": "Single contract value",        "limit": "NOK 8–10M",  "rationale": "≤25% of annual revenue (NOK 38M)"},
            {"threshold": "Annual spend (all contracts)", "limit": "NOK 15M",    "rationale": "Avoid >40% client dependency"},
            {"threshold": "Advance / mobilisation",       "limit": "NOK 500K",   "rationale": "Cash reserves ~3 months only"},
            {"threshold": "Unsecured exposure",           "limit": "NOK 5M",     "rationale": "Require performance bond above this"},
        ],
        "protections": [
            "Performance bond (5–10% of contract value)",
            "Stage-gate milestone payment — not lump sum advance",
            "5% retention, released 12 months post-completion",
            "Key-person clause for lead engineer",
        ]
    },
    "capacity_fit": [
        {"factor": "Facility / floor space",    "supplier": "2,400 m² workshop",  "requirement": "~800 m² laydown",  "fit": "✅"},
        {"factor": "Max. single order value",   "supplier": "NOK 20M proven",     "requirement": "NOK 12M scope",    "fit": "✅"},
        {"factor": "Machining capability",      "supplier": "CNC to 6m, class X", "requirement": "4m CNC + cert",    "fit": "✅"},
        {"factor": "Concurrent project load",   "supplier": "3 projects max",     "requirement": "1 project (med.)", "fit": "✅"},
        {"factor": "Workforce / surge capacity","supplier": "45 FTE + 20% flex",  "requirement": "~30 FTE equiv.",   "fit": "⚠️"},
    ],
    "technology_lockIn": [
        "SCADA: Siemens-only platform — cannot integrate ABB or Emerson without re-engineering",
        "Proprietary connector design — switching supplier requires hardware re-spec (est. 6–8 weeks)",
    ],
    "recommendation_rationale": "Supplier demonstrates the technical capability and track record required. Financial risk is manageable with appropriate contract protections. Concentration risk is the primary concern — recommend proceeding while actively developing an alternative supplier pipeline in parallel.",
    "next_steps": [
        {"action": "Request audited FY2024 financials before contract", "timeline": "Immediate"},
        {"action": "Reference call with 2 existing customers",          "timeline": "Within 30 days"},
        {"action": "Negotiate net-45 payment and 2-year frame",         "timeline": "Negotiation phase"},
        {"action": "Begin alternative supplier qualification",          "timeline": "Q3 2026"},
    ]
}


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate procurement intelligence report')
    parser.add_argument('--data',   help='Path to JSON data file (optional — uses sample if omitted)')
    parser.add_argument('--out',    help='Output directory', default='/Users/jonathonmilne/.openclaw/workspace/venture/reports')
    parser.add_argument('--name',   help='Output filename prefix', default='procurement_report')
    args = parser.parse_args()

    os.makedirs(args.out, exist_ok=True)

    if args.data:
        with open(args.data) as f:
            data = json.load(f)
    else:
        data = SAMPLE_DATA
        print("ℹ️  No data file provided — using sample data (Seabased demo)")

    ts   = datetime.now().strftime('%Y%m%d_%H%M')
    pptx_path = os.path.join(args.out, f"{args.name}_{ts}.pptx")

    generate_pptx(data, pptx_path)
    print(f"\n📊 Report generated in: {args.out}")
