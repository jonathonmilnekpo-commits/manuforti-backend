from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY    = RGBColor(0x00,0x21,0x47)
STEEL   = RGBColor(0x2B,0x6C,0xB0)
GREY    = RGBColor(0x71,0x80,0x96)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
AMBER   = RGBColor(0xD9,0x7F,0x06)
LTBLUE  = RGBColor(0xEB,0xF4,0xFF)
DARK    = RGBColor(0x2D,0x3A,0x4A)
GREEN   = RGBColor(0x27,0xAE,0x60)
RED     = RGBColor(0xC0,0x39,0x2B)
TEAL    = RGBColor(0x0B,0x7A,0x75)
DEEP    = RGBColor(0x04,0x12,0x2E)
PGREEN  = RGBColor(0xD4,0xED,0xDA)
PAMBER  = RGBColor(0xFF,0xF3,0xCD)
PRED    = RGBColor(0xF8,0xD7,0xDA)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def R(s,l,t,w,h,fill=None,line=None):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb=fill
    sh.line.fill.background()
    if line: sh.line.color.rgb=line; sh.line.width=Pt(0.75)
    return sh

def T(s,text,l,t,w,h,sz=10,bold=False,color=WHITE,align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.color.rgb=color; r.font.italic=italic
    return tb

def HDR(s,tag,title,tc=STEEL):
    R(s,0,0,13.333,1.05,fill=NAVY)
    T(s,tag,0.4,0.07,12,0.32,sz=10,bold=True,color=tc)
    T(s,title,0.4,0.35,12.5,0.58,sz=19,bold=True,color=WHITE)

def FTR(s,txt="Statkraft Procurement AI Strategy  |  Confidential  |  March 2026"):
    R(s,0,7.1,13.333,0.4,fill=NAVY)
    T(s,txt,0.3,7.12,12.7,0.3,sz=9,color=GREY,align=PP_ALIGN.CENTER)

def LINE(s,t,c=STEEL,l=0.4,w=12.5):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(0.04))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def BANNER(s,txt,l,t,w,h,bg=LTBLUE,fg=DARK,sz=12,accent=STEEL):
    R(s,l,t,0.07,h,fill=accent)
    R(s,l+0.07,t,w-0.07,h,fill=bg)
    T(s,txt,l+0.22,t+0.08,w-0.32,h-0.14,sz=sz,color=fg)

def STAT(s,stats,t,bw=2.3,bg=NAVY,nc=STEEL):
    x=0.3
    for num,lbl in stats:
        R(s,x,t,bw,1.3,fill=bg)
        T(s,num,x+0.1,t+0.1,bw-0.2,0.65,sz=24,bold=True,color=nc,align=PP_ALIGN.CENTER)
        T(s,lbl,x+0.1,t+0.72,bw-0.2,0.52,sz=9,color=WHITE,align=PP_ALIGN.CENTER)
        x+=bw+0.19

# ══ SLIDE 1 — TITLE ══════════════════════════════════════════════════════════
s=blank()
R(s,0,0,13.333,7.5,fill=NAVY)
R(s,0,5.6,13.333,1.9,fill=STEEL)
T(s,"PROCUREMENT AI STRATEGY",0.6,0.8,12,0.5,sz=14,bold=False,color=STEEL)
T(s,"Statkraft",0.6,1.3,12,1.1,sz=54,bold=True,color=WHITE)
T(s,"From Framework to Agentic Capability — v2",0.6,2.4,12,0.55,sz=22,color=GREY)
LINE(s,3.1,STEEL)
T(s,'"Every company in the world today needs to have an OpenClaw strategy, an agentic system strategy.\nThis is the new computer."  —  Jensen Huang, CEO Nvidia, GTC 2026',
  0.6,3.2,12,0.85,sz=13,color=WHITE,italic=True)
T(s,"Market Intelligence  ·  Peer Analysis  ·  Agentic Use Cases  ·  Decision Intelligence  ·  Cybersecurity  ·  Execution Roadmap",
  0.6,4.15,12,0.4,sz=11,color=GREY)
T(s,"Prepared by Aiden  |  Manu Forti Intelligence  |  March 17, 2026  |  Confidential",
  0.6,5.75,12,0.4,sz=11,color=WHITE,italic=True)

# ══ SLIDE 2 — NEMOCLAW / WHY NOW ══════════════════════════════════════════════
s=blank()
HDR(s,"THE MOMENT  |  Nvidia GTC 2026 — Yesterday","OpenClaw is the new computer. Nvidia just made it enterprise-ready.")
FTR(s)
# Left — big Jensen quote
R(s,0.3,1.15,7.2,5.55,fill=DEEP)
R(s,0.3,1.15,0.08,5.55,fill=STEEL)
T(s,'"OpenClaw has made it possible for us to create a world where everyone has their own AI agent."',
  0.55,1.35,6.7,1.1,sz=16,bold=True,color=WHITE,italic=True)
T(s,"— Jensen Huang, CEO Nvidia  |  GTC 2026, San Jose",0.55,2.48,6.7,0.38,sz=11,color=GREY,italic=True)
LINE(s,2.95,STEEL,0.55,6.7)
T(s,"What NemoClaw means for enterprise procurement:",0.55,3.1,6.7,0.35,sz=12,bold=True,color=STEEL)
pts=[
    "Enterprise-secure OpenClaw stack — protects sensitive corporate data",
    "Runs on Nvidia hardware: on-prem GPU servers compatible with H100/A100",
    "Any open-source model: Llama, Mistral, Qwen — no vendor lock-in",
    "Peter Steinberger (OpenClaw creator): 'We're building the claws and guardrails that let enterprises act confidently'",
    "Nvidia investing $26B in open-source AI — this ecosystem is not going away",
]
y=3.55
for pt in pts:
    T(s,f"▸  {pt}",0.55,y,6.7,0.48,sz=10.5,color=WHITE)
    y+=0.52
# Right — what this means
R(s,7.75,1.15,5.25,5.55,fill=LTBLUE)
R(s,7.75,1.15,0.07,5.55,fill=AMBER)
T(s,"What This Means for Statkraft",8.0,1.25,4.9,0.38,sz=13,bold=True,color=NAVY)
T(s,"The AI moment has a name — and it's now.",8.0,1.7,4.9,0.35,sz=11,color=STEEL,bold=True)
LINE(s,2.15,AMBER,8.0,4.9)
implications=[
    ("The platform is proven","OpenClaw runs 2.2M installs weekly. This isn't a pilot."),
    ("The security problem is solved","NemoClaw = enterprise-grade, on-premises, auditable."),
    ("The model cost is zero","Llama/Mistral/Qwen are free. Only compute costs money."),
    ("Procurement is next","Operations AI is mature. Procurement AI is the frontier."),
    ("First movers win","Equinor set the standard with $130M. Statkraft can set\nthe Nordic procurement standard."),
]
y=2.28
for title,desc in implications:
    T(s,f"► {title}:",8.0,y,4.9,0.3,sz=11,bold=True,color=NAVY)
    T(s,desc,8.0,y+0.3,4.9,0.42,sz=10,color=DARK)
    y+=0.85

# ══ SLIDE 3 — MARKET INTELLIGENCE (EXPANDED) ══════════════════════════════════
s=blank()
HDR(s,"MARKET INTELLIGENCE  |  2026 Research — 15+ Sources","The data is unambiguous: procurement AI is real, scaled, and delivering")
FTR(s)

BANNER(s,"McKinsey (Jan 2026): AI copilots improve procurement productivity 25–40%. Autonomous category agents deliver 15–30% efficiency gains on top. The next frontier is end-to-end workflow automation — not task tools.",
       0.3,1.15,12.73,0.72,sz=11)

# 5 headline stats
STAT(s,[
    ("94%","of procurement execs\nuse GenAI weekly\n(Art of Procurement 2026)"),
    ("4%","have large-scale\ndeployment — huge\nexecution gap"),
    ("25–40%","efficiency gain from\nagentic AI\n(McKinsey 2026)"),
    ("~2×","AI use in procurement\nyear-over-year\n(Hackett 2026)"),
    ("+8%","workloads rising\nwhile headcount\nfalls (Hackett 2026)"),
],t=1.98,bw=2.42)

# Two column deep dive
R(s,0.3,3.42,6.1,3.32,fill=DEEP)
T(s,"📊  Source Data Points",0.5,3.52,5.7,0.35,sz=12,bold=True,color=STEEL)
data_points=[
    "EY CPO Survey: 80% of CPOs plan GenAI in 3 years; only 36% have meaningful implementations",
    "Gartner: 74% of procurement leaders say data isn't AI-ready; GenAI can automate 50–80% of work",
    "McKinsey: Spend managed per FTE is 50% more today than 5 years ago — AI is what keeps scaling",
    "Hackett: 34% efficiency gains & 23% cost savings from AI-driven procurement outsourcing",
    "AWS/McKinsey: AI supply chains = 15% lower logistics costs, 35% less inventory, 65% better service",
    "Art of Procurement: Top CPO use cases — spend analytics (53%), RFQ (42%), contract (41%)",
]
y=3.95
for pt in data_points:
    T(s,f"▸  {pt}",0.5,y,5.7,0.46,sz=9.5,color=WHITE); y+=0.48

R(s,6.7,3.42,6.3,3.32,fill=LTBLUE)
R(s,6.7,3.42,0.07,3.32,fill=STEEL)
T(s,"⚡  The Adoption Curve — Where Are We?",6.95,3.52,5.9,0.35,sz=12,bold=True,color=NAVY)
phases=[
    ("2022–2023","Experimentation","ChatGPT moments. Individual users exploring.","GREY"),
    ("2024","Pilot Projects","Teams running structured pilots. Mostly text tasks.","AMBER"),
    ("2025","Production Tools","Contract review, spend analytics, RFQ drafting going live.","STEEL"),
    ("2026","Agentic Systems","Autonomous agents. End-to-end workflows. THIS IS NOW.","GREEN"),
    ("2027+","Autonomous Procurement","Self-executing sourcing cycles with human gates.","GREEN"),
]
y=3.95
for yr,phase,desc,col_name in phases:
    col={"GREY":GREY,"AMBER":AMBER,"STEEL":STEEL,"GREEN":GREEN}[col_name]
    R(s,6.87,y,0.9,0.46,fill=col)
    T(s,yr,6.88,y+0.07,0.88,0.32,sz=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,phase,7.85,y+0.03,1.5,0.38,sz=10,bold=True,color=col)
    T(s,desc,9.3,y+0.07,3.5,0.35,sz=9.5,color=DARK)
    y+=0.52

# ══ SLIDE 4 — ENERGY PEER ANALYSIS ════════════════════════════════════════════
s=blank()
HDR(s,"ENERGY SECTOR PEER ANALYSIS  |  What They're Running in Production","Your peers have moved. The gap is closing. Statkraft needs to move now.")
FTR(s)
peers=[
    ("Equinor 🇳🇴",DEEP,[
        "$130M AI savings in 2025 (official, Jan 2026)",
        "500M NOK/yr from low-code automation",
        "Automated spend classification at scale",
        "AI logistics: 13% inventory reduction",
        "24,000 sensors — predictive maintenance",
        "100+ new AI use cases identified",
        "NOK 1B+ projected from robotics",
        "Most relevant peer — watch this closely",
    ]),
    ("Shell / TotalEnergies",DEEP,[
        "Shell: AI contract authoring — first draft automated",
        "Shell: Supplier compliance monitoring, real-time",
        "Shell: Embedding AI in workflows, not dashboards",
        "TotalEnergies: 100K suppliers / $31B spend",
        "Contract NLP scanning: 80% faster review (live)",
        "Continuous supplier financial monitoring agents",
        "AI at $31B spend = $310M saving per 1% gain",
        "Both using on-prem + private cloud hybrid",
    ]),
    ("Vattenfall / Ørsted / BP",DEEP,[
        "Vattenfall: Predictive MRO/capex demand (6–12m)",
        "Vattenfall: Engaged Hackett Group — paying £50K+",
        "Ørsted: Autonomous RFQ from spec docs — saves 4–8h",
        "BP: Autonomous PO routing and compliance checks",
        "BP: AI workflow named as specific cost reduction lever",
        "Multiple majors: Contract NLP on full portfolio",
        "Ørsted: AI cost estimating for offshore wind BOS",
        "All: Moving from dashboards to autonomous agents",
    ]),
]
x=0.3
for title,bg,pts in peers:
    R(s,x,1.15,4.1,5.55,fill=bg)
    T(s,title,x+0.15,1.22,3.8,0.5,sz=13,bold=True,color=STEEL)
    LINE(s,1.82,c=STEEL,l=x+0.15,w=3.78)
    y=1.98
    for pt in pts:
        T(s,f"▸  {pt}",x+0.15,y,3.8,0.44,sz=9.5,color=WHITE); y+=0.46
    x+=4.37
R(s,0.3,6.78,12.73,0.5,fill=STEEL)
T(s,"⚡  The Statkraft Gap: Renewables-specific procurement AI (HV equipment, EPC, offshore wind BOS) is unaddressed by any major. First-mover advantage is real and available now.",
  0.5,6.83,12.4,0.42,sz=11,bold=True,color=WHITE)

# ══ SLIDE 5 — 8 AGENTIC WORKERS ═══════════════════════════════════════════════
s=blank()
HDR(s,"AGENTIC AI WORKERS — THE PROCUREMENT TEAM  |  8 Autonomous Agents","Not tools you prompt — agents that work 24/7, flag risks, draft documents, and feed decisions")
FTR(s)
BANNER(s,'McKinsey: "The real leap forward comes when AI moves beyond individual use cases to agents autonomously executing end-to-end workflows." Category agents alone deliver 15–30% efficiency improvement.',
       0.3,1.15,12.73,0.68,sz=11)
agents=[
    ("🔍","Supplier\nMonitor",STEEL,
     "Daily news + financial scan on 50+ suppliers. Flags risks before they become crises. "
     "Replaces annual review cycle with continuous early warning. Feeds project risk register."),
    ("📊","Market\nIntelligence",TEAL,
     "Commodity prices, steel, copper, cable, transformer indices. Daily briefings. "
     "Informs project estimates at DG2. Flags cost pressure before budget is locked."),
    ("📋","RFQ\nDrafter",STEEL,
     "Given a spec, outputs complete RFQ + evaluation criteria + supplier shortlist. "
     "Saves 4–8 hrs per event. Proven at Ørsted and BP. Works from Jaggaer templates."),
    ("⚖️","Contract\nReviewer",RED,
     "Reads contracts vs. standard playbook. Flags non-standard clauses. "
     "80% faster review benchmark (Luminance/Icertis). Covers the full active portfolio."),
    ("💰","Spend\nClassifier",GREEN,
     "Processes SAP invoices + POs continuously. 70–90% accuracy. "
     "Foundation for all other AI use cases. Cleans data for project cost intelligence."),
    ("🗺️","Category\nStrategist",AMBER,
     "Given a brief, researches market and drafts full category strategy. "
     "Weeks → days. This is what Manu Forti delivers commercially. Highest-value agent."),
    ("🤝","Negotiation\nPrep",TEAL,
     "Market benchmarks, BATNA analysis, walk-away price, supplier SWOT. "
     "Category managers walk into negotiations prepared — not improvising."),
    ("🌿","CSRD\nCompliance",GREY,
     "Monitors supplier ESG, sanctions, and regulatory status. "
     "Automated surveillance for CSRD/EUDR reporting. Audit-ready output."),
]
x,y=0.3,1.95; col=0
for icon,name,ac,desc in agents:
    R(s,x,y,3.08,1.45,fill=DEEP)
    R(s,x,y,0.38,1.45,fill=ac)
    T(s,icon,x+0.04,y+0.42,0.3,0.48,sz=16,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,name,x+0.48,y+0.08,2.52,0.42,sz=11,bold=True,color=ac)
    T(s,desc,x+0.48,y+0.5,2.52,0.9,sz=8.5,color=WHITE)
    col+=1
    if col==4: x=0.3; y+=1.55; col=0
    else: x+=3.27
R(s,0.3,5.12,12.73,0.58,fill=NAVY)
T(s,"⚙️  HUMAN-IN-THE-LOOP THRESHOLDS:  Draft doc → Always auto (logged)  |  Supplier comms → Auto <€50K / Human above  |  Create PO → Auto <€10K / Human above  |  Award contract → ALWAYS HUMAN  |  Confidential data → NEVER autonomous",
  0.5,5.16,12.4,0.5,sz=9.5,color=WHITE)

# ══ SLIDE 6 — AGENTIC USE CASES PART 2 ════════════════════════════════════════
s=blank()
HDR(s,"AGENTIC USE CASES — ADVANCED APPLICATIONS  |  Beyond the Basics","From individual task agents to multi-agent workflows that transform how procurement works")
FTR(s)
BANNER(s,"The shift from 'AI tools' to 'AI systems': Multiple agents working together, passing data between them, with procurement professionals as orchestrators not operators.",
       0.3,1.15,12.73,0.65,sz=11)

advanced=[
    ("🔄","Multi-Agent\nSourcing Workflow",STEEL,
     "Trigger: New capex project spec approved\n"
     "→ Market Intel agent pulls live prices\n"
     "→ RFQ Drafter creates tender document\n"
     "→ Supplier Monitor shortlists qualified vendors\n"
     "→ Negotiation Prep builds target/walk-away\n"
     "Result: Full sourcing pack in hours, not weeks"),
    ("📈","Project Cost\nIntelligence Agent",AMBER,
     "Mines SAP/Jaggaer historical data:\n"
     "→ Actual prices paid vs. original estimates\n"
     "→ Lead time actuals vs. vendor claims\n"
     "→ Cost overrun patterns by category\n"
     "→ Live commodity index overlays\n"
     "Result: DG2 estimates grounded in reality"),
    ("🏗️","EPC Contract\nAI Review Agent",RED,
     "Reads full EPC/EPCI contract packages:\n"
     "→ Flags deviations from Statkraft playbook\n"
     "→ Identifies risk allocation non-standards\n"
     "→ Compares to 5 most recent similar contracts\n"
     "→ Generates exception report for legal review\n"
     "Result: Legal cycle time cut 40–60%"),
    ("🌐","Supplier\nNetwork Intelligence",TEAL,
     "Maps supply chain dependencies:\n"
     "→ Which suppliers share sub-tier vendors?\n"
     "→ Geographic concentration risk (Taiwan, China)\n"
     "→ Single-source dependencies flagged\n"
     "→ Alternative supplier scoring\n"
     "Result: Resilience built into category strategy"),
    ("📉","Cost Overrun\nEarly Warning System",GREEN,
     "Cross-references live commitments vs. budget:\n"
     "→ Monitors committed spend vs. project forecast\n"
     "→ Flags if PO prices exceed DG2 assumption >5%\n"
     "→ Tracks commodity index vs. contract price\n"
     "→ Triggers review if project exposure breached\n"
     "Result: CFO sees risks before they land"),
    ("🤖","Autonomous\nTender Management",NAVY,
     "End-to-end low-value tender automation:\n"
     "→ Receives spec → builds RFQ → issues to panel\n"
     "→ Evaluates responses vs. criteria\n"
     "→ Recommends award with justification\n"
     "→ Human approval gate before award\n"
     "Result: Procurement bandwidth freed for strategy"),
]
x,y=0.3,1.9; col=0
for icon,name,ac,desc in advanced:
    R(s,x,y,4.1,2.25,fill=DEEP)
    R(s,x,y,0.4,2.25,fill=ac)
    T(s,icon,x+0.04,y+0.74,0.34,0.5,sz=18,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,name,x+0.55,y+0.1,3.42,0.5,sz=12,bold=True,color=ac)
    T(s,desc,x+0.55,y+0.65,3.42,1.55,sz=9,color=WHITE)
    col+=1
    if col==3: x=0.3; y+=2.35; col=0
    else: x+=4.35

# ══ SLIDE 7 — DECISION INTELLIGENCE ══════════════════════════════════════════
s=blank()
HDR(s,"DECISION INTELLIGENCE  |  Procurement Data → Better Project Decisions",
    "The most strategic use case: mining what Statkraft actually paid to inform what we should budget next")
FTR(s)

# Shocking stat banner
R(s,0.3,1.15,12.73,0.72,fill=RED)
T(s,"⚠️  Energy infrastructure projects cost on average 40% MORE than estimated and take 2 years longer. (ScienceDirect, 2025 — 662 projects, $1.358 trillion in investment, 83 countries)",
  0.5,1.22,12.35,0.6,sz=12,bold=True,color=WHITE)

# Left: the problem
R(s,0.3,2.0,6.0,4.72,fill=DEEP)
T(s,"The Root Cause",0.5,2.1,5.6,0.38,sz=13,bold=True,color=RED)
T(s,"Cost overruns don't start on site.\nThey start in the estimate.",0.5,2.55,5.6,0.58,sz=12,color=WHITE,bold=True)
LINE(s,3.2,RED,0.5,5.6)
causes=[
    ("Optimism bias","Planners anchor to best-case scenarios. Kahneman's reference class problem."),
    ("Stale benchmarks","Project teams use market reports, not Statkraft's actual paid prices."),
    ("Unlocked data","SAP and Jaggaer hold years of actual costs. Nobody mines them systematically."),
    ("Market movement","HV equipment prices have moved 25–40% since 2020. Old benchmarks are wrong."),
    ("No feedback loop","When a project overruns, the cost model doesn't update for next time."),
]
y=3.35
for title,desc in causes:
    T(s,f"► {title}:",0.5,y,1.8,0.36,sz=10,bold=True,color=RED)
    T(s,desc,2.25,y+0.03,4.0,0.36,sz=9.5,color=WHITE)
    y+=0.56

# Right: the solution
R(s,6.6,2.0,6.43,4.72,fill=LTBLUE)
R(s,6.6,2.0,0.07,4.72,fill=GREEN)
T(s,"The Procurement AI Solution",6.85,2.1,6.05,0.38,sz=13,bold=True,color=NAVY)
T(s,"Mine Statkraft's own data to build a live Cost Intelligence engine",6.85,2.52,6.05,0.5,sz=11,color=STEEL,bold=True)
LINE(s,3.1,GREEN,6.85,6.05)
solutions=[
    ("Cost Intelligence Agent","Queries SAP/Jaggaer: what did Statkraft actually pay for HV transformers in the last 5 years? Adjusts for spec, inflation, and FX."),
    ("DG2 Decision Support","At DG2 (T-24m), agent delivers: real price range, lead time estimate, commodity trend, Bear/Base/Bull cost scenario. Replaces assumptions with data."),
    ("Overrun Early Warning","Monitors committed PO prices vs. DG2 budget assumption. Flags if >5% variance. CFO and PM see it before it hits the P&L."),
    ("Feedback Loop","Post-project: actual vs. estimate auto-analysed. Cost model updates. Every project makes the next estimate more accurate."),
    ("HV Example","At Statkraft's scale: 2% improvement in estimate accuracy on €500M capex = €10M better capital allocation per year."),
]
y=3.18
for title,desc in solutions:
    T(s,f"► {title}:",6.85,y,2.2,0.34,sz=10,bold=True,color=GREEN)
    T(s,desc,9.0,y+0.03,4.0,0.44,sz=9.5,color=DARK)
    y+=0.58

# ══ SLIDE 8 — VALUE MEASUREMENT ════════════════════════════════════════════════
s=blank()
HDR(s,"VALUE MEASUREMENT METHODOLOGY  |  Every Initiative Declares Its ROI",
    "No go-live without a baseline. No scale without 90-day review. Kill what doesn't work.")
FTR(s)
BANNER(s,"Rule: Every AI initiative must declare — BEFORE go-live — its baseline metric, target metric, measurement method, review date, and owner. The Excel model calculates ROI automatically.",
       0.3,1.15,12.73,0.65,sz=11)

# 4-step cycle
steps=[
    ("1","DECLARE\nBASELINE",STEEL,
     "Before any work starts:\n• What is the current state? (hours, cost, accuracy)\n• What data proves this is the baseline?\n• Who owns the measurement?\n• What tool/system tracks it?"),
    ("2","SET\nTARGET",AMBER,
     "Use benchmarks to set targets:\n• RFQ: 75% time reduction (Ørsted)\n• Contract NLP: 80% faster (Luminance)\n• Spend classification: 85% accuracy\n• Supplier monitoring: continuous vs. annual"),
    ("3","GO LIVE\n+ MEASURE",GREEN,
     "At go-live and monthly:\n• Track KPI in initiative tracker (Excel)\n• Log actual savings vs. target\n• Update RAG status: 🟢🟡🔴\n• Flag blockers immediately"),
    ("4","90-DAY\nREVIEW",RED,
     "At 90 days — mandatory gate:\n• Achieved >50% of target? → Scale it\n• 25–50%? → Fix root cause, 90-day extension\n• <25%? → Kill it, document lesson\n• CPO signs off on every outcome"),
]
x=0.3
for num,title,ac,desc in steps:
    R(s,x,1.95,2.95,4.72,fill=DEEP)
    R(s,x,1.95,2.95,0.65,fill=ac)
    T(s,num,x+0.1,2.0,0.5,0.55,sz=22,bold=True,color=WHITE)
    T(s,title,x+0.6,2.04,2.25,0.56,sz=12,bold=True,color=WHITE)
    T(s,desc,x+0.15,2.68,2.68,3.92,sz=10,color=WHITE)
    # Arrow
    if x < 9.0:
        T(s,"▶",x+3.05,3.95,0.25,0.5,sz=18,color=STEEL,align=PP_ALIGN.CENTER)
    x+=3.27

# KPI examples bottom
R(s,0.3,6.75,12.73,0.53,fill=NAVY)
T(s,"📊  KPIs tracked in Value Model Excel:  "
   "Spend classification accuracy %  |  Hours per sourcing event  |  Contract review time (hrs)  |  "
   "Supplier monitoring coverage  |  Estimate accuracy (actual vs. DG2)  |  EUR saved (actual)",
  0.5,6.79,12.4,0.44,sz=9.5,color=WHITE)

# ══ SLIDE 9 — CYBERSECURITY ════════════════════════════════════════════════════
s=blank()
HDR(s,"7-LAYER CYBERSECURITY ARCHITECTURE  |  Hardware, Software & Expected Costs",
    "On-premises first. Every layer specified. Total cost of ownership transparent.")
FTR(s)

R(s,0.3,1.15,12.73,0.52,fill=RED)
T(s,"🔒  Core Principle: Procurement data must NEVER leave the corporate perimeter. "
   "Banned inputs: active tender pricing | negotiation BATNAs | non-public supplier financials | personal data",
  0.5,1.19,12.35,0.42,sz=10,bold=True,color=WHITE)

layers=[
    ("L1","MODEL DEPLOYMENT",STEEL,
     "Llama 3.3 70B / Mistral 3 / Qwen 3.5 on-premises via Ollama or vLLM",
     "Hardware: NVIDIA A100 80GB (×2) or DGX H100  |  Cost: €45–150K CAPEX  |  Running cost: €0/token"),
    ("L2","NETWORK ISOLATION",TEAL,
     "AI inference servers on dedicated VLAN. Zero outbound internet from inference layer.",
     "Hardware: 10GbE managed switch (€2,500)  |  Software: pfSense or Cisco  |  Config: IT Basis team"),
    ("L3","DATA CLASSIFICATION",AMBER,
     "Tag all data: PUBLIC / INTERNAL / CONFIDENTIAL / RESTRICTED. Agents access at or below clearance level.",
     "Software: Microsoft Purview or custom tagging  |  Cost: Existing M365 licence or €5–15K/yr"),
    ("L4","AGENT SANDBOXING",GREEN,
     "Each agent in isolated Docker container. No agent-to-agent comms — all via audited message broker.",
     "Software: Docker + Kubernetes (K3s)  |  Cost: Open source (free)  |  Orchestration: Rancher Desktop"),
    ("L5","PROMPT INJECTION DEFENCE",RED,
     "No.1 attack vector: malicious supplier embeds instructions in document. Input sanitisation + output validation.",
     "Software: Custom sanitisation pipeline (Python)  |  Separate parsing model (Llama 3.2 3B)  |  Cost: ~€5K build"),
    ("L6","AUDIT & MONITORING",GREY,
     "Every AI action logged: timestamp, agent, input hash, output, reviewer. Immutable signed audit log.",
     "Software: ELK Stack (Elasticsearch/Kibana)  |  Cost: €3,600/yr (Elastic Cloud basic)  |  NIS2 compliant"),
    ("L7","HUMAN-IN-THE-LOOP",STEEL,
     "Defined approval gates by action type and value. Draft: always auto. Award contract: always human.",
     "Software: OpenClaw approval workflow  |  Cost: Open source  |  Governance: documented HITL policy"),
]
y=1.78
for code,title,color,desc,hw in layers:
    R(s,0.3,y,0.62,0.62,fill=color)
    T(s,code,0.32,y+0.14,0.58,0.35,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    R(s,0.95,y,12.08,0.62,fill=DEEP)
    T(s,title+":",1.07,y+0.04,2.8,0.28,sz=9.5,bold=True,color=color)
    T(s,desc,1.07,y+0.3,5.5,0.28,sz=8.5,color=WHITE)
    T(s,hw,6.6,y+0.1,6.3,0.46,sz=8.5,color=GREY,italic=True)
    y+=0.68

# ══ SLIDE 10 — TECHNOLOGY STACK & COSTS ════════════════════════════════════════
s=blank()
HDR(s,"TECHNOLOGY STACK  |  3 Tiers + Total Cost of Ownership","On-prem for sensitive data. Private cloud for internal. Cloud tools for public research only.")
FTR(s)

# TCO summary boxes
tcos=[
    ("Year 1 CAPEX\n(Hardware)","€45–150K","A100 server + networking + storage"),
    ("Annual OPEX\n(Software + IT)","~€90K/yr","Elastic, Azure OpenAI, ElevenLabs, 0.5 FTE IT"),
    ("3-Year TCO","~€320K","Hardware + 3yr OPEX"),
    ("Annual Value\n(Value Model)","€1.3–2.5M","From Value Calculator Excel"),
    ("3-Year ROI","8–12×","Conservative estimate at 50% target achievement"),
]
x=0.3; bw=2.42
for label,val,note in tcos:
    R(s,x,1.15,bw,1.3,fill=NAVY)
    T(s,val,x+0.1,1.2,bw-0.2,0.65,sz=22,bold=True,color=STEEL,align=PP_ALIGN.CENTER)
    T(s,label,x+0.1,1.82,bw-0.2,0.4,sz=9,color=WHITE,align=PP_ALIGN.CENTER)
    x+=bw+0.19

# 3-tier stack
tiers=[
    ("TIER 1\nOn-Premises",RED,"RESTRICTED + CONFIDENTIAL",[
        "Models: Llama 3.3 70B / Mistral 3 / Qwen 3.5","Inference: Ollama (dev) / vLLM (prod)",
        "Orchestration: OpenClaw + LangChain","Vector DB: ChromaDB / Qdrant (on-prem)",
        "Hardware: NVIDIA A100 80GB ×2 (~€45K)","Storage: 4TB NVMe (~€6K)",
        "Network: Dedicated VLAN isolation","Audit: ELK Stack (~€3.6K/yr)",
        "Cost: €0/token for inference","Agents: Supplier Monitor, Contract Review,","  Spend Classifier, Negotiation Prep",
    ]),
    ("TIER 2\nPrivate Cloud",AMBER,"INTERNAL data — controlled VPC",[
        "Platform: Azure OpenAI (EU region) / AWS Bedrock","Data residency: EU (GDPR compliant)",
        "No training on Statkraft data (contractual)","Integration: VNET to on-prem",
        "Auth: Entra ID / SSO","Audit: Azure Monitor",
        "Cost: ~€2K/month (~€24K/yr)","Models: GPT-4o, Claude Sonnet (API)",
        "Agents: RFQ Drafter, Market Intel,","  Category Strategist",
    ]),
    ("TIER 3\nCloud Tools",GREEN,"PUBLIC data ONLY — no corporate data",[
        "Tools: Claude, ChatGPT, Gemini, Perplexity","Rule: NO corporate data input. Ever.",
        "Use: Market research, benchmarking, ideation","Manager approval + data handling agreement",
        "Usage logged centrally","Annual policy review",
        "Cost: ~€0–50/user/month","Individual track (Track A/B)",
        "Agents: Market research, learning,","  public benchmarking only",
    ]),
]
x=0.3
for title,color,sub,items in tiers:
    R(s,x,2.58,4.13,4.2,fill=DEEP)
    R(s,x,2.58,4.13,0.6,fill=color)
    T(s,title,x+0.15,2.61,3.85,0.55,sz=13,bold=True,color=WHITE)
    T(s,sub,x+0.15,3.22,3.85,0.28,sz=9,color=color,italic=True)
    y=3.55
    for item in items:
        T(s,f"· {item}",x+0.15,y,3.85,0.36,sz=8.5,color=WHITE); y+=0.36
    x+=4.37

# ══ SLIDE 11 — SAP & JAGGAER INTEGRATION ══════════════════════════════════════
s=blank()
HDR(s,"DATA INTEGRATION  |  SAP & Jaggaer → AI Agent Layer",
    "The agents are only as good as the data. Here's how we connect the systems.")
FTR(s)

# Architecture diagram (text-based)
R(s,0.3,1.15,12.73,0.6,fill=DEEP)
T(s,"DATA FLOW:  SAP S/4HANA  →  Jaggaer One  →  Integration Layer  →  On-Prem Vector DB  →  AI Agents  →  Decisions",
  0.5,1.25,12.35,0.42,sz=11,bold=True,color=STEEL,align=PP_ALIGN.CENTER)

# Two system columns
R(s,0.3,1.85,5.85,3.6,fill=DEEP)
T(s,"🏢  SAP S/4HANA",0.5,1.95,5.45,0.38,sz=13,bold=True,color=STEEL)
T(s,"Key APIs:",0.5,2.38,5.45,0.28,sz=10,bold=True,color=GREY)
sap_apis=[
    "PurchaseOrderAPI — All PO data (header, items, pricing)",
    "SupplierAPI — Supplier master data + payment history",
    "ContractAPI — Contract portfolio (terms, values, expiry)",
    "PurchaseRequisitionAPI — Demand signal for forecasting",
    "SupplierInvoiceAPI — Invoice data for spend classification",
    "SAP Joule — Native AI copilot, no custom integration needed",
]
y=2.72
for api in sap_apis:
    T(s,f"· {api}",0.5,y,5.45,0.36,sz=9.5,color=WHITE); y+=0.38
T(s,"Protocol: OData v4 / REST  |  Auth: OAuth 2.0  |  Access: Read-only service account",
  0.5,5.2,5.45,0.22,sz=8.5,color=GREY,italic=True)

R(s,6.48,1.85,6.55,3.6,fill=DEEP)
T(s,"🔗  Jaggaer One",6.68,1.95,6.15,0.38,sz=13,bold=True,color=AMBER)
T(s,"Key APIs + Webhooks:",6.68,2.38,6.15,0.28,sz=10,bold=True,color=GREY)
jag_apis=[
    "/rfx — All sourcing events and responses",
    "/contract — Portfolio (clauses, expiry, supplier)",
    "/supplier — Profiles, scores, qualification status",
    "/savings — Savings tracking and reporting",
    "Webhooks: Contract expiry → trigger Contract Reviewer",
    "Webhooks: New sourcing event → trigger RFQ Drafter",
]
y=2.72
for api in jag_apis:
    T(s,f"· {api}",6.68,y,6.15,0.36,sz=9.5,color=WHITE); y+=0.38
T(s,"Protocol: REST API v2  |  Auth: API key + OAuth  |  Certified: SAP S4HANA integration",
  6.68,5.2,6.15,0.22,sz=8.5,color=GREY,italic=True)

# MD file output + 30-day plan
R(s,0.3,5.55,7.65,1.22,fill=RGBColor(0x0A,0x1A,0x35))
T(s,"📄  Structured Output — MD Files for AI Consumption",0.5,5.63,7.3,0.3,sz=11,bold=True,color=STEEL)
md_files="spend_summary_YYYY-MM.md  |  active_contracts.md  |  supplier_master.md  |  open_pos.md  |  expiring_contracts_90d.md  |  mro_demand_history.md  |  active_sourcing_events.md"
T(s,md_files,0.5,5.97,7.3,0.72,sz=9,color=WHITE)

R(s,8.2,5.55,4.83,1.22,fill=RGBColor(0x0A,0x1A,0x35))
T(s,"⚡  30-Day Quick Start",8.35,5.63,4.55,0.3,sz=11,bold=True,color=AMBER)
steps30="Day 1-3: Request SAP read-only service user  |  Day 3-5: Request Jaggaer API credentials  |  Week 2: Run extraction scripts (non-prod)  |  Week 3: Validate & connect ChromaDB  |  Week 4: Test Supplier Monitor on live data"
T(s,steps30,8.35,5.97,4.55,0.72,sz=9,color=WHITE)

# ══ SLIDE 12 — EXECUTION ROADMAP ══════════════════════════════════════════════
s=blank()
HDR(s,"EXECUTION ROADMAP  |  90 Days to Pilots. 12 Months to Agentic Scale.",
    "Three fast wins in 90 days. Decision Intelligence by month 6. Full agent team by month 12.")
FTR(s)

phases=[
    ("PHASE 0\nWks 1–2\nFoundation",STEEL,[
        "✓ Add use case library (top 12)",
        "✓ Mandate value measurement template",
        "✓ Define data classification rules",
        "✓ Draft Track D governance framework",
        "✓ Request SAP + Jaggaer API access",
        "✓ Cybersecurity requirements defined",
    ]),
    ("PHASE 1\nWks 3–12\n3 Fast Pilots",GREEN,[
        "✓ Pilot: Spend classification (SAP data)",
        "✓ Pilot: Contract NLP scanning (Jaggaer)",
        "✓ Pilot: Supplier monitoring (top 50)",
        "✓ Baseline metrics captured for each",
        "✓ On-prem server procured + configured",
        "✓ 90-day review gate — scale or kill",
    ]),
    ("PHASE 2\nMths 4–6\nExpand",AMBER,[
        "✓ RFQ drafter live (target: 4–8 hrs saved)",
        "✓ Cost Intelligence Agent deployed",
        "✓ Demand forecasting (MRO/capex)",
        "✓ Track D governance live",
        "✓ AI Forum → product team with backlog",
        "✓ DG2 decision support piloted",
    ]),
    ("PHASE 3\nMths 7–12\nFull Agentic",RED,[
        "✓ All 8 agents + 6 advanced use cases",
        "✓ Category strategy agent live",
        "✓ Project cost intelligence at every DG2",
        "✓ CSRD compliance agent live",
        "✓ Full 7-layer security stack certified",
        "✓ Statkraft sets Nordic standard",
    ]),
]
x=0.3
for phase,color,items in phases:
    R(s,x,1.15,3.06,5.65,fill=DEEP)
    R(s,x,1.15,3.06,0.78,fill=color)
    T(s,phase,x+0.15,1.18,2.82,0.74,sz=11,bold=True,color=WHITE)
    y=2.02
    for item in items:
        T(s,item,x+0.15,y,2.82,0.5,sz=9.5,color=WHITE); y+=0.54
    x+=3.24

R(s,0.3,6.87,12.73,0.43,fill=NAVY)
targets="Phase 0: Framework improved  |  Phase 1: 3 pilots live, baselines set  |  Phase 2: 25%+ time saving, decision intelligence live  |  Phase 3: 40% efficiency gain, cost overruns reduced  |  Year 2: Nordic standard"
T(s,targets,0.5,6.9,12.4,0.38,sz=8.5,color=WHITE,align=PP_ALIGN.CENTER)

# ══ SLIDE 13 — CALL TO ACTION ══════════════════════════════════════════════════
s=blank()
R(s,0,0,13.333,7.5,fill=NAVY)
R(s,0,5.55,13.333,1.95,fill=STEEL)
T(s,"THE BOTTOM LINE",0.6,0.7,12,0.48,sz=14,color=STEEL,bold=True)
T(s,"Statkraft has years of procurement data sitting unused.\nThe tools to unlock it cost less than one project overrun.",
  0.6,1.2,12,1.35,sz=30,bold=True,color=WHITE)
LINE(s,2.7,STEEL)
calls=[
    ("This week",   "Use case library + value measurement template added to framework"),
    ("This month",  "3 pilots launched: spend classification, contract NLP, supplier monitoring"),
    ("Month 3",     "Decision Intelligence agent live — DG2 estimates grounded in actual SAP data"),
    ("Month 12",    "All 8 agents live. 40% efficiency gain. Cost overrun risk materially reduced."),
]
y=2.85
for when,action in calls:
    R(s,0.6,y,1.85,0.44,fill=STEEL)
    T(s,when,0.65,y+0.09,1.76,0.3,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,f"→  {action}",2.6,y+0.09,10.0,0.32,sz=12,color=WHITE)
    y+=0.56

T(s,"Sources: Hackett Group 2026 · McKinsey Jan 2026 · EY CPO Survey 2025 · Gartner · Equinor.com · ScienceDirect (662 energy projects) · Art of Procurement",
  0.6,5.1,12,0.35,sz=9,color=GREY,italic=True)
T(s,"Prepared by Aiden  |  Manu Forti Intelligence  |  For Statkraft Procurement  |  March 17, 2026",
  0.6,5.68,12,0.4,sz=14,bold=True,color=WHITE)
T(s,"manuforti.as@gmail.com  |  www.manuforti.no",
  0.6,6.12,12,0.35,sz=12,color=WHITE,italic=True)

out="/Users/jonathonmilne/.openclaw/workspace/venture/Statkraft_Procurement_AI_Strategy_v2_March2026.pptx"
prs.save(out)
print(f"Saved: {out}")
