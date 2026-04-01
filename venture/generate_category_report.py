"""
Category Market Intelligence Report Generator
Generates a professional PPTX for supply market / category mapping reports.
Usage: python3 generate_category_report.py --data data/cermaq_category_mapping.json
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json, os, argparse
from datetime import datetime

# ─── Brand colours ─────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
STEEL      = RGBColor(0x1C, 0x3A, 0x55)
ACCENT     = RGBColor(0x00, 0x9B, 0xD9)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY   = RGBColor(0x8A, 0x94, 0xA6)
DARK_GREY  = RGBColor(0x2D, 0x3A, 0x4A)
GREEN_RAG  = RGBColor(0x10, 0xB9, 0x81)
AMBER_RAG  = RGBColor(0xF5, 0x9E, 0x0B)
RED_RAG    = RGBColor(0xEF, 0x44, 0x44)
TEAL       = RGBColor(0x0F, 0x76, 0x6E)
SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)


def rag(r):
    r = str(r).upper()
    if any(x in r for x in ['LOW', 'GREEN', '🟢', 'STRONG']):   return GREEN_RAG
    if any(x in r for x in ['HIGH', 'RED', '🔴', 'CAUTION']):   return RED_RAG
    return AMBER_RAG


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else setattr(s.line.color, 'rgb', line)
    return s


def txt(slide, text, l, t, w, h, size=11, bold=False, colour=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = colour
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    rect(slide, 0, 1.05, 13.33, 0.07, ACCENT)
    txt(slide, title, 0.4, 0.14, 10, 0.62, size=20, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.66, 10, 0.38, size=11,
            colour=RGBColor(0xB0, 0xC4, 0xDE))


def footer(slide, client, date_str, page):
    rect(slide, 0, 7.15, 13.33, 0.35, NAVY)
    txt(slide, f"Prepared for: {client}", 0.3, 7.18, 5, 0.28, size=8, colour=MID_GREY)
    txt(slide, "CONFIDENTIAL", 5.5, 7.18, 3, 0.28, size=8, colour=MID_GREY, align=PP_ALIGN.CENTER)
    txt(slide, f"{date_str}  |  {page}", 9.5, 7.18, 3.5, 0.28, size=8, colour=MID_GREY, align=PP_ALIGN.RIGHT)


# ─── Slides ────────────────────────────────────────────────────────────────

def slide_cover(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    rect(s, 0, 0, 0.18, 7.5, ACCENT)
    rect(s, 0.18, 1.85, 13.15, 0.05, ACCENT)
    # Badge
    rect(s, 0.5, 1.1, 4.2, 0.52, STEEL)
    txt(s, "CATEGORY MARKET INTELLIGENCE REPORT", 0.6, 1.16, 4.0, 0.38,
        size=9, bold=True, colour=ACCENT)
    # Title
    txt(s, data['subject_name'], 0.5, 2.0, 10, 0.9, size=34, bold=True)
    txt(s, data['category'], 0.5, 2.9, 10, 0.5, size=16,
        colour=RGBColor(0xB0, 0xC4, 0xDE))
    rect(s, 0.5, 3.55, 6.5, 0.04, ACCENT)
    # Meta
    meta = [("Prepared for", data['client_name']),
            ("Prepared by",  "Jonathon Milne  |  AI-Powered Procurement Intelligence"),
            ("Date",         data.get('date', datetime.today().strftime('%B %Y'))),
            ("Classification", "Confidential")]
    my = 3.72
    for label, val in meta:
        txt(s, label.upper(), 0.5, my, 2.8, 0.35, size=8, colour=MID_GREY, bold=True)
        txt(s, val, 3.5, my, 7.5, 0.35, size=10)
        my += 0.42
    rect(s, 0, 6.9, 13.33, 0.6, STEEL)
    txt(s, "AI-POWERED PROCUREMENT INTELLIGENCE", 0.5, 6.95, 12, 0.45,
        size=10, bold=True, colour=ACCENT)


def slide_exec_summary(prs, data, client, date_str):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, LIGHT_GREY)
    header(s, "Executive Summary", data['client_name'] + " — " + data['category'])
    footer(s, client, date_str, "2")

    es = data['executive_summary']

    # Market snapshot boxes
    snap = [
        ("Market Size",    es.get('market_size', '—')),
        ("Growth (CAGR)",  es.get('growth_rate', '—')),
        ("Viable Suppliers", str(es.get('viable_suppliers', '—'))),
        ("Buyer Power",    es.get('buyer_power', '—')),
    ]
    bx = 0.4
    for label, val in snap:
        rect(s, bx, 1.25, 2.3, 1.0, NAVY)
        rect(s, bx, 1.25, 2.3, 0.08, ACCENT)
        txt(s, val, bx, 1.52, 2.3, 0.5, size=18, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)
        txt(s, label, bx, 1.98, 2.3, 0.3, size=8, colour=MID_GREY, align=PP_ALIGN.CENTER)
        bx += 2.4

    # Key findings
    rect(s, 0.4, 2.45, 8.5, 0.38, NAVY)
    txt(s, "KEY FINDINGS", 0.55, 2.5, 8, 0.28, size=10, bold=True, colour=ACCENT)
    for i, f in enumerate(es.get('key_findings', [])[:4]):
        y = 2.95 + i * 0.6
        rect(s, 0.4, y, 0.35, 0.35, ACCENT)
        txt(s, str(i+1), 0.4, y, 0.35, 0.35, size=11, bold=True, align=PP_ALIGN.CENTER)
        txt(s, f, 0.85, y, 8.0, 0.52, size=9.5, colour=DARK_GREY)

    # Recommendation
    rect(s, 9.2, 1.25, 3.9, 5.5, NAVY)
    txt(s, "STRATEGIC RECOMMENDATION", 9.3, 1.3, 3.7, 0.38, size=8, bold=True, colour=ACCENT)
    rec = es.get('recommendation', '')
    txt(s, rec, 9.3, 1.78, 3.6, 1.8, size=9.5, colour=WHITE)
    rect(s, 9.3, 3.7, 3.7, 0.04, ACCENT)
    txt(s, "QUICK WINS", 9.3, 3.82, 3.7, 0.28, size=8, bold=True, colour=AMBER_RAG)
    for i, qw in enumerate(data.get('quick_wins', [])[:4]):
        y = 4.18 + i * 0.6
        rect(s, 9.3, y, 0.22, 0.22, AMBER_RAG)
        txt(s, qw, 9.6, y - 0.04, 3.4, 0.55, size=8, colour=WHITE)


def slide_category_landscape(prs, data, client, date_str):
    """Overview of all categories on one slide — heat-map style."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, LIGHT_GREY)
    header(s, "Supply Category Landscape", "Risk & Concentration Overview — All Categories")
    footer(s, client, date_str, "3")

    cats = data.get('categories', [])

    # Column headers
    col_headers = ["Supply Category", "Spend %", "Supply Risk", "Concentration", "Lock-In Risk", "Key Suppliers"]
    col_x       = [0.35, 3.0, 4.0, 5.2, 6.4, 7.7]
    col_w       = [2.6, 0.9, 1.1, 1.1, 1.1, 5.4]
    ty = 1.22
    for h, x, w in zip(col_headers, col_x, col_w):
        txt(s, h, x, ty, w, 0.3, size=8, bold=True, colour=MID_GREY)
    rect(s, 0.35, 1.5, 12.65, 0.04, ACCENT)

    for i, cat in enumerate(cats[:7]):
        y = 1.62 + i * 0.73
        bg = NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A)
        rect(s, 0.35, y, 12.65, 0.7, bg)

        # Category name
        txt(s, cat['name'], 0.45, y + 0.13, 2.5, 0.45, size=9, bold=True)

        # Spend
        txt(s, cat.get('spend_estimate', '—').split(' ')[0], 3.0, y + 0.2, 0.9, 0.3,
            size=9, colour=ACCENT, align=PP_ALIGN.CENTER)

        # RAG chips
        for val, x in [(cat.get('risk','MEDIUM'), 4.02), (cat.get('concentration','MEDIUM'), 5.22)]:
            rc = rag(val)
            rect(s, x, y + 0.17, 1.05, 0.35, rc)
            lbl = val.replace('🟢','').replace('🟡','').replace('🔴','').strip()[:6]
            txt(s, lbl, x, y + 0.17, 1.05, 0.35, size=8, bold=True, align=PP_ALIGN.CENTER)

        # Lock-in flag
        lockin = cat.get('lock_in_flag', '')
        has_lockin = '⚠' in lockin or 'lock' in lockin.lower() or 'proprietary' in lockin.lower()
        rect(s, 6.42, y + 0.17, 1.05, 0.35, AMBER_RAG if has_lockin else GREEN_RAG)
        txt(s, "HIGH" if has_lockin else "LOW", 6.42, y + 0.17, 1.05, 0.35,
            size=8, bold=True, align=PP_ALIGN.CENTER)

        # Key suppliers (Tier 1 names)
        t1 = [sup['name'] for sup in cat.get('tier1', [])[:3]]
        txt(s, "  ·  ".join(t1), 7.75, y + 0.2, 5.2, 0.35, size=8.5, colour=RGBColor(0xB0, 0xC4, 0xDE))


def slide_category_deep(prs, data, client, date_str, cat_data, page):
    """One slide per high-risk/important category."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, LIGHT_GREY)
    header(s, cat_data['name'], cat_data.get('description', ''))
    footer(s, client, date_str, str(page))

    risk_col = rag(cat_data.get('risk', 'MEDIUM'))
    conc_col = rag(cat_data.get('concentration', 'MEDIUM'))

    # RAG badges row
    for label, val, col, x in [
        ("SUPPLY RISK",    cat_data.get('risk','—'),          risk_col, 0.4),
        ("CONCENTRATION",  cat_data.get('concentration','—'), conc_col, 2.35),
        ("SPEND (CAPEX)",  cat_data.get('spend_estimate','—'), STEEL,   4.3),
    ]:
        rect(s, x, 1.25, 1.75, 0.9, col if col != STEEL else NAVY)
        txt(s, label, x, 1.25, 1.75, 0.38, size=7, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        v = val.replace('🟢','').replace('🟡','').replace('🔴','').strip()[:20]
        txt(s, v, x, 1.6, 1.75, 0.5, size=10, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    # Tier 1 suppliers (left panel)
    rect(s, 0.4, 2.3, 6.1, 0.38, NAVY)
    txt(s, "TIER 1 — RECOMMENDED FOR ENGAGEMENT", 0.55, 2.35, 5.9, 0.28, size=9, bold=True, colour=ACCENT)

    for i, sup in enumerate(cat_data.get('tier1', [])[:3]):
        y = 2.78 + i * 1.2
        rect(s, 0.4, y, 6.1, 1.12, NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A))
        rect(s, 0.4, y, 0.08, 1.12, ACCENT)
        txt(s, sup['name'], 0.58, y + 0.06, 3.0, 0.38, size=11, bold=True)
        txt(s, sup.get('hq',''), 3.65, y + 0.09, 2.7, 0.28, size=8, colour=MID_GREY)
        txt(s, sup.get('strength',''), 0.58, y + 0.42, 4.5, 0.5, size=8.5, colour=RGBColor(0xB0,0xC4,0xDE))
        if sup.get('risk'):
            rect(s, 0.58, y + 0.82, 5.8, 0.24, RGBColor(0x3D, 0x1F, 0x1F))
            txt(s, "⚠  " + sup['risk'], 0.68, y + 0.83, 5.6, 0.22, size=7.5, colour=AMBER_RAG)

    # Tier 2 (right top)
    rect(s, 6.75, 2.3, 6.1, 0.38, NAVY)
    txt(s, "TIER 2 — MONITOR / CONDITIONAL", 6.9, 2.35, 5.9, 0.28, size=9, bold=True, colour=AMBER_RAG)
    for i, sup in enumerate(cat_data.get('tier2', [])[:3]):
        y = 2.78 + i * 0.75
        rect(s, 6.75, y, 6.1, 0.7, NAVY if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3A))
        txt(s, sup['name'], 6.9, y + 0.07, 2.5, 0.3, size=10, bold=True)
        txt(s, sup.get('hq',''), 9.45, y + 0.1, 3.3, 0.25, size=8, colour=MID_GREY)
        txt(s, sup.get('note',''), 6.9, y + 0.38, 5.8, 0.28, size=8, colour=RGBColor(0xB0,0xC4,0xDE))

    # Excluded (right mid)
    excluded = cat_data.get('excluded', [])
    if excluded:
        ey = 2.78 + len(cat_data.get('tier2', [])) * 0.75 + 0.2
        rect(s, 6.75, ey, 6.1, 0.35, RGBColor(0x2D, 0x10, 0x10))
        txt(s, "NOT RECOMMENDED / EXCLUDED", 6.9, ey + 0.06, 5.9, 0.24, size=8, bold=True, colour=RED_RAG)
        for i, ex in enumerate(excluded[:2]):
            y = ey + 0.42 + i * 0.48
            rect(s, 6.75, y, 6.1, 0.42, RGBColor(0x1F, 0x10, 0x10))
            txt(s, ex['name'], 6.9, y + 0.08, 2.2, 0.28, size=9, bold=True, colour=RED_RAG)
            txt(s, ex.get('reason',''), 9.15, y + 0.08, 3.6, 0.28, size=7.5, colour=MID_GREY)

    # Lock-in warning (bottom)
    lockin = cat_data.get('lock_in_flag','')
    if lockin:
        rect(s, 6.75, 5.8, 6.1, 0.05, AMBER_RAG)
        rect(s, 6.75, 5.88, 6.1, 1.02, RGBColor(0x2A, 0x1F, 0x08))
        txt(s, "⚠  LOCK-IN / DEPENDENCY WARNING", 6.9, 5.92, 5.9, 0.3, size=9, bold=True, colour=AMBER_RAG)
        txt(s, lockin, 6.9, 6.28, 5.9, 0.6, size=8, colour=WHITE)


def slide_strategy(prs, data, client, date_str, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    rect(s, 0, 0, 0.18, 7.5, ACCENT)
    rect(s, 0, 6.9, 13.33, 0.6, STEEL)

    txt(s, "STRATEGIC SOURCING RECOMMENDATION", 0.5, 0.3, 12, 0.5,
        size=13, bold=True, colour=MID_GREY, align=PP_ALIGN.CENTER)
    txt(s, data['subject_name'], 0.5, 0.85, 12, 0.7,
        size=28, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)
    rect(s, 2.5, 1.7, 8.33, 0.05, ACCENT)

    rec = data.get('strategic_recommendation', '')
    txt(s, rec, 0.8, 1.88, 11.7, 1.4, size=11, colour=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)

    txt(s, "IMMEDIATE ACTIONS (NEXT 30 DAYS)", 0.8, 3.45, 11.7, 0.38,
        size=11, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

    qw = data.get('quick_wins', [])
    sx = [0.5, 3.55, 6.6, 9.65]
    for i, win in enumerate(qw[:4]):
        x = sx[i]
        rect(s, x, 3.95, 2.8, 1.9, STEEL)
        rect(s, x, 3.95, 2.8, 0.12, ACCENT)
        txt(s, str(i+1), x + 0.05, 3.97, 0.4, 0.35, size=14, bold=True, colour=ACCENT)
        txt(s, win, x + 0.12, 4.18, 2.6, 1.4, size=9, colour=WHITE)

    txt(s, "AI-POWERED PROCUREMENT INTELLIGENCE  |  JONATHON MILNE",
        0.5, 6.95, 12, 0.4, size=9, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)
    footer(s, client, date_str, str(page))


# ─── Main ──────────────────────────────────────────────────────────────────

def generate(data, output_path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    client   = data.get('client_name', 'Client')
    date_str = data.get('date', datetime.today().strftime('%d %B %Y'))

    slide_cover(prs, data)
    slide_exec_summary(prs, data, client, date_str)
    slide_category_landscape(prs, data, client, date_str)

    # Deep-dive slides for highest-risk categories
    priority_cats = [c for c in data.get('categories', []) if c.get('risk') == 'HIGH']
    for i, cat in enumerate(priority_cats[:4]):
        slide_category_deep(prs, data, client, date_str, cat, page=4+i)

    slide_strategy(prs, data, client, date_str, page=4+len(priority_cats))

    prs.save(output_path)
    print(f"✅  Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--data', required=True, help='JSON data file')
    parser.add_argument('--out', default='/Users/jonathonmilne/.openclaw/workspace/venture/reports')
    parser.add_argument('--name', default='category_report')
    args = parser.parse_args()

    os.makedirs(args.out, exist_ok=True)
    with open(args.data) as f:
        data = json.load(f)

    ts   = datetime.now().strftime('%Y%m%d_%H%M')
    path = os.path.join(args.out, f"{args.name}_{ts}.pptx")
    generate(data, path)
