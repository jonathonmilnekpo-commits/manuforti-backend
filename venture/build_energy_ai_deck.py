from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
NAVY = RGBColor(0x00, 0x21, 0x47)
STEEL_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
MID_GREY = RGBColor(0x71, 0x80, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xD9, 0x7F, 0x06)
LIGHT_BLUE_BG = RGBColor(0xEB, 0xF4, 0xFF)
DARK_GREY = RGBColor(0x2D, 0x3A, 0x4A)
GREEN = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_divider(slide, top, color=STEEL_BLUE):
    shape = slide.shapes.add_shape(1, Inches(0.4), Inches(top), Inches(12.5), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def slide_footer(slide, text="Manu Forti Intelligence  |  Confidential  |  March 2026"):
    add_rect(slide, 0, 7.1, 13.333, 0.4, fill_color=NAVY)
    add_text(slide, text, 0.3, 7.12, 12.5, 0.3, font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)

# ── SLIDE 1: TITLE ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=NAVY)
add_rect(slide, 0, 5.8, 13.333, 1.7, fill_color=STEEL_BLUE)

add_text(slide, "ENERGY SECTOR", 0.6, 1.0, 12, 1.0, font_size=16, bold=False, color=STEEL_BLUE)
add_text(slide, "Procurement AI", 0.6, 1.7, 12, 1.6, font_size=52, bold=True, color=WHITE)
add_text(slide, "Competitive Intelligence Report", 0.6, 3.2, 12, 0.6, font_size=22, bold=False, color=MID_GREY)
add_divider(slide, 3.85, color=STEEL_BLUE)
add_text(slide, "What Equinor, Shell, Vattenfall, BP, TotalEnergies & RWE are doing —\nand where the gap is for Statkraft and Manu Forti", 0.6, 4.0, 12, 1.2, font_size=14, color=WHITE)
add_text(slide, "Prepared by Aiden  |  Manu Forti Intelligence  |  March 17, 2026", 0.6, 5.95, 12, 0.5, font_size=11, color=WHITE, italic=True)

# ── SLIDE 2: EXECUTIVE SUMMARY ─────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "EXECUTIVE SUMMARY", 0.4, 0.08, 9, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "Operational AI is generating real savings — Procurement AI is still early stage", 0.4, 0.35, 12.5, 0.55, font_size=20, bold=True, color=WHITE)
slide_footer(slide)

# Insight banner
add_rect(slide, 0.3, 1.15, 12.73, 0.85, fill_color=LIGHT_BLUE_BG)
shape = slide.shapes.add_shape(1, Inches(0.3), Inches(1.15), Inches(0.08), Inches(0.85))
shape.fill.solid(); shape.fill.fore_color.rgb = STEEL_BLUE; shape.line.fill.background()
add_text(slide, "Major energy companies are ahead on operational AI but procurement-specific capability lags 2-3 years. This is the window for Manu Forti.", 0.5, 1.2, 12.3, 0.75, font_size=13, color=DARK_GREY, bold=False)

# Key figures - 5 stat boxes
stats = [
    ("$130M", "Equinor AI savings\nin 2025"),
    ("NOK 1B+", "Equinor projected\nrobotics savings"),
    ("100,000+", "TotalEnergies\nsuppliers / $31B spend"),
    ("2×", "AI use in procurement\nyear-over-year (Hackett 2026)"),
    ("+8%", "Procurement workloads\nrising as headcount falls"),
]
box_w = 2.3
for i, (num, label) in enumerate(stats):
    x = 0.3 + i * (box_w + 0.18)
    add_rect(slide, x, 2.15, box_w, 1.55, fill_color=NAVY)
    add_text(slide, num, x + 0.1, 2.25, box_w - 0.2, 0.65, font_size=26, bold=True, color=STEEL_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 2.88, box_w - 0.2, 0.75, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Gap statement
add_rect(slide, 0.3, 3.9, 12.73, 1.0, fill_color=RGBColor(0x0F, 0x2D, 0x5A))
add_text(slide, "⚡  The Gap:", 0.5, 3.97, 2.5, 0.4, font_size=13, bold=True, color=AMBER)
add_text(slide, "All majors lagging on Category Strategy AI, Supplier Evaluation Intelligence, and Procurement Advisory AI — exactly where Manu Forti operates.", 2.9, 3.97, 10.0, 0.8, font_size=13, color=WHITE)

# Source line
add_text(slide, "Sources: Equinor.com (Jan 2026), Hackett Group Key Issues Study (Mar 2026), Klover.ai, Vattenfall/Hackett Case Study", 0.3, 5.1, 12.7, 0.4, font_size=9, color=MID_GREY, italic=True)

# ── SLIDE 3: EQUINOR ───────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "EQUINOR  |  Norway  ⭐ Most Relevant Peer", 0.4, 0.08, 10, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "$130M AI savings in 2025 — but procurement AI is still emerging", 0.4, 0.35, 12.5, 0.55, font_size=20, bold=True, color=WHITE)
slide_footer(slide)

# Left: savings breakdown
add_rect(slide, 0.3, 1.15, 5.8, 4.7, fill_color=RGBColor(0x04, 0x12, 0x2E))
add_text(slide, "2025 AI Savings Breakdown", 0.5, 1.25, 5.4, 0.4, font_size=13, bold=True, color=STEEL_BLUE)

items = [
    ("$120M", "Since 2020 — Predictive maintenance\n24,000 sensors, 700+ rotating machines"),
    ("$12M", "Johan Sverdrup Phase 3 well planning\nAI found solution humans missed"),
    ("2M km²", "Seismic data interpreted in 2025\n10× faster than manual process"),
    ("100+", "New AI use cases identified\nProcurement on the roadmap"),
]
y = 1.75
for val, desc in items:
    add_rect(slide, 0.45, y, 1.1, 0.6, fill_color=STEEL_BLUE)
    add_text(slide, val, 0.47, y + 0.05, 1.06, 0.5, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, 1.65, y + 0.05, 4.3, 0.55, font_size=10, color=WHITE)
    y += 0.82

# Right: strategic read
add_rect(slide, 6.4, 1.15, 6.6, 4.7, fill_color=LIGHT_BLUE_BG)
shape = slide.shapes.add_shape(1, Inches(6.4), Inches(1.15), Inches(0.08), Inches(4.7))
shape.fill.solid(); shape.fill.fore_color.rgb = STEEL_BLUE; shape.line.fill.background()

add_text(slide, "Strategic Read for Statkraft", 6.6, 1.25, 6.2, 0.4, font_size=13, bold=True, color=NAVY)
add_text(slide, 'Quote: "AI is a central part of our operations. Moving forward, AI will become even more important for solving industrial tasks safely, faster, more profitably, and at scale."\n\n— Hege Skryseth, EVP Technology, Digital & Innovation', 6.6, 1.75, 6.2, 1.4, font_size=11, color=DARK_GREY, italic=False)
add_divider(slide, 3.25, color=MID_GREY)
points = [
    "Equinor is 2-3 years ahead of Statkraft on operational AI",
    "Procurement-specific AI not yet public — on the roadmap",
    "Norwegian Continental Shelf 2035 ambition: AI is 'crucial'",
    "Statkraft can get ahead now before Equinor sets the standard",
]
y2 = 3.35
for pt in points:
    add_text(slide, f"▸  {pt}", 6.6, y2, 6.2, 0.4, font_size=11, color=DARK_GREY)
    y2 += 0.45

# ── SLIDE 4: SHELL & VATTENFALL ────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "SHELL  &  VATTENFALL", 0.4, 0.08, 10, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "Shell leads on contract AI; Vattenfall is paying Hackett Group to do what Manu Forti does", 0.4, 0.35, 12.5, 0.55, font_size=18, bold=True, color=WHITE)
slide_footer(slide)

# Shell box
add_rect(slide, 0.3, 1.15, 6.0, 4.7, fill_color=RGBColor(0x04, 0x12, 0x2E))
add_text(slide, "🐚  Shell", 0.5, 1.22, 5.6, 0.38, font_size=15, bold=True, color=WHITE)
add_text(slide, "Most advanced on procurement AI among majors", 0.5, 1.58, 5.6, 0.35, font_size=11, color=STEEL_BLUE)
add_divider(slide, 2.0, color=STEEL_BLUE)
shell_pts = [
    ("Contract Authoring", "AI writes first drafts of procurement contracts automatically"),
    ("Supplier Compliance", "AI monitors obligations, flags deviations in real-time"),
    ("Workflow Embedding", "AI inside core procurement workflows, not just analytics"),
    ("Goal", "Free up procurement professionals for strategic supplier relationships and diversity"),
]
y = 2.15
for title, desc in shell_pts:
    add_text(slide, f"► {title}:", 0.5, y, 2.0, 0.38, font_size=11, bold=True, color=STEEL_BLUE)
    add_text(slide, desc, 2.45, y, 3.7, 0.38, font_size=11, color=WHITE)
    y += 0.6

add_rect(slide, 0.4, 4.85, 5.8, 0.65, fill_color=STEEL_BLUE)
add_text(slide, "Key Insight: Shell is 2-3 years ahead on contract AI. Statkraft's HV equipment contracts are a direct analogue use case.", 0.55, 4.9, 5.55, 0.6, font_size=10, bold=True, color=WHITE)

# Vattenfall box
add_rect(slide, 6.7, 1.15, 6.3, 4.7, fill_color=LIGHT_BLUE_BG)
shape = slide.shapes.add_shape(1, Inches(6.7), Inches(1.15), Inches(0.08), Inches(4.7))
shape.fill.solid(); shape.fill.fore_color.rgb = AMBER; shape.line.fill.background()
add_text(slide, "⚡  Vattenfall", 6.9, 1.22, 5.9, 0.38, font_size=15, bold=True, color=NAVY)
add_text(slide, "Most comparable peer — European utility in active transformation", 6.9, 1.58, 5.9, 0.35, font_size=11, color=STEEL_BLUE)
add_divider(slide, 2.0, color=AMBER)
vattenfall_pts = [
    "Engaged The Hackett Group for AI-guided procurement transformation",
    "Led by Arjan De Jong (Head of Solutions & Insights)",
    "Hackett providing: insights, expertise, practical AI guidance",
    "Focus: Moving from transactional to strategic, AI-augmented category management",
    "Revenue: $26.47B (2024) — significant procurement spend at stake",
]
y = 2.15
for pt in vattenfall_pts:
    add_text(slide, f"▸  {pt}", 6.9, y, 5.9, 0.45, font_size=11, color=DARK_GREY)
    y += 0.52

add_rect(slide, 6.8, 4.85, 6.1, 0.65, fill_color=AMBER)
add_text(slide, "🎯  Vattenfall is paying Hackett Group (£50K+ retainer) for what Manu Forti delivers in days at a fraction of the cost.", 6.95, 4.9, 5.85, 0.6, font_size=10, bold=True, color=NAVY)

# ── SLIDE 5: BP, TOTALENERGIES, RWE ────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "BP  |  TOTALENERGIES  |  RWE", 0.4, 0.08, 10, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "Scale and tech investment are forcing procurement AI — whether they're ready or not", 0.4, 0.35, 12.5, 0.55, font_size=20, bold=True, color=WHITE)
slide_footer(slide)

panels = [
    ("BP", RGBColor(0x04, 0x12, 0x2E), WHITE, [
        "AI agents routing purchase requests and compliance checks",
        "Workflow optimisation: named as a specific AI workstream",
        "Cost reduction focus post-2020 restructuring",
        "SAP Ariba + custom AI layer",
        "Every efficiency gain matters under cost pressure",
    ]),
    ("TotalEnergies", LIGHT_BLUE_BG, DARK_GREY, [
        "100,000+ suppliers globally / $31B spend (2024)",
        "Positioning as 'data company evolving into AI company'",
        "1% AI efficiency = $310M annual savings at their scale",
        "Battery supply chain AI via Saft subsidiary",
        "Present in ~120 countries — complexity demands AI",
    ]),
    ("RWE", RGBColor(0x04, 0x12, 0x2E), WHITE, [
        "AWS, HPE, Infosys, FPT: major AI tech partnerships",
        "Microsoft for AI tools across operations",
        "Bought Vattenfall's Norfolk offshore assets (€1.1B)",
        "Scaling fast — procurement complexity increasing",
        "Less visible externally but significant investment",
    ]),
]
x = 0.3
for name, bg, fg, pts in panels:
    add_rect(slide, x, 1.15, 4.1, 4.7, fill_color=bg)
    title_color = STEEL_BLUE if bg != LIGHT_BLUE_BG else NAVY
    add_text(slide, name, x + 0.2, 1.22, 3.7, 0.4, font_size=15, bold=True, color=title_color)
    add_divider(slide, 1.7, color=STEEL_BLUE)
    y = 1.85
    for pt in pts:
        add_text(slide, f"▸  {pt}", x + 0.2, y, 3.7, 0.5, font_size=10, color=fg)
        y += 0.6
    x += 4.35

# ── SLIDE 6: MARKET CONTEXT ────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "MARKET CONTEXT  |  Hackett Group 2026 Study — Published Today", 0.4, 0.08, 12, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "AI use in procurement nearly doubled — workloads rising while headcount falls", 0.4, 0.35, 12.5, 0.55, font_size=20, bold=True, color=WHITE)
slide_footer(slide)

# 4 stat blocks top row
hackett_stats = [
    ("~2×", "AI use in procurement\nnearly doubled YoY"),
    ("+8%", "Procurement workloads\nin 2026"),
    ("34%", "Efficiency gains from\nAI-driven procurement"),
    ("23%", "Cost savings from\nAI procurement platforms"),
]
bw = 2.9
for i, (num, label) in enumerate(hackett_stats):
    x = 0.3 + i * (bw + 0.2)
    add_rect(slide, x, 1.15, bw, 1.3, fill_color=STEEL_BLUE)
    add_text(slide, num, x + 0.1, 1.22, bw - 0.2, 0.65, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 1.82, bw - 0.2, 0.55, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Gartner stats
add_text(slide, "Gartner 2025:", 0.3, 2.65, 2.5, 0.35, font_size=12, bold=True, color=NAVY)
gartner_pts = [
    ("74%", "of procurement leaders say their data isn't AI-ready"),
    ("50-80%", "of current procurement work can be automated by GenAI"),
]
x = 0.3
for stat, desc in gartner_pts:
    add_rect(slide, x, 3.05, 1.3, 0.65, fill_color=NAVY)
    add_text(slide, stat, x + 0.05, 3.1, 1.2, 0.55, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + 1.4, 3.15, 4.3, 0.5, font_size=11, color=DARK_GREY)
    x += 6.2

# Gap table
add_rect(slide, 0.3, 3.9, 12.73, 0.35, fill_color=NAVY)
add_text(slide, "WHERE ENERGY COMPANIES ARE  →  AND WHERE THE GAP IS", 0.5, 3.93, 12.3, 0.28, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Operational AI (maintenance, drilling, seismic)", "✅ Advanced", GREEN),
    ("Contract automation & supplier compliance AI", "⚠️ Developing", AMBER),
    ("Category Strategy AI & Supplier Evaluation Intelligence", "❌ Lagging — The Manu Forti Gap", RGBColor(0xC0, 0x39, 0x2B)),
]
y = 4.3
for area, status, status_color in rows:
    add_rect(slide, 0.3, y, 9.0, 0.45, fill_color=LIGHT_BLUE_BG)
    add_rect(slide, 9.35, y, 3.68, 0.45, fill_color=LIGHT_BLUE_BG)
    add_text(slide, area, 0.5, y + 0.07, 8.7, 0.35, font_size=11, color=DARK_GREY)
    add_text(slide, status, 9.45, y + 0.07, 3.5, 0.35, font_size=11, bold=True, color=status_color)
    y += 0.5

# ── SLIDE 7: COMPETITIVE POSITIONING ──────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "COMPETITIVE POSITIONING", 0.4, 0.08, 10, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "Manu Forti: faster, cheaper, and more procurement-specific than any alternative", 0.4, 0.35, 12.5, 0.55, font_size=20, bold=True, color=WHITE)
slide_footer(slide)

# Comparison table
headers = ["", "Manu Forti", "Equinor\n(Internal)", "Shell\n(Internal)", "Vattenfall\n(Hackett)", "Big Consulting"]
col_w = [2.5, 1.8, 1.8, 1.8, 1.8, 1.8]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.02)

# Header row
y = 1.15
for i, (h, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    bg = STEEL_BLUE if i == 1 else NAVY
    add_rect(slide, cx, y, cw, 0.5, fill_color=bg)
    add_text(slide, h, cx + 0.05, y + 0.05, cw - 0.1, 0.4, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Procurement AI Focus", ["✅ Core product", "⚠️ Emerging", "✅ Contract/\ncompliance", "✅ Transformation", "✅ Generic"]),
    ("Speed to Value", ["Days", "Years", "Months", "12+ months", "6+ months"]),
    ("Price Point", ["€99 – €3,999", "Internal\n(millions)", "Internal\n(millions)", "£50K+ retainer", "£100K+"]),
    ("Renewables Focus", ["✅ Wind, solar, HV", "O&G first", "O&G first", "✅ Utility", "Generic"]),
    ("AI-Native", ["✅ Built on agents", "Legacy + AI", "Legacy + AI", "Legacy + AI", "Human + AI"]),
    ("Supplier Analysis", ["✅ Product 1", "Internal", "Internal", "Internal", "Not offered"]),
]

y = 1.68
for row_label, row_vals in rows:
    bg = RGBColor(0xF0, 0xF7, 0xFF)
    add_rect(slide, col_x[0], y, col_w[0], 0.48, fill_color=NAVY)
    add_text(slide, row_label, col_x[0] + 0.1, y + 0.08, col_w[0] - 0.15, 0.38, font_size=10, bold=True, color=WHITE)
    for i, (val, cw, cx) in enumerate(zip(row_vals, col_w[1:], col_x[1:])):
        cell_bg = STEEL_BLUE if i == 0 else bg
        cell_fg = WHITE if i == 0 else DARK_GREY
        add_rect(slide, cx, y, cw, 0.48, fill_color=cell_bg)
        add_text(slide, val, cx + 0.05, y + 0.06, cw - 0.1, 0.38, font_size=9, color=cell_fg, align=PP_ALIGN.CENTER)
    y += 0.5

# ── SLIDE 8: STATKRAFT IMPLICATIONS ────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 1.05, fill_color=NAVY)
add_text(slide, "IMPLICATIONS FOR STATKRAFT", 0.4, 0.08, 10, 0.35, font_size=11, color=STEEL_BLUE, bold=True)
add_text(slide, "Act now — the procurement AI gap is closing, and early movers will set the standard", 0.4, 0.35, 12.5, 0.55, font_size=20, bold=True, color=WHITE)
slide_footer(slide)

actions = [
    ("1", "Don't wait for IT programmes", "Start with high-value categories where AI analysis already exists — HV equipment and EPC contractors are ready now. Use Manu Forti's Product 1 for vendor qualification today."),
    ("2", "Shell's contract AI is your roadmap", "Shell's contract authoring automation is the direction of travel. Statkraft's HV equipment contracts (€100M+ framework agreements) are exactly the right complexity level for AI-assisted drafting."),
    ("3", "Vattenfall is paying £50K+ for what you already have access to", "Vattenfall engaged Hackett Group to guide their AI procurement transformation. Statkraft can use Manu Forti's Category Strategy methodology at a fraction of the cost and 10× the speed."),
    ("4", "Build a Procurement AI roadmap now", "Use the Manu Forti methodology as the framework. Positions Statkraft ahead of Norwegian peers and ahead of where Equinor will be in 2-3 years."),
    ("5", "Monitor Equinor's public disclosures", "They're the leading benchmark for Norwegian energy AI. Their next milestone: NOK 1B+ annual savings from robotics. Watch the procurement-specific use cases emerging."),
]

y = 1.2
for num, title, desc in actions:
    add_rect(slide, 0.3, y, 0.55, 0.75, fill_color=STEEL_BLUE)
    add_text(slide, num, 0.3, y + 0.18, 0.55, 0.4, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.9, y, 12.1, 0.75, fill_color=LIGHT_BLUE_BG)
    add_text(slide, title, 1.05, y + 0.04, 3.5, 0.35, font_size=12, bold=True, color=NAVY)
    add_text(slide, desc, 4.5, y + 0.08, 8.4, 0.6, font_size=10, color=DARK_GREY)
    y += 0.83

# ── SLIDE 9: CLOSING ───────────────────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=NAVY)
add_rect(slide, 0, 5.5, 13.333, 2.0, fill_color=STEEL_BLUE)

add_text(slide, "THE BOTTOM LINE", 0.6, 1.0, 12, 0.45, font_size=14, color=STEEL_BLUE, bold=True)
add_text(slide, "Every major energy company is investing in AI.\nProcurement is the last frontier — and the biggest prize.", 0.6, 1.5, 12, 1.4, font_size=32, bold=True, color=WHITE)
add_divider(slide, 3.0, color=STEEL_BLUE)

bullets = [
    "Equinor proved AI delivers at scale ($130M in a year). Procurement is next.",
    "Shell proved contract AI works. Procurement cycle times can be cut dramatically.",
    "Vattenfall proved the demand. Companies will pay to get there — the question is who guides them.",
    "Manu Forti is already built for this. Product 1 + Category Strategy = the answer the market needs.",
]
y = 3.15
for b in bullets:
    add_text(slide, f"▸  {b}", 0.6, y, 12, 0.4, font_size=13, color=WHITE)
    y += 0.48

add_text(slide, "Manu Forti Intelligence  |  Procurement AI for the Energy Transition", 0.6, 5.65, 12, 0.45, font_size=16, bold=True, color=WHITE)
add_text(slide, "Contact: manuforti.as@gmail.com  |  www.manuforti.no", 0.6, 6.15, 12, 0.35, font_size=12, color=WHITE, italic=True)
add_text(slide, "Prepared by Aiden  |  March 17, 2026  |  Confidential", 0.6, 6.6, 12, 0.35, font_size=11, color=RGBColor(0xBB, 0xCC, 0xDD), italic=True)

# Save
output_path = "/Users/jonathonmilne/.openclaw/workspace/venture/Energy_Procurement_AI_Intelligence_March2026.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
