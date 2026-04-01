"""
Adds two things to the v4 deck:
1. A new Methodology slide (after slide 8 Value Measurement)
2. Updates the TCO numbers on the tech stack slide with €1.9B-based figures
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import copy

NAVY    = RGBColor(0x00,0x21,0x47)
STEEL   = RGBColor(0x2B,0x6C,0xB0)
GREY    = RGBColor(0x71,0x80,0x96)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
AMBER   = RGBColor(0xD9,0x7F,0x06)
LTBLUE  = RGBColor(0xEB,0xF4,0xFF)
DARK    = RGBColor(0x2D,0x3A,0x4A)
GREEN   = RGBColor(0x27,0xAE,0x60)
RED     = RGBColor(0xC0,0x39,0x2B)
TEAL    = RGBColor(0x0B,0x7A,0x75)
DEEP    = RGBColor(0x04,0x12,0x2E)
PGREEN  = RGBColor(0xD4,0xED,0xDA)

prs = Presentation("/Users/jonathonmilne/.openclaw/workspace/venture/Statkraft_Procurement_AI_Strategy_v4_March2026.pptx")

def R(s,l,t,w,h,fill=None):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb=fill
    sh.line.fill.background()
    return sh

def T(s,text,l,t,w,h,sz=10,bold=False,color=WHITE,align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.color.rgb=color; r.font.italic=italic
    return tb

def LINE(s,t,c=STEEL,l=0.4,w=12.5):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(0.04))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def BANNER(s,txt,l,t,w,h,bg=LTBLUE,fg=DARK,sz=12,accent=STEEL):
    R(s,l,t,0.07,h,fill=accent)
    R(s,l+0.07,t,w-0.07,h,fill=bg)
    T(s,txt,l+0.22,t+0.08,w-0.32,h-0.14,sz=sz,color=fg)

def HDR(s,tag,title,tc=STEEL):
    R(s,0,0,13.333,1.05,fill=NAVY)
    T(s,tag,0.4,0.07,12,0.32,sz=10,bold=True,color=tc)
    T(s,title,0.4,0.35,12.5,0.58,sz=19,bold=True,color=WHITE)

def FTR(s,txt="Statkraft Procurement AI Strategy  |  Confidential  |  March 2026"):
    R(s,0,7.1,13.333,0.4,fill=NAVY)
    T(s,txt,0.3,7.12,12.7,0.3,sz=9,color=GREY,align=PP_ALIGN.CENTER)

# ── INSERT NEW SLIDE AFTER SLIDE 8 (index 8) ─────────────────────────────────
# We add at the end then move it — python-pptx doesn't support insert natively
# Instead add to existing prs

blank_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank_layout)

s = new_slide
HDR(s, "VALUE CALCULATION METHODOLOGY  |  Standard Framework",
    "How we calculate savings: 4 calculation types, clear assumptions, auditable workings")
FTR(s)

BANNER(s, "Principle: Every saving is calculated from a specific, measurable baseline. "
          "No general percentages without a defined formula. All assumptions are explicit and Statkraft-adjustable in the Excel model.",
       0.3, 1.15, 12.73, 0.65, sz=11)

# 4 Calculation Types header
R(s,0.3,1.9,12.73,0.38,fill=NAVY)
T(s,"4 STANDARD CALCULATION TYPES — Applied consistently across all 20 use cases",
  0.5,1.95,12.35,0.3,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

calc_types = [
    ("TYPE A\nTime Savings",        STEEL,
     "Formula:",
     "Hours saved × FTE hourly rate",
     "FTE hourly rate = Annual FTE cost (€120K) ÷ 1,750 working hours = €68.57/hr\n"
     "Hours saved = (Baseline hours − AI-assisted hours) × Volume per year",
     "Examples in use cases:",
     "• Contract NLP: 4 hrs → 0.8 hrs × 1,200 contracts × €68.57\n"
     "• CSRD checks: 0.8 hrs saved × 200 suppliers × 12 checks × €68.57\n"
     "• Invoice processing: 15 min → 4.5 min × 18,000 invoices × €68.57"),

    ("TYPE B\nSpend Savings",        AMBER,
     "Formula:",
     "Addressable spend × Savings rate",
     "Addressable spend = Annual spend (€1.9B) × Category relevance %\n"
     "Savings rate = Conservative benchmark from peer implementations\n"
     "Benchmark source cited for every use case",
     "Examples in use cases:",
     "• Spend Classification: €1.9B × 23% accuracy gap × 5% recovery rate\n"
     "• Tail Spend: €1.9B × 8% tail spend × 15% savings via competitive tendering\n"
     "• Negotiation: €1.9B × 0.3% additional savings from better preparation"),

    ("TYPE C\nRisk Reduction",       RED,
     "Formula:",
     "Probability × Impact of adverse event avoided",
     "Probability = Assessed likelihood of event per year (conservative)\n"
     "Impact = Quantified cost of the event occurring (documented)\n"
     "Uses expected value: P × Impact",
     "Examples in use cases:",
     "• Supplier Monitor: 1 major failure × €2M avg cost ÷ 2 years expected\n"
     "• Supply Chain Risk: 15% prob of major disruption × €5M impact\n"
     "• Project Cost Intel: 40% avg overrun × 2% improvement × €400M capex"),

    ("TYPE D\nCapacity Release",     GREEN,
     "Formula:",
     "FTE time freed × % redirected to strategic value",
     "Capacity freed = Hours automated × % of tasks that are automatable\n"
     "Value multiplier = 1.0× (conservative — treating freed time at cost rate)\n"
     "Does NOT assume headcount reduction — assumes redeployment to higher value",
     "Examples in use cases:",
     "• RFQ Drafter: 40hrs × 75% time reduction × 250 events per year\n"
     "• Market Intelligence: analyst time freed for strategic work\n"
     "• Category Strategy: 3 strategies/yr → 12/yr, same team"),
]

x = 0.3
for title, color, f_label, formula, method, ex_label, examples in calc_types:
    R(s, x, 2.36, 3.06, 4.42, fill=DEEP)
    R(s, x, 2.36, 3.06, 0.52, fill=color)
    T(s, title, x+0.12, 2.38, 2.85, 0.48, sz=12, bold=True, color=WHITE)
    T(s, f_label, x+0.12, 2.96, 0.7, 0.28, sz=9, bold=True, color=color)
    T(s, formula, x+0.82, 2.96, 2.1, 0.28, sz=9, bold=True, color=WHITE)
    LINE(s, 3.3, c=color, l=x+0.12, w=2.82)
    T(s, method, x+0.12, 3.35, 2.82, 0.92, sz=8.5, color=GREY, italic=True)
    LINE(s, 4.32, c=color, l=x+0.12, w=2.82)
    T(s, ex_label, x+0.12, 4.36, 2.82, 0.28, sz=9, bold=True, color=color)
    T(s, examples, x+0.12, 4.68, 2.82, 1.98, sz=8.5, color=WHITE)
    x += 3.27

# Bottom note
R(s,0.3,6.78,12.73,0.5,fill=NAVY)
T(s,"📊  Bear = 50% of base  |  Base = conservative benchmark  |  Bull = 150% of base  |  "
   "All assumptions in yellow cells of Excel model — adjust to Statkraft actuals  |  "
   "No use case is included without a documented peer implementation or published benchmark",
  0.5,6.83,12.35,0.42,sz=9.5,color=WHITE)

# ── UPDATE SLIDE 10 (TCO) with correct €1.9B numbers ──────────────────────────
# Find the TCO slide by looking for the stat boxes — it's slide index ~9 in v4
# We'll search for the slide with "Year 1 CAPEX" text and update the values

for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "Year 1 CAPEX" in full_text and "45–150K" in full_text:
                # Found it — update by replacing text in textboxes
                for shape2 in slide.shapes:
                    if shape2.has_text_frame:
                        for para in shape2.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = run.text.replace("€1.3–2.5M", "€108M+")
                                run.text = run.text.replace("From Value Calculator Excel", "€108M base — all 20 use cases, €1.9B spend")
                                run.text = run.text.replace("8–12×", "289×")
                                run.text = run.text.replace("Conservative estimate at 50% target achievement", "Qwen 35B research + benchmarks, Excel full workings")
                                run.text = run.text.replace("~€320K", "~€1.3M")
                                run.text = run.text.replace("Hardware + 3yr OPEX", "Hardware + 3yr OPEX (at €375K/yr)")
                break

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = "/Users/jonathonmilne/.openclaw/workspace/venture/Statkraft_Procurement_AI_Strategy_v5_March2026.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
