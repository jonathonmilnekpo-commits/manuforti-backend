from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── COLOURS ────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x00, 0x21, 0x47)
STEEL      = RGBColor(0x2B, 0x6C, 0xB0)
MID_GREY   = RGBColor(0x71, 0x80, 0x96)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
AMBER      = RGBColor(0xD9, 0x7F, 0x06)
LT_BLUE    = RGBColor(0xEB, 0xF4, 0xFF)
DARK       = RGBColor(0x2D, 0x3A, 0x4A)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xC0, 0x39, 0x2B)
TEAL       = RGBColor(0x0B, 0x7A, 0x75)
DEEP_NAVY  = RGBColor(0x04, 0x12, 0x2E)
PALE_AMBER = RGBColor(0xFF, 0xF3, 0xCD)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── HELPERS ────────────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None, lw=1.0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    return s

def txt(slide, text, l, t, w, h, size=11, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def line(slide, t, color=STEEL, l=0.4, w=12.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def header(slide, tag, title, tag_color=STEEL):
    rect(slide, 0, 0, 13.333, 1.05, fill=NAVY)
    txt(slide, tag,   0.4, 0.08, 12, 0.32, size=10, bold=True,  color=tag_color)
    txt(slide, title, 0.4, 0.35, 12.5, 0.58, size=19, bold=True, color=WHITE)

def footer(slide, text="Statkraft Procurement AI Strategy  |  Confidential  |  March 2026"):
    rect(slide, 0, 7.1, 13.333, 0.4, fill=NAVY)
    txt(slide, text, 0.3, 7.12, 12.7, 0.3, size=9, color=MID_GREY, align=PP_ALIGN.CENTER)

def insight(slide, text, t=1.15, l=0.3, w=12.73, h=0.75):
    rect(slide, l, t, 0.07, h, fill=STEEL)
    rect(slide, l+0.07, t, w-0.07, h, fill=LT_BLUE)
    txt(slide, text, l+0.22, t+0.1, w-0.35, h-0.15, size=12, color=DARK)

def bullet_box(slide, title, pts, l, t, w, h, bg=DEEP_NAVY, fg=WHITE, title_color=STEEL, size=10):
    rect(slide, l, t, w, h, fill=bg)
    txt(slide, title, l+0.15, t+0.1, w-0.25, 0.35, size=12, bold=True, color=title_color)
    y = t + 0.5
    for pt in pts:
        txt(slide, f"▸  {pt}", l+0.15, y, w-0.25, 0.42, size=size, color=fg)
        y += 0.44

def stat_row(slide, stats, t, bw=2.3, bg=NAVY, num_color=STEEL):
    x = 0.3
    for num, label in stats:
        rect(slide, x, t, bw, 1.3, fill=bg)
        txt(slide, num,   x+0.1, t+0.1,  bw-0.2, 0.65, size=24, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        txt(slide, label, x+0.1, t+0.72, bw-0.2, 0.52, size=9,  color=WHITE, align=PP_ALIGN.CENTER)
        x += bw + 0.19

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 0, 5.7, 13.333, 1.8, fill=STEEL)
txt(s, "PROCUREMENT AI STRATEGY", 0.6, 0.85, 12, 0.5, size=15, bold=False, color=STEEL)
txt(s, "Statkraft", 0.6, 1.35, 12, 1.1, size=54, bold=True, color=WHITE)
txt(s, "From Framework to Agentic Capability", 0.6, 2.45, 12, 0.6, size=24, color=MID_GREY)
line(s, 3.2, color=STEEL)
txt(s, "Market Intelligence  ·  Peer Analysis  ·  Agentic Use Cases  ·  Cybersecurity Stack  ·  Execution Roadmap",
    0.6, 3.35, 12, 0.45, size=13, color=WHITE)
txt(s, "Prepared by Aiden  |  Manu Forti Intelligence  |  March 17, 2026  |  Confidential",
    0.6, 5.85, 12, 0.45, size=11, color=WHITE, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "EXECUTIVE SUMMARY", "The market has moved — Statkraft must move from pilots to agentic production")
footer(s)
insight(s, "94% of procurement executives now use GenAI weekly. Only 4% have achieved large-scale deployment. "
           "The gap between ambition and production is the defining challenge of 2026 — and Statkraft's opportunity to lead in Norway.")

stat_row(s, [
    ("94%",    "of procurement execs\nuse GenAI weekly"),
    ("4%",     "have large-scale\ndeployment (Hackett)"),
    ("25–40%", "efficiency gain from\nagentic AI (McKinsey)"),
    ("500M NOK", "Equinor annual savings\nfrom automation"),
    ("80%",    "of CPOs plan GenAI\nin 3 years (EY)"),
], t=2.05, bw=2.42)

# 3 column summary
cols = [
    ("🔍 Where We Are",   DEEP_NAVY,  WHITE,  [
        "Framework exists — governance structure in place",
        "Copilot Studio as primary tool (individual track)",
        "No use case library or ROI measurement",
        "No agentic AI track (Track D missing)",
        "Data governance vague on CONFIDENTIAL data",
    ]),
    ("🎯 Where We Need to Be",  STEEL, WHITE, [
        "AI-enabled category management at scale",
        "8 agentic workers running continuously",
        "Value measurement on every initiative",
        "On-prem deployment for sensitive data",
        "Statkraft sets the Nordic standard by 2027",
    ]),
    ("🚀 How We Get There",  RGBColor(0x0A,0x45,0x1A), WHITE, [
        "3 immediate pilots: spend classify, contract NLP, supplier monitor",
        "Add Track D: Agentic AI governance tier",
        "Deploy on-prem LLM stack (Llama 3 / Mistral)",
        "7-layer cybersecurity architecture",
        "Quarterly value reviews — scale or kill",
    ]),
]
x = 0.3
for title, bg, fg, pts in cols:
    rect(s, x, 3.55, 4.2, 3.22, fill=bg)
    txt(s, title, x+0.15, 3.62, 3.9, 0.38, size=12, bold=True, color=AMBER)
    y = 4.05
    for pt in pts:
        txt(s, f"▸  {pt}", x+0.15, y, 3.9, 0.44, size=9.5, color=fg)
        y += 0.46
    x += 4.37

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MARKET INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "MARKET INTELLIGENCE  |  2025/26 Research", "The data is in: procurement AI is real, proven, and scaling fast")
footer(s)
insight(s, "43% of organisations actively deploying AI in procurement (doubled YoY, Hackett 2026). "
           "Workloads up 8% while headcount falls — AI is the only bridge.")

# Top use cases
rect(s, 0.3, 2.0, 6.1, 4.7, fill=DEEP_NAVY)
txt(s, "Top CPO Use Cases (2025/26)", 0.5, 2.1, 5.7, 0.38, size=13, bold=True, color=STEEL)
use_cases = [
    ("53%", "Spend analytics & classification"),
    ("42%", "RFQ / RFP generation"),
    ("41%", "Contract summarisation & review"),
    ("38%", "Supplier risk monitoring"),
    ("31%", "Demand forecasting (MRO/capex)"),
    ("28%", "Negotiation preparation"),
    ("22%", "Compliance monitoring (CSRD/ESG)"),
]
y = 2.58
for pct, label in use_cases:
    rect(s, 0.45, y, 0.9, 0.46, fill=STEEL)
    txt(s, pct, 0.47, y+0.08, 0.86, 0.32, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, 1.45, y+0.1, 4.8, 0.32, size=11, color=WHITE)
    y += 0.54

# Right: key stats + McKinsey + EY
rect(s, 6.7, 2.0, 6.3, 4.7, fill=LT_BLUE)
rect(s, 6.7, 2.0, 0.07, 4.7, fill=STEEL)
txt(s, "Key Research Findings", 6.95, 2.1, 5.8, 0.38, size=13, bold=True, color=NAVY)
findings = [
    ("McKinsey", "25–40% efficiency improvement potential from agentic AI in procurement — not assistive, agentic."),
    ("Hackett 2026", "Procurement workloads +8% in 2026 while headcount and budgets fall. AI is no longer optional."),
    ("EY CPO Survey", "80% of CPOs plan to deploy GenAI in 3 years. Only 36% have meaningful implementations today."),
    ("Gartner", "74% of procurement leaders say their data isn't AI-ready. GenAI can automate 50–80% of procurement work."),
    ("Art of Procurement", "94% of procurement executives use GenAI at least weekly — up 44pp from 2023."),
]
y = 2.58
for source, finding in findings:
    txt(s, source+":", 6.95, y, 1.5, 0.38, size=10, bold=True, color=STEEL)
    txt(s, finding, 8.35, y, 4.5, 0.5, size=10, color=DARK)
    y += 0.7

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PEER ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "ENERGY SECTOR PEER ANALYSIS", "What your direct peers are actually running in production")
footer(s)

peers = [
    ("Equinor\n🇳🇴 Norway", DEEP_NAVY, [
        "$130M AI savings in 2025 (official disclosure Jan 2026)",
        "500M NOK/yr from low-code automation",
        "Automated spend classification at scale",
        "AI logistics: 13% inventory reduction",
        "100+ new AI use cases identified",
        "NOK 1B+ projected from robotics",
    ]),
    ("Shell / TotalEnergies", DEEP_NAVY, [
        "Continuous supplier financial monitoring agents",
        "Automated contract authoring (first draft)",
        "AI supplier compliance checks in real-time",
        "TotalEnergies: 100K suppliers / $31B spend",
        "Contract NLP: 80% faster review",
        "Embedding AI into core workflows — not dashboards",
    ]),
    ("Vattenfall / Ørsted / BP", DEEP_NAVY, [
        "Vattenfall: predictive MRO/capex demand (6-12m horizon)",
        "Ørsted + BP: autonomous RFQ generation from spec docs",
        "Saves 4–8 hrs per sourcing event",
        "Vattenfall engaged Hackett Group for AI transformation",
        "BP: AI agents routing POs and compliance checks",
        "Multiple majors: contract NLP on full portfolio",
    ]),
]
x = 0.3
for title, bg, pts in peers:
    rect(s, x, 1.15, 4.1, 5.55, fill=bg)
    txt(s, title, x+0.15, 1.22, 3.8, 0.55, size=13, bold=True, color=STEEL)
    line(s, 1.85, l=x+0.15, w=3.8)
    y = 2.0
    for pt in pts:
        txt(s, f"▸  {pt}", x+0.15, y, 3.8, 0.46, size=10, color=WHITE)
        y += 0.48
    x += 4.37

# Gap banner
rect(s, 0.3, 6.78, 12.73, 0.5, fill=STEEL)
txt(s, "⚡  The Statkraft Gap: Renewables-specific procurement AI (HV equipment, EPC, O&M) is unaddressed by any major. "
       "First mover advantage is available now.", 0.5, 6.83, 12.4, 0.42, size=11, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 6 WEAKNESSES
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "CRITICAL FRAMEWORK REVIEW  |  Current State Assessment",
       "6 weaknesses in the current Procurement AI framework — and why they matter")
footer(s)

weaknesses = [
    ("1", "No Value Measurement",
     "Framework describes how to build — not what success looks like. No KPIs, no ROI tracking, "
     "no 90-day post-implementation review. Initiatives exist but impact is unmeasured."),
    ("2", "Copilot Studio Is a Ceiling",
     "Restricting individuals to low-complexity chatbots prevents anything transformative. "
     "The 25–40% efficiency gain from McKinsey lives in agentic tools, not Copilot Lite."),
    ("3", "No Use Case Library",
     "Blank-sheet ideation without a framework produces low-value ideas. Teams need a ranked "
     "library of 10–15 high-impact use cases to start from, not a blank page."),
    ("4", "Citizen Dev Without Data Guardrails",
     "Local shared use path is vague. Who approves agent access to supplier financials, "
     "contract pricing, and negotiation positions? Current framework doesn't say."),
    ("5", "No Agentic AI Track",
     "Framework covers assistive AI only. Zero mention of autonomous agents that take actions. "
     "Track D is missing — this is where the real value sits. Peers are already there."),
    ("6", "Governance Is IT-Heavy, Procurement-Light",
     "Category managers know whether a tool solves a procurement problem. IT architects don't. "
     "The AI Forum needs procurement ownership, not just IT participation."),
]

x, y = 0.3, 1.15
col = 0
for num, title, desc in weaknesses:
    rect(s, x, y, 6.1, 1.55, fill=DEEP_NAVY)
    rect(s, x, y, 0.5, 1.55, fill=RED)
    txt(s, num, x+0.07, y+0.5, 0.38, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.65, y+0.12, 5.3, 0.35, size=12, bold=True, color=STEEL)
    txt(s, desc,  x+0.65, y+0.48, 5.3, 1.0,  size=9.5, color=WHITE)
    col += 1
    if col == 2:
        x = 0.3; y += 1.65; col = 0
    else:
        x += 6.55

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 10 IMPROVEMENT ACTIONS
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "10 CONCRETE IMPROVEMENT ACTIONS", "Immediate, structural, and pilot actions — in priority order")
footer(s)

rect(s, 0.3, 1.15, 3.5, 0.38, fill=RED)
txt(s, "IMMEDIATE (Weeks 1–4)", 0.45, 1.19, 3.3, 0.3, size=11, bold=True, color=WHITE)
rect(s, 4.05, 1.15, 4.15, 0.38, fill=AMBER)
txt(s, "STRUCTURAL (Months 1–3)", 4.2, 1.19, 3.9, 0.3, size=11, bold=True, color=WHITE)
rect(s, 8.45, 1.15, 4.55, 0.38, fill=GREEN)
txt(s, "PILOT FIRST (Months 1–2)", 8.6, 1.19, 4.3, 0.3, size=11, bold=True, color=WHITE)

immediate = [
    ("1", "ADD USE CASE LIBRARY",       "Top 10–15 high-value use cases ranked by effort vs. impact. Give teams a starting point."),
    ("2", "MANDATE VALUE MEASUREMENT",  "Every initiative declares: baseline metric, target, measurement method, review date. No go-live without this."),
    ("3", "ADD DATA GOVERNANCE RULES",  "Tag data by classification. Supplier pricing, contract terms, negotiation positions = RESTRICTED. Agents can only access at their clearance level."),
    ("4", "UPGRADE INDIVIDUAL TRACK",   "Allow Claude, Gemini, OpenAI with manager approval + data handling rules. Copilot Lite is not enough."),
]
structural = [
    ("5", "ADD TRACK D — AGENTIC AI",   "Autonomous agents need their own governance tier. Define HITL thresholds: auto below limit, human approval above."),
    ("6", "GIVE AI FORUM A BACKLOG",    "The forum should own a prioritised use case backlog. Make it a product team, not a discussion group."),
    ("7", "90-DAY POST GO-LIVE REVIEWS","Every initiative reviewed vs. declared success metrics. Kill what doesn't work. Scale what does."),
]
pilots = [
    ("8", "SPEND CLASSIFICATION",       "Fast win, 70–90% accuracy, high ROI, zero supplier-facing risk. Cleans up spend data for everything else."),
    ("9", "CONTRACT NLP SCANNING",      "Flag non-standard terms across existing contract portfolio. Immediate value, no supplier interaction required."),
    ("10","SUPPLIER MONITORING AGENTS", "Automated financial + media monitoring on top 50 suppliers. Continuous early warning vs. annual review."),
]

def action_col(slide, actions, x, y_start, w, bg, num_bg):
    y = y_start
    for num, title, desc in actions:
        rect(slide, x, y, w, 1.42, fill=bg)
        rect(slide, x, y, 0.45, 1.42, fill=num_bg)
        txt(slide, num, x+0.04, y+0.47, 0.38, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, title, x+0.55, y+0.1,  w-0.65, 0.32, size=10, bold=True, color=STEEL)
        txt(slide, desc,  x+0.55, y+0.45, w-0.65, 0.9,  size=9.5, color=WHITE)
        y += 1.5

action_col(s, immediate,   0.3,  1.6,  3.5,  DEEP_NAVY, RED)
action_col(s, structural,  4.05, 1.6,  4.15, DEEP_NAVY, AMBER)
action_col(s, pilots,      8.45, 1.6,  4.55, DEEP_NAVY, GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 8 AGENTIC WORKERS
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "AGENTIC AI WORKERS  |  The Procurement Agent Team",
       "8 autonomous AI workers that run continuously — not tools you prompt, agents that act")
footer(s)

insight(s, "Agentic AI is not ChatGPT with a procurement wrapper. These are autonomous workers that monitor, analyse, draft, "
           "and flag — 24/7 — with human approval thresholds governing action. McKinsey: 25–40% efficiency gain lives here.")

agents = [
    ("🔍", "Supplier Monitor",      STEEL,  "Daily news + financial scan on 50+ suppliers. Flags risks before they become crises. Replaces annual reviews with continuous early warning."),
    ("📊", "Market Intelligence",   TEAL,   "Commodity prices, indices, competitor moves, regulatory changes. Daily briefings for category managers. Never miss a market move."),
    ("📋", "RFQ Drafter",           STEEL,  "Given a spec, outputs a complete RFQ with evaluation criteria and supplier shortlist. Saves 4–8 hrs per sourcing event. Used by Ørsted and BP."),
    ("⚖️", "Contract Reviewer",     RED,    "Reads contracts, flags non-standard clauses vs. playbook. 80% faster review (industry benchmark). Immediate value on existing portfolio."),
    ("💰", "Spend Classifier",      GREEN,  "Processes invoices and POs continuously. Keeps spend data clean at 70–90% accuracy. Foundation for every other AI use case."),
    ("🗺️", "Category Strategist",   AMBER,  "Given a brief, researches the market and drafts a category strategy. This is what Manu Forti delivers. The highest-value agent."),
    ("🤝", "Negotiation Prep",      TEAL,   "Builds negotiation pack with market benchmarks, BATNA analysis, and walk-away price. Category managers walk in prepared, not improvising."),
    ("🌿", "Compliance Checker",    MID_GREY, "Monitors supplier ESG, sanctions, and regulatory status for CSRD/EUDR compliance. Automated surveillance vs. periodic manual review."),
]

x, y = 0.3, 2.02
col = 0
for icon, name, accent, desc in agents:
    rect(s, x, y, 3.08, 1.42, fill=DEEP_NAVY)
    rect(s, x, y, 0.38, 1.42, fill=accent)
    txt(s, icon, x+0.04, y+0.44, 0.3, 0.45, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name, x+0.48, y+0.1,  2.5, 0.35, size=11, bold=True, color=accent)
    txt(s, desc, x+0.48, y+0.48, 2.52, 0.9,  size=9, color=WHITE)
    col += 1
    if col == 4:
        x = 0.3; y += 1.5; col = 0
    else:
        x += 3.27

# HITL banner
rect(s, 0.3, 5.08, 12.73, 0.55, fill=NAVY)
txt(s, "⚙️  Human-in-the-Loop Thresholds:  "
       "Draft doc → Always autonomous (logged)  |  "
       "Supplier comms → Auto <€50K, Human above  |  "
       "Create PO → Auto <€10K, Human above  |  "
       "Award contract → Always Human  |  "
       "Share confidential data → Never autonomous",
    0.5, 5.13, 12.4, 0.45, size=9.5, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CYBERSECURITY STACK
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "CYBERSECURITY ARCHITECTURE  |  7-Layer On-Premises Security Stack",
       "Core principle: procurement data must never leave the corporate perimeter")
footer(s)

rect(s, 0.3, 1.15, 12.73, 0.55, fill=RED)
txt(s, "🔒  On-Premises First Policy: All CONFIDENTIAL data processed on-prem. "
       "Cloud LLMs prohibited for: tender pricing, negotiation BATNAs, non-public supplier financials, contract pricing.",
    0.5, 1.2, 12.4, 0.45, size=10, bold=True, color=WHITE)

layers = [
    ("L1", "MODEL DEPLOYMENT",         STEEL,
     "Run open-source LLMs on-premises: Llama 3.3, Mistral 3, Qwen 3.5. "
     "Performance comparable to GPT-4 for most procurement tasks. No data sent to any cloud provider. "
     "Deployment via Ollama (lightweight) or vLLM (production scale)."),
    ("L2", "NETWORK ISOLATION",        TEAL,
     "AI inference servers on dedicated VLAN. Zero outbound internet from inference layer. "
     "External data (market prices, news, supplier info) fetched by a sanitising data layer before reaching AI."),
    ("L3", "DATA CLASSIFICATION & ACCESS CONTROL",  AMBER,
     "Tag all data: PUBLIC / INTERNAL / CONFIDENTIAL / RESTRICTED. "
     "AI agents can only access data at or below their clearance level. "
     "Supplier pricing + contract terms + negotiation positions = RESTRICTED."),
    ("L4", "AGENT SANDBOXING",         GREEN,
     "Each agent runs in an isolated container (Docker/Kubernetes). No direct agent-to-agent comms — "
     "all via audited message broker. Read-only access to source systems. No agent can write to ERP without HITL approval."),
    ("L5", "PROMPT INJECTION DEFENCE", RED,
     "No. 1 attack vector: malicious supplier embeds instructions in a document (invoice, contract, email). "
     "Mitigations: input sanitisation pipeline, output validation layer, separate parsing model."),
    ("L6", "AUDIT & MONITORING",       MID_GREY,
     "Every AI action logged: timestamp, agent ID, input hash, output, reviewer. "
     "Immutable append-only signed audit log. Real-time anomaly detection. Required for CSRD/NIS2 compliance."),
    ("L7", "HUMAN-IN-THE-LOOP THRESHOLDS", STEEL,
     "Defined approval gates by action type and value. Auto-execute below threshold, human approval above. "
     "Award contract = always human. Share confidential data = never autonomous."),
]

y = 1.82
for code, title, color, desc in layers:
    rect(s, 0.3,  y, 0.65, 0.62, fill=color)
    txt(s, code, 0.32, y+0.15, 0.62, 0.35, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.98, y, 12.05, 0.62, fill=DEEP_NAVY)
    txt(s, title+":", 1.1, y+0.05, 3.0, 0.3, size=10, bold=True, color=color)
    txt(s, desc,  4.05, y+0.06, 8.85, 0.52, size=9.5, color=WHITE)
    y += 0.68

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "TECHNOLOGY STACK  |  Agentic Procurement Architecture",
       "Three deployment tiers — matched to data sensitivity and use case")
footer(s)

# Tier cards
tiers = [
    ("TIER 1\nOn-Premises", RED,
     "For RESTRICTED & CONFIDENTIAL data",
     ["LLM: Llama 3.3 / Mistral 3 / Qwen 3.5 (self-hosted)",
      "Inference: Ollama (dev) or vLLM (production)",
      "Orchestration: OpenClaw / LangChain / AutoGen",
      "Vector DB: ChromaDB or Qdrant (on-prem)",
      "Storage: SAP ERP / SharePoint (read-only API)",
      "Auth: Active Directory / SSO",
      "Monitoring: ELK Stack / Grafana",
      "Containers: Docker + Kubernetes"],
     ["Supplier Monitor", "Contract Reviewer", "Spend Classifier", "Negotiation Prep"]),
    ("TIER 2\nPrivate Cloud", AMBER,
     "For INTERNAL data — controlled VPC",
     ["LLM: Azure OpenAI (EU data residency) or AWS Bedrock",
      "Orchestration: OpenClaw agents / AutoGen",
      "Data: Statkraft-controlled Azure tenant",
      "No training on Statkraft data",
      "Data residency: EU (GDPR compliant)",
      "VNET integration with on-prem",
      "Azure Monitor for audit logging",
      "Entra ID for access control"],
     ["RFQ Drafter", "Market Intelligence", "Category Strategist"]),
    ("TIER 3\nCloud Tools", GREEN,
     "For PUBLIC data only — no corporate data",
     ["Tools: Claude, ChatGPT, Gemini, Perplexity",
      "Use case: market research, benchmarking",
      "Rule: NO corporate data input. Ever.",
      "Individual track with manager approval",
      "Data handling agreement required",
      "Output review before use in Statkraft docs",
      "Usage logged centrally",
      "Annual policy review"],
     ["Market Research", "Benchmarking", "Draft ideas", "Learning"]),
]

x = 0.3
for tier_title, color, sub, stack, use_cases in tiers:
    rect(s, x, 1.15, 4.15, 5.55, fill=DEEP_NAVY)
    rect(s, x, 1.15, 4.15, 0.55, fill=color)
    txt(s, tier_title, x+0.15, 1.18, 3.85, 0.5, size=14, bold=True, color=WHITE)
    txt(s, sub, x+0.15, 1.75, 3.85, 0.3, size=9.5, color=color, italic=True)
    txt(s, "Technology:", x+0.15, 2.12, 3.85, 0.28, size=10, bold=True, color=MID_GREY)
    y = 2.4
    for item in stack:
        txt(s, f"· {item}", x+0.15, y, 3.85, 0.38, size=8.5, color=WHITE)
        y += 0.36
    txt(s, "Key Agents:", x+0.15, y+0.05, 3.85, 0.28, size=10, bold=True, color=color)
    uc_txt = "  ·  ".join(use_cases)
    txt(s, uc_txt, x+0.15, y+0.33, 3.85, 0.4, size=9, color=WHITE)
    x += 4.37

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — USE CASE LIBRARY
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "USE CASE LIBRARY  |  Top 12 Ranked by Impact vs. Effort",
       "Start here — not a blank sheet. Ranked by value delivered and implementation speed.")
footer(s)

use_cases_lib = [
    ("SPEND CLASSIFICATION",         "High",   "Low",   "Fast Win",  "70–90% accuracy, cleans data for everything else. Foundation use case."),
    ("CONTRACT NLP SCANNING",        "High",   "Low",   "Fast Win",  "Flag non-standard clauses across portfolio. 80% faster review. No supplier risk."),
    ("SUPPLIER RISK MONITORING",     "High",   "Medium","Priority",  "Daily monitoring of top 50 suppliers. Continuous early warning system."),
    ("RFQ GENERATION",               "High",   "Medium","Priority",  "Auto-draft RFQs from spec docs. Saves 4–8 hrs per event. Proven at Ørsted/BP."),
    ("DEMAND FORECASTING (MRO)",     "High",   "Medium","Priority",  "6–12 month horizon for capex/MRO. Proven at Vattenfall. Reduces emergency buys."),
    ("NEGOTIATION PACK BUILDER",     "High",   "Medium","Priority",  "Market benchmarks, BATNA analysis, walk-away price. Category managers walk in prepared."),
    ("CATEGORY STRATEGY DRAFTING",   "V.High", "High",  "Strategic", "Given a brief, outputs full strategy. Manu Forti methodology. Highest value agent."),
    ("ESG / CSRD COMPLIANCE CHECK",  "High",   "Medium","Compliance","Monitor supplier ESG and sanctions status continuously. Audit-ready output."),
    ("INVOICE PROCESSING",           "Medium", "Low",   "Fast Win",  "3-way matching + anomaly detection. Reduces manual processing, catches errors."),
    ("MARKET PRICE INTELLIGENCE",    "Medium", "Low",   "Fast Win",  "Daily briefings on commodity prices and indices. Category managers always current."),
    ("SUPPLIER ONBOARDING CHECKS",   "Medium", "Medium","Structured","Automate due diligence checks for new suppliers. Faster onboarding, lower risk."),
    ("CONTRACT DRAFTING ASSIST",     "Medium", "High",  "Structured","AI-assisted first draft of standard contracts. Reduces legal cycle time by 40–60%."),
]

rect(s, 0.3, 1.15, 12.73, 0.38, fill=NAVY)
for col, label, w, x in [("USE CASE","",4.0,0.3),("IMPACT","",1.0,4.35),("EFFORT","",1.0,5.4),
                           ("TYPE","",1.1,6.45),("VALUE STATEMENT","",6.45,7.6)]:
    txt(s, label, x+0.1, 1.18, w-0.1, 0.28, size=9, bold=True, color=WHITE)

y = 1.58
for i, (name, impact, effort, typ, value) in enumerate(use_cases_lib):
    bg = RGBColor(0xF5,0xF9,0xFF) if i % 2 == 0 else WHITE
    rect(s, 0.3, y, 12.73, 0.43, fill=bg if i >= 0 else bg, line=LT_BLUE, lw=0.3)
    i_color = GREEN if impact=="High" else (STEEL if impact=="V.High" else AMBER)
    e_color = GREEN if effort=="Low" else (AMBER if effort=="Medium" else RED)
    t_color = STEEL if typ=="Fast Win" else (AMBER if typ=="Priority" else (TEAL if typ=="Strategic" else MID_GREY))
    txt(s, name, 0.4, y+0.07, 3.9, 0.3, size=9.5, bold=True, color=DARK)
    txt(s, impact, 4.45, y+0.07, 0.9, 0.3, size=9.5, bold=True, color=i_color, align=PP_ALIGN.CENTER)
    txt(s, effort, 5.5,  y+0.07, 0.9, 0.3, size=9.5, bold=True, color=e_color, align=PP_ALIGN.CENTER)
    txt(s, typ,    6.55, y+0.07, 1.0, 0.3, size=9.5, color=t_color, align=PP_ALIGN.CENTER)
    txt(s, value,  7.7,  y+0.07, 5.2, 0.3, size=9.5, color=DARK)
    y += 0.44

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GOVERNANCE & TRACK D
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "GOVERNANCE FRAMEWORK  |  Four-Track Model with Track D Added",
       "Revised governance model: from IT-heavy to procurement-led, with agentic AI tier")
footer(s)

tracks = [
    ("A", "INDIVIDUAL\nASSISTIVE", STEEL,
     ["Copilot, Claude, ChatGPT, Gemini", "Manager approval + data handling rules",
      "PUBLIC/INTERNAL data only", "Usage logged centrally", "Annual policy review"]),
    ("B", "LOCAL\nSHARED USE", TEAL,
     ["Team-level tools with shared prompt libraries", "Data governance rules mandatory",
      "No CONFIDENTIAL data without IT review", "Category manager leads deployment",
      "Quarterly review against declared KPIs"]),
    ("C", "ENTERPRISE\nPLATFORM", AMBER,
     ["Coupa AI, SAP Ariba, approved platforms", "Full IT + Procurement governance board",
      "RESTRICTED data permitted with controls", "6-monthly review vs. business case",
      "Change management mandatory"]),
    ("D\n⭐NEW", "AGENTIC AI\n(AUTONOMOUS)", RED,
     ["Autonomous agents with continuous operation", "Dedicated governance tier — not Track C",
      "HITL thresholds defined per action type", "Cybersecurity assessment mandatory",
      "On-prem deployment for CONFIDENTIAL data",
      "Real-time audit log of all agent actions",
      "CPO sign-off required for production"]),
]

x = 0.3
for track, title, color, items in tracks:
    w = 3.08 if track != "D\n⭐NEW" else 3.2
    rect(s, x, 1.15, w, 5.55, fill=DEEP_NAVY)
    rect(s, x, 1.15, w, 0.7,  fill=color)
    txt(s, track, x+0.1, 1.2, 0.55, 0.6, size=18, bold=True, color=WHITE)
    txt(s, title, x+0.65, 1.22, w-0.8, 0.62, size=12, bold=True, color=WHITE)
    y = 2.0
    for item in items:
        txt(s, f"▸  {item}", x+0.15, y, w-0.25, 0.5, size=9.5, color=WHITE)
        y += 0.55
    x += w + 0.18

# 3 additions footer
rect(s, 0.3, 6.78, 12.73, 0.5, fill=NAVY)
txt(s, "3 Mandatory Additions:  "
       "A) Cybersecurity assessment at every path selection review  |  "
       "B) On-premises first policy for CONFIDENTIAL data — no cloud LLMs  |  "
       "C) Banned inputs: active tender pricing, negotiation BATNAs, non-public supplier financials",
    0.5, 6.82, 12.4, 0.42, size=9.5, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — EXECUTION ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
header(s, "EXECUTION ROADMAP  |  90-Day Sprint + 12-Month Programme",
       "Three pilots in 90 days. Full agentic capability in 12 months. Measured at every gate.")
footer(s)

phases = [
    ("PHASE 0\nWeeks 1–2\nFoundation", STEEL, [
        "Add use case library (top 12)",
        "Mandate value measurement template",
        "Define data classification rules",
        "Draft Track D governance framework",
        "Identify cybersecurity requirements",
    ]),
    ("PHASE 1\nWeeks 3–12\nThree Pilots", GREEN, [
        "Pilot 1: Spend classification cleanup",
        "Pilot 2: Contract NLP scanning",
        "Pilot 3: Supplier monitoring agent (top 50)",
        "Baseline metrics captured for each",
        "90-day review gate — scale or kill",
    ]),
    ("PHASE 2\nMonths 4–6\nExpand", AMBER, [
        "Deploy RFQ drafter (target: 4–8 hrs saved/event)",
        "Launch demand forecasting (MRO/capex)",
        "Track D governance live",
        "On-prem LLM stack deployed",
        "AI Forum becomes product team w/ backlog",
    ]),
    ("PHASE 3\nMonths 7–12\nAgentic Scale", RED, [
        "All 8 agentic workers operational",
        "Category strategy agent live",
        "Negotiation pack builder deployed",
        "CSRD compliance agent live",
        "Full 7-layer security stack certified",
        "Nordic benchmark leadership",
    ]),
]

x = 0.3
for phase, color, items in phases:
    rect(s, x, 1.15, 3.06, 5.55, fill=DEEP_NAVY)
    rect(s, x, 1.15, 3.06, 0.75, fill=color)
    txt(s, phase, x+0.12, 1.18, 2.85, 0.7, size=11, bold=True, color=WHITE)
    y = 2.0
    for item in items:
        txt(s, f"✓  {item}", x+0.12, y, 2.85, 0.48, size=9.5, color=WHITE)
        y += 0.52
    x += 3.24

# Value target row
rect(s, 0.3, 6.78, 12.73, 0.5, fill=NAVY)
targets = [
    "Phase 0: Framework improved",
    "Phase 1: 3 pilots live, ROI baseline set",
    "Phase 2: 25%+ time saving on sourcing events",
    "Phase 3: 40% efficiency gain (McKinsey benchmark)",
    "Year 2: Statkraft is the Nordic standard",
]
x_t = 0.45
for t in targets:
    txt(s, f"📌 {t}", x_t, 6.82, 2.4, 0.42, size=8.5, color=WHITE)
    x_t += 2.56

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSING / CALL TO ACTION
# ═══════════════════════════════════════════════════════════════════════════════
s = blank()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 0, 5.5, 13.333, 2.0, fill=STEEL)

txt(s, "THE BOTTOM LINE", 0.6, 0.8, 12, 0.5, size=14, color=STEEL, bold=True)
txt(s, "Statkraft can lead in Procurement AI.\nThe gap is open. The tools are ready. The case is proven.", 0.6, 1.3, 12, 1.4, size=30, bold=True, color=WHITE)
line(s, 2.85, color=STEEL)

calls = [
    ("This week",   "Add use case library + value measurement template to the framework"),
    ("This month",  "Launch 3 pilots: spend classification, contract NLP, supplier monitoring"),
    ("This quarter","Add Track D governance, deploy on-prem LLM stack, upgrade individual track"),
    ("This year",   "All 8 agents operational, 40% efficiency gain, Statkraft sets the Nordic standard"),
]
y = 3.0
for when, action in calls:
    rect(s, 0.6, y, 1.8, 0.42, fill=STEEL)
    txt(s, when,   0.65, y+0.08, 1.7, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, f"→  {action}", 2.55, y+0.08, 10.0, 0.3, size=12, color=WHITE)
    y += 0.52

txt(s, "Prepared by Aiden  |  Manu Forti Intelligence  |  For Statkraft Procurement Leadership  |  March 17, 2026",
    0.6, 5.65, 12, 0.4, size=13, bold=True, color=WHITE)
txt(s, "Contact: manuforti.as@gmail.com  |  www.manuforti.no", 0.6, 6.1, 12, 0.35, size=12, color=WHITE, italic=True)
txt(s, "Sources: Hackett Group (2026), McKinsey, EY CPO Survey (2025), Gartner, Equinor.com, Art of Procurement, Spend Matters",
    0.6, 6.55, 12, 0.35, size=9, color=MID_GREY, italic=True)

# ── SAVE ───────────────────────────────────────────────────────────────────────
out = "/Users/jonathonmilne/.openclaw/workspace/venture/Statkraft_Procurement_AI_Strategy_March2026.pptx"
prs.save(out)
print(f"Saved: {out}")
