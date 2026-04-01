#!/usr/bin/env python3
"""
Product 1 v15 BULLETPROOF Generator for Prospex Energy
Creates PPTX with EXACT canonical positioning from template.
"""

import json
import sys
from pathlib import Path
from datetime import datetime

# Add common site-packages paths
for path in [
    '/opt/homebrew/lib/python3.14/site-packages',
    '/opt/homebrew/lib/python3.13/site-packages',
    '/usr/local/lib/python3.11/site-packages',
    str(Path.home() / '.local/lib/python3.11/site-packages'),
    str(Path.home() / 'Library/Python/3.11/lib/python/site-packages'),
]:
    if Path(path).exists():
        sys.path.insert(0, path)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# CANONICAL DIMENSIONS - HARD-CODED FROM v15 TEMPLATE
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

RISK_GAUGE_LEFT = Inches(0.3)
RISK_GAUGE_TOP = Inches(1.05)
RISK_GAUGE_WIDTH = Inches(7.2)
RISK_GAUGE_HEIGHT = Inches(3.1)

SNAPSHOT_LEFT = Inches(7.6)
SNAPSHOT_TOP = Inches(1.05)
SNAPSHOT_WIDTH = Inches(5.4)
SNAPSHOT_HEIGHT = Inches(5.6)

MANU_FORTI_LEFT = Inches(12.533)
MANU_FORTI_TOP = Inches(6.613)
MANU_FORTI_WIDTH = Inches(0.6)
MANU_FORTI_HEIGHT = Inches(0.737)

SOURCE_LEFT = Inches(0.3)
SOURCE_TOP = Inches(7.15)
SOURCE_WIDTH = Inches(12.8)
SOURCE_HEIGHT = Inches(0.3)

LOGO_BOX_LEFT = Inches(10.92)
LOGO_BOX_TOP = Inches(0.06)
LOGO_BOX_WIDTH = Inches(2.26)
LOGO_BOX_HEIGHT = Inches(0.87)

NAVY = RGBColor(0x00, 0x21, 0x47)
STEEL_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
MID_GREY = RGBColor(0x71, 0x80, 0x96)
GREEN = RGBColor(0x48, 0xBB, 0x78)
AMBER = RGBColor(0xD6, 0x9E, 0x2E)
RED = RGBColor(0xE5, 0x3E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_source_line(slide, month_year="March 2026"):
    textbox = slide.shapes.add_textbox(SOURCE_LEFT, SOURCE_TOP, SOURCE_WIDTH, SOURCE_HEIGHT)
    tf = textbox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"Source: Manu Forti Intelligence  |  Confidential  |  {month_year}"
    p.font.size = Pt(8)
    p.font.color.rgb = MID_GREY
    return textbox

def add_manu_forti_logo(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MANU_FORTI_LEFT, MANU_FORTI_TOP,
        MANU_FORTI_WIDTH, MANU_FORTI_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "MF"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_header(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
        tf = sub_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = MID_GREY

def add_supplier_logo_box(slide):
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LOGO_BOX_LEFT, LOGO_BOX_TOP,
        LOGO_BOX_WIDTH, LOGO_BOX_HEIGHT
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = MID_GREY
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PXEN"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    return box

def get_risk_color(rating):
    rating = rating.upper()
    if rating == "LOW":
        return GREEN
    elif rating == "MEDIUM":
        return AMBER
    else:
        return RED

def create_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SUPPLIER EVALUATION REPORT"
    p.font.size = Pt(20)
    p.font.color.rgb = MID_GREY
    
    name_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(12), Inches(1.2))
    tf = name_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['supplier_legal_name']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    sector_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(0.6))
    tf = sector_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Natural Gas  |  Investment Company  |  European Energy"
    p.font.size = Pt(16)
    p.font.color.rgb = STEEL_BLUE
    
    stats_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(12), Inches(0.4))
    tf = stats_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"£27.8M Market Cap  |  AIM-Listed  |  25 Years  |  Italy/Spain/Poland"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GREY
    
    conf_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(12), Inches(0.4))
    tf = conf_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Confidential  |  {data['report_date']}  |  Manu Forti Intelligence"
    p.font.size = Pt(12)
    p.font.color.rgb = MID_GREY
    
    add_manu_forti_logo(slide)
    return slide

def create_exec_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "EXECUTIVE SUMMARY", 
               f"{data['supplier_legal_name']} — Risk & Suitability Overview")
    add_supplier_logo_box(slide)
    
    gauge_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        RISK_GAUGE_LEFT, RISK_GAUGE_TOP,
        RISK_GAUGE_WIDTH, RISK_GAUGE_HEIGHT
    )
    gauge_bg.fill.solid()
    gauge_bg.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    gauge_bg.line.color.rgb = MID_GREY
    
    score = data['overall_risk_score']
    rating = data['overall_risk_rating']
    risk_color = get_risk_color(rating)
    
    score_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(3), Inches(1.2))
    tf = score_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{score}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = risk_color
    p.alignment = PP_ALIGN.CENTER
    
    rating_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(3), Inches(0.6))
    tf = rating_box.text_frame
    p = tf.paragraphs[0]
    p.text = rating
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = risk_color
    p.alignment = PP_ALIGN.CENTER
    
    label_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.0), Inches(3), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "OVERALL RISK"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GREY
    p.alignment = PP_ALIGN.CENTER
    
    findings_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(4.3),
        Inches(7.2), Inches(2.0)
    )
    findings_bg.fill.solid()
    findings_bg.fill.fore_color.rgb = RGBColor(0xEB, 0xF4, 0xFF)
    findings_bg.line.fill.background()
    
    findings_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(7), Inches(0.3))
    tf = findings_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Findings"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = STEEL_BLUE
    
    findings_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.75), Inches(7), Inches(1.5))
    tf = findings_text.text_frame
    tf.word_wrap = True
    findings = [
        f"• Micro-cap investment entity: £27.8M market cap, limited liquidity",
        f"• Diversified European gas portfolio: Italy (37%), Spain (100%), Poland",
        f"• No operating revenue — dependent on investee company distributions",
        f"• HIGH risk rating (72/100) due to investment structure and micro-cap status"
    ]
    for i, finding in enumerate(findings):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = finding
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
    
    snap_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SNAPSHOT_LEFT, SNAPSHOT_TOP,
        SNAPSHOT_WIDTH, SNAPSHOT_HEIGHT
    )
    snap_bg.fill.solid()
    snap_bg.fill.fore_color.rgb = WHITE
    snap_bg.line.color.rgb = MID_GREY
    
    snap_title = slide.shapes.add_textbox(Inches(7.8), Inches(1.22), Inches(5.2), Inches(0.32))
    tf = snap_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SUPPLIER SNAPSHOT"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    snapshot_data = [
        ("Supplier:", data['supplier_legal_name']),
        ("Ticker:", data['ticker']),
        ("Sector:", "Natural Gas / Investment"),
        ("HQ:", data['operational_capability']['headquarters']),
        ("Founded:", "2000 (25 years)"),
        ("Market Cap:", data['financial_health']['market_cap']),
        ("Structure:", "Investment Company (IFRS10)"),
        ("Focus:", "European Natural Gas"),
    ]
    
    y_pos = 1.64
    for label, value in snapshot_data:
        lbl_box = slide.shapes.add_textbox(Inches(7.8), Inches(y_pos), Inches(2.3), Inches(0.27))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = MID_GREY
        
        val_box = slide.shapes.add_textbox(Inches(10.1), Inches(y_pos), Inches(2.7), Inches(0.27))
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        y_pos += 0.3
    
    add_source_line(slide, data['report_date'])
    add_manu_forti_logo(slide)
    return slide

def create_recommendation_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "RECOMMENDATION", "Decision Summary & Commercial Conditions")
    add_supplier_logo_box(slide)
    
    banner_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(12.33), Inches(1.35)
    )
    banner_bg.fill.solid()
    banner_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xE6)
    banner_bg.line.color.rgb = AMBER
    
    verdict_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.18), Inches(12), Inches(0.55))
    tf = verdict_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"⚠  RECOMMENDATION: {data['recommendation']}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = AMBER
    
    desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.78), Inches(12), Inches(0.58))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{data['supplier_legal_name']} is recommended for natural gas sector engagement with strict conditions. Micro-cap investment structure and limited liquidity require enhanced due diligence and monitoring."
    p.font.size = Pt(12)
    p.font.color.rgb = NAVY
    
    comm_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.65), Inches(6), Inches(0.3))
    tf = comm_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Commercial Conditions"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    comm_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(6), Inches(2.2))
    tf = comm_box.text_frame
    tf.word_wrap = True
    for i, cond in enumerate(data['conditions'][:5]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {cond}"
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
    
    risk_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(5.35),
        Inches(12.33), Inches(1.5)
    )
    risk_bg.fill.solid()
    risk_bg.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    risk_bg.line.color.rgb = MID_GREY
    
    risk_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(0.3))
    tf = risk_title.text_frame
    p = tf.paragraphs[0]
    p.text = f"Overall Risk Summary — {data['overall_risk_rating']} ({data['overall_risk_score']}/100)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    risks = [
        ("Financial:", data['financial_health']['rating'], data['financial_health']['score']),
        ("Operational:", data['operational_capability']['rating'], data['operational_capability']['score']),
        ("Geopolitical:", data['geopolitical_risk']['rating'], data['geopolitical_risk']['score']),
        ("ESG:", data['esg_rating']['rating'], data['esg_rating']['score']),
    ]
    
    x_pos = 0.6
    for label, rating, score in risks:
        lbl_box = slide.shapes.add_textbox(Inches(x_pos), Inches(5.88), Inches(3.1), Inches(0.27))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        rating_color = get_risk_color(rating)
        rat_box = slide.shapes.add_textbox(Inches(x_pos), Inches(6.17), Inches(3.1), Inches(0.27))
        tf = rat_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{rating} ({score})"
        p.font.size = Pt(12)
        p.font.color.rgb = rating_color
        
        x_pos += 3.1
    
    add_source_line(slide, data['report_date'])
    add_manu_forti_logo(slide)
    return slide

def create_profile_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "SUPPLIER PROFILE", "Corporate Structure & Global Footprint")
    add_supplier_logo_box(slide)
    
    overview_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.05),
        Inches(6.5), Inches(3.5)
    )
    overview_bg.fill.solid()
    overview_bg.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    overview_bg.line.fill.background()
    
    overview_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(6.1), Inches(0.3))
    tf = overview_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Company Overview"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = STEEL_BLUE
    
    overview_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.1), Inches(3.0))
    tf = overview_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{data['supplier_legal_name']} is an AIM-listed investment company holding working interests in European natural gas projects. Founded in 2000, the company operates a low-capex, late-stage exploration and production strategy focused on the energy transition."
    p.font.size = Pt(12)
    p.font.color.rgb = NAVY
    
    p2 = tf.add_paragraph()
    p2.text = ""
    
    p3 = tf.add_paragraph()
    p3.text = "Unlike traditional E&P companies, Prospex does not operate assets directly. Instead, it holds passive working interests and relies on partner operators for production, creating a unique risk profile that requires careful evaluation."
    p3.font.size = Pt(12)
    p3.font.color.rgb = NAVY
    
    inv_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.0), Inches(1.05),
        Inches(6.0), Inches(3.5)
    )
    inv_bg.fill.solid()
    inv_bg.fill.fore_color.rgb = WHITE
    inv_bg.line.color.rgb = MID_GREY
    
    inv_title = slide.shapes.add_textbox(Inches(7.2), Inches(1.15), Inches(5.6), Inches(0.3))
    tf = inv_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Investments"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = STEEL_BLUE
    
    inv_text = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.6), Inches(3.0))
    tf = inv_text.text_frame
    tf.word_wrap = True
    
    investments = data['operational_capability']['core_investments']
    for i, inv in enumerate(investments):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {inv}"
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
    
    strengths_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(4.7),
        Inches(6.5), Inches(2.3)
    )
    strengths_bg.fill.solid()
    strengths_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF4)
    strengths_bg.line.fill.background()
    
    strengths_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(6.1), Inches(0.3))
    tf = strengths_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Strengths"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    strengths_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(6.1), Inches(1.8))
    tf = strengths_text.text_frame
    tf.word_wrap = True
    for i, strength in enumerate(data['key_strengths'][:4]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {strength}"
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
    
    risks_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.0), Inches(4.7),
        Inches(6.0), Inches(2.3)
    )
    risks_bg.fill.solid()
    risks_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
    risks_bg.line.fill.background()
    
    risks_title = slide.shapes.add_textbox(Inches(7.2), Inches(4.8), Inches(5.6), Inches(0.3))
    tf = risks_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Risks"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED
    
    risks_text = slide.shapes.add_textbox(Inches(7.2), Inches(5.15), Inches(5.6), Inches(1.8))
    tf = risks_text.text_frame
    tf.word_wrap = True
    for i, risk in enumerate(data['key_risks'][:4]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {risk}"
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
    
    add_source_line(slide, data['report_date'])
    add_manu_forti_logo(slide)
    return slide

def create_financial_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fin = data['financial_health']
    
    add_header(slide, "FINANCIAL HEALTH", 
               f"Investment Company Structure  |  Market Cap: {fin['market_cap']}  |  Minimal Debt")
    add_supplier_logo_box(slide)
    
    chart_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.05),
        Inches(7.5), Inches(4.55)
    )
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    chart_bg.line.color.rgb = MID_GREY
    
    chart_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(7.1), Inches(0.4))
    tf = chart_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Investment Entity Financial Profile"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    struct_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(7.1), Inches(3.5))
    tf = struct_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "IFRS 10 Investment Entity Structure"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = STEEL_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = ""
    
    p3 = tf.add_paragraph()
    p3.text = "Prospex Energy PLC operates as an investment company under IFRS 10, meaning it holds controlling interests in subsidiaries rather than consolidating operational revenue. This creates a unique financial profile:"
    p3.font.size = Pt(11)
    p3.font.color.rgb = NAVY
    
    p4 = tf.add_paragraph()
    p4.text = ""
    
    bullets = [
        "• Revenue not consolidated at parent level",
        "• Income derived from investee distributions",
        "• Asset value fluctuations drive NAV changes",
        "• Lower operational overhead vs. E&P operators"
    ]
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
    
    highlights_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.1), Inches(1.05),
        Inches(4.9), Inches(5.8)
    )
    highlights_bg.fill.solid()
    highlights_bg.fill.fore_color.rgb = WHITE
    highlights_bg.line.color.rgb = MID_GREY
    
    hl_title = slide.shapes.add_textbox(Inches(8.25), Inches(1.15), Inches(4.6), Inches(0.35))
    tf = hl_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Highlights"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = STEEL_BLUE
    
    metrics = [
        ("Market Cap:", fin['market_cap']),
        ("Shares Outstanding:", fin['shares_outstanding']),
        ("Stock Price:", fin['stock_price']),
        ("2024 Net Loss:", fin['net_loss_2024']),
        ("Ownership:", fin['ownership']),
        ("Liquidity:", fin['liquidity']),
        ("Debt:", fin['debt']),
    ]
    
    y_pos = 1.6
    for label, value in metrics:
        lbl_box = slide.shapes.add_textbox(Inches(8.25), Inches(y_pos), Inches(2.4), Inches(0.28))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = MID_GREY
        
        val_box = slide.shapes.add_textbox(Inches(10.65), Inches(y_pos), Inches(2.2), Inches(0.28))
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        y_pos += 0.35
    
    risk_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.1), Inches(4.05),
        Inches(4.8), Inches(0.44)
    )
    risk_bg.fill.solid()
    risk_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
    risk_bg.line.color.rgb = RED
    
    risk_text = slide.shapes.add_textbox(Inches(8.25), Inches(4.12), Inches(4.5), Inches(0.3))
    tf = risk_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Risk: HIGH ⚠"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED
    
    exposure_text = slide.shapes.add_textbox(Inches(8.1), Inches(4.6), Inches(4.8), Inches(2.3))
    tf = exposure_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Exposure Guidance"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p2 = tf.add_paragraph()
    p2.text = ""
    
    guidance = [
        "• Micro-cap: Limited liquidity, wide spreads",
        "• Investment entity: No operating cash flow",
        "• Enhanced due diligence required",
        "• Quarterly NAV monitoring recommended"
    ]
    for g in guidance:
        p = tf.add_paragraph()
        p.text = g
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
    
    add_source_line(slide, data['report_date'])
    add_manu_forti_logo(slide)
    return slide

def create_market_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_header(slide, "MARKET POSITION", "AIM-Listed Natural Gas Investment Companies")
    add_supplier_logo_box(slide)
    
    chart_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.05),
        Inches(8.0), Inches(4.46)
    )
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    chart_bg.line.color.rgb = MID_GREY
    
    chart_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(7.6), Inches(0.4))
    tf = chart_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Peer Comparison — European Gas Investment Companies"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    peers_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(7.6), Inches(3.5))
    tf = peers_text.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]