#!/usr/bin/env python3
"""
Nel ASA Product 1 — Rev X (Stress-Test Quality vs Boskalis v15)
"""

import sys
sys.path.insert(0, '/opt/homebrew/lib/python3.14/site-packages')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Circle
import numpy as np
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors
NAVY = RGBColor(0, 33, 71)
STEEL = RGBColor(43, 108, 176)
MID_GREY = RGBColor(113, 128, 150)
GREEN = RGBColor(72, 187, 120)
AMBER = RGBColor(214, 158, 46)
RED = RGBColor(229, 62, 62)
WHITE = RGBColor(255, 255, 255)
TEXT_PRIMARY = RGBColor(26, 32, 44)
TEXT_SECONDARY = RGBColor(74, 85, 104)
LIGHT_GREY = RGBColor(200, 210, 220)
LIGHT_BG = RGBColor(240, 245, 250)

M_NAVY = '#002147'
M_STEEL = '#2B6CB0'
M_GREY = '#718096'
M_GREEN = '#48BB78'
M_AMBER = '#D69E2E'
M_RED = '#E53E3E'
M_LIGHT = '#EBF4FF'

BASE = Path('/Users/jonathonmilne/.openclaw/workspace')
VISUALS = BASE / 'nel_revX_visuals'
VISUALS.mkdir(exist_ok=True)
LOGO = BASE / 'skills/product-1-generator/assets/manu_forti_logo.png'
OUTPUT = BASE / 'Nel_ASA_Product1_RevX_Final.pptx'

NEL_DATA = {
    'name': 'Nel ASA',
    'ticker': 'OSE: NEL',
    'sector': 'Hydrogen Electrolysers',
    'subsector': 'Alkaline & PEM Technology',
    'hq': 'Oslo, Norway',
    'founded': 1927,
    'employees': '~550',
    'ceo': 'Håkon Volldal',
    'ceo_since': 'July 2022',
    'revenue_2024': 'NOK 1,600M',
    'revenue_2023': 'NOK 1,253M',
    'revenue_growth': '+28%',
    'ebitda_2024': 'NOK -275M',
    'ebitda_margin': '-17.2%',
    'order_backlog': 'NOK 1,614M',
    'cash_balance': 'NOK 1,876M',
    'total_equity': 'NOK 4,800M',
    'total_debt': 'NOK 0M',
    'debt_equity': '0%',
    'capacity_heroya': '1,000 MW',
    'capacity_wallingford': '500 MW',
    'capacity_planned_us': '4,000 MW',
    'risk_scores': {'financial': 58, 'operational': 28, 'geopolitical': 15, 'esg': 22},
    'overall_risk': 48,
    'verdict': 'CONDITIONAL APPROVAL',
}

def add_header(slide, title, subtitle):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    h.fill.solid(); h.fill.fore_color.rgb = NAVY; h.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE
    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = subtitle; r2.font.size = Pt(12); r2.font.color.rgb = LIGHT_GREY

def add_source(slide, date='March 2026'):
    t = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = f'Source: Manu Forti Intelligence  |  Confidential  |  {date}'
    r.font.size = Pt(8); r.font.color.rgb = MID_GREY

def add_text(slide, l, t, w, h, text, sz=14, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return tb

def add_bullets(slide, l, t, w, h, title, items, tsz=14, isz=11, tc=NAVY, ic=TEXT_PRIMARY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True; r.font.color.rgb = tc
    for item in items:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = f'• {item}'; r.font.size = Pt(isz); r.font.color.rgb = ic

def add_rounded_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_logo(slide):
    try:
        cimg = Image.open(str(LOGO))
        ca = cimg.size[0] / cimg.size[1]
        cw = 0.6; ch = cw / ca
        cx = 13.333 - cw - 0.2
        cy = 7.5 - ch - 0.15
        slide.shapes.add_picture(str(LOGO), Inches(cx), Inches(cy), Inches(cw), Inches(ch))
    except:
        pass

def gen_risk_gauge():
    fig, ax = plt.subplots(figsize=(14, 6.5), facecolor='white')
    segments = [(0, 60, M_GREEN, 'LOW\n0-33'), (60, 120, M_AMBER, 'MEDIUM\n34-66'), (120, 180, M_RED, 'HIGH\n67-100')]
    for start, end, color, label in segments:
        theta = np.linspace(np.radians(start), np.radians(end), 100)
        x_o = np.cos(theta); y_o = np.sin(theta)
        x_i = np.cos(theta[::-1]) * 0.6; y_i = np.sin(theta[::-1]) * 0.6
        ax.fill(np.concatenate([x_o, x_i]), np.concatenate([y_o, y_i]), color=color, alpha=0.85)
        mid = np.radians((start+end)/2)
        ax.text(np.cos(mid)*0.8, np.sin(mid)*0.8, label, ha='center', va='center', fontsize=11, fontweight='bold', color='white')
    score = NEL_DATA['overall_risk']
    needle = np.radians(score * 180 / 100)
    ax.annotate('', xy=(np.cos(needle)*0.55, np.sin(needle)*0.55), xytext=(0,0), arrowprops=dict(arrowstyle='->', color=M_NAVY, lw=3.5))
    ax.add_patch(Circle((0,0), 0.08, color=M_NAVY, zorder=5))
    ax.text(0, -0.2, f'{score}/100', ha='center', fontsize=28, fontweight='bold', color=M_NAVY)
    ax.text(0, -0.35, 'OVERALL RISK SCORE', ha='center', fontsize=12, color=M_GREY)
    ax.text(0, 1.25, f'{NEL_DATA["name"]} - Overall Risk Assessment', ha='center', fontsize=18, fontweight='bold', color=M_NAVY)
    ax.text(0, 1.10, f'{NEL_DATA["subsector"]} | {NEL_DATA["capacity_heroya"]} Manufacturing Capacity', ha='center', fontsize=12, color=M_GREY)
    boxes = [('Financial', 'MEDIUM', M_AMBER, 58), ('Operational', 'LOW', M_GREEN, 28), ('Geopolitical', 'LOW', M_GREEN, 15), ('ESG', 'LOW', M_GREEN, 22)]
    for idx, (cat, rating, color, score_val) in enumerate(boxes):
        bx = -0.75 + idx * 0.5
        rect = FancyBboxPatch((bx-0.2, -0.67), 0.4, 0.22, boxstyle="round,pad=0.03", facecolor=color, alpha=0.15, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(bx, -0.54, cat, ha='center', fontsize=8, fontweight='bold', color=M_NAVY)
        ax.text(bx, -0.62, f'{rating} ({score_val})', ha='center', fontsize=8, fontweight='bold', color=color)
    ax.set_xlim(-1.4, 1.4); ax.set_ylim(-0.75, 1.4); ax.set_aspect('equal'); ax.axis('off')
    p = VISUALS / '01_risk_gauge.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_org_chart():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    parent = FancyBboxPatch((5.5, 7.5), 5, 1.2, boxstyle="round,pad=0.15", facecolor=M_NAVY, edgecolor='none')
    ax.add_patch(parent)
    ax.text(8, 8.1, f'{NEL_DATA["name"]} ({NEL_DATA["ticker"]})', ha='center', fontsize=18, fontweight='bold', color='white')
    ax.text(8, 7.75, f'{NEL_DATA["hq"]} | Public Company | Founded {NEL_DATA["founded"]}', ha='center', fontsize=10, color='#cccccc')
    divs = [('Alkaline\nElectrolyser', f'Heroya, Norway\n{NEL_DATA["capacity_heroya"]} Capacity', 1.0), ('PEM\nElectrolyser', 'Wallingford, CT\n500 MW + GM Partnership', 4.5), ('US\nGigafactory', f'{NEL_DATA["capacity_planned_us"]} Planned\nLocation TBD', 8.0), ('R&D &\nCorporate', 'Oslo HQ\nGlobal Sales', 11.5)]
    for label, desc, x in divs:
        box = FancyBboxPatch((x, 4.8), 3, 1.6, boxstyle="round,pad=0.1", facecolor=M_STEEL, edgecolor='none', alpha=0.9)
        ax.add_patch(box)
        ax.text(x+1.5, 5.85, label, ha='center', fontsize=12, fontweight='bold', color='white')
        ax.text(x+1.5, 5.2, desc, ha='center', fontsize=9, color='#e0e0e0')
        ax.plot([8, x+1.5], [7.5, 6.4], color=M_GREY, lw=1.5, alpha=0.6)
    facts = [(f'{NEL_DATA["order_backlog"]}\nOrder Backlog', 1.5, 2.8), (f'{NEL_DATA["employees"]}\nEmployees', 5.0, 2.8), (f'{NEL_DATA["cash_balance"]}\nCash', 8.5, 2.8), ('Listed on\nOslo Bors', 12.0, 2.8)]
    for text, x, y in facts:
        box = FancyBboxPatch((x, y), 2.5, 1.2, boxstyle="round,pad=0.08", facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1)
        ax.add_patch(box)
        ax.text(x+1.25, y+0.6, text, ha='center', fontsize=11, fontweight='bold', color=M_NAVY)
    ax.set_xlim(0, 16); ax.set_ylim(2, 9.5); ax.axis('off')
    p = VISUALS / '02_org_chart.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_financial():
    fig, ax1 = plt.subplots(figsize=(15, 9), facecolor='white')
    years = ['2021', '2022', '2023', '2024']
    revenue = [796, 1064, 1253, 1600]
    ebitda = [-120, -250, -310, -275]
    margin = [-15.1, -23.5, -24.7, -17.2]
    bars = ax1.bar(years, revenue, color=M_NAVY, width=0.5, zorder=3, label='Revenue (NOK M)')
    for bar, val in zip(bars, revenue):
        ax1.text(bar.get_x()+bar.get_width()/2, bar.get_height()+30, f'NOK {val}M', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax1.set_ylabel('Revenue (NOK Millions)', fontsize=13, color=M_NAVY, fontweight='bold')
    ax1.set_ylim(0, 2000); ax1.tick_params(axis='y', labelcolor=M_NAVY, labelsize=11)
    ax1.tick_params(axis='x', labelsize=13)
    ax2 = ax1.twinx()
    ax2.plot(years, ebitda, color=M_RED, marker='o', markersize=10, linewidth=3, zorder=4, label='EBITDA (NOK M)')
    for x, y, m in zip(years, ebitda, margin):
        color = M_GREEN if y > 0 else M_RED
        ax2.text(x, y-40, f'NOK {y}M\n({m:.1f}%)', ha='center', fontsize=11, fontweight='bold', color=color)
    ax2.set_ylabel('EBITDA (NOK Millions)', fontsize=13, color=M_RED, fontweight='bold')
    ax2.set_ylim(-450, 0); ax2.tick_params(axis='y', labelcolor=M_RED, labelsize=11)
    ax1.spines['top'].set_visible(False); ax2.spines['top'].set_visible(False)
    ax1.grid(axis='y', alpha=0.2, linestyle='--')
    ax1.set_title('Revenue & EBITDA Trajectory 2021-2024', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    props = dict(boxstyle='round,pad=0.5', facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1.5)
    ax1.text(0.02, 0.97, 'Revenue growing +28% YoY; EBITDA improving but still negative.\nAlkaline division now EBITDA positive (NOK 127M FY2024).\nPath to group profitability expected 2026-27.', transform=ax1.transAxes, fontsize=11, va='top', bbox=props, color=M_NAVY)
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, loc='upper right', fontsize=11)
    plt.tight_layout()
    p = VISUALS / '03_financial.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_market():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    companies = ['Thyssenkrupp Nucera', 'Siemens Energy', 'John Cockerill', 'Nel ASA', 'ITM Power', 'Plug Power', 'Bloom Energy']
    capacity = [3000, 2500, 2000, 1500, 800, 600, 400]
    colors = [M_GREY, M_GREY, M_GREY, M_NAVY, M_GREY, M_GREY, M_GREY]
    bars = ax.barh(range(len(companies)), capacity, color=colors, height=0.55, zorder=3)
    for bar, val, c in zip(bars, capacity, companies):
        lc = 'white' if c == 'Nel ASA' else M_NAVY
        ax.text(val-100, bar.get_y()+bar.get_height()/2, f'{val} MW', ha='right', va='center', fontsize=13, fontweight='bold', color=lc)
    ax.set_yticks(range(len(companies))); ax.set_yticklabels(companies, fontsize=13, fontweight='bold')
    ax.invert_yaxis(); ax.get_yticklabels()[3].set_color(M_NAVY)
    ax.set_xlabel('Estimated Manufacturing Capacity / Market Position (MW)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Hydrogen Electrolyser Market Position - Global Players', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.annotate(f'Nel: {NEL_DATA["capacity_heroya"]} Heroya + {NEL_DATA["capacity_wallingford"]} Wallingford\n+ {NEL_DATA["capacity_planned_us"]} US Gigafactory planned', xy=(1500, 3), xytext=(2200, 1.5), fontsize=11, fontweight='bold', color=M_STEEL, arrowprops=dict(arrowstyle='->', color=M_STEEL, lw=2))
    plt.tight_layout()
    p = VISUALS / '04_market.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_timeline():
    fig, ax = plt.subplots(figsize=(14, 6), facecolor='white')
    milestones = [('1927', 'Founded as\nNorsk Elektrisk', M_NAVY), ('2015', 'Spun out\nas Nel ASA', M_STEEL), ('2021', 'Heroya Line 1\n500 MW', M_GREEN), ('2022', 'Heroya Line 2\n1 GW Total', M_GREEN), ('2023', 'GM Partnership\nPEM Technology', M_STEEL), ('2024', 'US Gigafactory\n4 GW Planned', M_NAVY), ('2025', 'Cost Reduction\nPath to Profit', M_AMBER)]
    ax.plot([0, len(milestones)-1], [0, 0], color=M_NAVY, linewidth=3, zorder=1)
    for i, (year, desc, color) in enumerate(milestones):
        ax.plot(i, 0, 'o', markersize=16, color=color, zorder=3)
        ax.plot(i, 0, 'o', markersize=10, color='white', zorder=4)
        y_off = 0.6 if i % 2 == 0 else -0.6
        ax.text(i, y_off, f'{year}\n{desc}', ha='center', va='bottom' if i%2==0 else 'top', fontsize=10, fontweight='bold', color=M_NAVY, bbox=dict(boxstyle='round,pad=0.3', facecolor=M_LIGHT, edgecolor=color, linewidth=1.5))
        ax.plot([i, i], [0, y_off*0.4], color=color, linewidth=1.5, linestyle='--')
    ax.set_title('Strategic Investment Timeline', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(-0.5, len(milestones)-0.5); ax.set_ylim(-1.4, 1.4); ax.axis('off')
    plt.tight_layout()
    p = VISUALS / '05_timeline.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_risk_matrix():
    fig, ax = plt.subplots(figsize=(11, 9), facecolor='white')
    ax.fill_between([0,5], [5,5], [10,10], color=M_AMBER, alpha=0.10)
    ax.fill_between([5,10], [5,5], [10,10], color=M_RED, alpha=0.10)
    ax.fill_between([0,5], [0,0], [5,5], color=M_GREEN, alpha=0.10)
    ax.fill_between([5,10], [0,0], [5,5], color=M_AMBER, alpha=0.10)
    ax.text(2.5, 7.5, 'MONITOR', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    ax.text(7.5, 7.5, 'CRITICAL', ha='center', fontsize=14, fontweight='bold', color=M_RED, alpha=0.4)
    ax.text(2.5, 2.5, 'ACCEPT', ha='center', fontsize=14, fontweight='bold', color=M_GREEN, alpha=0.4)
    ax.text(7.5, 2.5, 'MITIGATE', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    risks = [('Negative\nEBITDA', 4.5, 6.5, 700, M_AMBER), ('Order Backlog\nDecline', 6.0, 5.5, 650, M_AMBER), ('Hydrogen Market\nTiming', 5.5, 7.0, 750, M_RED), ('Competition\n(China/EU)', 6.5, 4.5, 600, M_AMBER), ('Manufacturing\nScale-Up', 4.0, 5.0, 550, M_AMBER), ('Norwegian\nJurisdiction', 2.5, 2.5, 350, M_GREEN), ('Cash Position\nStrong', 2.0, 3.5, 400, M_GREEN)]
    for label, x, y, size, color in risks:
        ax.scatter(x, y, s=size, color=color, alpha=0.7, edgecolors=M_NAVY, linewidth=1.5, zorder=3)
        ax.text(x, y, label, ha='center', va='center', fontsize=8, fontweight='bold', color=M_NAVY)
    ax.set_xlabel('Probability ->', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_ylabel('Impact ->', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Risk Matrix - Nel ASA', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.set_xticks([2.5, 7.5]); ax.set_xticklabels(['Low', 'High'], fontsize=11)
    ax.set_yticks([2.5, 7.5]); ax.set_yticklabels(['Low', 'High'], fontsize=11)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    ax.axhline(y=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    ax.axvline(x=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    plt.tight_layout()
    p = VISUALS / '06_risk_matrix.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_radar():
    fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(projection='polar'), facecolor='white')
    cats = ['Price\nCompetitiveness', 'Manufacturing\nScale', 'Technology\nMaturity', 'Financial\nStability', 'Supply Chain\nSecurity', 'Market\nPosition']
    N = len(cats)
    nel = [7, 8, 8, 5, 7, 7]
    siemens = [6, 9, 9, 9, 8, 9]
    angles = np.linspace(0, 2*np.pi, N, endpoint=False).tolist()
    nel_c = nel + [nel[0]]; s_c = siemens + [siemens[0]]; ang_c = angles + [angles[0]]
    ax.fill(ang_c, nel_c, color=M_NAVY, alpha=0.2)
    ax.plot(ang_c, nel_c, color=M_NAVY, linewidth=2.5, marker='o', markersize=8, label='Nel ASA')
    ax.fill(ang_c, s_c, color=M_STEEL, alpha=0.1)
    ax.plot(ang_c, s_c, color=M_STEEL, linewidth=2, marker='s', markersize=7, linestyle='--', label='Siemens Energy (Benchmark)')
    ax.set_xticks(angles); ax.set_xticklabels(cats, fontsize=11, fontweight='bold', color=M_NAVY)
    ax.set_ylim(0, 10); ax.set_yticks([2,4,6,8,10])
    ax.set_yticklabels(['2','4','6','8','10'], fontsize=9, color=M_GREY)
    ax.grid(color=M_GREY, alpha=0.3)
    ax.set_title('Supplier Benchmarking - Nel vs Siemens Energy', fontsize=14, fontweight='bold', color=M_NAVY, pad=25)
    ax.legend(loc='lower right', fontsize=11, bbox_to_anchor=(1.15, -0.05))
    plt.tight_layout()
    p = VISUALS / '07_radar.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_peer_risk():
    fig, ax = plt.subplots(figsize=(14, 8), facecolor='white')
    companies = ['Siemens Energy', 'Thyssenkrupp', 'Nel ASA', 'Bloom Energy', 'ITM Power', 'Plug Power']
    scores = [32, 35, NEL_DATA['overall_risk'], 52, 55, 68]
    colors = [M_GREEN if s<=33 else M_AMBER if s<=66 else M_RED for s in scores]
    ax.axhspan(0, 33, facecolor=M_GREEN, alpha=0.07); ax.axhspan(33, 66, facecolor=M_AMBER, alpha=0.07)
    ax.axhspan(66, 100, facecolor=M_RED, alpha=0.07)
    ax.text(5.3, 16, 'LOW', fontsize=10, color=M_GREEN, fontweight='bold', alpha=0.6, ha='right')
    ax.text(5.3, 50, 'MEDIUM', fontsize=10, color=M_AMBER, fontweight='bold', alpha=0.6, ha='right')
    ax.text(5.3, 80, 'HIGH', fontsize=10, color=M_RED, fontweight='bold', alpha=0.6, ha='right')
    for i, (c, s, col) in enumerate(zip(companies, scores, colors)):
        ax.plot([i, i], [0, s], color=col, linewidth=3, zorder=2)
        ax.scatter(i, s, s=200, color=col, zorder=3, edgecolors=M_NAVY, linewidth=1.5)
        ax.text(i, s+3, f'{s}', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.scatter(2, 48, s=350, color=M_NAVY, zorder=4, edgecolors=M_NAVY, linewidth=2)
    ax.text(2, 51, '48', ha='center', fontsize=14, fontweight='bold', color=M_NAVY)
    ax.set_xticks(range(len(companies))); ax.set_xticklabels(companies, fontsize=12, fontweight='bold', rotation=15)
    ax.get_xticklabels()[2].set_color(M_NAVY)
    ax.set_ylabel('Risk Score (0-100)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Peer Risk Comparison - Hydrogen Electrolyser Manufacturers', fontsize=14, fontweight='bold', color=M_NAVY, pad=15)
    ax.set_ylim(0, 100); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.15)
    plt.tight_layout()
    p = VISUALS / '08_peer_risk.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_esg():
    fig, axes = plt.subplots(1, 3, figsize=(20, 12), facecolor='white')
    fig.suptitle('ESG Assessment - Nel ASA', fontsize=22, fontweight='bold', color=M_NAVY, y=0.98)
    fig.text(0.5, 0.94, 'Overall ESG Rating: LOW | Strong environmental profile as enabler of green hydrogen economy', ha='center', fontsize=13, color=M_GREY)
    pillars = [
        {'title': 'Environmental', 'rating': 'LOW', 'color': M_GREEN, 'items': [('Green Hydrogen Enabler', M_GREEN, '✓'), ('Renewable Energy Focus', M_GREEN, '✓'), ('Carbon Reduction Tech', M_GREEN, '✓'), ('Circular Economy', M_GREEN, '✓'), ('Manufacturing Emissions', M_AMBER, '⚠'), ('Supply Chain Scope 3', M_AMBER, '⚠')]},
        {'title': 'Social', 'rating': 'LOW', 'color': M_GREEN, 'items': [('Norwegian Labour Standards', M_GREEN, '✓'), ('Employee Safety (Industrial)', M_GREEN, '✓'), ('Green Jobs Creation', M_GREEN, '✓'), ('US Manufacturing Expansion', M_GREEN, '✓'), ('Diversity Reporting', M_AMBER, '⚠'), ('Community Engagement', M_AMBER, '⚠')]},
        {'title': 'Governance', 'rating': 'MEDIUM', 'color': M_AMBER, 'items': [('Public Company Disclosure', M_GREEN, '✓'), ('Oslo Bors Listed', M_GREEN, '✓'), ('Board Independence', M_GREEN, '✓'), ('Executive Compensation', M_AMBER, '⚠'), ('Shareholder Structure', M_AMBER, '⚠'), ('Cavendish Spin-off', M_AMBER, '⚠')]}
    ]
    for ax, pillar in zip(axes, pillars):
        ax.set_xlim(0, 10); ax.set_ylim(0, 12); ax.axis('off')
        title_box = FancyBboxPatch((0.5, 10.5), 9, 1.2, boxstyle="round,pad=0.15", facecolor=pillar['color'], edgecolor='none', alpha=0.9)
        ax.add_patch(title_box)
        ax.text(5, 11.1, pillar['title'], ha='center', fontsize=16, fontweight='bold', color='white')
        ax.text(5, 10.7, f"Rating: {pillar['rating']}", ha='center', fontsize=12, color='white', alpha=0.9)
        for i, (item, color, symbol) in enumerate(pillar['items']):
            y = 9.2 - i * 1.3
            ax.add_patch(Circle((1.5, y), 0.35, color=color, alpha=0.8))
            ax.text(1.5, y, symbol, ha='center', va='center', fontsize=14, fontweight='bold', color='white')
            ax.text(2.5, y, item, va='center', fontsize=12, fontweight='bold', color=M_NAVY)
    controversy_box = FancyBboxPatch((0.03, 0.015), 0.94, 0.105, boxstyle="round,pad=0.01", facecolor=M_RED, edgecolor=M_RED, linewidth=1.5, alpha=0.10, transform=fig.transFigure)
    fig.patches.append(controversy_box)
    fig.text(0.5, 0.098, '⚠  CONTROVERSY SCREENING', ha='center', fontsize=20, fontweight='bold', color=M_RED)
    fig.text(0.5, 0.068, '⚠  Negative EBITDA: Company is loss-making with path to profitability expected 2026-27', ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.042, '⚠  Order backlog decline: Q4 2024 backlog NOK 1.6B down from NOK 2.4B in Q1 2024', ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.016, '⚠  Hydrogen market timing: Green hydrogen adoption slower than projected; revenue uncertainty', ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    plt.subplots_adjust(wspace=0.15, top=0.92, bottom=0.15)
    p = VISUALS / '09_esg.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

print("Generating Nel ASA Product 1 Rev X...")
print("Step 1: Charts...")
visuals = {
    'gauge': gen_risk_gauge(), 'org': gen_org_chart(), 'fin': gen_financial(),
    'mkt': gen_market(), 'timeline': gen_timeline(), 'risk': gen_risk_matrix(),
    'radar': gen_radar(), 'peer': gen_peer_risk(), 'esg': gen_esg(),
}
print("Step 2: Building PPTX...")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# S1: TITLE
print("  Slide 1: Title")
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.55), Inches(13.33), Inches(0.08))
ln.fill.solid(); ln.fill.fore_color.rgb = STEEL; ln.line.fill.background()
add_text(s, 0.7, 1.4, 12, 0.5, 'SUPPLIER EVALUATION REPORT', 16, False, MID_GREY)
add_text(s, 0.7, 2.05, 12, 1.2, 'Nel ASA', 44, True, WHITE)
add_text(s, 0.7, 3.5, 12, 0.6, 'Hydrogen Electrolysers | Alkaline & PEM | 1 GW+ Manufacturing', 18, False, LIGHT_GREY)
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.35), Inches(4.5), Inches(0.05))
div.fill.solid(); div.fill.fore_color.rgb = STEEL; div.line.fill.background()
add_text(s, 0.7, 4.55, 12, 0.4, f'{NEL_DATA["order_backlog"]} Order Backlog | {NEL_DATA["employees"]} Employees | Oslo Bors Listed | Founded 1927', 14, False, LIGHT_GREY)
add_text(s, 0.7, 6.2, 12, 0.4, 'Confidential | March 2026 | Manu Forti Intelligence', 11, False, MID_GREY)
add_logo(s)

# S2: EXEC SUMMARY
print("  Slide 2: Executive Summary")
s = prs.slides.add_slide(blank)
add_header(s, 'EXECUTIVE SUMMARY', 'Nel ASA - Risk & Suitability Overview')
s.shapes.add_picture(str(visuals['gauge']), Inches(0.3), Inches(1.05), Inches(7.2), Inches(3.1))
add_rounded_rect(s, 7.6, 1.05, 5.4, 5.6, LIGHT_BG)
add_text(s, 7.8, 1.22, 5.2, 0.32, 'SUPPLIER SNAPSHOT', 14, True, NAVY)
snap = [('Supplier:', 'Nel ASA (OSE: NEL)'), ('Type:', 'Public Company'), ('Sector:', 'Hydrogen Electrolysers'), ('HQ:', 'Oslo, Norway'), ('Founded:', '1927'), ('Employees:', '~550'), ('Revenue:', 'NOK 1.6B (2024)'), ('Order Backlog:', 'NOK 1.6B'), ('Cash:', 'NOK 1.9B'), ('Manufacturing:', '1.5 GW Capacity')]
for idx, (label, val) in enumerate(snap):
    y = 1.64 + idx * 0.30
    add_text(s, 7.8, y, 2.3, 0.27, label, 11, True, TEXT_SECONDARY)
    add_text(s, 10.1, y, 2.7, 0.27, val, 11, False, TEXT_PRIMARY)
add_rounded_rect(s, 0.3, 5.22, 7.2, 1.43, LIGHT_BG)
add_text(s, 0.5, 5.33, 7, 0.3, 'Key Findings', 13, True, NAVY)
findings = ['• Leading Nordic hydrogen electrolyser manufacturer with 1+ GW capacity', '• Dual technology: Alkaline (Heroya, Norway) and PEM (Wallingford, US)', '• Negative EBITDA but strong cash position (NOK 1.9B) and order backlog', '• ESG leader: Enabling green hydrogen economy; strong environmental profile']
for idx, f in enumerate(findings):
    add_text(s, 0.5, 5.65+idx*0.27, 7, 0.26, f, 11, False, TEXT_PRIMARY)
add_source(s); add_logo(s)

# S3: RECOMMENDATION
print("  Slide 3: Recommendation")
s = prs.slides.add_slide(blank)
add_header(s, 'RECOMMENDATION', 'Decision Summary & Commercial Conditions')
add_rounded_rect(s, 0.5, 1.1, 12.33, 1.35, AMBER)
add_text(s, 0.7, 1.18, 12, 0.55, '⚠ RECOMMENDATION: CONDITIONAL APPROVAL', 20, True, WHITE)
add_text(s, 0.7, 1.78, 12, 0.58, 'Nel is a technically capable, established electrolyser manufacturer with strong ESG credentials and manufacturing scale. However, negative EBITDA and order backlog decline require enhanced due diligence and milestone-based engagement.', 11, False, WHITE)
add_bullets(s, 0.5, 2.65, 6, 2.5, 'Commercial Conditions', ['Payment: Milestone-based (no advance payment)', 'Performance bonds: 10-15% of contract value', 'Liquidated damages: For delivery delays', 'Parent guarantee: Not applicable (public company)', 'Local entity: Norwegian contract preferred'], 14, 11)
add_bullets(s, 6.7, 2.65, 6, 2.5, 'Enhanced Due Diligence (Mandatory)', ['Review latest quarterly financials (Q4 2024)', 'Verify order backlog quality and customer creditworthiness', 'Site visit: Heroya manufacturing facility', 'Technology validation: Independent electrolyser efficiency test', 'Reference checks: Existing utility customers'], 14, 11)
add_rounded_rect(s, 0.5, 5.35, 12.33, 1.5, LIGHT_BG)
add_text(s, 0.7, 5.5, 12, 0.3, 'Overall Risk Summary - MEDIUM (48/100)', 14, True, NAVY)
for x, t, d in [(0.6,'Financial: MEDIUM','Negative EBITDA; strong cash position'), (3.7,'Operational: LOW','Proven manufacturing; 1+ GW capacity'), (6.8,'Geopolitical: LOW','Norwegian-based; stable jurisdiction'), (9.9,'ESG: LOW','Green hydrogen enabler; strong profile')]:
    add_text(s, x, 5.88, 3.1, 0.27, t, 11, True, NAVY)
    add_text(s, x, 6.17, 3.1, 0.27, d, 9, False, TEXT_SECONDARY)
add_source(s); add_logo(s)

# S4: SUPPLIER PROFILE
print("  Slide 4: Supplier Profile")
s = prs.slides.add_slide(blank)
add_header(s, 'SUPPLIER PROFILE', 'Corporate Structure & Global Footprint')
s.shapes.add_picture(str(visuals['org']), Inches(0.3), Inches(1.05), Inches(8.2), Inches(4.48))
add_bullets(s, 8.7, 1.05, 4.3, 3.5, 'Company Overview', ['Nel ASA is a leading global hydrogen', 'company specializing in electrolyser', 'technology for production of green', 'hydrogen from renewable energy.', '', 'Founded 1927; spun out as dedicated', 'hydrogen company in 2015. Listed on', 'Oslo Bors (OSE: NEL).'], 14, 11)
add_bullets(s, 8.7, 4.5, 4.3, 2.4, 'Leadership & Operations', ['CEO: Håkon Volldal', '  (appointed July 2022)', '', 'HQ: Oslo, Norway', 'Manufacturing: Heroya (NO),', '  Wallingford (US)', '', 'R&D: GM partnership (PEM)'], 14, 11)
add_source(s); add_logo(s)

# S5: FINANCIAL HEALTH
print("  Slide 5: Financial Health")
s = prs.slides.add_slide(blank)
add_header(s, 'FINANCIAL HEALTH', 'Revenue & EBITDA Trajectory 2021-2024 | Listed on Oslo Bors')
s.shapes.add_picture(str(visuals['fin']), Inches(0.3), Inches(1.05), Inches(7.5), Inches(4.55))
add_text(s, 8.1, 1.1, 4.9, 0.35, 'Financial Highlights', 14, True, NAVY)
metrics = [('2024 Revenue:', 'NOK 1,600M'), ('2023 Revenue:', 'NOK 1,253M'), ('Revenue Growth:', '+28% YoY'), ('EBITDA 2024:', 'NOK -275M'), ('EBITDA Margin:', '-17.2%'), ('Order Backlog:', 'NOK 1,614M'), ('Cash Balance:', 'NOK 1,876M'), ('Total Debt:', 'NOK 0M (Debt-free)')]
for idx, (label, val) in enumerate(metrics):
    y = 1.56 + idx * 0.30
    add_text(s, 8.1, y, 2.55, 0.29, label, 11, True, TEXT_SECONDARY)
    vc = AMBER if 'EBITDA' in label else GREEN if 'Debt' in label else TEXT_PRIMARY
    add_text(s, 10.65, y, 2.25, 0.29, val, 11, False, vc)
add_rounded_rect(s, 8.1, 4.05, 4.8, 0.44, AMBER)
add_text(s, 8.25, 4.12, 4.5, 0.3, 'Financial Risk: MEDIUM ⚠', 12, True, WHITE, PP_ALIGN.CENTER)
add_bullets(s, 8.1, 4.6, 4.8, 2.3, 'Exposure Guidance', ['Negative EBITDA but strong cash buffer', 'Public company - full financial disclosure', 'Milestone-based payments recommended', 'Monitor quarterly results for improvement'], 12, 10)
add_source(s); add_logo(s)

# S6: MARKET POSITION
print("  Slide 6: Market Position")
s = prs.slides.add_slide(blank)
add_header(s, 'MARKET POSITION', 'Hydrogen Electrolyser Market Position - Global Players')
s.shapes.add_picture(str(visuals['mkt']), Inches(0.3), Inches(1.05), Inches(8.0), Inches(4.46))
add_bullets(s, 8.5, 1.1, 4.5, 2.0, 'Competitive Landscape', ['Nel is a top-4 global electrolyser', 'manufacturer with established', 'technology and manufacturing scale.', 'Competes with German industrial', 'giants and emerging players.'], 14, 11)
add_bullets(s, 8.5, 3.2, 4.5, 2.6, 'Nel Competitive Advantages', ['1+ GW manufacturing capacity', 'Dual technology (Alkaline + PEM)', 'Norwegian/Nordic green energy access', 'US expansion (4 GW Gigafactory planned)', 'GM partnership for PEM technology', 'Public company transparency'], 14, 11)
add_text(s, 0.5, 5.5, 12, 0.3, 'Key Competitors', 13, True, NAVY)
for idx, c in enumerate(['• Thyssenkrupp Nucera - Germany - 3 GW capacity - Industrial giant', '• Siemens Energy - Germany - 2.5 GW - Energy infrastructure leader', '• John Cockerill - Belgium - 2 GW - Industrial equipment focus', '• ITM Power - UK - 800 MW - PEM specialist', '• Plug Power - USA - 600 MW - Integrated hydrogen solutions']):
    add_text(s, 0.6, 5.84+idx*0.29, 12, 0.27, c, 10, False, TEXT_PRIMARY)
add_source(s); add_logo(s)

# S7: OPS & RISK
print("  Slide 7: Ops & Risk")
s = prs.slides.add_slide(blank)
add_header(s, 'OPERATIONAL CAPABILITY & RISK ASSESSMENT', 'Milestones & Strategic Investment | Risk Matrix | Capability Summary')
s.shapes.add_picture(str(visuals['timeline']), Inches(0.3), Inches(1.05), Inches(7.0), Inches(2.92))
s.shapes.add_picture(str(visuals['risk']), Inches(7.65), Inches(1.05), Inches(5.5), Inches(4.46))
add_text(s, 0.3, 5.58, 7, 0.28, 'Risk Summary', 13, True, NAVY)
for idx, (cat, rating, desc) in enumerate([('Financial:', 'MEDIUM', 'Negative EBITDA; strong cash position'), ('Operational:', 'LOW', 'Proven manufacturing; 1+ GW capacity'), ('Technology:', 'LOW', 'Mature alkaline + GM PEM partnership'), ('Competition:', 'MEDIUM', 'Strong German/Chinese competition'), ('Market Timing:', 'MEDIUM', 'Green hydrogen adoption slower than expected')]):
    y = 5.89 + idx * 0.25
    add_text(s, 0.3, y, 2.1, 0.24, cat, 10, True, TEXT_PRIMARY)
    add_text(s, 2.4, y, 1.4, 0.24, rating, 10, True, TEXT_PRIMARY)
    add_text(s, 3.8, y, 3.7, 0.24, desc, 9, False, TEXT_SECONDARY)
add_bullets(s, 7.65, 5.58, 5.5, 1.85, 'Key Capabilities', ['Alkaline electrolyser manufacturing (1 GW)', 'PEM electrolyser technology (GM partnership)', 'Modular, scalable system design', 'Global service and support network'], 12, 10)
add_source(s); add_logo(s)

# S8: COMMERCIAL INTEL
print("  Slide 8: Commercial Intelligence")
s = prs.slides.add_slide(blank)
add_header(s, 'COMMERCIAL INTELLIGENCE & PEER RISK COMPARISON', 'Benchmarking Radar | Peer Risk Profile | Commercial Terms')
s.shapes.add_picture(str(visuals['radar']), Inches(0.2), Inches(1.0), Inches(5.8), Inches(4.74))
s.shapes.add_picture(str(visuals['peer']), Inches(6.2), Inches(1.0), Inches(6.9), Inches(3.85))
add_rounded_rect(s, 0.3, 5.8, 12.7, 1.5, LIGHT_BG)
add_bullets(s, 0.5, 5.85, 6.5, 1.4, 'Commercial Terms & Negotiation', ['Pricing: Competitive vs Siemens/Thyssenkrupp', 'Lead time: 12-18 months for large-scale projects', 'Payment: Milestone-based recommended', 'Warranty: Standard 2-year equipment warranty'], 12, 10)
add_bullets(s, 7.0, 5.85, 5.8, 1.4, 'Key Watch Points', ['⚠ Monitor quarterly financial results for EBITDA improvement', '⚠ Track order backlog trend (declined from NOK 2.4B to 1.6B)', '⚠ Verify US Gigafactory construction timeline', '⚠ Assess competition from Chinese electrolyser manufacturers'], 12, 10)
add_source(s); add_logo(s)

# S9: ESG
print("  Slide 9: ESG Assessment")
s = prs.slides.add_slide(blank)
add_header(s, 'ESG ASSESSMENT', 'Environmental, Social & Governance Screening | Overall Rating: LOW')
s.shapes.add_picture(str(visuals['esg']), Inches(1.67), Inches(1.05), Inches(10.0), Inches(5.93))
add_source(s); add_logo(s)

prs.save(str(OUTPUT))
print(f"Saved: {OUTPUT}")
import os
print(f"Size: {os.path.getsize(str(OUTPUT))/1024:.1f} KB")
