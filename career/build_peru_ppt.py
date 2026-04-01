from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colour Palette (Wood Mac / Bloomberg NEF style) ──────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
NAVY_LIGHT  = RGBColor(0x16, 0x2A, 0x44)   # card/panel
TEAL        = RGBColor(0x00, 0xB4, 0xD8)   # primary accent
AMBER       = RGBColor(0xF5, 0xA6, 0x23)   # secondary accent
GREEN       = RGBColor(0x2E, 0xCC, 0x71)   # positive / high
RED         = RGBColor(0xE7, 0x4C, 0x3C)   # risk / low
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT  = RGBColor(0xB0, 0xBE, 0xC5)
GREY_MID    = RGBColor(0x54, 0x6E, 0x7A)
DIVIDER     = RGBColor(0x1E, 0x3A, 0x5F)

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=NAVY):
    r = slide.shapes.add_shape(1,0,0,Inches(13.33),Inches(7.5))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def rect(slide, l,t,w,h, fill=None, line=None, lw=Pt(0.5)):
    s = slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       size=14, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf  = box.text_frame; tf.word_wrap = wrap
    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name   = 'Calibri'
    return box

def header_bar(slide, title, subtitle=''):
    rect(slide, 0, 0, 13.33, 1.05, fill=NAVY_LIGHT)
    rect(slide, 0, 0, 0.22,  1.05, fill=TEAL)
    tb(slide, title,    0.35, 0.12, 10, 0.55, size=24, bold=True,  color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.35, 0.65, 10, 0.38, size=11, italic=True, color=GREY_LIGHT)
    tb(slide, 'STATKRAFT  ·  CONFIDENTIAL', 11.0, 0.38, 2.1, 0.35,
       size=8, color=GREY_MID, align=PP_ALIGN.RIGHT)

def footer(slide, page):
    rect(slide, 0, 7.22, 13.33, 0.28, fill=NAVY_LIGHT)
    tb(slide, 'Peru BOS/BOP Market Analysis  ·  Statkraft Procurement  ·  February 2026',
       0.35, 7.25, 10, 0.22, size=8, color=GREY_MID)
    tb(slide, str(page), 12.9, 7.25, 0.3, 0.22, size=8, color=GREY_MID, align=PP_ALIGN.RIGHT)

def divider(slide, y):
    rect(slide, 0.35, y, 12.63, 0.02, fill=DIVIDER)

def kpi_box(slide, l, t, w, h, value, label, color=TEAL):
    rect(slide, l, t, w, h, fill=NAVY_LIGHT, line=color, lw=Pt(1))
    rect(slide, l, t, w, 0.06, fill=color)
    tb(slide, value, l+0.1, t+0.15, w-0.2, 0.55, size=22, bold=True, color=color)
    tb(slide, label, l+0.1, t+0.72, w-0.2, 0.4,  size=9,  color=GREY_LIGHT, wrap=True)

def table_header(slide, cols, l, t, col_widths, row_h=0.38):
    x = l
    for i, (col, cw) in enumerate(zip(cols, col_widths)):
        fill = TEAL if i == 0 else NAVY_LIGHT
        rect(slide, x, t, cw, row_h, fill=fill)
        tb(slide, col, x+0.07, t+0.07, cw-0.1, row_h-0.1,
           size=9, bold=True, color=WHITE)
        x += cw

def table_row(slide, cells, l, t, col_widths, row_h=0.32, shade=False):
    x = l
    fill = RGBColor(0x10,0x22,0x38) if shade else NAVY_LIGHT
    for i, (cell, cw) in enumerate(zip(cells, col_widths)):
        rect(slide, x, t, cw-0.01, row_h, fill=fill, line=DIVIDER, lw=Pt(0.3))
        c = TEAL if i == 0 else GREY_LIGHT
        tb(slide, str(cell), x+0.07, t+0.04, cw-0.15, row_h-0.08, size=8.5, color=c)
        x += cw

# ─────────────────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────
prs = new_prs()

# ── SLIDE 1: COVER ────────────────────────────────────────────────────────────
s = blank(prs); bg(s)
rect(s, 0, 0, 0.35, 7.5,   fill=TEAL)
rect(s, 0.35, 2.8, 12.98, 0.04, fill=TEAL)
rect(s, 0.35, 5.2, 12.98, 0.04, fill=AMBER)

tb(s, 'PERU RENEWABLE ENERGY',  0.6, 1.1, 12, 0.8,  size=13, color=TEAL, bold=True)
tb(s, 'BOS & BOP\nMarket Analysis', 0.6, 1.7, 12, 2.2, size=52, bold=True, color=WHITE)
tb(s, 'Balance of System  ·  Balance of Plant  ·  Wind & Solar', 0.6, 3.0, 11, 0.5,
   size=16, color=GREY_LIGHT, italic=True)
tb(s, 'Statkraft Procurement International  ·  February 2026  ·  CONFIDENTIAL',
   0.6, 5.4, 11, 0.45, size=12, color=AMBER)
tb(s, 'Prepared by: Aiden  ·  For: Jonathon Milne, VP Procurement International',
   0.6, 5.95, 11, 0.4,  size=11, color=GREY_LIGHT)
tb(s, 'INTERNAL USE ONLY', 0.6, 6.6, 4, 0.5, size=9, color=GREY_MID)

# ── SLIDE 2: EXECUTIVE SUMMARY ────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Executive Summary', 'Peru BOS/BOP — Key Findings at a Glance')

kpi_box(s, 0.35, 1.25, 2.9,  1.15, '20,000+ MW', 'Renewable pipeline (solar+wind)\nin pre-approval studies')
kpi_box(s, 3.40, 1.25, 2.9,  1.15, '19 / 114',   'Projects with final MINEM concession\nto break ground', color=AMBER)
kpi_box(s, 6.45, 1.25, 2.9,  1.15, '$700–900/kW','Utility-scale solar all-in CAPEX\n(Peru, 2025)')
kpi_box(s,9.50,  1.25, 3.48, 1.15, '$1,100–1,400/kW','Onshore wind all-in CAPEX\n(Peru, 2025)', color=GREEN)

divider(s, 2.6)

findings = [
    ('🌞  Solar BOS', '~40–55% of project CAPEX. Inverters, trackers and transformers are 70–95% imported (China-dominant). Local value concentrated in civil works and DC cabling.'),
    ('💨  Wind BOP',  '~30–45% of project CAPEX. Turbines 100% imported. Foundations, roads and crane pads 85–90% local — strong domestic civil contractor base.'),
    ('⚡  Grid Risk',  'Transmission bottleneck is the #1 procurement risk. $430M national investment underway but insufficient for the 20 GW pipeline. Build 15–20% cost contingency.'),
    ('📋  Bankability','No PPA = no financing. Law 32249 (Jan 2026) improves the framework. Procurement timelines gate-kept by offtake availability.'),
    ('🏗  Strategy',   'If Statkraft develops 1–2 GW in Peru: $400M–$900M in BOS/BOP contracts. Window to lock contractor capacity and frame agreements is NOW.'),
]
y = 2.75
for icon_title, body in findings:
    rect(s, 0.35, y, 12.63, 0.72, fill=NAVY_LIGHT)
    rect(s, 0.35, y, 0.08, 0.72, fill=TEAL)
    tb(s, icon_title, 0.55, y+0.08, 2.8,  0.55, size=10, bold=True,  color=TEAL)
    tb(s, body,       3.35, y+0.08, 9.45, 0.55, size=9.5, color=GREY_LIGHT, wrap=True)
    y += 0.8

footer(s, 2)

# ── SLIDE 3: MARKET OVERVIEW ──────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Peru Renewable Energy Market', 'Installed Capacity & Pipeline — 2025 to 2030')

# Capacity table
cols = ['Scenario', 'Solar PV', 'Wind', 'Total RE', 'Commentary']
cw   = [3.2, 1.8, 1.8, 1.8, 4.55]
table_header(s, cols, 0.35, 1.18, cw)
rows = [
    ['Current (2025)',         '938 MW',   '1,021 MW', '~1,960 MW', 'Operational today'],
    ['Base case 2026',         '2,362 MW', '1,021 MW', '3,383 MW',  'Final concessions only'],
    ['Base case 2027',         '3,242 MW', '1,867 MW', '5,109 MW',  'Sustained growth trajectory'],
    ['Base case 2028',         '3,337 MW', '1,867 MW', '5,203 MW',  'Slight solar increase'],
    ['Full pipeline 2026',     '9,838 MW', '7,281 MW', '17,120 MW', 'All COES pre-approved projects'],
    ['Full potential 2030',    '15,185 MW','9,344 MW', '24,529 MW', '56% solar / 44% wind'],
]
for i, row in enumerate(rows):
    table_row(s, row, 0.35, 1.56 + i*0.35, cw, shade=i%2==0)

divider(s, 3.7)

# Key regions
tb(s, 'KEY DEVELOPMENT REGIONS', 0.35, 3.78, 6, 0.3, size=9, bold=True, color=TEAL)
tb(s, 'CRITICAL MARKET DYNAMIC', 7.0, 3.78, 6, 0.3, size=9, bold=True, color=AMBER)

regions = [
    ('SOLAR',  'Ica — Prime irradiance, flat terrain, grid access\nArequipa/Moquegua — Atacama fringe, GHI >6.5 kWh/m²/day\nLima Periphery — Smaller projects, demand proximity'),
    ('WIND',   'Piura — Strongest resource, existing wind hub\nIca/Lambayeque — Secondary corridors\nAncash/La Libertad — Emerging coastal/highland sites'),
]
x = 0.35
for tech, desc in regions:
    rect(s, x, 4.12, 3.1, 1.5, fill=NAVY_LIGHT)
    tb(s, tech, x+0.1, 4.18, 2.8, 0.3, size=10, bold=True, color=TEAL)
    tb(s, desc, x+0.1, 4.52, 2.9, 1.0, size=9,  color=GREY_LIGHT, wrap=True)
    x += 3.25

rect(s, 7.0, 4.12, 5.98, 1.5, fill=NAVY_LIGHT)
tb(s, '114 projects approved by COES  ·  Only 19 have MINEM final concession\n\n95 projects = 21,142 MW still awaiting authorisation to break ground\n\nPermitting backlog is the primary constraint on procurement volume activation',
   7.1, 4.18, 5.8, 1.4, size=9.5, color=GREY_LIGHT, wrap=True)

footer(s, 3)

# ── SLIDE 4: BOS SOLAR ────────────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Balance of System — Solar PV', 'Component Breakdown, Cost Structure & Supply Origin')

# Left: component table
cols2 = ['BOS Component', '% of CAPEX', 'Local / Import', 'Key Suppliers']
cw2   = [3.0, 1.4, 1.7, 3.1]
table_header(s, cols2, 0.35, 1.18, cw2)
bos_rows = [
    ['PV Modules (reference)',     '30–35%', '100% Import', 'LONGi, JA Solar, Trina, JinkoSolar'],
    ['Inverters (string/central)', '7–10%',  '95% Import',  'Huawei, SMA, Sungrow, ABB'],
    ['Mounting / Trackers',        '10–14%', '70% Import',  'Nextracker, Array Tech, Soltec'],
    ['DC & AC Cabling',            '5–7%',   '40–50% Local','Indeco, CEP, Phelps Dodge'],
    ['MV/HV Transformers',         '6–8%',   '65% Import',  'ABB, Hitachi Energy, Siemens'],
    ['Switchgear & Protection',    '4–6%',   '70% Import',  'Schneider, ABB, Eaton'],
    ['SCADA & Monitoring',         '1–2%',   '90% Import',  'Siemens, GE, vendor-specific'],
    ['Civil Works',                '8–12%',  '85–90% Local','Cosapi, JJC, Mota-Engil Peru'],
    ['Grid Connection (to POI)',   '5–15%',  'Mixed',       'REP, ISA, Cobra (ACS)'],
]
for i, row in enumerate(bos_rows):
    table_row(s, row, 0.35, 1.56 + i*0.33, cw2, row_h=0.33, shade=i%2==0)

# Right: cost waterfall bars
tb(s, 'TYPICAL BOS COST SPLIT  (% of total project CAPEX)', 9.7, 1.18, 3.3, 0.3,
   size=8, bold=True, color=TEAL)

bar_data = [
    ('PV Modules',    32, TEAL),
    ('Inverters',      9, AMBER),
    ('Trackers/Mount',12, GREEN),
    ('Cabling',        6, TEAL),
    ('Transformers',   7, AMBER),
    ('Civil Works',   10, GREEN),
    ('Grid/Other',    10, GREY_MID),
    ('Switchgear+Mon', 6, RGBColor(0x9B,0x59,0xB6)),
]
max_w = 2.8; y0 = 1.6
for label, pct, col in bar_data:
    bw = max_w * pct / 35
    rect(s, 9.7, y0, bw, 0.25, fill=col)
    tb(s, f'{label} {pct}%', 9.7+bw+0.05, y0, 2.5, 0.25, size=8, color=GREY_LIGHT)
    y0 += 0.32

divider(s, 4.6)
tb(s, '⚠  FEOC / ESG NOTE', 0.35, 4.65, 4, 0.28, size=9, bold=True, color=AMBER)
tb(s, 'Huawei dominates the inverter market (~40% share in Peru). European lenders applying IFC/EBRD standards increasingly scrutinise Xinjiang-linked module supply chains. FEOC rules (US-specific) not applicable for Statkraft Latam financing structures, but reputational risk warrants supply chain due diligence documentation.',
   0.35, 4.97, 12.63, 0.9, size=9, color=GREY_LIGHT, wrap=True)

tb(s, 'All-in solar CAPEX: $700–900/kW  ·  BOS = ~40–55% of CAPEX  ·  Grid connection highly variable ($30–150+/kW)',
   0.35, 5.9, 12.63, 0.3, size=9, bold=True, color=TEAL)

footer(s, 4)

# ── SLIDE 5: BOP WIND ─────────────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Balance of Plant — Wind', 'Component Breakdown, Turbine OEM Landscape & Cost Structure')

# Turbine OEM table (top right)
tb(s, 'TURBINE OEM LANDSCAPE (Peru)', 7.1, 1.18, 5.9, 0.3, size=9, bold=True, color=TEAL)
oem_cols = ['OEM', 'Origin', 'Peru Presence', 'FEOC Risk']
oem_cw   = [2.1, 1.2, 1.8, 1.0]
table_header(s, oem_cols, 7.1, 1.5, oem_cw, row_h=0.32)
oem_rows = [
    ['Vestas',           'Denmark', 'Active — existing projects', 'Low'],
    ['Siemens Gamesa',   'Spain/DE', 'Active',                   'Low'],
    ['GE Vernova',       'US',       'Limited (via Brazil)',      'Low'],
    ['Goldwind',         'China',    'Growing',                   'HIGH'],
    ['Envision',         'China',    'Emerging',                  'HIGH'],
]
for i, row in enumerate(oem_rows):
    color_override = RED if row[3]=='HIGH' else None
    table_row(s, row, 7.1, 1.82+i*0.32, oem_cw, row_h=0.32, shade=i%2==0)

# BOP components table
cols3 = ['BOP Component', '% of CAPEX', 'Local / Import', 'Key Suppliers']
cw3   = [3.0, 1.4, 1.7, 2.7]
table_header(s, cols3, 0.35, 1.18, cw3)
bop_rows = [
    ['Wind Turbines (reference)',    '40–50%', '100% Import',   'Vestas, SGRE, GE, Goldwind'],
    ['Civil (roads, pads, founds)',  '10–15%', '85–90% Local',  'Cosapi, Graña y Montero, JJC'],
    ['Internal Collection Grid',     '5–8%',   '50% Local',     'Indeco, Phelps Dodge'],
    ['Onsite Substation',            '4–6%',   '65% Import',    'ABB/Hitachi, Siemens, Schneider'],
    ['Grid Connection (to POI)',     '5–20%',  'Mixed',         'REP, ISA, Cobra (ACS)'],
    ['SCADA & Communications',       '1–2%',   '90% Import',    'OEM-integrated (Vestas/SGRE)'],
    ['O&M Facility & Access',        '1–3%',   '100% Local',    'Local civil contractors'],
]
for i, row in enumerate(bop_rows):
    table_row(s, row, 0.35, 1.56+i*0.38, cw3, row_h=0.38, shade=i%2==0)

divider(s, 4.35)
tb(s, 'KEY BOP PROCUREMENT PRINCIPLES', 0.35, 4.42, 6, 0.28, size=9, bold=True, color=TEAL)
principles = [
    '→  Split turbine supply from BOP civil — local BOP contracts reduce FX exposure and support social licence',
    '→  SCADA: keep OEM-integrated — attempting to separate creates interface risk and voids turbine warranties',
    '→  Grid connection: owner-procure directly with REP/ISA — EPC wrapping adds premium with no value',
    '→  Secure Cosapi/JJC early — dominant contractors face capacity constraints as pipeline accelerates',
]
y = 4.75
for p in principles:
    tb(s, p, 0.45, y, 12.5, 0.28, size=9, color=GREY_LIGHT)
    y += 0.33

tb(s, 'All-in wind CAPEX: $1,100–1,400/kW  ·  BOP = ~30–45% of CAPEX  ·  Coastal/flat sites at lower end',
   0.35, 6.05, 12.63, 0.28, size=9, bold=True, color=TEAL)

footer(s, 5)

# ── SLIDE 6: SUPPLY CHAIN & LOGISTICS ────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Supply Chain Analysis', 'Import Dependency, Logistics & Lead Times')

# Import dependency visual
tb(s, 'EQUIPMENT IMPORT DEPENDENCY', 0.35, 1.18, 6, 0.3, size=9, bold=True, color=TEAL)

import_items = [
    ('PV Modules',         100, RGBColor(0xE7,0x4C,0x3C)),
    ('Wind Turbines',      100, RGBColor(0xE7,0x4C,0x3C)),
    ('Inverters',           95, RGBColor(0xE7,0x4C,0x3C)),
    ('MV/HV Transformers',  70, AMBER),
    ('Tracker Systems',     70, AMBER),
    ('Switchgear',          70, AMBER),
    ('HV Cabling',          60, AMBER),
    ('DC/AC Cabling',       55, RGBColor(0xF3,0xD5,0x00)),
    ('Civil Works',         12, GREEN),
    ('Concrete/Aggregate',   5, GREEN),
]
bar_max = 5.5; y0 = 1.55
for item, pct, col in import_items:
    bw = bar_max * pct / 100
    rect(s, 2.4, y0, bw, 0.25, fill=col)
    tb(s, item,    0.35, y0, 2.0, 0.25, size=8.5, color=GREY_LIGHT)
    tb(s, f'{pct}%', 2.42+bw, y0, 0.5, 0.25, size=8, bold=True, color=col)
    y0 += 0.31

# Legend
for col, label in [(RGBColor(0xE7,0x4C,0x3C),'Critical import risk'), (AMBER,'Moderate import'), (GREEN,'Primarily local')]:
    rect(s, 2.4 + [0, 2.3, 4.0][[(RGBColor(0xE7,0x4C,0x3C),'Critical import risk'), (AMBER,'Moderate import'), (GREEN,'Primarily local')].index((col,label))],
         4.7, 0.18, 0.18, fill=col)

tb(s, '■ Critical import     ■ Moderate import     ■ Primarily local',
   2.35, 4.68, 4.5, 0.28, size=8, color=GREY_LIGHT)

divider(s, 4.98)

# Right panel: logistics
tb(s, 'LOGISTICS & LEAD TIME MATRIX', 8.3, 1.18, 4.7, 0.3, size=9, bold=True, color=TEAL)
log_cols = ['Item', 'Lead Time', 'Key Risk']
log_cw   = [2.0, 1.5, 2.0]
table_header(s, log_cols, 8.3, 1.5, log_cw, row_h=0.32)
log_rows = [
    ['Turbines (blades/nacelle)', '12–18 mths', 'Port handling, convoy'],
    ['HV Transformers',           '10–14 mths', 'Global shortage 2024–26'],
    ['MV Switchgear',             '6–10 mths',  'Tariff disruption'],
    ['Tracker Systems',           '4–6 mths',   'Low'],
    ['Inverters',                 '3–5 mths',   'Low — high supply'],
    ['PV Modules',                '2–4 mths',   'Low'],
    ['Cables (HV)',               '4–8 mths',   'Capacity constraints'],
    ['Civil Equipment',           '1–3 mths',   'Generally available'],
]
for i, row in enumerate(log_rows):
    table_row(s, row, 8.3, 1.82+i*0.35, log_cw, row_h=0.35, shade=i%2==0)

# Callout boxes
boxes = [
    (0.35, 5.08, 4.1, 1.0, AMBER, 'CUSTOMS', 'Peru FTAs with China, US & EU — duties 0–6% on RE equipment. Clearance time variance 30–90 days. Use specialist customs broker.'),
    (4.6,  5.08, 4.1, 1.0, TEAL,  'CURRENCY', 'Equipment contracts dollar-denominated. Civil/labour in PEN. Natural hedge available — structure BOP civil in PEN with CPI adjustment clause.'),
    (8.85, 5.08, 4.1, 1.0, RED,   'GRID CONGESTION', 'Callao port handles oversized cargo. HV transformer shortage is GLOBAL — secure early. Grid POI congestion adds 6–12 month delay risk.'),
]
for l, t, w, h, col, title, body in boxes:
    rect(s, l, t, w, h, fill=NAVY_LIGHT, line=col, lw=Pt(1))
    rect(s, l, t, w, 0.06, fill=col)
    tb(s, title, l+0.1, t+0.12, w-0.2, 0.28, size=9, bold=True, color=col)
    tb(s, body,  l+0.1, t+0.44, w-0.2, 0.52, size=8.5, color=GREY_LIGHT, wrap=True)

footer(s, 6)

# ── SLIDE 7: PROCUREMENT RISKS ────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Procurement Risk Register', 'BOS & BOP — Wind & Solar Projects, Peru')

risk_cols = ['Risk', 'Severity', 'Likelihood', 'Procurement Impact', 'Mitigation']
risk_cw   = [2.8, 1.1, 1.2, 3.2, 4.7]
table_header(s, risk_cols, 0.35, 1.18, risk_cw)
risks = [
    ['Grid connection delay',    'HIGH',   'HIGH',   'Cost overrun 15–25%; schedule slip 6–12 mths', 'Owner-procure; early COES/REP engagement; 20% contingency'],
    ['HV transformer shortage',  'HIGH',   'HIGH',   'Critical path item; 10–14 mth lead time',       'Order at FID; frame agreements with ABB/Hitachi/Siemens'],
    ['Civil contractor capacity','MEDIUM', 'HIGH',   'Price escalation; schedule risk',                'Pre-qualify Cosapi/JJC/Mota-Engil; early works contract'],
    ['PPA non-availability',     'HIGH',   'MEDIUM', 'Project non-bankable; procurement spend at risk','Gate procurement to PPA milestone; use NTP structure'],
    ['Currency (PEN) inflation', 'MEDIUM', 'MEDIUM', 'Civil BOP cost overrun',                         'PEN contracts with CPI adjustment; FX hedging for equipment'],
    ['Customs clearance delay',  'MEDIUM', 'MEDIUM', '30–90 day variance impacts schedule',            'Specialist broker; buffer lead times; staged delivery'],
    ['FEOC/ESG scrutiny',        'LOW',    'LOW',    'Lender pushback on Chinese equipment',           'Supply chain due diligence; Xinjiang attestations'],
    ['Permitting delay (MINEM)', 'HIGH',   'HIGH',   'Procurement wasted if concession withheld',      'Do not commit spend until final concession confirmed'],
    ['Single-source turbine OEM','MEDIUM', 'LOW',    'No competitive tension; warranty leverage lost', 'Dual-source short-list; avoid sole-source at FID'],
]
sev_colors = {'HIGH': RED, 'MEDIUM': AMBER, 'LOW': GREEN}
for i, row in enumerate(risks):
    y = 1.56 + i*0.53
    shade = i%2==0
    # Draw row
    x = 0.35
    for j, (cell, cw) in enumerate(zip(row, risk_cw)):
        fill = RGBColor(0x10,0x22,0x38) if shade else NAVY_LIGHT
        rect(s, x, y, cw-0.01, 0.52, fill=fill, line=DIVIDER, lw=Pt(0.3))
        if j in [1,2] and cell in sev_colors:
            rect(s, x+0.05, y+0.1, cw-0.15, 0.32, fill=sev_colors[cell])
            tb(s, cell, x+0.05, y+0.1, cw-0.15, 0.32, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        else:
            c = TEAL if j==0 else GREY_LIGHT
            tb(s, cell, x+0.07, y+0.06, cw-0.15, 0.42, size=8, color=c, wrap=True)
        x += cw

footer(s, 7)

# ── SLIDE 8: RECOMMENDATIONS ─────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Strategic Procurement Recommendations', 'Statkraft Peru — BOS/BOP Positioning')

recs = [
    (TEAL,  '01', 'Establish Market Intelligence Now',
     'Map active EPC contractors, equipment suppliers, and workload commitments across Peru\'s 20 GW pipeline. The market will tighten rapidly as MINEM concessions accelerate. Information advantage = commercial leverage.'),
    (AMBER, '02', 'Pre-Qualify & Reserve Civil Contractor Capacity',
     'Cosapi, JJC Contratistas and Mota-Engil Peru are the tier-1 civil contractors. Issue pre-qualification RFQs now. Consider early works contracts to secure scheduling priority before the pipeline heats up.'),
    (GREEN, '03', 'Frame Agreements: Transformers & Trackers First',
     'HV transformers are 10–14 month critical path items with global supply constraints. Nextracker/Array Tech have strong Latam coverage. Early engagement yields better pricing and delivery priority.'),
    (TEAL,  '04', 'Owner-Procure Grid Connection',
     'Do not wrap grid connection in the EPC package. Peru\'s transmission bottleneck means EPC contractors will add risk premium. Engage REP/ISA/COES directly. Budget 15–20% contingency.'),
    (AMBER, '05', 'Gate Procurement to Concession Milestones',
     'With 95 of 114 projects still awaiting final MINEM concession, committed procurement spend before concession is capital at risk. Structure NTP-linked procurement with minimal pre-FID commitments.'),
    (RED,   '06', 'Address FEOC/ESG Supply Chain Risk Proactively',
     'Chinese OEMs dominate PV modules and inverters. European lenders require supply chain due diligence. Document Xinjiang attestations, prepare origin-of-manufacture evidence — build this into supplier pre-qualification.'),
]

y = 1.18
for i, (col, num, title, body) in enumerate(recs):
    lx = 0.35 if i%2==0 else 6.74
    if i%2==0 and i>0: y += 1.5
    rect(s, lx, y, 6.24, 1.38, fill=NAVY_LIGHT)
    rect(s, lx, y, 0.07, 1.38, fill=col)
    rect(s, lx+0.15, y+0.1, 0.55, 0.55, fill=col)
    tb(s, num, lx+0.15, y+0.1, 0.55, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, lx+0.85, y+0.1,  5.2, 0.38, size=11, bold=True, color=col)
    tb(s, body,  lx+0.85, y+0.52, 5.2, 0.82, size=9,  color=GREY_LIGHT, wrap=True)

footer(s, 8)

# ── SLIDE 9: SCORECARD ────────────────────────────────────────────────────────
s = blank(prs); bg(s)
header_bar(s, 'Peru BOS/BOP — Market Scorecard', 'Summary Assessment for Statkraft Procurement Strategy')

scores = [
    ('Market Growth Potential', 5, 5, GREEN,
     '20 GW pipeline; government commitment strong; CAGR 15%+'),
    ('BOS Supply Chain Maturity', 3, 5, AMBER,
     'Import-dependent but established import channels; local civil strong'),
    ('BOP Contractor Capability', 4, 5, TEAL,
     'Strong civil base (Cosapi, JJC); electrical contractors thinner'),
    ('Grid Infrastructure', 2, 5, RED,
     'Transmission bottleneck; $430M plan underway but insufficient'),
    ('Offtake/PPA Maturity', 2, 5, RED,
     'Underdeveloped; Law 32249 improving but regulation still being written'),
    ('Regulatory Environment', 3, 5, AMBER,
     'Improving; MINEM concession queue a constraint; direction positive'),
    ('Logistics Accessibility', 3, 5, AMBER,
     'Callao port functional; remote site access variable; lead times manageable'),
    ('Procurement Complexity', 1, 5, RED,
     'HIGH complexity — import dependency, grid risk, contractor capacity, FX'),
]

y0 = 1.22
for label, score, max_score, col, comment in scores:
    rect(s, 0.35, y0, 3.8, 0.5, fill=NAVY_LIGHT)
    tb(s, label, 0.45, y0+0.09, 3.6, 0.35, size=10, bold=True, color=WHITE)
    # Score dots
    for d in range(max_score):
        fc = col if d < score else DIVIDER
        rect(s, 4.25+d*0.42, y0+0.12, 0.32, 0.28, fill=fc)
    # Score label
    stars = '●'*score + '○'*(max_score-score)
    tb(s, f'{score}/{max_score}', 6.45, y0+0.09, 0.7, 0.32, size=10, bold=True, color=col)
    tb(s, comment, 7.2, y0+0.09, 5.78, 0.35, size=8.5, color=GREY_LIGHT, wrap=True)
    y0 += 0.62

divider(s, 6.22)
rect(s, 0.35, 6.3, 12.63, 0.65, fill=NAVY_LIGHT)
rect(s, 0.35, 6.3, 0.08,  0.65, fill=TEAL)
tb(s, 'OVERALL VERDICT', 0.55, 6.36, 2.5, 0.28, size=10, bold=True, color=TEAL)
tb(s, 'Strong market opportunity with HIGH procurement complexity. Peru rewards early movers who establish supplier relationships, pre-qualify contractors and secure grid connections before the 20 GW pipeline unlocks. The window is now.',
   3.1, 6.34, 10.0, 0.55, size=9.5, color=GREY_LIGHT, wrap=True)

footer(s, 9)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = '/Users/jonathonmilne/Documents/Aiden/Peru_BOS_BOP_Market_Analysis.pptx'
prs.save(out)
print(f'Saved → {out}')
