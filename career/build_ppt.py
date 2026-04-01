from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Statkraft brand colours
DARK_GREEN = RGBColor(0x00, 0x4B, 0x2D)
LIGHT_GREEN = RGBColor(0x78, 0xBE, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF4, 0xF4)
MID_GREY = RGBColor(0x55, 0x55, 0x55)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x00, 0x8E, 0x5A)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

# ─────────────────────────────────────────────
# SLIDE 1: Title / Vision
# ─────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Background
add_rect(slide1, 0, 0, 13.33, 7.5, fill_color=DARK_GREEN)

# Green accent bar left
add_rect(slide1, 0, 0, 0.18, 7.5, fill_color=LIGHT_GREEN)

# Candidate + role (top right)
add_text_box(slide1, "Jonathon Milne  |  SVP Procurement Candidate",
             1.0, 0.25, 11.5, 0.5, font_size=11, color=LIGHT_GREEN, align=PP_ALIGN.RIGHT)

# Main title
add_text_box(slide1, "Procurement Vision",
             1.0, 1.4, 11.0, 1.2, font_size=52, bold=True, color=WHITE)
add_text_box(slide1, "Statkraft  ·  All Spend & Categories",
             1.0, 2.55, 9.0, 0.6, font_size=24, color=LIGHT_GREEN)

# Divider line
add_rect(slide1, 1.0, 3.25, 11.5, 0.03, fill_color=LIGHT_GREEN)

# Tagline
add_text_box(slide1, "Cost-Conscious  ·  Strategically Structured  ·  Built to Last",
             1.0, 3.45, 11.5, 0.55, font_size=16, italic=True, color=RGBColor(0xCC,0xFF,0xCC))

# Footer
add_text_box(slide1, "Confidential  ·  February 2026",
             1.0, 6.9, 11.5, 0.4, font_size=10, color=RGBColor(0x88,0xAA,0x88), align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 2: Cost-Conscious Metrics
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_rect(slide2, 0, 0, 13.33, 7.5, fill_color=LIGHT_GREY)
add_rect(slide2, 0, 0, 13.33, 1.1, fill_color=DARK_GREEN)
add_rect(slide2, 0, 0, 0.18, 7.5, fill_color=LIGHT_GREEN)

add_text_box(slide2, "Cost-Conscious Procurement", 0.4, 0.2, 10.0, 0.75,
             font_size=28, bold=True, color=WHITE)
add_text_box(slide2, "Measuring what matters across every level", 0.4, 0.72, 10.0, 0.4,
             font_size=13, color=LIGHT_GREEN)

# Three columns
cols = [
    {
        "title": "⚙  Operations",
        "metrics": ["$/MWh", "OEE / Availability"],
        "desc": "Drive day-to-day efficiency in operational spend",
        "x": 0.4
    },
    {
        "title": "🏗  Projects",
        "metrics": ["$/MW Installed", "IRR / NPV"],
        "desc": "Maximise value on capital projects and investments",
        "x": 4.7
    },
    {
        "title": "🏢  Corporate",
        "metrics": ["IT $/FTE", "$/MWh Benchmark"],
        "desc": "Benchmark overheads against best-in-class peers",
        "x": 9.0
    },
]

col_w = 3.9
for col in cols:
    x = col["x"]
    # Card background
    card = add_rect(slide2, x, 1.4, col_w, 5.5, fill_color=WHITE)
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    # Top accent
    add_rect(slide2, x, 1.4, col_w, 0.12, fill_color=LIGHT_GREEN)
    # Title
    add_text_box(slide2, col["title"], x + 0.15, 1.6, col_w - 0.3, 0.55,
                 font_size=16, bold=True, color=DARK_GREEN)
    # Metrics
    y_m = 2.3
    for m in col["metrics"]:
        add_rect(slide2, x + 0.15, y_m, col_w - 0.3, 0.5, fill_color=LIGHT_GREY)
        add_text_box(slide2, m, x + 0.25, y_m + 0.05, col_w - 0.5, 0.42,
                     font_size=18, bold=True, color=DARK_GREEN)
        y_m += 0.65
    # Description
    add_text_box(slide2, col["desc"], x + 0.15, y_m + 0.2, col_w - 0.3, 1.2,
                 font_size=12, color=MID_GREY, wrap=True)

add_text_box(slide2, "Statkraft  ·  Procurement Vision  ·  Jonathon Milne",
             0.4, 7.1, 12.5, 0.35, font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 3: Strategic Priorities + 5-Year Plan
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, 13.33, 7.5, fill_color=LIGHT_GREY)
add_rect(slide3, 0, 0, 13.33, 1.1, fill_color=DARK_GREEN)
add_rect(slide3, 0, 0, 0.18, 7.5, fill_color=LIGHT_GREEN)

add_text_box(slide3, "Strategic Priorities", 0.4, 0.2, 10.0, 0.75,
             font_size=28, bold=True, color=WHITE)
add_text_box(slide3, "A structured path to procurement excellence", 0.4, 0.72, 10.0, 0.4,
             font_size=13, color=LIGHT_GREEN)

priorities = [
    {
        "num": "01",
        "title": "Strategic Positioning",
        "body": "Elevate procurement from a cost function to a strategic enabler — embedded in business decisions from day one.",
    },
    {
        "num": "02",
        "title": "Staircase Planning",
        "body": "Build capability in structured steps: establish the baseline, identify gaps, set targets, and execute with discipline.",
    },
    {
        "num": "03",
        "title": "Controls & Compliance",
        "body": "Strengthen governance frameworks, improve transparency, and ensure procurement operates to the highest standards.",
    },
]

y_start = 1.35
row_h = 1.65
for i, p in enumerate(priorities):
    y = y_start + i * row_h
    # Row background
    card = add_rect(slide3, 0.4, y, 8.5, row_h - 0.1, fill_color=WHITE)
    card.line.color.rgb = RGBColor(0xDD,0xDD,0xDD)
    # Number badge
    badge = add_rect(slide3, 0.4, y, 0.9, row_h - 0.1, fill_color=DARK_GREEN)
    add_text_box(slide3, p["num"], 0.4, y + 0.35, 0.9, 0.7,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title + body
    add_text_box(slide3, p["title"], 1.45, y + 0.12, 7.3, 0.45,
                 font_size=15, bold=True, color=DARK_GREEN)
    add_text_box(slide3, p["body"], 1.45, y + 0.58, 7.3, 0.95,
                 font_size=11.5, color=MID_GREY, wrap=True)

# 5-Year panel (right column)
add_rect(slide3, 9.3, 1.35, 3.8, 4.85, fill_color=DARK_GREEN)
add_rect(slide3, 9.3, 1.35, 3.8, 0.08, fill_color=LIGHT_GREEN)
add_text_box(slide3, "5-Year Vision", 9.45, 1.5, 3.5, 0.5,
             font_size=16, bold=True, color=WHITE)
add_text_box(slide3, "Learning Organisation", 9.45, 2.05, 3.5, 0.45,
             font_size=13, bold=True, color=LIGHT_GREEN)

milestones = [
    "Y1  ·  Baseline & Quick Wins",
    "Y2  ·  Capability Building",
    "Y3  ·  Strategic Integration",
    "Y4  ·  Performance Culture",
    "Y5  ·  Best-in-Class",
]
y_m = 2.65
for m in milestones:
    add_rect(slide3, 9.45, y_m, 3.5, 0.5, fill_color=ACCENT)
    add_text_box(slide3, m, 9.55, y_m + 0.07, 3.3, 0.38,
                 font_size=11, color=WHITE, bold=False)
    y_m += 0.62

add_text_box(slide3, "Statkraft  ·  Procurement Vision  ·  Jonathon Milne",
             0.4, 7.1, 12.5, 0.35, font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)


# Save
out_path = "/Users/jonathonmilne/.openclaw/workspace/career/Procurement_Vision_JM.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
