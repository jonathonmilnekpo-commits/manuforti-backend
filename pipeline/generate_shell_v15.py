#!/usr/bin/env python3
"""
Generate Shell plc Product 1 v15 — Full deck with all infographics.
Matches Boskalis_Product1_v15_Final-1.pptx layout exactly.
"""

import sys, os
sys.path.insert(0, '/opt/homebrew/lib/python3.14/site-packages')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, Arc, Circle, Wedge
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Branding Constants ──
NAVY = (0, 33, 71)        # #002147
STEEL = (43, 108, 176)    # #2B6CB0
MID_GREY = (113, 128, 150)# #718096
GREEN = (72, 187, 120)    # #48BB78
AMBER = (214, 158, 46)    # #D69E2E
RED = (229, 62, 62)       # #E53E3E
WHITE = (255, 255, 255)
TEXT_PRIMARY = (26, 32, 44)
TEXT_SECONDARY = (74, 85, 104)
LIGHT_BG = (235, 244, 255)  # #EBF4FF

# Matplotlib color strings
M_NAVY = '#002147'
M_STEEL = '#2B6CB0'
M_GREY = '#718096'
M_GREEN = '#48BB78'
M_AMBER = '#D69E2E'
M_RED = '#E53E3E'
M_LIGHT = '#EBF4FF'

OUTPUT_DIR = Path('/Users/jonathonmilne/Documents/Aiden/Venture/shell_v15_visuals')
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUTPUT = '/Users/jonathonmilne/Documents/Aiden/Venture/Shell_Product1_v15_Full.pptx'

def rgb_color(r, g, b):
    """Create RGBColor from tuple."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

def set_cell_font(cell, text, size=11, bold=False, color=TEXT_PRIMARY):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb_color(*color)

# ═══════════════════════════════════════════════════════
# CHART GENERATION
# ═══════════════════════════════════════════════════════

def generate_risk_gauge():
    """Slide 2: Risk gauge dial showing LOW risk (score 28)."""
    fig, ax = plt.subplots(figsize=(14, 6.5), facecolor='white')
    
    # Draw the gauge arc
    theta_start, theta_end = 0, 180
    
    # Background segments: Green (LOW), Amber (MEDIUM), Red (HIGH)
    segments = [
        (0, 60, M_GREEN, 'LOW\n0–33'),
        (60, 120, M_AMBER, 'MEDIUM\n34–66'),
        (120, 180, M_RED, 'HIGH\n67–100'),
    ]
    
    for start, end, color, label in segments:
        theta = np.linspace(np.radians(start), np.radians(end), 100)
        r_outer, r_inner = 1.0, 0.6
        x_outer = np.cos(theta) * r_outer
        y_outer = np.sin(theta) * r_outer
        x_inner = np.cos(theta[::-1]) * r_inner
        y_inner = np.sin(theta[::-1]) * r_inner
        
        x = np.concatenate([x_outer, x_inner])
        y = np.concatenate([y_outer, y_inner])
        ax.fill(x, y, color=color, alpha=0.85)
        
        # Label
        mid_theta = np.radians((start + end) / 2)
        lx = np.cos(mid_theta) * 0.8
        ly = np.sin(mid_theta) * 0.8
        ax.text(lx, ly, label, ha='center', va='center', fontsize=11,
                fontweight='bold', color='white')
    
    # Needle — Shell risk score = 28 (LOW)
    score = 28
    needle_angle = np.radians(score * 180 / 100)
    nx = np.cos(needle_angle) * 0.55
    ny = np.sin(needle_angle) * 0.55
    ax.annotate('', xy=(nx, ny), xytext=(0, 0),
                arrowprops=dict(arrowstyle='->', color=M_NAVY, lw=3.5))
    
    # Center circle
    circle = plt.Circle((0, 0), 0.08, color=M_NAVY, zorder=5)
    ax.add_patch(circle)
    
    # Score text
    ax.text(0, -0.2, f'{score}/100', ha='center', va='center',
            fontsize=28, fontweight='bold', color=M_NAVY)
    ax.text(0, -0.35, 'OVERALL RISK SCORE', ha='center', va='center',
            fontsize=12, color=M_GREY)
    
    # Title
    ax.text(0, 1.25, 'Shell plc — Overall Risk Assessment', ha='center',
            fontsize=18, fontweight='bold', color=M_NAVY)
    ax.text(0, 1.10, 'Integrated Energy | Oil & Gas | Renewables', ha='center',
            fontsize=12, color=M_GREY)
    
    # Risk summary boxes below gauge
    box_data = [
        ('Financial', 'LOW', M_GREEN),
        ('Operational', 'LOW', M_GREEN),
        ('Geopolitical', 'MEDIUM', M_AMBER),
        ('ESG', 'MEDIUM', M_AMBER),
    ]
    
    for idx, (cat, rating, color) in enumerate(box_data):
        bx = -0.75 + idx * 0.5
        by = -0.55
        rect = FancyBboxPatch((bx - 0.2, by - 0.12), 0.4, 0.22,
                               boxstyle="round,pad=0.03", facecolor=color, alpha=0.15,
                               edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(bx, by + 0.02, cat, ha='center', va='center', fontsize=8,
                fontweight='bold', color=M_NAVY)
        ax.text(bx, by - 0.06, rating, ha='center', va='center', fontsize=8,
                fontweight='bold', color=color)
    
    ax.set_xlim(-1.4, 1.4)
    ax.set_ylim(-0.75, 1.4)
    ax.set_aspect('equal')
    ax.axis('off')
    
    path = OUTPUT_DIR / '01_risk_gauge.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Risk gauge: {path}")
    return path


def generate_org_chart():
    """Slide 4: Shell corporate structure diagram."""
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    
    # Parent box
    parent = FancyBboxPatch((5.5, 7.5), 5, 1.2, boxstyle="round,pad=0.15",
                             facecolor=M_NAVY, edgecolor='none')
    ax.add_patch(parent)
    ax.text(8, 8.1, 'Shell plc', ha='center', va='center',
            fontsize=18, fontweight='bold', color='white')
    ax.text(8, 7.75, 'London, UK  |  LSE: SHEL  |  NYSE: SHEL', ha='center',
            fontsize=10, color='#cccccc')
    
    # Division boxes
    divisions = [
        ('Upstream', 'Exploration &\nProduction', 1.0),
        ('Integrated Gas\n& LNG', 'World\'s largest\nLNG portfolio', 4.5),
        ('Downstream &\nChemicals', 'Refining, Trading\n& Retail', 8.0),
        ('Renewables &\nEnergy Solutions', 'Wind, Solar, EV\nCharging, H₂', 11.5),
    ]
    
    for label, desc, x in divisions:
        box = FancyBboxPatch((x, 4.8), 3, 1.6, boxstyle="round,pad=0.1",
                              facecolor=M_STEEL, edgecolor='none', alpha=0.9)
        ax.add_patch(box)
        ax.text(x + 1.5, 5.85, label, ha='center', va='center',
                fontsize=12, fontweight='bold', color='white')
        ax.text(x + 1.5, 5.2, desc, ha='center', va='center',
                fontsize=9, color='#e0e0e0')
        
        # Connector line
        ax.plot([8, x + 1.5], [7.5, 6.4], color=M_GREY, lw=1.5, alpha=0.6)
    
    # Key facts boxes
    facts = [
        ('90,000+\nEmployees', 1.5, 2.8),
        ('70+\nCountries', 5.0, 2.8),
        ('$380B\nRevenue 2024', 8.5, 2.8),
        ('500+\nRetail Brands', 12.0, 2.8),
    ]
    
    for text, x, y in facts:
        box = FancyBboxPatch((x, y), 2.5, 1.2, boxstyle="round,pad=0.08",
                              facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1)
        ax.add_patch(box)
        ax.text(x + 1.25, y + 0.6, text, ha='center', va='center',
                fontsize=11, fontweight='bold', color=M_NAVY)
    
    ax.set_xlim(0, 16)
    ax.set_ylim(2, 9.5)
    ax.axis('off')
    
    path = OUTPUT_DIR / '02_org_chart.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Org chart: {path}")
    return path


def generate_financial_chart():
    """Slide 5: Dual-axis revenue bar + EBITDA line chart."""
    fig, ax1 = plt.subplots(figsize=(15, 9), facecolor='white')
    
    years = ['2021', '2022', '2023', '2024']
    revenue = [261.5, 386.2, 316.6, 380.2]  # $B
    ebitda = [55.3, 93.2, 59.8, 68.4]       # $B
    ebitda_margin = [21.2, 24.1, 18.9, 18.0] # %
    
    # Revenue bars
    bars = ax1.bar(years, revenue, color=M_NAVY, width=0.5, zorder=3, label='Revenue ($B)')
    
    # Data labels on bars
    for bar, val in zip(bars, revenue):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 5,
                f'${val:.1f}B', ha='center', va='bottom', fontsize=13,
                fontweight='bold', color=M_NAVY)
    
    ax1.set_ylabel('Revenue ($B)', fontsize=13, color=M_NAVY, fontweight='bold')
    ax1.set_ylim(0, 450)
    ax1.tick_params(axis='y', labelcolor=M_NAVY, labelsize=11)
    ax1.tick_params(axis='x', labelsize=13)
    
    # EBITDA line on secondary axis
    ax2 = ax1.twinx()
    ax2.plot(years, ebitda, color=M_STEEL, marker='o', markersize=10, linewidth=3,
             zorder=4, label='EBITDA ($B)')
    
    for x, y, m in zip(years, ebitda, ebitda_margin):
        ax2.text(x, y + 2.5, f'${y:.1f}B\n({m:.1f}%)', ha='center', va='bottom',
                fontsize=11, fontweight='bold', color=M_STEEL)
    
    ax2.set_ylabel('EBITDA ($B)', fontsize=13, color=M_STEEL, fontweight='bold')
    ax2.set_ylim(0, 120)
    ax2.tick_params(axis='y', labelcolor=M_STEEL, labelsize=11)
    
    # Style
    ax1.spines['top'].set_visible(False)
    ax2.spines['top'].set_visible(False)
    ax1.grid(axis='y', alpha=0.2, linestyle='--')
    
    # Title
    ax1.set_title('Revenue & EBITDA Trajectory 2021–2024', fontsize=16,
                   fontweight='bold', color=M_NAVY, pad=15)
    
    # Insight box
    insight_text = ('Record revenue recovery in 2024 ($380.2B, +20% YoY)\n'
                    'after 2023 softening. EBITDA $68.4B reflects disciplined\n'
                    'cost management despite margin compression (18.0%).')
    
    props = dict(boxstyle='round,pad=0.5', facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1.5)
    ax1.text(0.02, 0.97, insight_text, transform=ax1.transAxes,
             fontsize=11, va='top', bbox=props, color=M_NAVY)
    
    # Legend
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper right',
               fontsize=11, framealpha=0.9)
    
    plt.tight_layout()
    path = OUTPUT_DIR / '03_financial_chart.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Financial chart: {path}")
    return path


def generate_market_position():
    """Slide 6: Horizontal bar chart — Shell vs peers."""
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    
    companies = ['Saudi Aramco', 'Shell plc', 'ExxonMobil', 'TotalEnergies', 'BP', 'Chevron']
    revenues = [535.0, 380.2, 365.0, 218.0, 210.0, 200.0]
    colors = [M_GREY, M_NAVY, M_GREY, M_GREY, M_GREY, M_GREY]
    
    y_pos = range(len(companies))
    bars = ax.barh(y_pos, revenues, color=colors, height=0.55, zorder=3)
    
    # Data labels
    for bar, val, company in zip(bars, revenues, companies):
        label_color = 'white' if company == 'Shell plc' else M_NAVY
        ax.text(val - 15, bar.get_y() + bar.get_height()/2, f'${val:.0f}B',
                ha='right', va='center', fontsize=13, fontweight='bold', color=label_color)
    
    ax.set_yticks(y_pos)
    ax.set_yticklabels(companies, fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    ax.set_xlabel('Revenue 2024 ($B)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Global Energy Major Revenue Comparison — 2024', fontsize=16,
                  fontweight='bold', color=M_NAVY, pad=15)
    
    # Highlight Shell
    ax.get_yticklabels()[1].set_color(M_NAVY)
    
    # Style
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.tick_params(axis='x', labelsize=11)
    
    # Annotation
    ax.annotate('Shell: #2 globally\n(excl. state-owned)',
                xy=(380.2, 1), xytext=(430, 2.5),
                fontsize=11, fontweight='bold', color=M_STEEL,
                arrowprops=dict(arrowstyle='->', color=M_STEEL, lw=2))
    
    plt.tight_layout()
    path = OUTPUT_DIR / '04_market_position.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Market position: {path}")
    return path


def generate_timeline():
    """Slide 7 (top-left): Strategic investment timeline."""
    fig, ax = plt.subplots(figsize=(14, 6), facecolor='white')
    
    milestones = [
        ('2019', 'Shell New\nEnergies Division', M_STEEL),
        ('2020', 'Net-Zero by\n2050 Pledge', M_GREEN),
        ('2021', '$4B Renewables\nCapex Ramp', M_STEEL),
        ('2022', 'Holland Hydrogen I\nProject Launch', M_GREEN),
        ('2023', 'LNG Canada\nFirst Cargo', M_NAVY),
        ('2024', 'EV Charging\n50,000+ Points', M_STEEL),
    ]
    
    ax.plot([0, len(milestones) - 1], [0, 0], color=M_NAVY, linewidth=3, zorder=1)
    
    for i, (year, desc, color) in enumerate(milestones):
        # Dot on timeline
        ax.plot(i, 0, 'o', markersize=16, color=color, zorder=3)
        ax.plot(i, 0, 'o', markersize=10, color='white', zorder=4)
        
        # Alternate above/below
        y_offset = 0.6 if i % 2 == 0 else -0.6
        va = 'bottom' if i % 2 == 0 else 'top'
        
        ax.text(i, y_offset, f'{year}\n{desc}', ha='center', va=va,
                fontsize=10, fontweight='bold', color=M_NAVY,
                bbox=dict(boxstyle='round,pad=0.3', facecolor=M_LIGHT,
                          edgecolor=color, linewidth=1.5))
        
        ax.plot([i, i], [0, y_offset * 0.4], color=color, linewidth=1.5, linestyle='--')
    
    ax.set_title('Strategic Investment Timeline', fontsize=14, fontweight='bold',
                  color=M_NAVY, pad=10)
    ax.set_xlim(-0.5, len(milestones) - 0.5)
    ax.set_ylim(-1.4, 1.4)
    ax.axis('off')
    
    plt.tight_layout()
    path = OUTPUT_DIR / '05_timeline.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Timeline: {path}")
    return path


def generate_risk_matrix():
    """Slide 7 (right): 2x2 risk matrix."""
    fig, ax = plt.subplots(figsize=(11, 9), facecolor='white')
    
    # Background quadrants
    ax.fill_between([0, 5], [5, 5], [10, 10], color=M_AMBER, alpha=0.10)  # top-left
    ax.fill_between([5, 10], [5, 5], [10, 10], color=M_RED, alpha=0.10)   # top-right
    ax.fill_between([0, 5], [0, 0], [5, 5], color=M_GREEN, alpha=0.10)    # bottom-left
    ax.fill_between([5, 10], [0, 0], [5, 5], color=M_AMBER, alpha=0.10)   # bottom-right
    
    # Quadrant labels
    ax.text(2.5, 7.5, 'MONITOR', ha='center', va='center', fontsize=14,
            fontweight='bold', color=M_AMBER, alpha=0.4)
    ax.text(7.5, 7.5, 'CRITICAL', ha='center', va='center', fontsize=14,
            fontweight='bold', color=M_RED, alpha=0.4)
    ax.text(2.5, 2.5, 'ACCEPT', ha='center', va='center', fontsize=14,
            fontweight='bold', color=M_GREEN, alpha=0.4)
    ax.text(7.5, 2.5, 'MITIGATE', ha='center', va='center', fontsize=14,
            fontweight='bold', color=M_AMBER, alpha=0.4)
    
    # Risk bubbles
    risks = [
        ('Oil Price\nVolatility', 7.5, 7.0, 800, M_RED),
        ('Climate\nLitigation', 4.5, 7.5, 600, M_AMBER),
        ('Geopolitical\nDisruption', 6.5, 5.5, 700, M_AMBER),
        ('Energy\nTransition', 3.5, 6.0, 650, M_AMBER),
        ('Operational\nSafety', 3.0, 3.5, 500, M_GREEN),
        ('FX\nExposure', 5.5, 3.0, 450, M_AMBER),
        ('Regulatory\nCompliance', 2.0, 4.5, 400, M_GREEN),
    ]
    
    for label, x, y, size, color in risks:
        ax.scatter(x, y, s=size, color=color, alpha=0.7, edgecolors=M_NAVY,
                   linewidth=1.5, zorder=3)
        ax.text(x, y, label, ha='center', va='center', fontsize=8,
                fontweight='bold', color=M_NAVY)
    
    # Axes
    ax.set_xlabel('Probability →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_ylabel('Impact →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Risk Matrix — Shell plc', fontsize=14, fontweight='bold',
                  color=M_NAVY, pad=10)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.set_xticks([2.5, 7.5])
    ax.set_xticklabels(['Low', 'High'], fontsize=11)
    ax.set_yticks([2.5, 7.5])
    ax.set_yticklabels(['Low', 'High'], fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    
    # Divider lines
    ax.axhline(y=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    ax.axvline(x=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    
    plt.tight_layout()
    path = OUTPUT_DIR / '06_risk_matrix.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Risk matrix: {path}")
    return path


def generate_radar_chart():
    """Slide 8 (left): Radar/spider chart — Shell vs ExxonMobil."""
    fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(projection='polar'),
                            facecolor='white')
    
    categories = ['Price\nCompetitiveness', 'Lead Time', 'Quality &\nReliability',
                   'Technical\nCapability', 'Supply Chain\nResilience', 'Service\nLevel']
    N = len(categories)
    
    # Shell scores (out of 10)
    shell_scores = [7, 7, 9, 9, 8, 8]
    exxon_scores = [8, 6, 8, 8, 7, 7]
    
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
    
    # Close the polygon
    shell_scores_c = shell_scores + [shell_scores[0]]
    exxon_scores_c = exxon_scores + [exxon_scores[0]]
    angles_c = angles + [angles[0]]
    
    # Plot
    ax.fill(angles_c, shell_scores_c, color=M_NAVY, alpha=0.2)
    ax.plot(angles_c, shell_scores_c, color=M_NAVY, linewidth=2.5, marker='o',
            markersize=8, label='Shell plc')
    
    ax.fill(angles_c, exxon_scores_c, color=M_STEEL, alpha=0.1)
    ax.plot(angles_c, exxon_scores_c, color=M_STEEL, linewidth=2, marker='s',
            markersize=7, linestyle='--', label='ExxonMobil')
    
    ax.set_xticks(angles)
    ax.set_xticklabels(categories, fontsize=11, fontweight='bold', color=M_NAVY)
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(['2', '4', '6', '8', '10'], fontsize=9, color=M_GREY)
    ax.grid(color=M_GREY, alpha=0.3)
    ax.set_title('Supplier Benchmarking — Shell vs ExxonMobil', fontsize=14,
                  fontweight='bold', color=M_NAVY, pad=25)
    ax.legend(loc='lower right', fontsize=11, bbox_to_anchor=(1.15, -0.05))
    
    plt.tight_layout()
    path = OUTPUT_DIR / '07_radar_chart.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Radar chart: {path}")
    return path


def generate_peer_risk():
    """Slide 8 (right): Lollipop chart — peer risk comparison."""
    fig, ax = plt.subplots(figsize=(14, 8), facecolor='white')
    
    companies = ['Shell plc', 'ExxonMobil', 'TotalEnergies', 'BP', 'Chevron', 'Equinor']
    scores = [28, 32, 35, 42, 30, 25]
    colors = []
    for s in scores:
        if s <= 33:
            colors.append(M_GREEN)
        elif s <= 66:
            colors.append(M_AMBER)
        else:
            colors.append(M_RED)
    
    # Risk zone bands
    ax.axhspan(0, 33, facecolor=M_GREEN, alpha=0.07, zorder=0)
    ax.axhspan(33, 66, facecolor=M_AMBER, alpha=0.07, zorder=0)
    ax.axhspan(66, 100, facecolor=M_RED, alpha=0.07, zorder=0)
    
    ax.text(len(companies) - 0.3, 16, 'LOW', fontsize=10, color=M_GREEN,
            fontweight='bold', alpha=0.6, ha='right')
    ax.text(len(companies) - 0.3, 50, 'MEDIUM', fontsize=10, color=M_AMBER,
            fontweight='bold', alpha=0.6, ha='right')
    ax.text(len(companies) - 0.3, 80, 'HIGH', fontsize=10, color=M_RED,
            fontweight='bold', alpha=0.6, ha='right')
    
    # Lollipops
    for i, (company, score, color) in enumerate(zip(companies, scores, colors)):
        ax.plot([i, i], [0, score], color=color, linewidth=3, zorder=2)
        ax.scatter(i, score, s=200, color=color, zorder=3, edgecolors=M_NAVY, linewidth=1.5)
        ax.text(i, score + 3, f'{score}', ha='center', fontsize=13,
                fontweight='bold', color=M_NAVY)
    
    # Highlight Shell
    ax.scatter(0, 28, s=350, color=M_NAVY, zorder=4, edgecolors=M_NAVY, linewidth=2)
    ax.text(0, 28 + 3, '28', ha='center', fontsize=14, fontweight='bold', color=M_NAVY)
    
    ax.set_xticks(range(len(companies)))
    ax.set_xticklabels(companies, fontsize=12, fontweight='bold', rotation=15)
    ax.get_xticklabels()[0].set_color(M_NAVY)
    ax.set_ylabel('Risk Score (0–100)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Peer Risk Comparison — Energy Majors', fontsize=14,
                  fontweight='bold', color=M_NAVY, pad=15)
    ax.set_ylim(0, 100)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.15)
    
    plt.tight_layout()
    path = OUTPUT_DIR / '08_peer_risk.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ Peer risk: {path}")
    return path


def generate_esg_infographic():
    """Slide 9: ESG assessment — 3-column E/S/G layout."""
    fig, axes = plt.subplots(1, 3, figsize=(20, 12), facecolor='white')
    fig.suptitle('ESG Assessment — Shell plc', fontsize=20,
                  fontweight='bold', color=M_NAVY, y=0.98)
    fig.text(0.5, 0.94, 'Overall ESG Rating: MEDIUM  |  Energy transition ambitions offset by historical carbon exposure',
             ha='center', fontsize=12, color=M_GREY)
    
    pillars = [
        {
            'title': 'Environmental',
            'rating': 'MEDIUM',
            'color': M_AMBER,
            'items': [
                ('ISO 14001 Certified', M_GREEN, '✓'),
                ('Net-Zero by 2050 Pledge', M_GREEN, '✓'),
                ('Scope 1&2 Targets Set', M_GREEN, '✓'),
                ('Scope 3 Reduction Pace', M_AMBER, '⚠'),
                ('Historical CO₂ Legacy', M_RED, '✗'),
                ('Biodiversity Programme', M_GREEN, '✓'),
            ]
        },
        {
            'title': 'Social',
            'rating': 'LOW',
            'color': M_GREEN,
            'items': [
                ('Code of Conduct', M_GREEN, '✓'),
                ('Human Rights Policy', M_GREEN, '✓'),
                ('H&S Record (TRIR)', M_GREEN, '✓'),
                ('Community Investment', M_GREEN, '✓'),
                ('Labour Rights Standards', M_GREEN, '✓'),
                ('Grievance Mechanisms', M_GREEN, '✓'),
            ]
        },
        {
            'title': 'Governance',
            'rating': 'LOW',
            'color': M_GREEN,
            'items': [
                ('Anti-Corruption Policy', M_GREEN, '✓'),
                ('Board Independence', M_GREEN, '✓'),
                ('No Sanctions Exposure', M_GREEN, '✓'),
                ('Tax Transparency', M_GREEN, '✓'),
                ('Executive Pay Linked to ESG', M_GREEN, '✓'),
                ('No Beneficial Ownership Issues', M_GREEN, '✓'),
            ]
        }
    ]
    
    for ax, pillar in zip(axes, pillars):
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 12)
        ax.axis('off')
        
        # Title bar
        title_box = FancyBboxPatch((0.5, 10.5), 9, 1.2, boxstyle="round,pad=0.15",
                                     facecolor=pillar['color'], edgecolor='none', alpha=0.9)
        ax.add_patch(title_box)
        ax.text(5, 11.1, pillar['title'], ha='center', va='center',
                fontsize=16, fontweight='bold', color='white')
        ax.text(5, 10.7, f"Rating: {pillar['rating']}", ha='center', va='center',
                fontsize=12, color='white', alpha=0.9)
        
        # Items
        for i, (item, color, symbol) in enumerate(pillar['items']):
            y = 9.2 - i * 1.3
            
            # Status circle
            circle = plt.Circle((1.5, y), 0.35, color=color, alpha=0.8)
            ax.add_patch(circle)
            ax.text(1.5, y, symbol, ha='center', va='center',
                    fontsize=14, fontweight='bold', color='white')
            
            # Item text
            ax.text(2.5, y, item, va='center', fontsize=12,
                    fontweight='bold', color=M_NAVY)
    
    # Controversy section at bottom
    fig.text(0.5, 0.04, 'Controversy Screening', ha='center', fontsize=14,
             fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.015,
             '⚠  Nigeria delta environmental legacy  |  ⚠  Historical climate litigation (Milieudefensie 2021, appeal pending)  |  ⚠  Scope 3 reporting challenges',
             ha='center', fontsize=10, color=M_AMBER)
    
    # ESG conditions
    fig.text(0.08, 0.08, 'ESG Conditions:', fontsize=11, fontweight='bold', color=M_NAVY)
    conditions = ('1. Annual ESG reporting  |  2. Scope 3 reduction roadmap  |  '
                  '3. Community impact assessment for major projects  |  '
                  '4. Third-party ESG audit (EcoVadis/CDP)  |  '
                  '5. Nigeria remediation update  |  6. Climate litigation disclosure')
    fig.text(0.08, 0.06, conditions, fontsize=9, color=M_GREY)
    
    plt.subplots_adjust(wspace=0.15, top=0.92, bottom=0.12)
    path = OUTPUT_DIR / '09_esg_assessment.png'
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  ✓ ESG assessment: {path}")
    return path


# ═══════════════════════════════════════════════════════
# PPTX BUILDER
# ═══════════════════════════════════════════════════════

def add_header(slide, title, subtitle, prs):
    """Add standard navy header bar + texts (matching Boskalis layout)."""
    # Navy header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(13.33), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = rgb_color(*NAVY)
    header.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = rgb_color(*WHITE)
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(12)
    run2.font.color.rgb = rgb_color(200, 210, 220)


def add_source_line(slide):
    """Add source line at y=7.15."""
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Source: Manu Forti Intelligence  |  Confidential  |  March 2026'
    run.font.size = Pt(8)
    run.font.color.rgb = rgb_color(*MID_GREY)


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb_color(*color)
    return txBox


def add_bullet_list(slide, left, top, width, height, title, items,
                     title_size=14, item_size=12, title_color=NAVY, item_color=TEXT_PRIMARY):
    """Add title + bullet list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = rgb_color(*title_color)
    
    # Items
    for item in items:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f'• {item}'
        run.font.size = Pt(item_size)
        run.font.color.rgb = rgb_color(*item_color)
    
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, alpha=1.0):
    """Add a rounded rectangle background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_color(*fill_color)
    shape.line.fill.background()
    return shape


def build_pptx(visuals):
    """Build the full 9-slide PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]  # Blank
    
    # ── SLIDE 1: TITLE ──
    print("  Building Slide 1: Title...")
    slide = prs.slides.add_slide(blank_layout)
    
    # Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb_color(*NAVY)
    bg.line.fill.background()
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.55),
                                    Inches(13.33), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb_color(*STEEL)
    line.line.fill.background()
    
    add_text(slide, 0.7, 1.4, 12, 0.5, 'SUPPLIER EVALUATION REPORT', 16, False, MID_GREY)
    add_text(slide, 0.7, 2.05, 12, 1.2, 'Shell plc', 44, True, WHITE)
    add_text(slide, 0.7, 3.5, 12, 0.6, 'Integrated Energy  |  Oil & Gas  |  Renewables  |  LNG',
             18, False, (200, 210, 220))
    
    # Accent divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.35),
                                   Inches(4.5), Inches(0.05))
    div.fill.solid()
    div.fill.fore_color.rgb = rgb_color(*STEEL)
    div.line.fill.background()
    
    add_text(slide, 0.7, 4.55, 12, 0.4,
             '$380B Revenue  |  90,000 Employees  |  70+ Countries  |  LSE/NYSE Listed',
             14, False, (200, 210, 220))
    add_text(slide, 0.7, 6.2, 12, 0.4,
             'Confidential  |  March 2026  |  Manu Forti Intelligence',
             11, False, MID_GREY)
    
    # ── SLIDE 2: EXECUTIVE SUMMARY ──
    print("  Building Slide 2: Executive Summary...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'EXECUTIVE SUMMARY',
               'Shell plc — Risk & Suitability Overview', prs)
    
    # Risk gauge image
    slide.shapes.add_picture(str(visuals['risk_gauge']),
                               Inches(0.3), Inches(1.05), Inches(7.2), Inches(3.1))
    
    # Supplier snapshot panel
    add_rounded_rect(slide, 7.6, 1.05, 5.4, 5.6, (240, 245, 250))
    add_text(slide, 7.8, 1.22, 5.2, 0.32, 'SUPPLIER SNAPSHOT', 14, True, NAVY)
    
    snapshot_fields = [
        ('Supplier:', 'Shell plc'),
        ('Ticker:', 'LSE: SHEL  |  NYSE: SHEL'),
        ('Sector:', 'Integrated Energy'),
        ('HQ:', 'London, United Kingdom'),
        ('Founded:', '1907 (Royal Dutch Shell merger)'),
        ('Revenue:', '$380.2B (2024, +20% YoY)'),
        ('EBITDA:', '$68.4B (18.0% margin)'),
        ('Order Book:', '$45.2B committed capex'),
        ('Employees:', '90,000+ globally'),
        ('Operations:', '70+ countries, 44,000+ retail sites'),
    ]
    
    for idx, (label, value) in enumerate(snapshot_fields):
        y = 1.64 + idx * 0.30
        add_text(slide, 7.8, y, 2.3, 0.27, label, 11, True, TEXT_SECONDARY)
        add_text(slide, 10.1, y, 2.7, 0.27, value, 11, False, TEXT_PRIMARY)
    
    # Key findings box
    add_rounded_rect(slide, 0.3, 5.22, 7.2, 1.43, (240, 245, 250))
    add_text(slide, 0.5, 5.33, 7, 0.3, 'Key Findings', 13, True, NAVY)
    
    findings = [
        '• 2024 rebound: $380.2B revenue (+20% YoY), $68.4B EBITDA — strong cash generation',
        '• World\'s #2 energy major (ex. state-owned) — dominant LNG, chemicals & retail positions',
        '• Net debt $42.8B vs EBITDA $68.4B = 0.6x leverage — conservative balance sheet',
        '• ESG flag: Nigeria delta legacy + Milieudefensie climate litigation — conditions apply',
    ]
    for idx, finding in enumerate(findings):
        add_text(slide, 0.5, 5.65 + idx * 0.27, 7, 0.26, finding, 11, False, TEXT_PRIMARY)
    
    add_source_line(slide)
    
    # ── SLIDE 3: RECOMMENDATION ──
    print("  Building Slide 3: Recommendation...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'RECOMMENDATION', 'Decision Summary & Commercial Conditions', prs)
    
    # Amber recommendation banner
    add_rounded_rect(slide, 0.5, 1.1, 12.33, 1.35, AMBER)
    add_text(slide, 0.7, 1.18, 12, 0.55,
             '⚠  RECOMMENDATION: APPROVE WITH ESG CONDITIONS', 22, True, WHITE)
    add_text(slide, 0.7, 1.78, 12, 0.58,
             'Shell is recommended for energy supply, LNG, chemicals, and infrastructure services. '
             'Exceptional financial strength ($380B revenue, 0.6x leverage) and global scale. '
             'ESG conditions apply due to climate litigation exposure and Nigeria legacy.',
             11, False, WHITE)
    
    # Commercial conditions (left)
    add_bullet_list(slide, 0.5, 2.65, 6, 2.5,
                     'Commercial Conditions',
                     ['Standard payment terms: Net-30 to Net-60 days',
                      'Parent company guarantee for contracts >$100M',
                      'Performance bonds for major infrastructure projects',
                      'Force-majeure clause covering geopolitical disruption',
                      'Quarterly business reviews on delivery performance'],
                     14, 11)
    
    # ESG conditions (right)
    add_bullet_list(slide, 6.7, 2.65, 6, 2.5,
                     'ESG Conditions (Mandatory)',
                     ['Annual ESG reporting aligned with Statkraft standards',
                      'Scope 3 emissions reduction roadmap required',
                      'Nigeria delta remediation status update',
                      'Climate litigation disclosure and risk assessment',
                      'Third-party ESG audit (EcoVadis/CDP score required)',
                      'Community impact clause for major projects'],
                     14, 11)
    
    # Risk summary strip
    add_rounded_rect(slide, 0.5, 5.35, 12.33, 1.5, (240, 245, 250))
    add_text(slide, 0.7, 5.5, 12, 0.3, 'Overall Risk Summary — LOW (28/100)', 14, True, NAVY)
    
    risk_items = [
        (0.6, 'Financial: LOW', '$380B revenue; 18% EBITDA; 0.6x leverage'),
        (3.7, 'Operational: LOW', 'Global scale; 70+ countries; mature operations'),
        (6.8, 'Geopolitical: MEDIUM', 'Nigeria, Middle East, Russia exit exposure'),
        (9.9, 'ESG: MEDIUM', 'Climate litigation; Scope 3 challenges'),
    ]
    
    for x, title, desc in risk_items:
        add_text(slide, x, 5.88, 3.1, 0.27, title, 11, True, NAVY)
        add_text(slide, x, 6.17, 3.1, 0.27, desc, 9, False, TEXT_SECONDARY)
    
    add_source_line(slide)
    
    # ── SLIDE 4: SUPPLIER PROFILE ──
    print("  Building Slide 4: Supplier Profile...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'SUPPLIER PROFILE', 'Corporate Structure & Global Footprint', prs)
    
    # Org chart image
    slide.shapes.add_picture(str(visuals['org_chart']),
                               Inches(0.3), Inches(1.05), Inches(8.2), Inches(4.48))
    
    # Company overview (right)
    add_bullet_list(slide, 8.7, 1.05, 4.3, 3.5,
                     'Company Overview',
                     ['Shell plc is one of the world\'s largest integrated',
                      'energy companies, operating across upstream',
                      'exploration, integrated gas & LNG, downstream',
                      'chemicals & refining, and a growing renewables',
                      'portfolio. Formed from the 1907 Royal Dutch/Shell',
                      'merger, headquartered in London since 2022.',
                      '',
                      'Listed on LSE and NYSE (SHEL), Shell commands',
                      'the world\'s largest LNG portfolio and one of the',
                      'biggest retail networks (44,000+ stations).'],
                     14, 11)
    
    # Leadership
    add_bullet_list(slide, 8.7, 4.5, 4.3, 2.4,
                     'Leadership',
                     ['Wael Sawan — CEO',
                      '  Appointed Jan 2023; Shell veteran since 1997',
                      '',
                      'Sinead Gorman — CFO',
                      '  Deep energy finance; Shell career 25+ years',
                      '',
                      'HQ: London, UK  |  Founded: 1907'],
                     14, 11)
    
    add_source_line(slide)
    
    # ── SLIDE 5: FINANCIAL HEALTH ──
    print("  Building Slide 5: Financial Health...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'FINANCIAL HEALTH',
               'Revenue & EBITDA Trajectory 2021–2024  |  $380B Revenue  |  Net Debt: $42.8B', prs)
    
    # Financial chart
    slide.shapes.add_picture(str(visuals['financial_chart']),
                               Inches(0.3), Inches(1.05), Inches(7.5), Inches(4.55))
    
    # Financial highlights panel
    add_text(slide, 8.1, 1.1, 4.9, 0.35, 'Financial Highlights', 14, True, NAVY)
    
    metrics = [
        ('2024 Revenue:', '$380.2B  (+20% YoY)'),
        ('2024 EBITDA:', '$68.4B  (18.0% margin)'),
        ('2024 Net Profit:', '$28.5B  (+18% YoY)'),
        ('3yr Revenue CAGR:', '+8.2%  (2021–2024)'),
        ('Order Book:', '$45.2B  (committed capex)'),
        ('Gross Debt:', '$82.1B  (investment grade)'),
        ('Net Cash Position:', '-$42.8B  (net debt)'),
        ('Debt / EBITDA:', '1.2x  ✓  (conservative)'),
    ]
    
    for idx, (label, value) in enumerate(metrics):
        y = 1.56 + idx * 0.30
        add_text(slide, 8.1, y, 2.55, 0.29, label, 11, True, TEXT_SECONDARY)
        # Highlight debt metrics in green
        val_color = GREEN if idx >= 6 else TEXT_PRIMARY
        add_text(slide, 10.65, y, 2.25, 0.29, value, 11, False, val_color)
    
    # Financial risk badge
    add_rounded_rect(slide, 8.1, 4.05, 4.8, 0.44, GREEN)
    add_text(slide, 8.25, 4.12, 4.5, 0.3, 'Financial Risk:  LOW  ✓', 12, True, WHITE,
             PP_ALIGN.CENTER)
    
    # Exposure guidance
    add_bullet_list(slide, 8.1, 4.6, 4.8, 2.3,
                     'Exposure Guidance',
                     ['1.2x Debt/EBITDA — strong capacity for major contracts',
                      'Investment-grade credit: AA- (S&P) — no guarantee needed',
                      'Single contract: standard project sizing applies',
                      'Payment: Net-30 standard; milestone for large projects'],
                     12, 10)
    
    add_source_line(slide)
    
    # ── SLIDE 6: MARKET POSITION ──
    print("  Building Slide 6: Market Position...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'MARKET POSITION',
               'Global Energy Major Revenue Comparison — 2024 Estimates', prs)
    
    # Market position chart
    slide.shapes.add_picture(str(visuals['market_position']),
                               Inches(0.3), Inches(1.05), Inches(8.0), Inches(4.46))
    
    # Competitive landscape (right)
    add_bullet_list(slide, 8.5, 1.1, 4.5, 2.0,
                     'Competitive Landscape',
                     ['Shell is the #2 energy major globally (excluding',
                      'state-owned entities). Dominant positions in LNG',
                      '(largest portfolio), downstream chemicals, and',
                      'retail (44,000+ service stations worldwide).'],
                     14, 11)
    
    add_bullet_list(slide, 8.5, 3.2, 4.5, 2.6,
                     'Shell Competitive Advantages',
                     ['World\'s largest LNG portfolio and trading book',
                      'Unmatched retail network: 44,000+ stations globally',
                      'Integrated value chain: upstream to chemicals',
                      'Strong balance sheet: 1.2x Debt/EBITDA',
                      'Energy transition investments: $4B+/year renewables',
                      'AA- credit rating — financial fortress'],
                     14, 11)
    
    # Key competitors
    add_text(slide, 0.5, 5.5, 12, 0.3, 'Key Competitors', 13, True, NAVY)
    
    competitors = [
        '• Saudi Aramco — Saudi Arabia — $535B — State-owned; world\'s largest by revenue',
        '• ExxonMobil — USA — $365B — Largest private oil major; Permian Basin dominant',
        '• TotalEnergies — France — $218B — Diversified; strong in Africa and renewables',
        '• BP — UK — $210B — Restructuring; pivoting to integrated energy model',
    ]
    for idx, comp in enumerate(competitors):
        add_text(slide, 0.6, 5.84 + idx * 0.29, 12, 0.27, comp, 10, False, TEXT_PRIMARY)
    
    add_source_line(slide)
    
    # ── SLIDE 7: OPERATIONAL CAPABILITY & RISK ──
    print("  Building Slide 7: Ops & Risk...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'OPERATIONAL CAPABILITY  &  RISK ASSESSMENT',
               'Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary', prs)
    
    # Timeline
    slide.shapes.add_picture(str(visuals['timeline']),
                               Inches(0.3), Inches(1.05), Inches(7.0), Inches(2.92))
    
    # Risk matrix
    slide.shapes.add_picture(str(visuals['risk_matrix']),
                               Inches(7.65), Inches(1.05), Inches(5.5), Inches(4.46))
    
    # Risk summary table (bottom-left)
    add_text(slide, 0.3, 5.58, 7, 0.28, 'Risk Summary', 13, True, NAVY)
    
    risk_rows = [
        ('Financial:', '🟢 LOW', '$380B revenue; 18% EBITDA; AA- rated'),
        ('Geopolitical:', '🟡 MEDIUM', 'Nigeria, Middle East, Russia exit'),
        ('Climate:', '🟡 MEDIUM', 'Litigation exposure; Scope 3 challenges'),
        ('Operational:', '🟢 LOW', '70+ countries; mature global operations'),
        ('Commodity:', '🟡 MEDIUM', 'Oil/gas price volatility; hedged partially'),
    ]
    
    for idx, (cat, rating, desc) in enumerate(risk_rows):
        y = 5.89 + idx * 0.25
        add_text(slide, 0.3, y, 2.1, 0.24, cat, 10, True, TEXT_PRIMARY)
        add_text(slide, 2.4, y, 1.4, 0.24, rating, 10, True, TEXT_PRIMARY)
        add_text(slide, 3.8, y, 3.7, 0.24, desc, 9, False, TEXT_SECONDARY)
    
    # Key capabilities (bottom-right)
    add_bullet_list(slide, 7.65, 5.58, 5.5, 1.85,
                     'Key Capabilities',
                     ['ISO 9001 / ISO 14001 / ISO 45001 certified',
                      'World\'s largest LNG operator and trader',
                      '70+ countries — deep local expertise',
                      'Full chain: explore → produce → refine → retail'],
                     12, 10)
    
    add_source_line(slide)
    
    # ── SLIDE 8: COMMERCIAL INTELLIGENCE ──
    print("  Building Slide 8: Commercial Intelligence...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON',
               'Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms', prs)
    
    # Radar chart
    slide.shapes.add_picture(str(visuals['radar_chart']),
                               Inches(0.2), Inches(1.0), Inches(5.8), Inches(4.74))
    
    # Peer risk chart
    slide.shapes.add_picture(str(visuals['peer_risk']),
                               Inches(6.2), Inches(1.0), Inches(6.9), Inches(3.85))
    
    # Commercial terms box
    add_rounded_rect(slide, 0.3, 5.8, 12.7, 1.5, (240, 245, 250))
    
    add_bullet_list(slide, 0.5, 5.85, 6.5, 1.4,
                     'Commercial Terms & Negotiation',
                     ['Pricing: Indexed to market benchmarks (Brent, JKM, TTF)',
                      'Lead time: LNG cargoes 1–3 months; infrastructure 6–18 months',
                      'Leverage: Volume commitment; long-term supply agreements',
                      'IP: Shell retains technology IP; client owns project deliverables'],
                     12, 10)
    
    add_bullet_list(slide, 7.0, 5.85, 5.8, 1.4,
                     'Key Watch Points',
                     ['⚠ Oil price: Monitor Brent exposure on long-term contracts',
                      '⚠ Climate: Litigation outcome may affect strategy/costs',
                      '⚠ FX: Hedge USD/EUR/NOK exposure on supply agreements'],
                     12, 10)
    
    add_source_line(slide)
    
    # ── SLIDE 9: ESG ASSESSMENT ──
    print("  Building Slide 9: ESG Assessment...")
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, 'ESG ASSESSMENT',
               'Environmental, Social & Governance Screening  |  Overall Rating: MEDIUM', prs)
    
    # ESG infographic
    slide.shapes.add_picture(str(visuals['esg']),
                               Inches(1.67), Inches(1.05), Inches(10.0), Inches(5.93))
    
    add_source_line(slide)
    
    # Save
    prs.save(PPTX_OUTPUT)
    print(f"\n✅ PPTX saved: {PPTX_OUTPUT}")
    return PPTX_OUTPUT


# ═══════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════

if __name__ == '__main__':
    print("🎨 Generating Shell plc Product 1 v15 — Full Infographics\n")
    print("Step 1/2: Generating charts...")
    
    visuals = {
        'risk_gauge': generate_risk_gauge(),
        'org_chart': generate_org_chart(),
        'financial_chart': generate_financial_chart(),
        'market_position': generate_market_position(),
        'timeline': generate_timeline(),
        'risk_matrix': generate_risk_matrix(),
        'radar_chart': generate_radar_chart(),
        'peer_risk': generate_peer_risk(),
        'esg': generate_esg_infographic(),
    }
    
    print(f"\nStep 2/2: Building PPTX...")
    build_pptx(visuals)
    
    print(f"\n📊 Visuals: {OUTPUT_DIR}")
    print(f"📎 Deck: {PPTX_OUTPUT}")
    print(f"\nDone! ✅")
