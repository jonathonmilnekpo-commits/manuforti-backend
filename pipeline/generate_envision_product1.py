#!/usr/bin/env python3
"""
Generate Envision Energy Product 1 — Full deck with all infographics.
Follows ALL locked Product 1 standards:
- Risk gauge on Exec Summary (LOCKED)
- Supplier logo auto-resized in top-right white box (LOCKED)
- Manu Forti clan logo bottom-right all slides (LOCKED)
- Controversy screening at 16pt (LOCKED)
- 8 mandatory financial metrics (LOCKED)
- 9-slide structure (LOCKED)
"""

import sys
sys.path.insert(0, '/opt/homebrew/lib/python3.14/site-packages')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── Branding ──
NAVY = (0, 33, 71)
STEEL = (43, 108, 176)
MID_GREY = (113, 128, 150)
GREEN = (72, 187, 120)
AMBER = (214, 158, 46)
RED = (229, 62, 62)
WHITE = (255, 255, 255)
TEXT_PRIMARY = (26, 32, 44)
TEXT_SECONDARY = (74, 85, 104)

M_NAVY = '#002147'
M_STEEL = '#2B6CB0'
M_GREY = '#718096'
M_GREEN = '#48BB78'
M_AMBER = '#D69E2E'
M_RED = '#E53E3E'
M_LIGHT = '#EBF4FF'

# ── Paths ──
BASE = Path('/Users/jonathonmilne/Documents/Aiden/Venture')
VISUALS = BASE / 'envision_visuals'
VISUALS.mkdir(exist_ok=True)
SUPPLIER_LOGO = BASE / 'envision_logo_transparent.png'
CLAN_LOGO = Path('/Users/jonathonmilne/.openclaw/workspace/skills/product-1-generator/assets/manu_forti_logo.png')
OUTPUT = BASE / 'Envision_Product1.pptx'

def rgb(r, g, b):
    return RGBColor(r, g, b)

def auto_resize_logo(logo_path, max_w=2.0, max_h=0.75):
    img = Image.open(logo_path)
    aspect = img.size[0] / img.size[1]
    w, h = max_w, max_w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    return w, h

def add_header(slide, title, subtitle):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1))
    h.fill.solid(); h.fill.fore_color.rgb = rgb(*NAVY); h.line.fill.background()
    
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = rgb(*WHITE)
    
    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = subtitle; r2.font.size = Pt(12); r2.font.color.rgb = rgb(200, 210, 220)

def add_source(slide):
    t = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = 'Source: Manu Forti Intelligence  |  Confidential  |  March 2026'
    r.font.size = Pt(8); r.font.color.rgb = rgb(*MID_GREY)

def add_text(slide, l, t, w, h, text, sz=14, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = rgb(*color)
    return tb

def add_bullets(slide, l, t, w, h, title, items, tsz=14, isz=12, tc=NAVY, ic=TEXT_PRIMARY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True; r.font.color.rgb = rgb(*tc)
    for item in items:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = f'• {item}'; r.font.size = Pt(isz); r.font.color.rgb = rgb(*ic)

def add_rounded_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(*color); s.line.fill.background()
    return s

def add_logos(slide, is_title=False):
    """Add supplier logo (top-right, auto-resized) + Manu Forti (bottom-right)."""
    # Supplier logo
    lw, lh = auto_resize_logo(str(SUPPLIER_LOGO))
    bx, bw, bh = 10.92, 2.26, 0.87
    by = 0.5 if is_title else 0.06
    backing = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(by), Inches(bw), Inches(bh))
    backing.fill.solid(); backing.fill.fore_color.rgb = rgb(*WHITE); backing.line.fill.background()
    lx = bx + (bw - lw) / 2
    ly = by + (bh - lh) / 2
    slide.shapes.add_picture(str(SUPPLIER_LOGO), Inches(lx), Inches(ly), Inches(lw), Inches(lh))
    
    # Clan logo
    cimg = Image.open(str(CLAN_LOGO))
    ca = cimg.size[0] / cimg.size[1]
    cw = 0.6; ch = cw / ca
    cx = 13.333 - cw - 0.2
    cy = 7.5 - ch - 0.15
    slide.shapes.add_picture(str(CLAN_LOGO), Inches(cx), Inches(cy), Inches(cw), Inches(ch))


# ═══════════════════════════════════════════════════════
# CHARTS
# ═══════════════════════════════════════════════════════

def gen_risk_gauge():
    fig, ax = plt.subplots(figsize=(14, 6.5), facecolor='white')
    
    segments = [(0, 60, M_GREEN, 'LOW\n0–33'), (60, 120, M_AMBER, 'MEDIUM\n34–66'), (120, 180, M_RED, 'HIGH\n67–100')]
    for start, end, color, label in segments:
        theta = np.linspace(np.radians(start), np.radians(end), 100)
        x_o = np.cos(theta); y_o = np.sin(theta)
        x_i = np.cos(theta[::-1]) * 0.6; y_i = np.sin(theta[::-1]) * 0.6
        ax.fill(np.concatenate([x_o, x_i]), np.concatenate([y_o, y_i]), color=color, alpha=0.85)
        mid = np.radians((start+end)/2)
        ax.text(np.cos(mid)*0.8, np.sin(mid)*0.8, label, ha='center', va='center', fontsize=11, fontweight='bold', color='white')
    
    # Score = 55 (MEDIUM — Chinese private company, geopolitical, ESG)
    score = 55
    needle = np.radians(score * 180 / 100)
    ax.annotate('', xy=(np.cos(needle)*0.55, np.sin(needle)*0.55), xytext=(0,0),
                arrowprops=dict(arrowstyle='->', color=M_NAVY, lw=3.5))
    ax.add_patch(plt.Circle((0,0), 0.08, color=M_NAVY, zorder=5))
    ax.text(0, -0.2, f'{score}/100', ha='center', fontsize=28, fontweight='bold', color=M_NAVY)
    ax.text(0, -0.35, 'OVERALL RISK SCORE', ha='center', fontsize=12, color=M_GREY)
    ax.text(0, 1.25, 'Envision Energy — Overall Risk Assessment', ha='center', fontsize=18, fontweight='bold', color=M_NAVY)
    ax.text(0, 1.10, 'Wind Turbines | Battery Storage | Energy Software', ha='center', fontsize=12, color=M_GREY)
    
    boxes = [('Financial', 'MEDIUM', M_AMBER), ('Operational', 'LOW', M_GREEN), ('Geopolitical', 'HIGH', M_RED), ('ESG', 'MEDIUM', M_AMBER)]
    for idx, (cat, rating, color) in enumerate(boxes):
        bx = -0.75 + idx * 0.5
        rect = FancyBboxPatch((bx-0.2, -0.67), 0.4, 0.22, boxstyle="round,pad=0.03", facecolor=color, alpha=0.15, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(bx, -0.54, cat, ha='center', fontsize=8, fontweight='bold', color=M_NAVY)
        ax.text(bx, -0.62, rating, ha='center', fontsize=8, fontweight='bold', color=color)
    
    ax.set_xlim(-1.4, 1.4); ax.set_ylim(-0.75, 1.4); ax.set_aspect('equal'); ax.axis('off')
    p = VISUALS / '01_risk_gauge.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Risk gauge"); return p

def gen_org_chart():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    
    parent = FancyBboxPatch((5.5, 7.5), 5, 1.2, boxstyle="round,pad=0.15", facecolor=M_NAVY, edgecolor='none')
    ax.add_patch(parent)
    ax.text(8, 8.1, 'Envision Group', ha='center', fontsize=18, fontweight='bold', color='white')
    ax.text(8, 7.75, 'Shanghai, China  |  Private  |  Founded 2007', ha='center', fontsize=10, color='#cccccc')
    
    divs = [
        ('Envision\nEnergy', 'Wind Turbines\n& Smart Grid', 1.0),
        ('Envision\nAESC', 'EV Batteries\n(ex-Nissan)', 4.5),
        ('Envision\nDigital', 'EnOS™ AIoT\nPlatform', 8.0),
        ('Envision\nGreen H₂', 'Green Hydrogen\n& Net Zero', 11.5),
    ]
    for label, desc, x in divs:
        box = FancyBboxPatch((x, 4.8), 3, 1.6, boxstyle="round,pad=0.1", facecolor=M_STEEL, edgecolor='none', alpha=0.9)
        ax.add_patch(box)
        ax.text(x+1.5, 5.85, label, ha='center', fontsize=12, fontweight='bold', color='white')
        ax.text(x+1.5, 5.2, desc, ha='center', fontsize=9, color='#e0e0e0')
        ax.plot([8, x+1.5], [7.5, 6.4], color=M_GREY, lw=1.5, alpha=0.6)
    
    facts = [('30+ GW\nInstalled Wind', 1.5, 2.8), ('6,000+\nEmployees', 5.0, 2.8),
             ('~$8B Est.\nRevenue 2024', 8.5, 2.8), ('20+\nCountries', 12.0, 2.8)]
    for text, x, y in facts:
        box = FancyBboxPatch((x, y), 2.5, 1.2, boxstyle="round,pad=0.08", facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1)
        ax.add_patch(box)
        ax.text(x+1.25, y+0.6, text, ha='center', fontsize=11, fontweight='bold', color=M_NAVY)
    
    ax.set_xlim(0, 16); ax.set_ylim(2, 9.5); ax.axis('off')
    p = VISUALS / '02_org_chart.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Org chart"); return p

def gen_financial():
    fig, ax1 = plt.subplots(figsize=(15, 9), facecolor='white')
    years = ['2021', '2022', '2023', '2024']
    revenue = [4.2, 6.8, 7.1, 8.0]  # $B estimated
    ebitda = [0.42, 0.75, 0.78, 0.96]
    margin = [10.0, 11.0, 11.0, 12.0]
    
    bars = ax1.bar(years, revenue, color=M_NAVY, width=0.5, zorder=3, label='Revenue ($B est.)')
    for bar, val in zip(bars, revenue):
        ax1.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.2, f'${val:.1f}B', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax1.set_ylabel('Revenue ($B est.)', fontsize=13, color=M_NAVY, fontweight='bold')
    ax1.set_ylim(0, 12); ax1.tick_params(axis='y', labelcolor=M_NAVY, labelsize=11)
    ax1.tick_params(axis='x', labelsize=13)
    
    ax2 = ax1.twinx()
    ax2.plot(years, ebitda, color=M_STEEL, marker='o', markersize=10, linewidth=3, zorder=4, label='EBITDA ($B est.)')
    for x, y, m in zip(years, ebitda, margin):
        ax2.text(x, y+0.06, f'${y:.2f}B\n({m:.0f}%)', ha='center', fontsize=11, fontweight='bold', color=M_STEEL)
    ax2.set_ylabel('EBITDA ($B est.)', fontsize=13, color=M_STEEL, fontweight='bold')
    ax2.set_ylim(0, 1.5); ax2.tick_params(axis='y', labelcolor=M_STEEL, labelsize=11)
    
    ax1.spines['top'].set_visible(False); ax2.spines['top'].set_visible(False)
    ax1.grid(axis='y', alpha=0.2, linestyle='--')
    ax1.set_title('Revenue & EBITDA Trajectory 2021–2024 (Estimated)', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    
    props = dict(boxstyle='round,pad=0.5', facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1.5)
    ax1.text(0.02, 0.97, 'Rapid growth driven by China wind boom and AESC\nbattery expansion. Privately held — financials estimated\nfrom industry sources and bond filings.',
             transform=ax1.transAxes, fontsize=11, va='top', bbox=props, color=M_NAVY)
    
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, loc='upper right', fontsize=11)
    plt.tight_layout()
    p = VISUALS / '03_financial.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Financial chart"); return p

def gen_market():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    companies = ['Vestas', 'Goldwind', 'GE Vernova', 'Envision Energy', 'Siemens Gamesa', 'Mingyang']
    gw = [17.5, 16.0, 8.5, 7.8, 7.2, 6.5]  # GW installed 2024 est.
    colors = [M_GREY, M_GREY, M_GREY, M_NAVY, M_GREY, M_GREY]
    
    bars = ax.barh(range(len(companies)), gw, color=colors, height=0.55, zorder=3)
    for bar, val, c in zip(bars, gw, companies):
        lc = 'white' if c == 'Envision Energy' else M_NAVY
        ax.text(val-0.3, bar.get_y()+bar.get_height()/2, f'{val:.1f} GW', ha='right', va='center', fontsize=13, fontweight='bold', color=lc)
    
    ax.set_yticks(range(len(companies))); ax.set_yticklabels(companies, fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    ax.set_xlabel('New Wind Capacity Installed 2024 (GW est.)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Global Wind Turbine OEM Ranking — 2024 Installations', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    ax.get_yticklabels()[3].set_color(M_NAVY)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.annotate('Envision: #4 globally\nChina\'s #2 after Goldwind', xy=(7.8, 3), xytext=(10, 1.5),
                fontsize=11, fontweight='bold', color=M_STEEL, arrowprops=dict(arrowstyle='->', color=M_STEEL, lw=2))
    plt.tight_layout()
    p = VISUALS / '04_market.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Market position"); return p

def gen_timeline():
    fig, ax = plt.subplots(figsize=(14, 6), facecolor='white')
    milestones = [
        ('2007', 'Founded in\nJiangyin', M_NAVY),
        ('2014', 'Global Smart\nEnergy Vision', M_STEEL),
        ('2018', 'Acquired Nissan\nAESC Battery', M_GREEN),
        ('2020', '#4 Global WTG\nOEM Ranking', M_NAVY),
        ('2022', 'Spain Giga-\nfactory Announced', M_STEEL),
        ('2024', 'EnOS™ Manages\n620+ GWh', M_GREEN),
    ]
    ax.plot([0, len(milestones)-1], [0, 0], color=M_NAVY, linewidth=3, zorder=1)
    for i, (year, desc, color) in enumerate(milestones):
        ax.plot(i, 0, 'o', markersize=16, color=color, zorder=3)
        ax.plot(i, 0, 'o', markersize=10, color='white', zorder=4)
        y_off = 0.6 if i % 2 == 0 else -0.6
        ax.text(i, y_off, f'{year}\n{desc}', ha='center', va='bottom' if i%2==0 else 'top',
                fontsize=10, fontweight='bold', color=M_NAVY,
                bbox=dict(boxstyle='round,pad=0.3', facecolor=M_LIGHT, edgecolor=color, linewidth=1.5))
        ax.plot([i, i], [0, y_off*0.4], color=color, linewidth=1.5, linestyle='--')
    ax.set_title('Strategic Investment Timeline', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(-0.5, len(milestones)-0.5); ax.set_ylim(-1.4, 1.4); ax.axis('off')
    plt.tight_layout()
    p = VISUALS / '05_timeline.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Timeline"); return p

def gen_risk_matrix():
    fig, ax = plt.subplots(figsize=(11, 9), facecolor='white')
    ax.fill_between([0,5], [5,5], [10,10], color=M_AMBER, alpha=0.10)
    ax.fill_between([5,10], [5,5], [10,10], color=M_RED, alpha=0.10)
    ax.fill_between([0,5], [0,0], [5,5], color=M_GREEN, alpha=0.10)
    ax.fill_between([5,10], [0,0], [5,5], color=M_AMBER, alpha=0.10)
    ax.text(2.5, 7.5, 'MONITOR', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    ax.text(7.5, 7.5, 'CRITICAL', ha='center', fontsize=14, fontweight='bold', color=M_RED, alpha=0.4)
    ax.text(2.5, 2.5, 'ACCEPT', ha='center', fontsize=14, fontweight='bold', color=M_GREEN, alpha=0.4)
    ax.text(7.5, 2.5, 'MITIGATE', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    
    risks = [
        ('US-China\nTrade War', 7.0, 8.0, 900, M_RED),
        ('IP Transfer\nConcerns', 6.0, 6.5, 650, M_RED),
        ('Financial\nTransparency', 5.5, 5.0, 600, M_AMBER),
        ('NDAA\nRestrictions', 7.5, 6.0, 750, M_RED),
        ('Quality\nConsistency', 4.0, 4.5, 500, M_AMBER),
        ('FX / RMB\nExposure', 5.0, 3.0, 450, M_AMBER),
        ('Supply Chain\nResilience', 3.0, 3.0, 400, M_GREEN),
    ]
    for label, x, y, size, color in risks:
        ax.scatter(x, y, s=size, color=color, alpha=0.7, edgecolors=M_NAVY, linewidth=1.5, zorder=3)
        ax.text(x, y, label, ha='center', va='center', fontsize=8, fontweight='bold', color=M_NAVY)
    
    ax.set_xlabel('Probability →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_ylabel('Impact →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Risk Matrix — Envision Energy', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.set_xticks([2.5, 7.5]); ax.set_xticklabels(['Low', 'High'], fontsize=11)
    ax.set_yticks([2.5, 7.5]); ax.set_yticklabels(['Low', 'High'], fontsize=11)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    ax.axhline(y=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    ax.axvline(x=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    plt.tight_layout()
    p = VISUALS / '06_risk_matrix.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Risk matrix"); return p

def gen_radar():
    fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(projection='polar'), facecolor='white')
    cats = ['Price\nCompetitiveness', 'Lead Time', 'Quality &\nReliability', 'Technical\nCapability', 'Supply Chain\nResilience', 'Service\nLevel']
    N = len(cats)
    envision = [9, 6, 7, 8, 6, 5]
    vestas = [6, 8, 9, 8, 8, 9]
    angles = np.linspace(0, 2*np.pi, N, endpoint=False).tolist()
    env_c = envision + [envision[0]]; ves_c = vestas + [vestas[0]]; ang_c = angles + [angles[0]]
    ax.fill(ang_c, env_c, color=M_NAVY, alpha=0.2)
    ax.plot(ang_c, env_c, color=M_NAVY, linewidth=2.5, marker='o', markersize=8, label='Envision Energy')
    ax.fill(ang_c, ves_c, color=M_STEEL, alpha=0.1)
    ax.plot(ang_c, ves_c, color=M_STEEL, linewidth=2, marker='s', markersize=7, linestyle='--', label='Vestas')
    ax.set_xticks(angles); ax.set_xticklabels(cats, fontsize=11, fontweight='bold', color=M_NAVY)
    ax.set_ylim(0, 10); ax.set_yticks([2,4,6,8,10])
    ax.set_yticklabels(['2','4','6','8','10'], fontsize=9, color=M_GREY)
    ax.grid(color=M_GREY, alpha=0.3)
    ax.set_title('Supplier Benchmarking — Envision vs Vestas', fontsize=14, fontweight='bold', color=M_NAVY, pad=25)
    ax.legend(loc='lower right', fontsize=11, bbox_to_anchor=(1.15, -0.05))
    plt.tight_layout()
    p = VISUALS / '07_radar.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Radar chart"); return p

def gen_peer_risk():
    fig, ax = plt.subplots(figsize=(14, 8), facecolor='white')
    companies = ['Envision Energy', 'Goldwind', 'Vestas', 'GE Vernova', 'Siemens Gamesa', 'Mingyang']
    scores = [55, 52, 22, 28, 35, 58]
    colors = [M_GREEN if s<=33 else M_AMBER if s<=66 else M_RED for s in scores]
    
    ax.axhspan(0, 33, facecolor=M_GREEN, alpha=0.07); ax.axhspan(33, 66, facecolor=M_AMBER, alpha=0.07)
    ax.axhspan(66, 100, facecolor=M_RED, alpha=0.07)
    ax.text(5.3, 16, 'LOW', fontsize=10, color=M_GREEN, fontweight='bold', alpha=0.6, ha='right')
    ax.text(5.3, 50, 'MEDIUM', fontsize=10, color=M_AMBER, fontweight='bold', alpha=0.6, ha='right')
    ax.text(5.3, 80, 'HIGH', fontsize=10, color=M_RED, fontweight='bold', alpha=0.6, ha='right')
    
    for i, (c, s, col) in enumerate(zip(companies, scores, colors)):
        ax.plot([i, i], [0, s], color=col, linewidth=3, zorder=2)
        ax.scatter(i, s, s=200, color=col, zorder=3, edgecolors=M_NAVY, linewidth=1.5)
        ax.text(i, s+3, f'{s}', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.scatter(0, 55, s=350, color=M_NAVY, zorder=4, edgecolors=M_NAVY, linewidth=2)
    ax.text(0, 58, '55', ha='center', fontsize=14, fontweight='bold', color=M_NAVY)
    
    ax.set_xticks(range(len(companies))); ax.set_xticklabels(companies, fontsize=12, fontweight='bold', rotation=15)
    ax.get_xticklabels()[0].set_color(M_NAVY)
    ax.set_ylabel('Risk Score (0–100)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Peer Risk Comparison — Wind Turbine OEMs', fontsize=14, fontweight='bold', color=M_NAVY, pad=15)
    ax.set_ylim(0, 100); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.15)
    plt.tight_layout()
    p = VISUALS / '08_peer_risk.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ Peer risk"); return p

def gen_esg():
    """ESG with 16pt controversy screening (LOCKED)."""
    fig, axes = plt.subplots(1, 3, figsize=(20, 12), facecolor='white')
    fig.suptitle('ESG Assessment — Envision Energy', fontsize=22, fontweight='bold', color=M_NAVY, y=0.98)
    fig.text(0.5, 0.94, 'Overall ESG Rating: MEDIUM  |  Strong green credentials offset by governance opacity and geopolitical risk',
             ha='center', fontsize=13, color=M_GREY)
    
    pillars = [
        {'title': 'Environmental', 'rating': 'LOW', 'color': M_GREEN, 'items': [
            ('Renewable Energy Core', M_GREEN, '✓'), ('Net-Zero Commitment', M_GREEN, '✓'),
            ('Green H₂ Investment', M_GREEN, '✓'), ('Circular Economy Prog.', M_GREEN, '✓'),
            ('Scope 1&2 Reporting', M_AMBER, '⚠'), ('Blade Recyclability', M_AMBER, '⚠')]},
        {'title': 'Social', 'rating': 'MEDIUM', 'color': M_AMBER, 'items': [
            ('Employee Safety Record', M_GREEN, '✓'), ('Community Programmes', M_GREEN, '✓'),
            ('Labour Standards', M_AMBER, '⚠'), ('Supply Chain Audits', M_AMBER, '⚠'),
            ('Forced Labour Risk', M_RED, '✗'), ('Grievance Mechanisms', M_AMBER, '⚠')]},
        {'title': 'Governance', 'rating': 'HIGH', 'color': M_RED, 'items': [
            ('Private — No Public Filings', M_RED, '✗'), ('Board Independence', M_RED, '✗'),
            ('Founder-Controlled', M_AMBER, '⚠'), ('CCP Influence Risk', M_RED, '✗'),
            ('Anti-Corruption Policy', M_AMBER, '⚠'), ('Beneficial Ownership', M_RED, '✗')]}
    ]
    
    for ax, pillar in zip(axes, pillars):
        ax.set_xlim(0, 10); ax.set_ylim(0, 12); ax.axis('off')
        title_box = FancyBboxPatch((0.5, 10.5), 9, 1.2, boxstyle="round,pad=0.15", facecolor=pillar['color'], edgecolor='none', alpha=0.9)
        ax.add_patch(title_box)
        ax.text(5, 11.1, pillar['title'], ha='center', fontsize=16, fontweight='bold', color='white')
        ax.text(5, 10.7, f"Rating: {pillar['rating']}", ha='center', fontsize=12, color='white', alpha=0.9)
        for i, (item, color, symbol) in enumerate(pillar['items']):
            y = 9.2 - i * 1.3
            ax.add_patch(plt.Circle((1.5, y), 0.35, color=color, alpha=0.8))
            ax.text(1.5, y, symbol, ha='center', va='center', fontsize=14, fontweight='bold', color='white')
            ax.text(2.5, y, item, va='center', fontsize=12, fontweight='bold', color=M_NAVY)
    
    # ── CONTROVERSY SCREENING AT 16pt (LOCKED) ──
    controversy_box = FancyBboxPatch((0.03, 0.015), 0.94, 0.105, boxstyle="round,pad=0.01",
                                       facecolor=M_RED, edgecolor=M_RED, linewidth=1.5, alpha=0.10, transform=fig.transFigure)
    fig.patches.append(controversy_box)
    fig.text(0.5, 0.098, '⚠  CONTROVERSY SCREENING', ha='center', fontsize=20, fontweight='bold', color=M_RED)
    fig.text(0.5, 0.068, '⚠  US NDAA FY2024: Proposed ban on Envision batteries in US defence systems — national security grounds',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.042, '⚠  Xinjiang supply chain: Polysilicon and rare earth sourcing linked to Uyghur forced labour regions',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.016, '⚠  Governance opacity: Private company — no public financial filings, limited board disclosure',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    
    plt.subplots_adjust(wspace=0.15, top=0.92, bottom=0.15)
    p = VISUALS / '09_esg.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ ESG assessment"); return p


# ═══════════════════════════════════════════════════════
# BUILD PPTX
# ═══════════════════════════════════════════════════════

def build(v):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    
    # ── S1: TITLE ──
    print("  Slide 1: Title")
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(*NAVY); bg.line.fill.background()
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.55), Inches(13.33), Inches(0.08))
    ln.fill.solid(); ln.fill.fore_color.rgb = rgb(*STEEL); ln.line.fill.background()
    add_text(s, 0.7, 1.4, 12, 0.5, 'SUPPLIER EVALUATION REPORT', 16, False, MID_GREY)
    add_text(s, 0.7, 2.05, 12, 1.2, 'Envision Energy', 44, True, WHITE)
    add_text(s, 0.7, 3.5, 12, 0.6, 'Wind Turbines  |  Battery Storage (AESC)  |  Energy Software (EnOS™)', 18, False, (200,210,220))
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.35), Inches(4.5), Inches(0.05))
    div.fill.solid(); div.fill.fore_color.rgb = rgb(*STEEL); div.line.fill.background()
    add_text(s, 0.7, 4.55, 12, 0.4, '~$8B Revenue (est.)  |  6,000+ Employees  |  20+ Countries  |  30+ GW Installed', 14, False, (200,210,220))
    add_text(s, 0.7, 6.2, 12, 0.4, 'Confidential  |  March 2026  |  Manu Forti Intelligence', 11, False, MID_GREY)
    add_logos(s, is_title=True)
    
    # ── S2: EXEC SUMMARY ──
    print("  Slide 2: Executive Summary")
    s = prs.slides.add_slide(blank)
    add_header(s, 'EXECUTIVE SUMMARY', 'Envision Energy — Risk & Suitability Overview')
    s.shapes.add_picture(str(v['gauge']), Inches(0.3), Inches(1.05), Inches(7.2), Inches(3.1))
    add_rounded_rect(s, 7.6, 1.05, 5.4, 5.6, (240,245,250))
    add_text(s, 7.8, 1.22, 5.2, 0.32, 'SUPPLIER SNAPSHOT', 14, True, NAVY)
    snap = [('Supplier:', 'Envision Energy (远景能源)'), ('Type:', 'Private (founder-controlled)'),
            ('Sector:', 'Wind Turbines / EV Batteries / Software'), ('HQ:', 'Shanghai, China'),
            ('Founded:', '2007 by Lei Zhang'), ('Revenue:', '~$8B est. (2024)'),
            ('EBITDA:', '~$960M est. (12% margin)'), ('Order Book:', '~$6.5B est. wind+battery'),
            ('Employees:', '6,000+ globally'), ('Operations:', '20+ countries, R&D in 6 countries')]
    for idx, (label, val) in enumerate(snap):
        y = 1.64 + idx * 0.30
        add_text(s, 7.8, y, 2.3, 0.27, label, 11, True, TEXT_SECONDARY)
        add_text(s, 10.1, y, 2.7, 0.27, val, 11, False, TEXT_PRIMARY)
    add_rounded_rect(s, 0.3, 5.22, 7.2, 1.43, (240,245,250))
    add_text(s, 0.5, 5.33, 7, 0.3, 'Key Findings', 13, True, NAVY)
    findings = [
        '• #4 global WTG OEM (2024) — strong China position, expanding internationally',
        '• Envision AESC battery division: partnerships with Renault, Nissan, BMW — gigafactory pipeline',
        '• EnOS™ AIoT platform manages 620+ GWh globally — significant digital energy moat',
        '• CRITICAL: US NDAA restrictions, Xinjiang supply chain risk, zero public financial disclosure']
    for idx, f in enumerate(findings):
        add_text(s, 0.5, 5.65+idx*0.27, 7, 0.26, f, 11, False, TEXT_PRIMARY)
    add_source(s); add_logos(s)
    
    # ── S3: RECOMMENDATION ──
    print("  Slide 3: Recommendation")
    s = prs.slides.add_slide(blank)
    add_header(s, 'RECOMMENDATION', 'Decision Summary & Commercial Conditions')
    add_rounded_rect(s, 0.5, 1.1, 12.33, 1.35, AMBER)
    add_text(s, 0.7, 1.18, 12, 0.55, '⚠  RECOMMENDATION: CONDITIONAL APPROVAL — ENHANCED DUE DILIGENCE REQUIRED', 20, True, WHITE)
    add_text(s, 0.7, 1.78, 12, 0.58,
             'Envision is a technically capable WTG and battery supplier with competitive pricing. However, significant geopolitical risk (US-China tensions, '
             'NDAA restrictions), governance opacity (private, no public filings), and Xinjiang supply chain concerns require enhanced due diligence before engagement.',
             11, False, WHITE)
    add_bullets(s, 0.5, 2.65, 6, 2.5, 'Commercial Conditions',
                ['Payment: Letter of Credit required (no open account)',
                 'Performance bonds: 15% of contract value (above standard)',
                 'Force-majeure: Explicit geopolitical/sanctions clause',
                 'Independent quality inspection at factory (pre-shipment)',
                 'Local entity contracting where possible (not PRC parent)'], 14, 11)
    add_bullets(s, 6.7, 2.65, 6, 2.5, 'Enhanced Due Diligence (Mandatory)',
                ['Supply chain audit: Xinjiang/forced labour attestation',
                 'Sanctions screening: OFAC, EU, UK — ongoing monitoring',
                 'IP protection: Ring-fenced technology sharing agreements',
                 'NDAA compliance check for any US-funded projects',
                 'Third-party ESG audit (EcoVadis or equivalent required)',
                 'Annual governance disclosure as contract condition'], 14, 11)
    add_rounded_rect(s, 0.5, 5.35, 12.33, 1.5, (240,245,250))
    add_text(s, 0.7, 5.5, 12, 0.3, 'Overall Risk Summary — MEDIUM (55/100)', 14, True, NAVY)
    for x, t, d in [(0.6,'Financial: MEDIUM','Private; no audited accounts; estimated ~$8B revenue'),
                     (3.7,'Operational: LOW','30+ GW installed; proven technology; global R&D'),
                     (6.8,'Geopolitical: HIGH','US-China tensions; NDAA ban; sanctions risk'),
                     (9.9,'ESG: MEDIUM','Green core but Xinjiang supply chain + governance gaps')]:
        add_text(s, x, 5.88, 3.1, 0.27, t, 11, True, NAVY)
        add_text(s, x, 6.17, 3.1, 0.27, d, 9, False, TEXT_SECONDARY)
    add_source(s); add_logos(s)
    
    # ── S4: SUPPLIER PROFILE ──
    print("  Slide 4: Supplier Profile")
    s = prs.slides.add_slide(blank)
    add_header(s, 'SUPPLIER PROFILE', 'Corporate Structure & Global Footprint')
    s.shapes.add_picture(str(v['org']), Inches(0.3), Inches(1.05), Inches(8.2), Inches(4.48))
    add_bullets(s, 8.7, 1.05, 4.3, 3.5, 'Company Overview',
                ['Envision Energy is a Chinese clean energy company',
                 'founded in 2007 by Lei Zhang in Jiangyin, Jiangsu.',
                 'Core business: onshore/offshore wind turbines, with',
                 'expansion into EV batteries (Envision AESC, acquired',
                 'from Nissan 2018) and digital energy management',
                 '(EnOS™ AIoT platform managing 620+ GWh).',
                 '',
                 'Privately held — Lei Zhang retains majority control.',
                 'R&D centres in Shanghai, Denmark, Japan, USA.',
                 'Battery gigafactories in China, Japan, UK, France.'], 14, 11)
    add_bullets(s, 8.7, 4.5, 4.3, 2.4, 'Leadership',
                ['Lei Zhang — Founder & CEO',
                 '  Founded 2007; visionary in smart energy',
                 '',
                 'Key: Single founder-CEO controls strategy',
                 '  No independent board disclosure',
                 '',
                 'HQ: Shanghai, CN  |  Founded: 2007'], 14, 11)
    add_source(s); add_logos(s)
    
    # ── S5: FINANCIAL HEALTH ──
    print("  Slide 5: Financial Health")
    s = prs.slides.add_slide(blank)
    add_header(s, 'FINANCIAL HEALTH', 'Revenue & EBITDA Trajectory 2021–2024 (Estimated)  |  Private Company — No Public Filings')
    s.shapes.add_picture(str(v['fin']), Inches(0.3), Inches(1.05), Inches(7.5), Inches(4.55))
    add_text(s, 8.1, 1.1, 4.9, 0.35, 'Financial Highlights', 14, True, NAVY)
    metrics = [('2024 Revenue (est.):', '~$8.0B  (+13% YoY)'), ('2024 EBITDA (est.):', '~$960M  (12% margin)'),
               ('2024 Net Profit (est.):', '~$480M'), ('3yr Revenue CAGR:', '~24%  (2021–2024)'),
               ('Order Book (est.):', '~$6.5B  (wind + battery)'), ('Gross Debt (est.):', '~$3.2B  (bond market data)'),
               ('Net Cash Position:', 'UNKNOWN  (no disclosure)'), ('Debt / EBITDA (est.):', '~3.3x  ⚠  (elevated)')]
    for idx, (label, val) in enumerate(metrics):
        y = 1.56 + idx * 0.30
        add_text(s, 8.1, y, 2.55, 0.29, label, 11, True, TEXT_SECONDARY)
        vc = RED if idx >= 6 else TEXT_PRIMARY
        add_text(s, 10.65, y, 2.25, 0.29, val, 11, False, vc)
    add_rounded_rect(s, 8.1, 4.05, 4.8, 0.44, AMBER)
    add_text(s, 8.25, 4.12, 4.5, 0.3, 'Financial Risk:  MEDIUM  ⚠', 12, True, WHITE, PP_ALIGN.CENTER)
    add_bullets(s, 8.1, 4.6, 4.8, 2.3, 'Exposure Guidance',
                ['Private company — ALL financials are estimates',
                 'No audited accounts available for verification',
                 'Debt/EBITDA ~3.3x — above comfortable threshold',
                 'Letter of Credit mandatory; no open account terms'], 12, 10)
    add_source(s); add_logos(s)
    
    # ── S6: MARKET POSITION ──
    print("  Slide 6: Market Position")
    s = prs.slides.add_slide(blank)
    add_header(s, 'MARKET POSITION', 'Global Wind Turbine OEM Ranking — 2024 Installations (GW)')
    s.shapes.add_picture(str(v['mkt']), Inches(0.3), Inches(1.05), Inches(8.0), Inches(4.46))
    add_bullets(s, 8.5, 1.1, 4.5, 2.0, 'Competitive Landscape',
                ['Envision is the #4 global WTG OEM and China\'s',
                 '#2 after Goldwind. Strong in onshore; growing',
                 'offshore capability. AESC battery division adds',
                 'a unique energy storage + automotive dimension.'], 14, 11)
    add_bullets(s, 8.5, 3.2, 4.5, 2.6, 'Envision Competitive Advantages',
                ['Price: 15-30% below Western OEMs on onshore WTGs',
                 'EnOS™ platform: 620+ GWh managed — digital moat',
                 'AESC battery partnerships: Renault, Nissan, BMW',
                 'Vertical integration: turbines + batteries + software',
                 'Green hydrogen R&D: Boston centre (2022)',
                 'Aggressive international expansion pipeline'], 14, 11)
    add_text(s, 0.5, 5.5, 12, 0.3, 'Key Competitors', 13, True, NAVY)
    for idx, c in enumerate([
        '• Vestas — Denmark — 17.5 GW — Global #1; premium pricing; strongest service network',
        '• Goldwind — China — 16.0 GW — China\'s #1; similar geopolitical risk profile',
        '• GE Vernova — USA — 8.5 GW — Restructured; Haliade-X offshore leader',
        '• Siemens Gamesa — Spain/Germany — 7.2 GW — Offshore dominant; quality issues recently']):
        add_text(s, 0.6, 5.84+idx*0.29, 12, 0.27, c, 10, False, TEXT_PRIMARY)
    add_source(s); add_logos(s)
    
    # ── S7: OPS & RISK ──
    print("  Slide 7: Ops & Risk")
    s = prs.slides.add_slide(blank)
    add_header(s, 'OPERATIONAL CAPABILITY  &  RISK ASSESSMENT', 'Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary')
    s.shapes.add_picture(str(v['timeline']), Inches(0.3), Inches(1.05), Inches(7.0), Inches(2.92))
    s.shapes.add_picture(str(v['risk']), Inches(7.65), Inches(1.05), Inches(5.5), Inches(4.46))
    add_text(s, 0.3, 5.58, 7, 0.28, 'Risk Summary', 13, True, NAVY)
    for idx, (cat, rating, desc) in enumerate([
        ('Financial:', '🟡 MEDIUM', 'Private; no audited accounts; ~3.3x leverage'),
        ('Geopolitical:', '🔴 HIGH', 'US-China tensions; NDAA battery ban proposed'),
        ('IP / Technology:', '🟡 MEDIUM', 'Competitive tech but IP transfer concerns'),
        ('Supply Chain:', '🟡 MEDIUM', 'Xinjiang rare earth / polysilicon exposure'),
        ('Delivery:', '🟢 LOW', '30+ GW installed; proven execution track record')]):
        y = 5.89 + idx * 0.25
        add_text(s, 0.3, y, 2.1, 0.24, cat, 10, True, TEXT_PRIMARY)
        add_text(s, 2.4, y, 1.4, 0.24, rating, 10, True, TEXT_PRIMARY)
        add_text(s, 3.8, y, 3.7, 0.24, desc, 9, False, TEXT_SECONDARY)
    add_bullets(s, 7.65, 5.58, 5.5, 1.85, 'Key Capabilities',
                ['30+ GW wind installed globally (onshore + offshore)',
                 'EnOS™ AIoT: world\'s largest energy IoT platform',
                 'AESC gigafactories: China, Japan, UK, France, Spain',
                 'R&D: Shanghai, Silkeborg (DK), Boulder (US), Osaka (JP)'], 12, 10)
    add_source(s); add_logos(s)
    
    # ── S8: COMMERCIAL INTEL ──
    print("  Slide 8: Commercial Intelligence")
    s = prs.slides.add_slide(blank)
    add_header(s, 'COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON', 'Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms')
    s.shapes.add_picture(str(v['radar']), Inches(0.2), Inches(1.0), Inches(5.8), Inches(4.74))
    s.shapes.add_picture(str(v['peer']), Inches(6.2), Inches(1.0), Inches(6.9), Inches(3.85))
    add_rounded_rect(s, 0.3, 5.8, 12.7, 1.5, (240,245,250))
    add_bullets(s, 0.5, 5.85, 6.5, 1.4, 'Commercial Terms & Negotiation',
                ['Pricing: 15-30% below Vestas/Siemens on comparable capacity',
                 'Lead time: 6–12 months (onshore); 12–24 months (offshore)',
                 'Payment: Letter of Credit mandatory — no open account',
                 'Leverage: Multi-turbine order discounts; bundled AESC battery deals'], 12, 10)
    add_bullets(s, 7.0, 5.85, 5.8, 1.4, 'Key Watch Points',
                ['⚠ CRITICAL: Sanctions / NDAA — check project-by-project',
                 '⚠ Quality: Independent inspection mandatory pre-shipment',
                 '⚠ IP: Ring-fence any shared technical specifications',
                 '⚠ FX: RMB exposure — hedge on contracts >$10M'], 12, 10)
    add_source(s); add_logos(s)
    
    # ── S9: ESG ──
    print("  Slide 9: ESG Assessment")
    s = prs.slides.add_slide(blank)
    add_header(s, 'ESG ASSESSMENT', 'Environmental, Social & Governance Screening  |  Overall Rating: MEDIUM')
    s.shapes.add_picture(str(v['esg']), Inches(1.67), Inches(1.05), Inches(10.0), Inches(5.93))
    add_source(s); add_logos(s)
    
    prs.save(str(OUTPUT))
    print(f"\n✅ Saved: {OUTPUT}")
    import os
    print(f"📎 Size: {os.path.getsize(str(OUTPUT))/1024/1024:.1f} MB")


# ═══════════════════════════════════════════════════════
if __name__ == '__main__':
    print("🎨 Generating Envision Energy Product 1 — Full Infographics\n")
    print("Step 1: Charts...")
    visuals = {
        'gauge': gen_risk_gauge(), 'org': gen_org_chart(), 'fin': gen_financial(),
        'mkt': gen_market(), 'timeline': gen_timeline(), 'risk': gen_risk_matrix(),
        'radar': gen_radar(), 'peer': gen_peer_risk(), 'esg': gen_esg(),
    }
    print("\nStep 2: Building PPTX...")
    build(visuals)
