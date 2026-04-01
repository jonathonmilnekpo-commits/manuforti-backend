#!/usr/bin/env python3
"""Generate ControlPartner AS Product 1 v15 — Second Draft using Jarotech as reference."""

import sys
sys.path.insert(0, '/opt/homebrew/lib/python3.14/site-packages')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Branding per Jarotech reference
NAVY = (0, 33, 71)
STEEL = (43, 108, 176)
MID_GREY = (113, 128, 150)
GREEN = (72, 187, 120)
AMBER = (214, 158, 46)
RED = (229, 62, 62)
WHITE = (255, 255, 255)
TEXT_PRIMARY = (26, 32, 44)
TEXT_SECONDARY = (74, 85, 104)

M_NAVY = '#002147'
M_STEEL = '#2B6CB0'
M_GREY = '#718096'
M_GREEN = '#48BB78'
M_AMBER = '#D69E2E'
M_RED = '#E53E3E'
M_LIGHT = '#EBF4FF'

BASE = Path('/Users/jonathonmilne/.openclaw/workspace')
VISUALS = BASE / 'controlpartner_v2_visuals'
VISUALS.mkdir(exist_ok=True)
CLAN_LOGO = BASE / 'skills/product-1-generator/assets/manu_forti_logo.png'
OUTPUT = BASE / 'ControlPartner_Product1_v15_SecondDraft.pptx'

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_header(slide, title, subtitle):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    h.fill.solid(); h.fill.fore_color.rgb = rgb(*NAVY); h.line.fill.background()
    
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = rgb(*WHITE)
    
    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = subtitle; r2.font.size = Pt(12); r2.font.color.rgb = rgb(200, 210, 220)

def add_source(slide):
    t = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = 'Source: Manu Forti Intelligence  |  Confidential  |  March 2026'
    r.font.size = Pt(8); r.font.color.rgb = rgb(*MID_GREY)

def add_text(slide, l, t, w, h, text, sz=14, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = rgb(*color)
    return tb

def add_bullets(slide, l, t, w, h, title, items, tsz=14, isz=11, tc=NAVY, ic=TEXT_PRIMARY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True; r.font.color.rgb = rgb(*tc)
    for item in items:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = f'• {item}'; r.font.size = Pt(isz); r.font.color.rgb = rgb(*ic)

def add_rounded_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(*color); s.line.fill.background()
    return s

def add_logos(slide, is_title=False):
    cimg = Image.open(str(CLAN_LOGO))
    ca = cimg.size[0] / cimg.size[1]
    cw = 0.6; ch = cw / ca
    cx = 13.333 - cw - 0.2
    cy = 7.5 - ch - 0.15
    slide.shapes.add_picture(str(CLAN_LOGO), Inches(cx), Inches(cy), Inches(cw), Inches(ch))

def gen_risk_gauge():
    fig, ax = plt.subplots(figsize=(14, 6.5), facecolor='white')
    segments = [(0, 60, M_GREEN, 'LOW\n0–33'), (60, 120, M_AMBER, 'MEDIUM\n34–66'), (120, 180, M_RED, 'HIGH\n67–100')]
    for start, end, color, label in segments:
        theta = np.linspace(np.radians(start), np.radians(end), 100)
        x_o = np.cos(theta); y_o = np.sin(theta)
        x_i = np.cos(theta[::-1]) * 0.6; y_i = np.sin(theta[::-1]) * 0.6
        ax.fill(np.concatenate([x_o, x_i]), np.concatenate([y_o, y_i]), color=color, alpha=0.85)
        mid = np.radians((start+end)/2)
        ax.text(np.cos(mid)*0.8, np.sin(mid)*0.8, label, ha='center', va='center', fontsize=11, fontweight='bold', color='white')
    
    score = 52
    needle = np.radians(score * 180 / 100)
    ax.annotate('', xy=(np.cos(needle)*0.55, np.sin(needle)*0.55), xytext=(0,0),
                arrowprops=dict(arrowstyle='->', color=M_NAVY, lw=3.5))
    ax.add_patch(plt.Circle((0,0), 0.08, color=M_NAVY, zorder=5))
    ax.text(0, -0.2, f'{score}/100', ha='center', fontsize=28, fontweight='bold', color=M_NAVY)
    ax.text(0, -0.35, 'OVERALL RISK SCORE', ha='center', fontsize=12, color=M_GREY)
    ax.text(0, 1.25, 'ControlPartner AS — Overall Risk Assessment', ha='center', fontsize=18, fontweight='bold', color=M_NAVY)
    ax.text(0, 1.10, 'Industrial Automation | SCADA Systems | Water/Wastewater', ha='center', fontsize=12, color=M_GREY)
    
    boxes = [('Financial', 'MEDIUM', M_AMBER), ('Operational', 'LOW', M_GREEN), ('Geopolitical', 'LOW', M_GREEN), ('ESG', 'MEDIUM', M_AMBER)]
    for idx, (cat, rating, color) in enumerate(boxes):
        bx = -0.75 + idx * 0.5
        rect = FancyBboxPatch((bx-0.2, -0.67), 0.4, 0.22, boxstyle="round,pad=0.03", facecolor=color, alpha=0.15, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(bx, -0.54, cat, ha='center', fontsize=8, fontweight='bold', color=M_NAVY)
        ax.text(bx, -0.62, rating, ha='center', fontsize=8, fontweight='bold', color=color)
    
    ax.set_xlim(-1.4, 1.4); ax.set_ylim(-0.75, 1.4); ax.set_aspect('equal'); ax.axis('off')
    p = VISUALS / '01_risk_gauge.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_org_chart():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    parent = FancyBboxPatch((5.5, 7.5), 5, 1.2, boxstyle="round,pad=0.15", facecolor=M_NAVY, edgecolor='none')
    ax.add_patch(parent)
    ax.text(8, 8.1, 'Uniwater AB', ha='center', fontsize=18, fontweight='bold', color='white')
    ax.text(8, 7.75, 'Sweden | Private Equity-backed | Water Tech Focus', ha='center', fontsize=10, color='#cccccc')
    
    divs = [('ControlPartner\nAS', 'Industrial\nAutomation', 1.0), ('Other Uniwater\nSubsidiaries', 'Water Tech\nPortfolio', 4.5),
            ('Uniwater\nGroup', 'Nordic Water\nSolutions', 8.0), ('Future\nAcquisitions', 'Growth via\nM&A', 11.5)]
    for label, desc, x in divs:
        box = FancyBboxPatch((x, 4.8), 3, 1.6, boxstyle="round,pad=0.1", facecolor=M_STEEL, edgecolor='none', alpha=0.9)
        ax.add_patch(box)
        ax.text(x+1.5, 5.85, label, ha='center', fontsize=12, fontweight='bold', color='white')
        ax.text(x+1.5, 5.2, desc, ha='center', fontsize=9, color='#e0e0e0')
        ax.plot([8, x+1.5], [7.5, 6.4], color=M_GREY, lw=1.5, alpha=0.6)
    
    facts = [('NOK 57.3M\nRevenue 2023', 1.5, 2.8), ('40-46\nEmployees', 5.0, 2.8),
             ('Porsgrunn\nHQ Norway', 8.5, 2.8), ('Oct 2025\nUniwater Acq.', 12.0, 2.8)]
    for text, x, y in facts:
        box = FancyBboxPatch((x, y), 2.5, 1.2, boxstyle="round,pad=0.08", facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1)
        ax.add_patch(box)
        ax.text(x+1.25, y+0.6, text, ha='center', fontsize=11, fontweight='bold', color=M_NAVY)
    
    ax.set_xlim(0, 16); ax.set_ylim(2, 9.5); ax.axis('off')
    p = VISUALS / '02_org_chart.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_financial():
    fig, ax1 = plt.subplots(figsize=(15, 9), facecolor='white')
    years = ['2021', '2022', '2023']
    revenue = [52.1, 54.8, 57.3]
    ebitda = [4.8, 5.0, 5.2]
    margin = [9.2, 9.1, 9.1]
    
    bars = ax1.bar(years, revenue, color=M_NAVY, width=0.5, zorder=3, label='Revenue (NOK M)')
    for bar, val in zip(bars, revenue):
        ax1.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1, f'NOK {val:.1f}M', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax1.set_ylabel('Revenue (NOK Millions)', fontsize=13, color=M_NAVY, fontweight='bold')
    ax1.set_ylim(0, 80); ax1.tick_params(axis='y', labelcolor=M_NAVY, labelsize=11)
    ax1.tick_params(axis='x', labelsize=13)
    
    ax2 = ax1.twinx()
    ax2.plot(years, ebitda, color=M_STEEL, marker='o', markersize=10, linewidth=3, zorder=4, label='EBITDA (NOK M)')
    for x, y, m in zip(years, ebitda, margin):
        ax2.text(x, y+0.3, f'NOK {y:.1f}M\n({m:.1f}%)', ha='center', fontsize=11, fontweight='bold', color=M_STEEL)
    ax2.set_ylabel('EBITDA (NOK Millions)', fontsize=13, color=M_STEEL, fontweight='bold')
    ax2.set_ylim(0, 8); ax2.tick_params(axis='y', labelcolor=M_STEEL, labelsize=11)
    
    ax1.spines['top'].set_visible(False); ax2.spines['top'].set_visible(False)
    ax1.grid(axis='y', alpha=0.2, linestyle='--')
    ax1.set_title('Revenue & EBITDA Trajectory 2021–2023 (Pre-Acquisition)', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    
    props = dict(boxstyle='round,pad=0.5', facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1.5)
    ax1.text(0.02, 0.97, 'Steady growth trajectory pre-Uniwater acquisition.\nOct 2025: Acquired by Uniwater AB — now part of\nNordic water technology group. Limited public data.',
             transform=ax1.transAxes, fontsize=11, va='top', bbox=props, color=M_NAVY)
    
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, loc='upper right', fontsize=11)
    plt.tight_layout()
    p = VISUALS / '03_financial.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_market():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    companies = ['Siemens Energy', 'ABB', 'Schneider Electric', 'ControlPartner AS', 'Emerson', 'Rockwell']
    revenue = [2800, 3200, 3500, 57.3, 15000, 7500]
    colors = [M_GREY, M_GREY, M_GREY, M_NAVY, M_GREY, M_GREY]
    
    bars = ax.barh(range(len(companies)), revenue, color=colors, height=0.55, zorder=3)
    for bar, val, c in zip(bars, revenue, companies):
        lc = 'white' if c == 'ControlPartner AS' else M_NAVY
        label = f'NOK {val:.1f}M' if val < 1000 else f'{val/1000:.1f}B'
        ax.text(val-50, bar.get_y()+bar.get_height()/2, label, ha='right', va='center', fontsize=13, fontweight='bold', color=lc)
    
    ax.set_yticks(range(len(companies))); ax.set_yticklabels(companies, fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    ax.set_xlabel('Estimated Revenue (NOK Millions)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Industrial Automation & SCADA Market Position — 2023', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    ax.get_yticklabels()[3].set_color(M_NAVY)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.annotate('ControlPartner: Niche player\nFocused on water/wastewater', xy=(57.3, 3), xytext=(800, 1.5),
                fontsize=11, fontweight='bold', color=M_STEEL, arrowprops=dict(arrowstyle='->', color=M_STEEL, lw=2))
    plt.tight_layout()
    p = VISUALS / '04_market.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_timeline():
    fig, ax = plt.subplots(figsize=(14, 6), facecolor='white')
    milestones = [('2015', 'Founded in\nPorsgrunn', M_NAVY), ('2017', 'First Major\nSCADA Project', M_STEEL),
                  ('2019', 'Bergen Office\nOpened', M_GREEN), ('2021', 'Oslo Office\nExpansion', M_NAVY),
                  ('2023', 'NOK 57.3M\nRevenue Peak', M_STEEL), ('Oct 2025', 'Uniwater\nAcquisition', M_GREEN)]
    ax.plot([0, len(milestones)-1], [0, 0], color=M_NAVY, linewidth=3, zorder=1)
    for i, (year, desc, color) in enumerate(milestones):
        ax.plot(i, 0, 'o', markersize=16, color=color, zorder=3)
        ax.plot(i, 0, 'o', markersize=10, color='white', zorder=4)
        y_off = 0.6 if i % 2 == 0 else -0.6
        ax.text(i, y_off, f'{year}\n{desc}', ha='center', va='bottom' if i%2==0 else 'top',
                fontsize=10, fontweight='bold', color=M_NAVY,
                bbox=dict(boxstyle='round,pad=0.3', facecolor=M_LIGHT, edgecolor=color, linewidth=1.5))
        ax.plot([i, i], [0, y_off*0.4], color=color, linewidth=1.5, linestyle='--')
    ax.set_title('Strategic Investment Timeline', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(-0.5, len(milestones)-0.5); ax.set_ylim(-1.4, 1.4); ax.axis('off')
    plt.tight_layout()
    p = VISUALS / '05_timeline.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_risk_matrix():
    fig, ax = plt.subplots(figsize=(11, 9), facecolor='white')
    ax.fill_between([0,5], [5,5], [10,10], color=M_AMBER, alpha=0.10)
    ax.fill_between([5,10], [5,5], [10,10], color=M_RED, alpha=0.10)
    ax.fill_between([0,5], [0,0], [5,5], color=M_GREEN, alpha=0.10)
    ax.fill_between([5,10], [0,0], [5,5], color=M_AMBER, alpha=0.10)
    ax.text(2.5, 7.5, 'MONITOR', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    ax.text(7.5, 7.5, 'CRITICAL', ha='center', fontsize=14, fontweight='bold', color=M_RED, alpha=0.4)
    ax.text(2.5, 2.5, 'ACCEPT', ha='center', fontsize=14, fontweight='bold', color=M_GREEN, alpha=0.4)
    ax.text(7.5, 2.5, 'MITIGATE', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    
    risks = [('Uniwater\nIntegration', 4.5, 5.5, 700, M_AMBER), ('Limited\nPublic Data', 6.0, 4.0, 600, M_AMBER),
             ('Customer\nConcentration', 5.0, 5.5, 650, M_AMBER), ('Small Scale\n(40-46 emp)', 3.5, 4.0, 500, M_GREEN),
             ('Norwegian\nMarket Focus', 2.5, 3.0, 400, M_GREEN)]
    for label, x, y, size, color in risks:
        ax.scatter(x, y, s=size, color=color, alpha=0.7, edgecolors=M_NAVY, linewidth=1.5, zorder=3)
        ax.text(x, y, label, ha='center', va='center', fontsize=8, fontweight='bold', color=M_NAVY)
    
    ax.set_xlabel('Probability →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_ylabel('Impact →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Risk Matrix — ControlPartner AS', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.set_xticks([2.5, 7.5]); ax.set_xticklabels(['Low', 'High'], fontsize=11)
    ax.set_yticks([2.5, 7.5]); ax.set_yticklabels(['Low', 'High'], fontsize=11)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    ax.axhline(y=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    ax.axvline(x=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    plt.tight_layout()
    p = VISUALS / '06_risk_matrix.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_radar():
    fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(projection='polar'), facecolor='white')
    cats = ['Price\nCompetitiveness', 'Lead Time', 'Quality &\nReliability', 'Technical\nCapability', 'Supply Chain\nResilience', 'Service\nLevel']
    N = len(cats)
    controlpartner = [7, 7, 8, 8, 6, 8]
    abb = [6, 8, 9, 9, 9, 7]
    angles = np.linspace(0, 2*np.pi, N, endpoint=False).tolist()
    cp_c = controlpartner + [controlpartner[0]]; abb_c = abb + [abb[0]]; ang_c = angles + [angles[0]]
    ax.fill(ang_c, cp_c, color=M_NAVY, alpha=0.2)
    ax.plot(ang_c, cp_c, color=M_NAVY, linewidth=2.5, marker='o', markersize=8, label='ControlPartner AS')
    ax.fill(ang_c, abb_c, color=M_STEEL, alpha=0.1)
    ax.plot(ang_c, abb_c, color=M_STEEL, linewidth=2, marker='s', markersize=7, linestyle='--', label='ABB (Benchmark)')
    ax.set_xticks(angles); ax.set_xticklabels(cats, fontsize=11, fontweight='bold', color=M_NAVY)
    ax.set_ylim(0, 10); ax.set_yticks([2,4,6,8,10])
    ax.set_yticklabels(['2','4','6','8','10'], fontsize=9, color=M_GREY)
    ax.grid(color=M_GREY, alpha=0.3)
    ax.set_title('Supplier Benchmarking — ControlPartner vs ABB', fontsize=14, fontweight='bold', color=M_NAVY, pad=25)
    ax.legend(loc='lower right', fontsize=11, bbox_to_anchor=(1.15, -0.05))
    plt.tight_layout()
    p = VISUALS / '07_radar.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_peer_risk():
    fig, ax = plt.subplots(figsize=(14, 8), facecolor='white')
    companies = ['ControlPartner AS', 'ABB Norway', 'Siemens Norway', 'Schneider NO', 'Emerson NO']
    scores = [52, 28, 25, 30, 35]
    colors = [M_GREEN if s<=33 else M_AMBER if s<=66 else M_RED for s in scores]
    
    ax.axhspan(0, 33, facecolor=M_GREEN, alpha=0.07); ax.axhspan(33, 66, facecolor=M_AMBER, alpha=0.07)
    ax.axhspan(66, 100, facecolor=M_RED, alpha=0.07)
    ax.text(4.3, 16, 'LOW', fontsize=10, color=M_GREEN, fontweight='bold', alpha=0.6, ha='right')
    ax.text(4.3, 50, 'MEDIUM', fontsize=10, color=M_AMBER, fontweight='bold', alpha=0.6, ha='right')
    ax.text(4.3, 80, 'HIGH', fontsize=10, color=M_RED, fontweight='bold', alpha=0.6, ha='right')
    
    for i, (c, s, col) in enumerate(zip(companies, scores, colors)):
        ax.plot([i, i], [0, s], color=col, linewidth=3, zorder=2)
        ax.scatter(i, s, s=200, color=col, zorder=3, edgecolors=M_NAVY, linewidth=1.5)
        ax.text(i, s+3, f'{s}', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.scatter(0, 52, s=350, color=M_NAVY, zorder=4, edgecolors=M_NAVY, linewidth=2)
    ax.text(0, 55, '52', ha='center', fontsize=14, fontweight='bold', color=M_NAVY)
    
    ax.set_xticks(range(len(companies))); ax.set_xticklabels(companies, fontsize=12, fontweight='bold', rotation=15)
    ax.get_xticklabels()[0].set_color(M_NAVY)
    ax.set_ylabel('Risk Score (0–100)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Peer Risk Comparison — Industrial Automation (Norway)', fontsize=14, fontweight='bold', color=M_NAVY, pad=15)
    ax.set_ylim(0, 100); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.15)
    plt.tight_layout()
    p = VISUALS / '08_peer_risk.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_esg():
    fig, axes = plt.subplots(1, 3, figsize=(20, 12), facecolor='white')
    fig.suptitle('ESG Assessment — ControlPartner AS', fontsize=22, fontweight='bold', color=M_NAVY, y=0.98)
    fig.text(0.5, 0.94, 'Overall ESG Rating: MEDIUM | Strong operational ESG offset by limited disclosure as private subsidiary',
             ha='center', fontsize=13, color=M_GREY)
    
    pillars = [
        {'title': 'Environmental', 'rating': 'LOW', 'color': M_GREEN, 'items': [
            ('Water Efficiency Focus', M_GREEN, '✓'), ('Energy Optimization', M_GREEN, '✓'),
            ('SCADA for Sustainability', M_GREEN, '✓'), ('ISO 14001 Certified', M_GREEN, '✓'),
            ('Scope 1&2 Tracking', M_AMBER, '⚠'), ('Net Zero Commitment', M_AMBER, '⚠')]},
        {'title': 'Social', 'rating': 'MEDIUM', 'color': M_AMBER, 'items': [
            ('Norwegian Labour Standards', M_GREEN, '✓'), ('Employee Safety', M_GREEN, '✓'),
            ('Local Hiring (Porsgrunn)', M_GREEN, '✓'), ('Uniwater Group Policies', M_AMBER, '⚠'),
            ('Limited Public Disclosure', M_AMBER, '⚠'), ('Grievance Mechanisms', M_AMBER, '⚠')]},
        {'title': 'Governance', 'rating': 'MEDIUM', 'color': M_AMBER, 'items': [
            ('Private Subsidiary', M_AMBER, '⚠'), ('No Public Filings', M_AMBER, '⚠'),
            ('Uniwater Oversight', M_GREEN, '✓'), ('Board Transparency', M_AMBER, '⚠'),
            ('Anti-Corruption Policy', M_GREEN, '✓'), ('Beneficial Ownership', M_AMBER, '⚠')]}
    ]
    
    for ax, pillar in zip(axes, pillars):
        ax.set_xlim(0, 10); ax.set_ylim(0, 12); ax.axis('off')
        title_box = FancyBboxPatch((0.5, 10.5), 9, 1.2, boxstyle="round,pad=0.15", facecolor=pillar['color'], edgecolor='none', alpha=0.9)
        ax.add_patch(title_box)
        ax.text(5, 11.1, pillar['title'], ha='center', fontsize=16, fontweight='bold', color='white')
        ax.text(5, 10.7, f"Rating: {pillar['rating']}", ha='center', fontsize=12, color='white', alpha=0.9)
        for i, (item, color, symbol) in enumerate(pillar['items']):
            y = 9.2 - i * 1.3
            ax.add_patch(plt.Circle((1.5, y), 0.35, color=color, alpha=0.8))
            ax.text(1.5, y, symbol, ha='center', va='center', fontsize=14, fontweight='bold', color='white')
            ax.text(2.5, y, item, va='center', fontsize=12, fontweight='bold', color=M_NAVY)
    
    controversy_box = FancyBboxPatch((0.03, 0.015), 0.94, 0.105, boxstyle="round,pad=0.01",
                                       facecolor=M_RED, edgecolor=M_RED, linewidth=1.5, alpha=0.10, transform=fig.transFigure)
    fig.patches.append(controversy_box)
    fig.text(0.5, 0.098, '⚠  CONTROVERSY SCREENING', ha='center', fontsize=20, fontweight='bold', color=M_RED)
    fig.text(0.5, 0.068, '⚠  Limited public disclosure: Private company — no audited financials publicly available',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.042, '⚠  Uniwater acquisition: Recent ownership change (Oct 2025) — governance integration ongoing',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.016, '⚠  Data limitations: 2023 financials only — pre-acquisition baseline; limited forward visibility',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    
    plt.subplots_adjust(wspace=0.15, top=0.92, bottom=0.15)
    p = VISUALS / '09_esg.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

print("🎨 Generating ControlPartner AS Product 1 v15 — Second Draft\n")
print("Step 1: Charts...")
visuals = {
    'gauge': gen_risk_gauge(), 'org': gen_org_chart(), 'fin': gen_financial(),
    'mkt': gen_market(), 'timeline': gen_timeline(), 'risk': gen_risk_matrix(),
    'radar': gen_radar(), 'peer': gen_peer_risk(), 'esg': gen_esg(),
}
print("\nStep 2: Building PPTX...")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# S1: TITLE
print("  Slide 1: Title")
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = rgb(*NAVY); bg.line.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.55), Inches(13.33), Inches(0.08))
ln.fill.solid(); ln.fill.fore_color.rgb = rgb(*STEEL); ln.line.fill.background()
add_text(s, 0.7, 1.4, 12, 0.5, 'SUPPLIER EVALUATION REPORT', 16, False, MID_GREY)
add_text(s, 0.7, 2.05, 12, 1.2, 'ControlPartner AS', 44, True, WHITE)
add_text(s, 0.7, 3.5, 12, 0.6, 'Industrial Automation  |  SCADA Systems  |  Water/Wastewater', 18, False, (200,210,220))
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.35), Inches(4.5), Inches(0.05))
div.fill.solid(); div.fill.fore_color.rgb = rgb(*STEEL); div.line.fill.background()
add_text(s, 0.7, 4.55, 12, 0.4, 'NOK 57.3M Revenue  |  40-46 Employees  |  Norway  |  Uniwater Subsidiary', 14, False, (200,210,220))
add_text(s, 0.7, 6.2, 12, 0.4, 'Confidential  |  March 2026  |  Manu Forti Intelligence', 11, False, MID_GREY)
add_logos(s, is_title=True)

# S2: EXEC SUMMARY
print("  Slide 2: Executive Summary")
s = prs.slides.add_slide(blank)
add_header(s, 'EXECUTIVE SUMMARY', 'ControlPartner AS — Risk & Suitability Overview')
s.shapes.add_picture(str(visuals['gauge']), Inches(0.3), Inches(1.05), Inches(7.2), Inches(3.1))
add_rounded_rect(s, 7.6, 1.05, 5.4, 5.6, (240,245,250))
add_text(s, 7.8, 1.22, 5.2, 0.32, 'SUPPLIER SNAPSHOT', 14, True, NAVY)
snap = [('Supplier:', 'ControlPartner AS'), ('Type:', 'Private (Uniwater subsidiary)'),
        ('Sector:', 'Industrial Automation / SCADA'), ('HQ:', 'Porsgrunn, Norway'),
        ('Founded:', '2015'), ('Revenue:', 'NOK 57.3M (2023)'),
        ('EBITDA:', 'NOK 5.2M (9.1% margin)'), ('Order Book:', 'N/A (private)'),
        ('Employees:', '40-46'), ('Parent:', 'Uniwater AB (Sweden)')]
for idx, (label, val) in enumerate(snap):
    y = 1.64 + idx * 0.30
    add_text(s, 7.8, y, 2.3, 0.27, label, 11, True, TEXT_SECONDARY)
    add_text(s, 10.1, y, 2.7, 0.27, val, 11, False, TEXT_PRIMARY)
add_rounded_rect(s, 0.3, 5.22, 7.2, 1.43, (240,245,250))
add_text(s, 0.5, 5.33, 7, 0.3, 'Key Findings', 13, True, NAVY)
findings = ['• Niche Norwegian automation specialist with steady 2021-2023 growth trajectory',
            '• Acquired by Uniwater AB in October 2025 — now part of Nordic water tech group',
            '• Limited public financial data as private subsidiary — governance integration ongoing',
            '• ESG flag: MEDIUM rating due to limited disclosure and recent ownership change']
for idx, f in enumerate(findings):
    add_text(s, 0.5, 5.65+idx*0.27, 7, 0.26, f, 11, False, TEXT_PRIMARY)
add_source(s); add_logos(s)

# S3: RECOMMENDATION
print("  Slide 3: Recommendation")
s = prs.slides.add_slide(blank)
add_header(s, 'RECOMMENDATION', 'Decision Summary & Commercial Conditions')
add_rounded_rect(s, 0.5, 1.1, 12.33, 1.35, AMBER)
add_text(s, 0.7, 1.18, 12, 0.55, '⚠  RECOMMENDATION: CONDITIONAL APPROVAL — ENHANCED DUE DILIGENCE REQUIRED', 20, True, WHITE)
add_text(s, 0.7, 1.78, 12, 0.58,
         'ControlPartner is a technically capable Norwegian automation specialist with water/wastewater expertise. '
         'However, significant data limitations due to private subsidiary status and recent Uniwater acquisition '
         '(October 2025) require enhanced due diligence before engagement.',
         11, False, WHITE)
add_bullets(s, 0.5, 2.65, 6, 2.5, 'Commercial Conditions',
            ['Payment: Letter of Credit required (no open account)',
             'Performance bonds: 15% of contract value (above standard)',
             'Force-majeure: Explicit geopolitical/sanctions clause',
             'Uniwater parent guarantee for contracts >NOK 5M',
             'Local entity contracting preferred'], 14, 11)
add_bullets(s, 6.7, 2.65, 6, 2.5, 'Enhanced Due Diligence (Mandatory)',
            ['Request 2024 financials post-Uniwater acquisition',
             'Verify Uniwater AB creditworthiness and governance',
             'Supply chain audit: Norwegian labour standards',
             'Confirm key personnel retention post-acquisition',
             'Request latest ISO certifications',
             'Annual financial disclosure as contract condition'], 14, 11)
add_rounded_rect(s, 0.5, 5.35, 12.33, 1.5, (240,245,250))
add_text(s, 0.7, 5.5, 12, 0.3, 'Overall Risk Summary — MEDIUM (52/100)', 14, True, NAVY)
for x, t, d in [(0.6,'Financial: MEDIUM','Private; limited disclosure; Uniwater backing'),
                 (3.7,'Operational: LOW','Proven SCADA expertise; Norwegian standards'),
                 (6.8,'Geopolitical: LOW','Norwegian-based; stable jurisdiction'),
                 (9.9,'ESG: MEDIUM','Limited public disclosure; recent acquisition')]:
    add_text(s, x, 5.88, 3.1, 0.27, t, 11, True, NAVY)
    add_text(s, x, 6.17, 3.1, 0.27, d, 9, False, TEXT_SECONDARY)
add_source(s); add_logos(s)

# S4: SUPPLIER PROFILE
print("  Slide 4: Supplier Profile")
s = prs.slides.add_slide(blank)
add_header(s, 'SUPPLIER PROFILE', 'Corporate Structure & Global Footprint')
s.shapes.add_picture(str(visuals['org']), Inches(0.3), Inches(1.05), Inches(8.2), Inches(4.48))
add_bullets(s, 8.7, 1.05, 4.3, 3.5, 'Company Overview',
            ['ControlPartner AS is a Norwegian industrial',
             'automation company founded in 2015 in',
             'Porsgrunn. Specializes in SCADA systems,',
             'process control, and automation solutions',
             'for water/wastewater and industrial clients.',
             '',
             'Acquired by Uniwater AB (Sweden) in',
             'October 2025 — 5th Uniwater acquisition',
             'in 2025. Now part of Nordic water tech group.'], 14, 11)
add_bullets(s, 8.7, 4.5, 4.3, 2.4, 'Leadership',
            ['Management team retained post-acquisition',
             '  Uniwater integration ongoing',
             '',
             'Key: Norwegian operations with Swedish',
             '  parent company oversight',
             '',
             'HQ: Porsgrunn, NO  |  Founded: 2015'], 14, 11)
add_source(s); add_logos(s)

# S5: FINANCIAL HEALTH
print("  Slide 5: Financial Health")
s = prs.slides.add_slide(blank)
add_header(s, 'FINANCIAL HEALTH', 'Revenue & EBITDA Trajectory 2021–2023 (Pre-Acquisition)  |  Limited Post-Acquisition Data')
s.shapes.add_picture(str(visuals['fin']), Inches(0.3), Inches(1.05), Inches(7.5), Inches(4.55))
add_text(s, 8.1, 1.1, 4.9, 0.35, 'Financial Highlights', 14, True, NAVY)
metrics = [('2023 Revenue:', 'NOK 57.3M  (+4.6% YoY)'), ('2023 EBITDA:', 'NOK 5.2M  (9.1% margin)'),
           ('2023 Operating:', 'NOK 5.2M  (pre-acquisition)'), ('3yr Revenue CAGR:', '+5.0%  (2021–2023)'),
           ('Order Book:', 'N/A  (private subsidiary)'), ('Gross Debt:', 'N/A  (Uniwater consolidated)'),
           ('Net Cash Position:', 'N/A  (no disclosure)'), ('Debt / EBITDA:', 'N/A  (no data)')]
for idx, (label, val) in enumerate(metrics):
    y = 1.56 + idx * 0.30
    add_text(s, 8.1, y, 2.55, 0.29, label, 11, True, TEXT_SECONDARY)
    add_text(s, 10.65, y, 2.25, 0.29, val, 11, False, TEXT_PRIMARY)
add_rounded_rect(s, 8.1, 4.05, 4.8, 0.44, AMBER)
add_text(s, 8.25, 4.12, 4.5, 0.3, 'Financial Risk:  MEDIUM  ⚠', 12, True, WHITE, PP_ALIGN.CENTER)
add_bullets(s, 8.1, 4.6, 4.8, 2.3, 'Exposure Guidance',
            ['Private subsidiary — limited financial disclosure',
             'Uniwater AB backing provides parent support',
             'Letter of Credit mandatory; no open account',
             'Request 2024 audited financials before award'], 12, 10)
add_source(s); add_logos(s)

# S6: MARKET POSITION
print("  Slide 6: Market Position")
s = prs.slides.add_slide(blank)
add_header(s, 'MARKET POSITION', 'Industrial Automation & SCADA Market Position — 2023')
s.shapes.add_picture(str(visuals['mkt']), Inches(0.3), Inches(1.05), Inches(8.0), Inches(4.46))
add_bullets(s, 8.5, 1.1, 4.5, 2.0, 'Competitive Landscape',
            ['ControlPartner is a niche player in the',
             'Norwegian industrial automation market.',
             'Competes with global giants through',
             'specialization in water/wastewater SCADA',
             'and local Norwegian expertise.'], 14, 11)
add_bullets(s, 8.5, 3.2, 4.5, 2.6, 'ControlPartner Competitive Advantages',
            ['Water/wastewater SCADA specialization',
             'Norwegian market knowledge and relationships',
             'Local service and support capability',
             'Integration with Nordic water infrastructure',
             'Uniwater group synergies (post-acquisition)',
             'Agile compared to global competitors'], 14, 11)
add_text(s, 0.5, 5.5, 12, 0.3, 'Key Competitors', 13, True, NAVY)
for idx, c in enumerate([
    '• ABB — Switzerland — NOK 3.2B Norway — Global automation leader',
    '• Siemens Energy — Germany — NOK 2.8B Norway — Diversified industrial',
    '• Schneider Electric — France — NOK 3.5B Norway — Energy management',
    '• Emerson — USA — NOK 15B global — Process automation']):
    add_text(s, 0.6, 5.84+idx*0.29, 12, 0.27, c, 10, False, TEXT_PRIMARY)
add_source(s); add_logos(s)

# S7: OPS & RISK
print("  Slide 7: Ops & Risk")
s = prs.slides.add_slide(blank)
add_header(s, 'OPERATIONAL CAPABILITY  &  RISK ASSESSMENT', 'Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary')
s.shapes.add_picture(str(visuals['timeline']), Inches(0.3), Inches(1.05), Inches(7.0), Inches(2.92))
s.shapes.add_picture(str(visuals['risk']), Inches(7.65), Inches(1.05), Inches(5.5), Inches(4.46))
add_text(s, 0.3, 5.58, 7, 0.28, 'Risk Summary', 13, True, NAVY)
for idx, (cat, rating, desc) in enumerate([
    ('Financial:', '🟡 MEDIUM', 'Private; limited disclosure; Uniwater backing'),
    ('Geopolitical:', '🟢 LOW', 'Norway-based; stable jurisdiction'),
    ('Uniwater Integration:', '🟡 MEDIUM', 'Recent acquisition; governance alignment'),
    ('Customer Concentration:', '🟡 MEDIUM', 'Niche market; limited diversification'),
    ('Delivery:', '🟢 LOW', 'Proven execution; Norwegian standards')]):
    y = 5.89 + idx * 0.25
    add_text(s, 0.3, y, 2.1, 0.24, cat, 10, True, TEXT_PRIMARY)
    add_text(s, 2.4, y, 1.4, 0.24, rating, 10, True, TEXT_PRIMARY)
    add_text(s, 3.8, y, 3.7, 0.24, desc, 9, False, TEXT_SECONDARY)
add_bullets(s, 7.65, 5.58, 5.5, 1.85, 'Key Capabilities',
            ['SCADA system design and implementation',
             'Process control and automation',
             'Water/wastewater sector expertise',
             'Norwegian regulatory compliance'], 12, 10)
add_source(s); add_logos(s)

# S8: COMMERCIAL INTEL
print("  Slide 8: Commercial Intelligence")
s = prs.slides.add_slide(blank)
add_header(s, 'COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON', 'Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms')
s.shapes.add_picture(str(visuals['radar']), Inches(0.2), Inches(1.0), Inches(5.8), Inches(4.74))
s.shapes.add_picture(str(visuals['peer']), Inches(6.2), Inches(1.0), Inches(6.9), Inches(3.85))
add_rounded_rect(s, 0.3, 5.8, 12.7, 1.5, (240,245,250))
add_bullets(s, 0.5, 5.85, 6.5, 1.4, 'Commercial Terms & Negotiation',
            ['Pricing: Competitive vs global players on local projects',
             'Lead time: 4–8 weeks for standard SCADA components',
             'Payment: Letter of Credit mandatory — no open account',
             'Leverage: Multi-year service agreements for discount'], 12, 10)
add_bullets(s, 7.0, 5.85, 5.8, 1.4, 'Key Watch Points',
            ['⚠ Request 2024 financials post-Uniwater acquisition',
             '⚠ Confirm key personnel retention post-acquisition',
             '⚠ Verify Uniwater AB creditworthiness',
             '⚠ Monitor integration progress with parent company'], 12, 10)
add_source(s); add_logos(s)

# S9: ESG
print("  Slide 9: ESG Assessment")
s = prs.slides.add_slide(blank)
add_header(s, 'ESG ASSESSMENT', 'Environmental, Social & Governance Screening  |  Overall Rating: MEDIUM')
s.shapes.add_picture(str(visuals['esg']), Inches(1.67), Inches(1.05), Inches(10.0), Inches(5.93))
add_source(s); add_logos(s)

prs.save(str(OUTPUT))
print(f"\n✅ Saved: {OUTPUT}")
import os
print(f"📎 Size: {os.path.getsize(str(OUTPUT))/1024:.1f} KB")
