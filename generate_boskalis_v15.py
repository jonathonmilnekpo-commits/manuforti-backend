#!/usr/bin/env python3
"""
Boskalis Product 1 — v14 Final
v13 layout fixes +
  Slide 5: Gross Debt (EUR 247M), Net Cash Position (EUR +518M), Debt/EBITDA (0.19x)
           added to Financial Highlights panel — real 2024 figures from annual results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Palette (v11 locked) ──────────────────────────────────────────────────────
NAVY      = RGBColor(  0,  51, 102)
COBALT    = RGBColor( 27,  94, 143)
GREEN     = RGBColor( 39, 174,  96)
AMBER     = RGBColor(243, 156,  18)
RED       = RGBColor(192,   0,   0)
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(127, 140, 141)
LGRAY     = RGBColor(236, 240, 241)
BLACK     = RGBColor(  0,   0,   0)
PALE_BLUE = RGBColor(160, 200, 240)
PALE_GRAY = RGBColor(200, 200, 200)

VISUALS = "/Users/jonathonmilne/.openclaw/workspace/boskalis_v11_visuals"
LOGO    = "/Users/jonathonmilne/.openclaw/workspace/boskalis_logo_transparent.png"

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def header(slide, title, subtitle=None):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = WHITE
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
        sp = s.text_frame.paragraphs[0]
        sp.text = subtitle; sp.font.size = Pt(10); sp.font.italic = True
        sp.font.color.rgb = PALE_BLUE

def logo_tag(slide):
    """Add company logo — white backing box + transparent logo PNG, top-right of header."""
    # White rounded backing
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(10.92), Inches(0.06), Inches(2.26), Inches(0.87))
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    # Logo image (292×110px, ratio 2.655 → at w=2.1": h=0.79")
    slide.shapes.add_picture(LOGO, Inches(11.0), Inches(0.10), width=Inches(2.1))

def source(slide):
    s = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    p = s.text_frame.paragraphs[0]
    p.text = "Source: Manu Forti Intelligence  |  Confidential  |  February 2026"
    p.font.size = Pt(8); p.font.color.rgb = GRAY

def txt(slide, t, x, y, w, h, size=11, bold=False, col=None, wrap=True, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.text = t; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = col if col else BLACK
    if align: p.alignment = align

def blist(slide, title, items, x, y, w, h, title_col=None, size=10):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title; p.font.bold = True; p.font.size = Pt(12)
    p.font.color.rgb = title_col or NAVY
    for item in items:
        p = tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = BLACK
        p.space_after = Pt(2)

def box(slide, x, y, w, h, fill, lc=None, rnd=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def img(slide, fname, x, y, w):
    slide.shapes.add_picture(f'{VISUALS}/{fname}', Inches(x), Inches(y), width=Inches(w))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
box(slide, 0, 0, 13.333, 7.5, NAVY)
box(slide, 0, 5.55, 13.333, 0.08, GREEN)

txt(slide, "SUPPLIER EVALUATION REPORT", 0.7, 1.4, 12, 0.5, size=15, col=PALE_BLUE)
txt(slide, "Royal Boskalis Westminster N.V.", 0.7, 2.05, 12, 1.2, size=42, bold=True, col=WHITE)
txt(slide, "Dredging  |  Offshore Energy  |  Towage & Salvage",
    0.7, 3.5, 12, 0.6, size=19, col=PALE_GRAY)
box(slide, 0.7, 4.35, 4.5, 0.05, COBALT)
txt(slide, "EUR 4.4B Revenue  |  11,683 Employees  |  94 Countries  |  500+ Vessels",
    0.7, 4.55, 12, 0.4, size=12, col=PALE_GRAY)
txt(slide, "Confidential  |  February 2026  |  Manu Forti Intelligence",
    0.7, 6.2, 12, 0.4, size=11, col=RGBColor(150, 150, 150))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY  (MEDIUM 42/100 — gauge PNG updated)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "EXECUTIVE SUMMARY",
       "Royal Boskalis Westminster N.V. — Risk & Suitability Overview")
logo_tag(slide)

# Risk gauge — 01_risk_gauge.png already reflects MEDIUM/amber/42
img(slide, "01_risk_gauge.png", 0.3, 1.05, 7.2)

# Supplier snapshot panel
box(slide, 7.6, 1.05, 5.4, 5.6, LGRAY, rnd=True)
txt(slide, "SUPPLIER SNAPSHOT", 7.8, 1.22, 5.2, 0.32, size=11, bold=True, col=NAVY)

snapshot = [
    ("Supplier",   "Royal Boskalis Westminster N.V."),
    ("Ticker",     "Euronext Amsterdam: BOKA"),
    ("Sector",     "Marine Infrastructure & Offshore Energy"),
    ("HQ",         "Papendrecht, Netherlands"),
    ("Founded",    "1910 (Royal designation 1978)"),
    ("Revenue",    "EUR 4.4B (2024, record year)"),
    ("EBITDA",     "EUR 1.3B (30% margin, +28% YoY)"),
    ("Order Book", "EUR 7.0B (16 months coverage)"),
    ("Employees",  "11,683 globally"),
    ("Fleet",      "500+ vessels, 94 countries"),
]
y = 1.64
for label, val in snapshot:
    txt(slide, label + ":", 7.8, y, 2.3, 0.27, size=9.5, bold=True, col=NAVY)
    txt(slide, val, 10.1, y, 2.7, 0.27, size=9.5)
    y += 0.30

# Key findings strip
box(slide, 0.3, 5.22, 7.2, 1.43, LGRAY, rnd=True)
txt(slide, "Key Findings", 0.5, 5.33, 7.0, 0.30, size=11, bold=True, col=NAVY)
findings = [
    "• Record 2024: EUR 4.4B revenue (+3% YoY), EUR 1.3B EBITDA (+28%) — margin at 30%",
    "• 115-year heritage — Royal designation, Euronext-listed, world-class governance",
    "• EUR 7.0B order book = 16 months of revenue locked in — exceptional pipeline visibility",
    "• ESG flag: Both ENDS legal action (2020) raises social impact exposure — conditions apply",
]
y = 5.65
for f in findings:
    txt(slide, f, 0.5, y, 7.0, 0.26, size=9.5)
    y += 0.27

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: RECOMMENDATION  (AMBER / APPROVE WITH ESG CONDITIONS)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "RECOMMENDATION", "Decision Summary & Commercial Conditions")
logo_tag(slide)

# Amber recommendation banner
box(slide, 0.5, 1.1, 12.333, 1.35, AMBER, rnd=True)
txt(slide, "⚠  RECOMMENDATION: APPROVE WITH ESG CONDITIONS",
    0.7, 1.18, 12, 0.55, size=22, bold=True, col=WHITE)
txt(slide,
    "Boskalis is recommended for marine infrastructure, dredging, offshore energy, and heavy maritime "
    "transport. Record financials and EUR 7.0B order book support approval. Conditions apply due to "
    "ESG flag: Both ENDS legal action (2020, South Sulawesi) must be verified as resolved.",
    0.7, 1.78, 12, 0.58, size=11, col=WHITE)

# Commercial Conditions
blist(slide, "Commercial Conditions", [
    "• Standard payment terms: Net-30 to Net-60 days",
    "• Joint venture structures recommended for projects >EUR 100M",
    "• Performance bonds for contracts >EUR 50M",
    "• Force-majeure clause covering geopolitical vessel access",
    "• Confirm key vessel availability before contract award",
], 0.5, 2.65, 6.0, 2.5)

# ESG Conditions
blist(slide, "ESG Conditions (Mandatory)", [
    "• Written confirmation Both ENDS case is resolved",
    "• Supplier Code of Conduct: community impact standards",
    "• Social Impact clause embedded in contract",
    "• Request latest EcoVadis or equivalent ESG rating",
    "• Annual ESG reporting milestone in framework contracts",
], 6.7, 2.65, 6.0, 2.5, title_col=AMBER)

# Risk summary strip — 4 pillars
box(slide, 0.5, 5.35, 12.333, 1.5, LGRAY, rnd=True)
txt(slide, "Overall Risk Summary — MEDIUM (42/100)", 0.7, 5.50, 12, 0.30, size=11, bold=True, col=NAVY)
risk_items = [
    ("Financial",    "LOW",    "EUR 4.4B revenue, 30% EBITDA, Euronext-listed governance"),
    ("Operational",  "LOW",    "500+ vessels, 94 countries — world-class fleet"),
    ("Geopolitical", "MEDIUM", "Complex jurisdictions — Middle East, West Africa"),
    ("ESG",          "MEDIUM", "Both ENDS action (2020) — social flag, conditions apply"),
]
x = 0.6
for cat, rating, note in risk_items:
    col = GREEN if rating == "LOW" else AMBER
    txt(slide, f"{cat}: {rating}", x, 5.88, 3.1, 0.27, size=9.5, bold=True, col=col)
    txt(slide, note,               x, 6.17, 3.1, 0.27, size=8.5, col=GRAY)
    x += 3.1

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: SUPPLIER PROFILE  (unchanged)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "SUPPLIER PROFILE", "Corporate Structure & Global Footprint")
logo_tag(slide)

# 02_org_structure.png: ratio=1.832 → at w=8.2": h=4.47", ends y=5.52 — fine
img(slide, "02_org_structure.png", 0.3, 1.05, 8.2)

blist(slide, "Company Overview", [
    "Royal Boskalis Westminster N.V. is a leading global",
    "contractor in dredging, marine infrastructure, and",
    "offshore energy services. With 115+ years of experience,",
    "Boskalis operates one of the world's largest specialised",
    "vessel fleets — 500+ units across 94 countries.",
    "",
    "Listed on Euronext Amsterdam since 1978 (BOKA), the",
    "company holds Royal designation from the Dutch Crown,",
    "reflecting its strategic national importance.",
], 8.7, 1.05, 4.3, 3.5)

blist(slide, "Leadership", [
    "Peter Berdowski — CEO",
    "  25+ years at Boskalis; industry veteran",
    "",
    "Michiel Borgers — CFO",
    "  Deep offshore energy finance background",
    "",
    "HQ: Papendrecht, NL  |  Founded: 1910",
], 8.7, 4.5, 4.3, 2.4)

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: FINANCIAL HEALTH  (v14: debt + D/E added)
# 03_financial_trend.png ratio=1.648 → w=7.5" → h=4.55" → ends y=5.60 ✓
# Right panel: 8 metrics incl. Gross Debt, Net Cash Position, Debt/EBITDA
# GREEN LOW box anchored below metrics; Exposure Guidance fits to footer
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "FINANCIAL HEALTH",
       "Revenue & EBITDA Trajectory 2021–2024  |  Record Year: EUR 4.4B  |  Net Cash: EUR +518M")
logo_tag(slide)

img(slide, "03_financial_trend.png", 0.3, 1.05, 7.5)

R = 8.1   # right panel x-start
txt(slide, "Financial Highlights", R, 1.10, 4.9, 0.35, size=12, bold=True, col=NAVY)

metrics = [
    ("2024 Revenue",     "EUR 4.4B  (+3% YoY)"),
    ("2024 EBITDA",      "EUR 1.3B  (30% margin, +28%)"),
    ("2024 Net Profit",  "EUR 781M  (+30% YoY)"),
    ("3yr Revenue CAGR", "+9%  (2021–2024)"),
    ("Order Book",       "EUR 7.0B  (16 months)"),
    ("Gross Debt",       "EUR 247M  (low absolute debt)"),
    ("Net Cash Position","EUR +518M  (cash EUR 765M)"),
    ("Debt / EBITDA",    "0.19x  ✓  (near debt-free)"),
]
y = 1.56
for label, val in metrics:
    # Highlight debt metrics in green (positive signal)
    val_col = GREEN if label in ("Net Cash Position", "Debt / EBITDA") else BLACK
    txt(slide, label + ":", R,     y, 2.55, 0.29, size=9.5, bold=True, col=NAVY)
    txt(slide, val,          R+2.55, y, 2.25, 0.29, size=9.5, col=val_col)
    y += 0.30

# GREEN LOW box — anchored cleanly below metrics
box(slide, R, 4.05, 4.8, 0.44, GREEN, rnd=True)
txt(slide, "Financial Risk:  LOW  ✓", R+0.15, 4.12, 4.5, 0.30, size=12, bold=True, col=WHITE)

blist(slide, "Exposure Guidance", [
    "• Near debt-free: EUR 0.19x Debt/EBITDA — exceptional headroom",
    "• No parent guarantee required — standalone credit",
    "• Single contract: standard project sizing applies",
    "• Payment: Net-30 standard; milestone for large projects",
], R, 4.60, 4.8, 2.30)

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: MARKET POSITION  (unchanged)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "MARKET POSITION",
       "Marine Contractor Revenue Comparison — 2024 Estimates")
logo_tag(slide)

# 04_market_position.png ratio=1.792 → w=8.0": h=4.46", ends y=5.51 — fine
img(slide, "04_market_position.png", 0.3, 1.05, 8.0)

blist(slide, "Competitive Landscape", [
    "Boskalis is one of only 3 global marine contractors",
    "capable of full-scope dredging + offshore energy +",
    "heavy transport simultaneously. Owned fleet and",
    "94-country presence creates a defensible moat.",
], 8.5, 1.1, 4.5, 2.0)

blist(slide, "Boskalis Competitive Advantages", [
    "• World's largest dredging fleet by capacity",
    "• Smit Lamnalco: 80+ global ports — unmatched",
    "• Royal designation — strategic Dutch asset",
    "• Fully integrated: vessels, crews, logistics",
    "• EUR 7.0B order book — market confidence signal",
    "• Euronext-listed: governance & financial discipline",
], 8.5, 3.2, 4.5, 2.6)

txt(slide, "Key Competitors", 0.5, 5.5, 12, 0.30, size=12, bold=True, col=NAVY)
comps = [
    "DEME (GeoSea) — Belgium — EUR 3.5B — Offshore wind specialist, subsea cables",
    "Jan De Nul — Luxembourg — EUR 2.8B — Privately held; strong in land reclamation",
    "Van Oord — Netherlands — EUR 2.0B — Dutch rival; wind farm focus",
    "Penta-Ocean — Japan — EUR 1.1B — Asia-Pacific dredging; less global",
]
y = 5.84
for c in comps:
    txt(slide, "• " + c, 0.6, y, 12, 0.27, size=9.5)
    y += 0.29

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: OPERATIONAL CAPABILITY + RISK ASSESSMENT  (layout rebalanced)
# 05_investment_timeline ratio=2.400 → w=7.0": h=2.92", ends y=3.97
# 06_risk_matrix ratio=1.233 → w=5.5": h=4.46", ends y=5.51
# Bottom section: y=5.58, 5 rows @ 0.25" spacing (fits before footer)
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "OPERATIONAL CAPABILITY  &  RISK ASSESSMENT",
       "Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary")
logo_tag(slide)

img(slide, "05_investment_timeline.png", 0.3,  1.05, 7.0)   # h=2.92, bottom=3.97
img(slide, "06_risk_matrix.png",         7.65, 1.05, 5.5)   # h=4.46, bottom=5.51

# Bottom section — starts at y=5.58 (0.07" below matrix bottom)
BOTY = 5.58
txt(slide, "Risk Summary", 0.3, BOTY, 7.0, 0.28, size=11, bold=True, col=NAVY)

risk_rows = [
    ("Financial",          "🟢 LOW",    "EUR 4.4B revenue; 30% EBITDA; Euronext-listed"),
    ("Geopolitical",       "🟡 MEDIUM", "Complex jurisdictions — Middle East, West Africa"),
    ("Fleet Availability", "🟢 LOW",    "500+ vessels; globally diversified"),
    ("FX / Currency",      "🟢 LOW",    "EUR-denominated; hedged globally"),
    ("Delivery Risk",      "🟡 MEDIUM", "Large projects carry execution risk"),
]
y = BOTY + 0.31
for cat, rating, note in risk_rows:
    col = GREEN if "🟢" in rating else AMBER
    txt(slide, cat + ":",  0.3, y, 2.1, 0.24, size=8.5, bold=True, col=NAVY)
    txt(slide, rating,     2.4, y, 1.4, 0.24, size=8.5, bold=True, col=col)
    txt(slide, note,       3.8, y, 3.7, 0.24, size=8.0, col=GRAY)
    y += 0.25

# Key Capabilities — aligned at same BOTY
blist(slide, "Key Capabilities", [
    "• ISO 9001 / ISO 14001 / OHSAS 18001 certified",
    "• 115 years continuous marine operations",
    "• 94 countries — local expertise globally",
    "• Full chain: design → execute → maintain",
], 7.65, BOTY, 5.5, 1.85)

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: COMMERCIAL INTELLIGENCE + PEER RISK
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON",
       "Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms")
logo_tag(slide)

# 07_benchmarking_radar ratio=1.223 → w=5.8": h=4.74", ends y=5.74
# 08_peer_risk ratio=1.791 → w=6.9": h=3.85", ends y=4.85
img(slide, "07_benchmarking_radar.png", 0.2, 1.0, 5.8)
img(slide, "08_peer_risk.png",          6.2, 1.0, 6.9)

box(slide, 0.3, 5.8, 12.7, 1.50, LGRAY, rnd=True)

blist(slide, "Commercial Terms & Negotiation", [
    "• Pricing: Milestone/progress payments standard; lump-sum for fixed-scope",
    "• Lead time: 3–9 months to mobilisation depending on vessel availability",
    "• Leverage: Volume commitment; long-term framework agreements",
    "• IP: Client owns design deliverables; Boskalis retains method IP",
], 0.5, 5.85, 6.5, 1.4, size=9)

blist(slide, "Key Watch Points", [
    "⚠ Geopolitical: Monitor project-country risk individually",
    "⚠ Order book: High demand may constrain vessel availability",
    "⚠ FX: Hedge EUR/USD/NOK exposure on long-term contracts",
], 7.0, 5.85, 5.8, 1.4, title_col=AMBER, size=9)

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: ESG ASSESSMENT
# 09_esg_assessment.png ratio=1.685 → w=10.0": h=5.94", ends y=6.99 ✓
# Centred: x = (13.333 - 10.0) / 2 = 1.67
# ══════════════════════════════════════════════════════════════════════════════
slide = blank()
header(slide, "ESG ASSESSMENT",
       "Environmental, Social & Governance Screening  |  Overall Rating: MEDIUM")
logo_tag(slide)

img(slide, "09_esg_assessment.png", 1.67, 1.05, 10.0)

source(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/jonathonmilne/.openclaw/workspace/Boskalis_Product1_v15_Final.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print("  9 slides | v15: Boskalis logo top-right every slide | MEDIUM (42/100) | Manu Forti branding")
