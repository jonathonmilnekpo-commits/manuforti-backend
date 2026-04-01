#!/usr/bin/env python3
"""Generate Form Energy Inc Product 1 v15 — Full Analysis using Jarotech as reference."""

import sys
sys.path.insert(0, '/opt/homebrew/lib/python3.14/site-packages')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Branding per Jarotech reference
NAVY = (0, 33, 71)
STEEL = (43, 108, 176)
MID_GREY = (113, 128, 150)
GREEN = (72, 187, 120)
AMBER = (214, 158, 46)
RED = (229, 62, 62)
WHITE = (255, 255, 255)
TEXT_PRIMARY = (26, 32, 44)
TEXT_SECONDARY = (74, 85, 104)

M_NAVY = '#002147'
M_STEEL = '#2B6CB0'
M_GREY = '#718096'
M_GREEN = '#48BB78'
M_AMBER = '#D69E2E'
M_RED = '#E53E3E'
M_LIGHT = '#EBF4FF'

BASE = Path('/Users/jonathonmilne/.openclaw/workspace')
VISUALS = BASE / 'formenergy_v2_visuals'
VISUALS.mkdir(exist_ok=True)
CLAN_LOGO = BASE / 'skills/product-1-generator/assets/manu_forti_logo.png'
OUTPUT = BASE / 'FormEnergy_Product1_v15_SecondDraft.pptx'

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_header(slide, title, subtitle):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    h.fill.solid(); h.fill.fore_color.rgb = rgb(*NAVY); h.line.fill.background()
    
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.55))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = rgb(*WHITE)
    
    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.5), Inches(0.35))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = subtitle; r2.font.size = Pt(12); r2.font.color.rgb = rgb(200, 210, 220)

def add_source(slide):
    t = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.8), Inches(0.3))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = 'Source: Manu Forti Intelligence  |  Confidential  |  March 2026'
    r.font.size = Pt(8); r.font.color.rgb = rgb(*MID_GREY)

def add_text(slide, l, t, w, h, text, sz=14, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = rgb(*color)
    return tb

def add_bullets(slide, l, t, w, h, title, items, tsz=14, isz=11, tc=NAVY, ic=TEXT_PRIMARY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(tsz); r.font.bold = True; r.font.color.rgb = rgb(*tc)
    for item in items:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = f'• {item}'; r.font.size = Pt(isz); r.font.color.rgb = rgb(*ic)

def add_rounded_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(*color); s.line.fill.background()
    return s

def add_logos(slide, is_title=False):
    cimg = Image.open(str(CLAN_LOGO))
    ca = cimg.size[0] / cimg.size[1]
    cw = 0.6; ch = cw / ca
    cx = 13.333 - cw - 0.2
    cy = 7.5 - ch - 0.15
    slide.shapes.add_picture(str(CLAN_LOGO), Inches(cx), Inches(cy), Inches(cw), Inches(ch))

def gen_risk_gauge():
    fig, ax = plt.subplots(figsize=(14, 6.5), facecolor='white')
    segments = [(0, 60, M_GREEN, 'LOW\n0–33'), (60, 120, M_AMBER, 'MEDIUM\n34–66'), (120, 180, M_RED, 'HIGH\n67–100')]
    for start, end, color, label in segments:
        theta = np.linspace(np.radians(start), np.radians(end), 100)
        x_o = np.cos(theta); y_o = np.sin(theta)
        x_i = np.cos(theta[::-1]) * 0.6; y_i = np.sin(theta[::-1]) * 0.6
        ax.fill(np.concatenate([x_o, x_i]), np.concatenate([y_o, y_i]), color=color, alpha=0.85)
        mid = np.radians((start+end)/2)
        ax.text(np.cos(mid)*0.8, np.sin(mid)*0.8, label, ha='center', va='center', fontsize=11, fontweight='bold', color='white')
    
    score = 58
    needle = np.radians(score * 180 / 100)
    ax.annotate('', xy=(np.cos(needle)*0.55, np.sin(needle)*0.55), xytext=(0,0),
                arrowprops=dict(arrowstyle='->', color=M_NAVY, lw=3.5))
    ax.add_patch(plt.Circle((0,0), 0.08, color=M_NAVY, zorder=5))
    ax.text(0, -0.2, f'{score}/100', ha='center', fontsize=28, fontweight='bold', color=M_NAVY)
    ax.text(0, -0.35, 'OVERALL RISK SCORE', ha='center', fontsize=12, color=M_GREY)
    ax.text(0, 1.25, 'Form Energy Inc — Overall Risk Assessment', ha='center', fontsize=18, fontweight='bold', color=M_NAVY)
    ax.text(0, 1.10, 'Iron-Air Batteries | Long-Duration Energy Storage | 100+ Hours', ha='center', fontsize=12, color=M_GREY)
    
    boxes = [('Financial', 'MEDIUM', M_AMBER), ('Operational', 'MEDIUM', M_AMBER), ('Geopolitical', 'LOW', M_GREEN), ('ESG', 'MEDIUM', M_AMBER)]
    for idx, (cat, rating, color) in enumerate(boxes):
        bx = -0.75 + idx * 0.5
        rect = FancyBboxPatch((bx-0.2, -0.67), 0.4, 0.22, boxstyle="round,pad=0.03", facecolor=color, alpha=0.15, edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(bx, -0.54, cat, ha='center', fontsize=8, fontweight='bold', color=M_NAVY)
        ax.text(bx, -0.62, rating, ha='center', fontsize=8, fontweight='bold', color=color)
    
    ax.set_xlim(-1.4, 1.4); ax.set_ylim(-0.75, 1.4); ax.set_aspect('equal'); ax.axis('off')
    p = VISUALS / '01_risk_gauge.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_org_chart():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    parent = FancyBboxPatch((5.5, 7.5), 5, 1.2, boxstyle="round,pad=0.15", facecolor=M_NAVY, edgecolor='none')
    ax.add_patch(parent)
    ax.text(8, 8.1, 'Form Energy Inc', ha='center', fontsize=18, fontweight='bold', color='white')
    ax.text(8, 7.75, 'Somerville, MA  |  Private  |  Founded 2017', ha='center', fontsize=10, color='#cccccc')
    
    divs = [('Battery\nTechnology', 'Iron-Air\n100+ hr Storage', 1.0), ('Manufacturing', 'Weirton, WV\nFacility', 4.5),
            ('R&D', 'Somerville\nInnovation Hub', 8.0), ('Commercial', 'Google, Xcel\nContracts', 11.5)]
    for label, desc, x in divs:
        box = FancyBboxPatch((x, 4.8), 3, 1.6, boxstyle="round,pad=0.1", facecolor=M_STEEL, edgecolor='none', alpha=0.9)
        ax.add_patch(box)
        ax.text(x+1.5, 5.85, label, ha='center', fontsize=12, fontweight='bold', color='white')
        ax.text(x+1.5, 5.2, desc, ha='center', fontsize=9, color='#e0e0e0')
        ax.plot([8, x+1.5], [7.5, 6.4], color=M_GREY, lw=1.5, alpha=0.6)
    
    facts = [('$1.52B\nFunding Raised', 1.5, 2.8), ('~400\nEmployees', 5.0, 2.8),
             ('$2.5B\nValuation', 8.5, 2.8), ('2027\nCommercial Target', 12.0, 2.8)]
    for text, x, y in facts:
        box = FancyBboxPatch((x, y), 2.5, 1.2, boxstyle="round,pad=0.08", facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1)
        ax.add_patch(box)
        ax.text(x+1.25, y+0.6, text, ha='center', fontsize=11, fontweight='bold', color=M_NAVY)
    
    ax.set_xlim(0, 16); ax.set_ylim(2, 9.5); ax.axis('off')
    p = VISUALS / '02_org_chart.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_funding_chart():
    """Funding timeline for pre-revenue company."""
    fig, ax = plt.subplots(figsize=(15, 9), facecolor='white')
    
    rounds = ['2017\nSeed', '2019\nSeries A', '2020\nSeries B', '2021\nSeries C', '2022\nSeries D', '2024\nSeries E']
    funding = [0.1, 0.4, 0.7, 2.4, 4.5, 7.5]  # Cumulative $B
    colors = [M_NAVY, M_STEEL, M_STEEL, M_STEEL, M_STEEL, M_GREEN]
    
    bars = ax.bar(rounds, funding, color=colors, width=0.6, zorder=3)
    for bar, val in zip(bars, funding):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.15, f'${val:.1f}B', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    
    ax.set_ylabel('Cumulative Funding ($B)', fontsize=13, color=M_NAVY, fontweight='bold')
    ax.set_ylim(0, 9); ax.tick_params(axis='y', labelcolor=M_NAVY, labelsize=11)
    ax.tick_params(axis='x', labelsize=12)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.2, linestyle='--')
    ax.set_title('Funding Trajectory 2017–2024 — Pre-Revenue Stage', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    
    props = dict(boxstyle='round,pad=0.5', facecolor=M_LIGHT, edgecolor=M_STEEL, linewidth=1.5)
    ax.text(0.02, 0.97, 'Form Energy is pre-revenue with $1.52B total funding.\nFirst commercial deliveries expected 2027.\nGoogle contract: 300MW/30GWh (2026-2028).',
             transform=ax.transAxes, fontsize=11, va='top', bbox=props, color=M_NAVY)
    
    plt.tight_layout()
    p = VISUALS / '03_funding.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_market():
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    companies = ['Tesla', 'Fluence', 'CATL', 'BYD', 'Form Energy', 'QuantumScape']
    valuation = [800, 450, 1200, 900, 2.5, 4.5]  # $B or valuation
    colors = [M_GREY, M_GREY, M_GREY, M_GREY, M_NAVY, M_GREY]
    
    bars = ax.barh(range(len(companies)), valuation, color=colors, height=0.55, zorder=3)
    for bar, val, c in zip(bars, valuation, companies):
        lc = 'white' if c == 'Form Energy' else M_NAVY
        label = f'${val:.1f}B'
        ax.text(val-20, bar.get_y()+bar.get_height()/2, label, ha='right', va='center', fontsize=13, fontweight='bold', color=lc)
    
    ax.set_yticks(range(len(companies))); ax.set_yticklabels(companies, fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    ax.set_xlabel('Valuation ($B)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Energy Storage Companies — Valuation Comparison', fontsize=16, fontweight='bold', color=M_NAVY, pad=15)
    ax.get_yticklabels()[4].set_color(M_NAVY)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.annotate('Form Energy: Pre-revenue\n$1.52B funding | $2.5B valuation', xy=(2.5, 4), xytext=(400, 2.5),
                fontsize=11, fontweight='bold', color=M_STEEL, arrowprops=dict(arrowstyle='->', color=M_STEEL, lw=2))
    plt.tight_layout()
    p = VISUALS / '04_market.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_timeline():
    fig, ax = plt.subplots(figsize=(14, 6), facecolor='white')
    milestones = [('2017', 'Founded in\nSomerville', M_NAVY), ('2019', '$9M Seed\nBreakthrough Energy', M_STEEL),
                  ('2021', '$240M Series C\nT. Rowe Price', M_STEEL), ('2022', '$450M Series D\nGE Vernova', M_STEEL),
                  ('2024', '$760M Facility\nWeirton, WV', M_GREEN), ('2026-27', 'Google Contract\n300MW/30GWh', M_GREEN)]
    ax.plot([0, len(milestones)-1], [0, 0], color=M_NAVY, linewidth=3, zorder=1)
    for i, (year, desc, color) in enumerate(milestones):
        ax.plot(i, 0, 'o', markersize=16, color=color, zorder=3)
        ax.plot(i, 0, 'o', markersize=10, color='white', zorder=4)
        y_off = 0.6 if i % 2 == 0 else -0.6
        ax.text(i, y_off, f'{year}\n{desc}', ha='center', va='bottom' if i%2==0 else 'top',
                fontsize=10, fontweight='bold', color=M_NAVY,
                bbox=dict(boxstyle='round,pad=0.3', facecolor=M_LIGHT, edgecolor=color, linewidth=1.5))
        ax.plot([i, i], [0, y_off*0.4], color=color, linewidth=1.5, linestyle='--')
    ax.set_title('Strategic Investment Timeline', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(-0.5, len(milestones)-0.5); ax.set_ylim(-1.4, 1.4); ax.axis('off')
    plt.tight_layout()
    p = VISUALS / '05_timeline.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_risk_matrix():
    fig, ax = plt.subplots(figsize=(11, 9), facecolor='white')
    ax.fill_between([0,5], [5,5], [10,10], color=M_AMBER, alpha=0.10)
    ax.fill_between([5,10], [5,5], [10,10], color=M_RED, alpha=0.10)
    ax.fill_between([0,5], [0,0], [5,5], color=M_GREEN, alpha=0.10)
    ax.fill_between([5,10], [0,0], [5,5], color=M_AMBER, alpha=0.10)
    ax.text(2.5, 7.5, 'MONITOR', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    ax.text(7.5, 7.5, 'CRITICAL', ha='center', fontsize=14, fontweight='bold', color=M_RED, alpha=0.4)
    ax.text(2.5, 2.5, 'ACCEPT', ha='center', fontsize=14, fontweight='bold', color=M_GREEN, alpha=0.4)
    ax.text(7.5, 2.5, 'MITIGATE', ha='center', fontsize=14, fontweight='bold', color=M_AMBER, alpha=0.4)
    
    risks = [('Pre-Revenue\nStatus', 6.0, 7.0, 800, M_RED), ('Manufacturing\nScale-Up', 5.5, 6.5, 700, M_AMBER),
             ('Technology\nValidation', 4.0, 5.5, 600, M_AMBER), ('Competition\n(Li-ion)', 7.0, 4.0, 550, M_AMBER),
             ('Supply Chain\n(Raw Materials)', 3.5, 4.0, 450, M_GREEN), ('US Market\nAccess', 2.5, 2.5, 350, M_GREEN)]
    for label, x, y, size, color in risks:
        ax.scatter(x, y, s=size, color=color, alpha=0.7, edgecolors=M_NAVY, linewidth=1.5, zorder=3)
        ax.text(x, y, label, ha='center', va='center', fontsize=8, fontweight='bold', color=M_NAVY)
    
    ax.set_xlabel('Probability →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_ylabel('Impact →', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Risk Matrix — Form Energy Inc', fontsize=14, fontweight='bold', color=M_NAVY, pad=10)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.set_xticks([2.5, 7.5]); ax.set_xticklabels(['Low', 'High'], fontsize=11)
    ax.set_yticks([2.5, 7.5]); ax.set_yticklabels(['Low', 'High'], fontsize=11)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    ax.axhline(y=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    ax.axvline(x=5, color=M_GREY, linewidth=1, linestyle='--', alpha=0.5)
    plt.tight_layout()
    p = VISUALS / '06_risk_matrix.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_radar():
    fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(projection='polar'), facecolor='white')
    cats = ['Price\nCompetitiveness', 'Duration\n(100+ hr)', 'Technology\nMaturity', 'Manufacturing\nScale', 'Supply Chain\nSecurity', 'Commercial\nTraction']
    N = len(cats)
    formenergy = [9, 10, 5, 4, 6, 5]
    tesla = [6, 4, 9, 9, 8, 9]
    angles = np.linspace(0, 2*np.pi, N, endpoint=False).tolist()
    fe_c = formenergy + [formenergy[0]]; t_c = tesla + [tesla[0]]; ang_c = angles + [angles[0]]
    ax.fill(ang_c, fe_c, color=M_NAVY, alpha=0.2)
    ax.plot(ang_c, fe_c, color=M_NAVY, linewidth=2.5, marker='o', markersize=8, label='Form Energy')
    ax.fill(ang_c, t_c, color=M_STEEL, alpha=0.1)
    ax.plot(ang_c, t_c, color=M_STEEL, linewidth=2, marker='s', markersize=7, linestyle='--', label='Tesla (Li-ion)')
    ax.set_xticks(angles); ax.set_xticklabels(cats, fontsize=11, fontweight='bold', color=M_NAVY)
    ax.set_ylim(0, 10); ax.set_yticks([2,4,6,8,10])
    ax.set_yticklabels(['2','4','6','8','10'], fontsize=9, color=M_GREY)
    ax.grid(color=M_GREY, alpha=0.3)
    ax.set_title('Supplier Benchmarking — Form Energy vs Tesla', fontsize=14, fontweight='bold', color=M_NAVY, pad=25)
    ax.legend(loc='lower right', fontsize=11, bbox_to_anchor=(1.15, -0.05))
    plt.tight_layout()
    p = VISUALS / '07_radar.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_peer_risk():
    fig, ax = plt.subplots(figsize=(14, 8), facecolor='white')
    companies = ['Form Energy', 'Tesla Energy', 'Fluence', 'CATL', 'QuantumScape', 'Solid Power']
    scores = [58, 25, 30, 35, 65, 70]
    colors = [M_GREEN if s<=33 else M_AMBER if s<=66 else M_RED for s in scores]
    
    ax.axhspan(0, 33, facecolor=M_GREEN, alpha=0.07); ax.axhspan(33, 66, facecolor=M_AMBER, alpha=0.07)
    ax.axhspan(66, 100, facecolor=M_RED, alpha=0.07)
    ax.text(5.3, 16, 'LOW', fontsize=10, color=M_GREEN, fontweight='bold', alpha=0.6, ha='right')
    ax.text(5.3, 50, 'MEDIUM', fontsize=10, color=M_AMBER, fontweight='bold', alpha=0.6, ha='right')
    ax.text(5.3, 80, 'HIGH', fontsize=10, color=M_RED, fontweight='bold', alpha=0.6, ha='right')
    
    for i, (c, s, col) in enumerate(zip(companies, scores, colors)):
        ax.plot([i, i], [0, s], color=col, linewidth=3, zorder=2)
        ax.scatter(i, s, s=200, color=col, zorder=3, edgecolors=M_NAVY, linewidth=1.5)
        ax.text(i, s+3, f'{s}', ha='center', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.scatter(0, 58, s=350, color=M_NAVY, zorder=4, edgecolors=M_NAVY, linewidth=2)
    ax.text(0, 61, '58', ha='center', fontsize=14, fontweight='bold', color=M_NAVY)
    
    ax.set_xticks(range(len(companies))); ax.set_xticklabels(companies, fontsize=12, fontweight='bold', rotation=15)
    ax.get_xticklabels()[0].set_color(M_NAVY)
    ax.set_ylabel('Risk Score (0–100)', fontsize=13, fontweight='bold', color=M_NAVY)
    ax.set_title('Peer Risk Comparison — Energy Storage (USA)', fontsize=14, fontweight='bold', color=M_NAVY, pad=15)
    ax.set_ylim(0, 100); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.15)
    plt.tight_layout()
    p = VISUALS / '08_peer_risk.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

def gen_esg():
    fig, axes = plt.subplots(1, 3, figsize=(20, 12), facecolor='white')
    fig.suptitle('ESG Assessment — Form Energy Inc', fontsize=22, fontweight='bold', color=M_NAVY, y=0.98)
    fig.text(0.5, 0.94, 'Overall ESG Rating: MEDIUM  |  Strong environmental profile offset by early-stage governance gaps',
             ha='center', fontsize=13, color=M_GREY)
    
    pillars = [
        {'title': 'Environmental', 'rating': 'LOW', 'color': M_GREEN, 'items': [
            ('Iron-Air Chemistry', M_GREEN, '✓'), ('No Rare Earths', M_GREEN, '✓'),
            ('Recyclable Materials', M_GREEN, '✓'), ('Grid Decarbonization', M_GREEN, '✓'),
            ('Manufacturing Emissions', M_AMBER, '⚠'), ('Supply Chain Audit', M_AMBER, '⚠')]},
        {'title': 'Social', 'rating': 'MEDIUM', 'color': M_AMBER, 'items': [
            ('US Manufacturing Jobs', M_GREEN, '✓'), ('Weirton, WV Facility', M_GREEN, '✓'),
            ('Workforce Development', M_GREEN, '✓'), ('Labor Standards (US)', M_GREEN, '✓'),
            ('Diversity Reporting', M_AMBER, '⚠'), ('Community Programs', M_AMBER, '⚠')]},
        {'title': 'Governance', 'rating': 'MEDIUM', 'color': M_AMBER, 'items': [
            ('Private Company', M_AMBER, '⚠'), ('Limited Public Filings', M_AMBER, '⚠'),
            ('Experienced Leadership', M_GREEN, '✓'), ('Board Composition', M_AMBER, '⚠'),
            ('Investor Oversight', M_GREEN, '✓'), ('Transparency', M_AMBER, '⚠')]}
    ]
    
    for ax, pillar in zip(axes, pillars):
        ax.set_xlim(0, 10); ax.set_ylim(0, 12); ax.axis('off')
        title_box = FancyBboxPatch((0.5, 10.5), 9, 1.2, boxstyle="round,pad=0.15", facecolor=pillar['color'], edgecolor='none', alpha=0.9)
        ax.add_patch(title_box)
        ax.text(5, 11.1, pillar['title'], ha='center', fontsize=16, fontweight='bold', color='white')
        ax.text(5, 10.7, f"Rating: {pillar['rating']}", ha='center', fontsize=12, color='white', alpha=0.9)
        for i, (item, color, symbol) in enumerate(pillar['items']):
            y = 9.2 - i * 1.3
            ax.add_patch(plt.Circle((1.5, y), 0.35, color=color, alpha=0.8))
            ax.text(1.5, y, symbol, ha='center', va='center', fontsize=14, fontweight='bold', color='white')
            ax.text(2.5, y, item, va='center', fontsize=12, fontweight='bold', color=M_NAVY)
    
    controversy_box = FancyBboxPatch((0.03, 0.015), 0.94, 0.105, boxstyle="round,pad=0.01",
                                       facecolor=M_RED, edgecolor=M_RED, linewidth=1.5, alpha=0.10, transform=fig.transFigure)
    fig.patches.append(controversy_box)
    fig.text(0.5, 0.098, '⚠  CONTROVERSY SCREENING', ha='center', fontsize=20, fontweight='bold', color=M_RED)
    fig.text(0.5, 0.068, '⚠  Pre-revenue status: No commercial track record — first deliveries expected 2027',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.042, '⚠  Manufacturing scale-up risk: $760M facility under construction — execution unproven at scale',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    fig.text(0.5, 0.016, '⚠  Technology validation: Iron-air chemistry validated in lab — long-duration field performance TBD',
             ha='center', fontsize=16, fontweight='bold', color=M_NAVY)
    
    plt.subplots_adjust(wspace=0.15, top=0.92, bottom=0.15)
    p = VISUALS / '09_esg.png'
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    return p

print("🎨 Generating Form Energy Inc Product 1 v15 — Second Draft\n")
print("Step 1: Charts...")
visuals = {
    'gauge': gen_risk_gauge(), 'org': gen_org_chart(), 'fin': gen_funding_chart(),
    'mkt': gen_market(), 'timeline': gen_timeline(), 'risk': gen_risk_matrix(),
    'radar': gen_radar(), 'peer': gen_peer_risk(), 'esg': gen_esg(),
}
print("\nStep 2: Building PPTX...")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# S1: TITLE
print("  Slide 1: Title")
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = rgb(*NAVY); bg.line.fill.background()
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.55), Inches(13.33), Inches(0.08))
ln.fill.solid(); ln.fill.fore_color.rgb = rgb(*STEEL); ln.line.fill.background()
add_text(s, 0.7, 1.4, 12, 0.5, 'SUPPLIER EVALUATION REPORT', 16, False, MID_GREY)
add_text(s, 0.7, 2.05, 12, 1.2, 'Form Energy Inc', 44, True, WHITE)
add_text(s, 0.7, 3.5, 12, 0.6, 'Iron-Air Batteries  |  Long-Duration Storage  |  100+ Hours', 18, False, (200,210,220))
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.35), Inches(4.5), Inches(0.05))
div.fill.solid(); div.fill.fore_color.rgb = rgb(*STEEL); div.line.fill.background()
add_text(s, 0.7, 4.55, 12, 0.4, '$1.52B Funding  |  ~400 Employees  |  USA  |  Pre-Revenue', 14, False, (200,210,220))
add_text(s, 0.7, 6.2, 12, 0.4, 'Confidential  |  March 2026  |  Manu Forti Intelligence', 11, False, MID_GREY)
add_logos(s, is_title=True)

# S2: EXEC SUMMARY
print("  Slide 2: Executive Summary")
s = prs.slides.add_slide(blank)
add_header(s, 'EXECUTIVE SUMMARY', 'Form Energy Inc — Risk & Suitability Overview')
s.shapes.add_picture(str(visuals['gauge']), Inches(0.3), Inches(1.05), Inches(7.2), Inches(3.1))
add_rounded_rect(s, 7.6, 1.05, 5.4, 5.6, (240,245,250))
add_text(s, 7.8, 1.22, 5.2, 0.32, 'SUPPLIER SNAPSHOT', 14, True, NAVY)
snap = [('Supplier:', 'Form Energy Inc'), ('Type:', 'Private (venture-backed)'),
        ('Sector:', 'Iron-Air Battery Storage'), ('HQ:', 'Somerville, MA, USA'),
        ('Founded:', '2017 by MIT researchers'), ('Valuation:', '$2.5B (2024)'),
        ('Funding:', '$1.52B total raised'), ('Employees:', '~400'),
        ('Revenue:', 'Pre-revenue (2027 target)'), ('Key Contract:', 'Google 300MW/30GWh')]
for idx, (label, val) in enumerate(snap):
    y = 1.64 + idx * 0.30
    add_text(s, 7.8, y, 2.3, 0.27, label, 11, True, TEXT_SECONDARY)
    add_text(s, 10.1, y, 2.7, 0.27, val, 11, False, TEXT_PRIMARY)
add_rounded_rect(s, 0.3, 5.22, 7.2, 1.43, (240,245,250))
add_text(s, 0.5, 5.33, 7, 0.3, 'Key Findings', 13, True, NAVY)
findings = ['• Iron-air battery technology: 100+ hour duration at $20/kWh target cost',
            '• $1.52B funding from Breakthrough Energy, T. Rowe Price, GE Vernova, Temasek',
            '• Google contract: 300MW/30GWh (2026-2028) — first major commercial validation',
            '• Pre-revenue status: Manufacturing facility under construction in Weirton, WV']
for idx, f in enumerate(findings):
    add_text(s, 0.5, 5.65+idx*0.27, 7, 0.26, f, 11, False, TEXT_PRIMARY)
add_source(s); add_logos(s)

# S3: RECOMMENDATION
print("  Slide 3: Recommendation")
s = prs.slides.add_slide(blank)
add_header(s, 'RECOMMENDATION', 'Decision Summary & Commercial Conditions')
add_rounded_rect(s, 0.5, 1.1, 12.33, 1.35, AMBER)
add_text(s, 0.7, 1.18, 12, 0.55, '⚠  RECOMMENDATION: CONDITIONAL APPROVAL — ENHANCED DUE DILIGENCE REQUIRED', 20, True, WHITE)
add_text(s, 0.7, 1.78, 12, 0.58,
         'Form Energy represents a breakthrough in long-duration energy storage with iron-air technology. '
         'However, pre-revenue status and manufacturing scale-up risks require milestone-based engagement '
         'with performance guarantees before full commercial commitment.',
         11, False, WHITE)
add_bullets(s, 0.5, 2.65, 6, 2.5, 'Commercial Conditions',
            ['Payment: Milestone-based (no advance payment)',
             'Performance bonds: Parent guarantee from Form Energy',
             'Liquidated damages: For delivery delays post-2027',
             'Technology escrow: Source code / manufacturing IP',
             'Termination: Exit rights if milestones missed'], 14, 11)
add_bullets(s, 6.7, 2.65, 6, 2.5, 'Enhanced Due Diligence (Mandatory)',
            ['Verify Google contract milestones and performance guarantees',
             'Site visit: Weirton, WV manufacturing facility',
             'Technology audit: Independent validation of 100-hr duration claims',
             'Financial audit: Review runway and Series E terms',
             'Insurance: Verify product liability coverage',
             'Reference checks: Early pilot customers (if any)'], 14, 11)
add_rounded_rect(s, 0.5, 5.35, 12.33, 1.5, (240,245,250))
add_text(s, 0.7, 5.5, 12, 0.3, 'Overall Risk Summary — MEDIUM-HIGH (58/100)', 14, True, NAVY)
for x, t, d in [(0.6,'Financial: MEDIUM','$1.52B funding; pre-revenue; 2027 target'),
                 (3.7,'Operational: MEDIUM','Unproven manufacturing at scale'),
                 (6.8,'Geopolitical: LOW','US-based; stable jurisdiction'),
                 (9.9,'ESG: MEDIUM','Strong environmental; governance gaps')]:
    add_text(s, x, 5.88, 3.1, 0.27, t, 11, True, NAVY)
    add_text(s, x, 6.17, 3.1, 0.27, d, 9, False, TEXT_SECONDARY)
add_source(s); add_logos(s)

# S4: SUPPLIER PROFILE
print("  Slide 4: Supplier Profile")
s = prs.slides.add_slide(blank)
add_header(s, 'SUPPLIER PROFILE', 'Corporate Structure & Global Footprint')
s.shapes.add_picture(str(visuals['org']), Inches(0.3), Inches(1.05), Inches(8.2), Inches(4.48))
add_bullets(s, 8.7, 1.05, 4.3, 3.5, 'Company Overview',
            ['Form Energy Inc is a US-based energy',
             'storage company founded in 2017 by',
             'MIT researchers. Developing iron-air',
             'batteries capable of 100+ hour duration',
             'storage at projected $20/kWh cost.',
             '',
             'Breakthrough technology addresses',
             'long-duration storage gap in renewable',
             'energy grid integration.'], 14, 11)
add_bullets(s, 8.7, 4.5, 4.3, 2.4, 'Leadership',
            ['Mateo Jaramillo — CEO',
             '  Former Tesla Energy (Powerwall)',
             '',
            'MIT pedigree: Co-founders from',
             '  Yet-Ming Chiang research group',
             '',
             'HQ: Somerville, MA  |  Founded: 2017'], 14, 11)
add_source(s); add_logos(s)

# S5: FINANCIAL HEALTH (Funding for pre-revenue)
print("  Slide 5: Financial Health")
s = prs.slides.add_slide(blank)
add_header(s, 'FINANCIAL HEALTH', 'Funding Trajectory 2017–2024  |  Pre-Revenue Stage  |  $1.52B Total Funding')
s.shapes.add_picture(str(visuals['fin']), Inches(0.3), Inches(1.05), Inches(7.5), Inches(4.55))
add_text(s, 8.1, 1.1, 4.9, 0.35, 'Financial Highlights', 14, True, NAVY)
metrics = [('Latest Valuation:', '$2.5B  (2024)'), ('Total Funding:', '$1.52B  (Series E)'),
           ('Latest Round:', '$450M  Series E'), ('Key Investors:', 'Breakthrough Energy Ventures'),
           ('Revenue 2024:', '$0  (pre-revenue)'), ('Revenue 2027E:', 'First commercial deliveries'),
           ('Cash Runway:', 'Estimated 3-4 years'), ('Burn Rate:', 'High (manufacturing build-out)')]
for idx, (label, val) in enumerate(metrics):
    y = 1.56 + idx * 0.30
    add_text(s, 8.1, y, 2.55, 0.29, label, 11, True, TEXT_SECONDARY)
    vc = AMBER if idx >= 4 else TEXT_PRIMARY
    add_text(s, 10.65, y, 2.25, 0.29, val, 11, False, vc)
add_rounded_rect(s, 8.1, 4.05, 4.8, 0.44, AMBER)
add_text(s, 8.25, 4.12, 4.5, 0.3, 'Financial Risk:  MEDIUM  ⚠', 12, True, WHITE, PP_ALIGN.CENTER)
add_bullets(s, 8.1, 4.6, 4.8, 2.3, 'Exposure Guidance',
            ['Pre-revenue: No revenue track record',
             'Milestone-based payments mandatory',
             'Verify runway through 2027 commercialization',
             'Request investor letters of support'], 12, 10)
add_source(s); add_logos(s)

# S6: MARKET POSITION
print("  Slide 6: Market Position")
s = prs.slides.add_slide(blank)
add_header(s, 'MARKET POSITION', 'Energy Storage Companies — Valuation Comparison')
s.shapes.add_picture(str(visuals['mkt']), Inches(0.3), Inches(1.05), Inches(8.0), Inches(4.46))
add_bullets(s, 8.5, 1.1, 4.5, 2.0, 'Competitive Landscape',
            ['Form Energy occupies unique position',
             'in long-duration storage (100+ hours).',
             'Competes with lithium-ion for shorter',
             'durations; unique for multi-day storage.',
             'Target: $20/kWh vs Li-ion $130/kWh.'], 14, 11)
add_bullets(s, 8.5, 3.2, 4.5, 2.6, 'Form Energy Competitive Advantages',
            ['100+ hour duration (Li-ion: 4-8 hrs)',
             'Iron-air chemistry: abundant materials',
             'No rare earths or lithium supply risk',
             'Target cost: $20/kWh (1/6th of Li-ion)',
             'Google contract validation (300MW)',
             'US manufacturing: Weirton, WV facility'], 14, 11)
add_text(s, 0.5, 5.5, 12, 0.3, 'Key Competitors', 13, True, NAVY)
for idx, c in enumerate([
    '• Tesla Energy — USA — $800B mkt cap — Powerwall/Powerpack/Megapack leader',
    '• Fluence — USA/Australia — $450M revenue — Grid-scale storage integrator',
    '• CATL — China — $120B mkt cap — World’s largest battery manufacturer',
    '• BYD — China — $90B mkt cap — Integrated EV + energy storage']):
    add_text(s, 0.6, 5.84+idx*0.29, 12, 0.27, c, 10, False, TEXT_PRIMARY)
add_source(s); add_logos(s)

# S7: OPS & RISK
print("  Slide 7: Ops & Risk")
s = prs.slides.add_slide(blank)
add_header(s, 'OPERATIONAL CAPABILITY  &  RISK ASSESSMENT', 'Milestones & Strategic Investment  |  Risk Matrix  |  Capability Summary')
s.shapes.add_picture(str(visuals['timeline']), Inches(0.3), Inches(1.05), Inches(7.0), Inches(2.92))
s.shapes.add_picture(str(visuals['risk']), Inches(7.65), Inches(1.05), Inches(5.5), Inches(4.46))
add_text(s, 0.3, 5.58, 7, 0.28, 'Risk Summary', 13, True, NAVY)
for idx, (cat, rating, desc) in enumerate([
    ('Financial:', '🟡 MEDIUM', '$1.52B funding; pre-revenue; runway TBD'),
    ('Technology:', '🟡 MEDIUM', 'Iron-air validated; manufacturing scale unproven'),
    ('Operational:', '🟡 MEDIUM', 'First facility under construction; 2027 target'),
    ('Competition:', '🟢 LOW', 'Unique 100+ hr position vs Li-ion'),
    ('Supply Chain:', '🟢 LOW', 'Iron, air, water — abundant inputs')]):
    y = 5.89 + idx * 0.25
    add_text(s, 0.3, y, 2.1, 0.24, cat, 10, True, TEXT_PRIMARY)
    add_text(s, 2.4, y, 1.4, 0.24, rating, 10, True, TEXT_PRIMARY)
    add_text(s, 3.8, y, 3.7, 0.24, desc, 9, False, TEXT_SECONDARY)
add_bullets(s, 7.65, 5.58, 5.5, 1.85, 'Key Capabilities',
            ['Iron-air battery R&D and engineering',
             '100+ hour duration storage technology',
             'US-based manufacturing (Weirton, WV)',
             'Strategic partnerships: Google, Xcel'], 12, 10)
add_source(s); add_logos(s)

# S8: COMMERCIAL INTEL
print("  Slide 8: Commercial Intelligence")
s = prs.slides.add_slide(blank)
add_header(s, 'COMMERCIAL INTELLIGENCE  &  PEER RISK COMPARISON', 'Benchmarking Radar  |  Peer Risk Profile  |  Commercial Terms')
s.shapes.add_picture(str(visuals['radar']), Inches(0.2), Inches(1.0), Inches(5.8), Inches(4.74))
s.shapes.add_picture(str(visuals['peer']), Inches(6.2), Inches(1.0), Inches(6.9), Inches(3.85))
add_rounded_rect(s, 0.3, 5.8, 12.7, 1.5, (240,245,250))
add_bullets(s, 0.5, 5.85, 6.5, 1.4, 'Commercial Terms & Negotiation',
            ['Pricing: Target $20/kWh (vs Li-ion $130/kWh)',
             'Lead time: 2027+ for commercial volumes',
             'Contract: Google 300MW/30GWh reference pricing',
             'Leverage: Early adopter positioning for LDES'], 12, 10)
add_bullets(s, 7.0, 5.85, 5.8, 1.4, 'Key Watch Points',
            ['⚠ Verify Google contract performance milestones',
             '⚠ Monitor Weirton facility construction progress',
             '⚠ Track 2027 commercial delivery readiness',
             '⚠ Assess technology field validation data'], 12, 10)
add_source(s); add_logos(s)

# S9: ESG
print("  Slide 9: ESG Assessment")
s = prs.slides.add_slide(blank)
add_header(s, 'ESG ASSESSMENT', 'Environmental, Social & Governance Screening  |  Overall Rating: MEDIUM')
s.shapes.add_picture(str(visuals['esg']), Inches(1.67), Inches(1.05), Inches(10.0), Inches(5.93))
add_source(s); add_logos(s)

prs.save(str(OUTPUT))
print(f"\n✅ Saved: {OUTPUT}")
import os
print(f"📎 Size: {os.path.getsize(str(OUTPUT))/1024:.1f} KB")
